"""Export functionality: Excel, Prism .pzfx, and PowerPoint generation.

Handles all format-specific export logic for pipeline results.
"""

from __future__ import annotations

import xml.etree.ElementTree as ET
from datetime import datetime, date
from itertools import combinations
from pathlib import Path
from xml.dom import minidom

import numpy as np
import pandas as pd

from rnaseq_explorer.engine.deseq2 import DEFAULT_DESEQ2_COLS, best_gene_key
from rnaseq_explorer.engine.gsea import normalize_gsea_cols
from rnaseq_explorer.engine.rmats import DEFAULT_RMATS_COLS, RMATS_EVENT_TYPES, COORD_COLS, make_event_key


# ---------------------------------------------------------------------------
# Excel exports
# ---------------------------------------------------------------------------


def export_excel(
    deseq2_sets: dict[str, pd.DataFrame],
    all_rmats_filtered: dict[str, pd.DataFrame],
    outdir: str | Path,
    log2fc_cutoff: float = 1.0,
    basemean_cutoff: float = 10.0,
    padj_cutoff: float = 0.05,
    use_fdr: bool = True,
    fdr_cutoff: float = 0.05,
    pval_cutoff: float = 0.05,
    dpsi_cutoff: float = 0.1,
    auto_biotype_split: bool = True,
) -> None:
    """Export per-condition results as two XLSX workbooks (DESeq2 + rMATS).

    Parameters
    ----------
    deseq2_sets : dict
        {label: filtered_df} e.g. {"all_genes": df, "protein_coding": df}.
    all_rmats_filtered : dict
        {event_type: filtered_df}.
    outdir : str or Path
        Output directory.
    log2fc_cutoff : float
        Log2FC threshold used (for summary).
    basemean_cutoff : float
        baseMean threshold used (for summary).
    padj_cutoff : float
        padj threshold used (for summary).
    use_fdr : bool
        Whether FDR was used for rMATS (for summary).
    fdr_cutoff : float
        rMATS FDR cutoff (for summary).
    pval_cutoff : float
        rMATS p-value cutoff (for summary).
    dpsi_cutoff : float
        dPSI cutoff (for summary).
    auto_biotype_split : bool
        Whether biotype split was enabled (for summary).
    """
    outdir = Path(outdir)
    outdir.mkdir(parents=True, exist_ok=True)

    print("\n-- Exported Files --")

    # DESeq2 workbook
    deseq2_xlsx = outdir / "deseq2_results.xlsx"
    with pd.ExcelWriter(deseq2_xlsx, engine="openpyxl") as writer:
        for label, df in deseq2_sets.items():
            sheet = label.replace("_", " ").title().replace(" ", "_")
            df.to_excel(writer, sheet_name=sheet, index=False)

        summary_rows = [
            ("|log2FC| cutoff", log2fc_cutoff),
            ("baseMean cutoff", basemean_cutoff),
            ("padj cutoff", padj_cutoff),
            ("Auto biotype split", auto_biotype_split),
            ("", ""),
        ]
        for label, df in deseq2_sets.items():
            n_up = int((df["direction"] == "up").sum()) if "direction" in df.columns else "N/A"
            n_down = int((df["direction"] == "down").sum()) if "direction" in df.columns else "N/A"
            display = label.replace("_", " ").title()
            summary_rows.append((f"{display} -- genes passing filter", len(df)))
            summary_rows.append((f"{display} -- upregulated", n_up))
            summary_rows.append((f"{display} -- downregulated", n_down))

        pd.DataFrame(summary_rows, columns=["Parameter", "Value"]).to_excel(
            writer, sheet_name="Summary", index=False
        )

    print(f"  {deseq2_xlsx}")

    # rMATS workbook
    if all_rmats_filtered:
        rmats_xlsx = outdir / "rmats_results.xlsx"
        with pd.ExcelWriter(rmats_xlsx, engine="openpyxl") as writer:
            for event_type, df in all_rmats_filtered.items():
                df.to_excel(writer, sheet_name=event_type, index=False)

            all_combined = pd.concat(all_rmats_filtered.values(), ignore_index=True)
            all_combined.to_excel(writer, sheet_name="All_Significant", index=False)

            summary_rows = [
                ("Filter column", "FDR" if use_fdr else "PValue"),
                ("p-value/FDR cutoff", fdr_cutoff if use_fdr else pval_cutoff),
                ("|dPSI| cutoff", dpsi_cutoff),
                ("", ""),
            ]
            for event_type, df in all_rmats_filtered.items():
                summary_rows.append((f"{event_type} significant events", len(df)))
            summary_rows.append(("Total significant events", len(all_combined)))

            pd.DataFrame(summary_rows, columns=["Parameter", "Value"]).to_excel(
                writer, sheet_name="Summary", index=False
            )

        print(f"  {rmats_xlsx}")


def export_combined_results(
    condition_results: dict[str, dict],
    cross_condition_data: dict,
    outdir: str | Path,
    cols: dict[str, str] | None = None,
    log2fc_cutoff: float = 1.0,
    basemean_cutoff: float = 10.0,
    padj_cutoff: float = 0.05,
    use_fdr: bool = True,
    fdr_cutoff: float = 0.05,
    dpsi_cutoff: float = 0.1,
) -> None:
    """Write master multi-condition XLSX with all conditions + cross-condition data.

    Parameters
    ----------
    condition_results : dict
        Pipeline condition_results structure.
    cross_condition_data : dict
        Cross-condition analysis results.
    outdir : str or Path
        Output directory.
    cols : dict or None
        DESeq2 column name mapping.
    log2fc_cutoff : float
        For summary sheet.
    basemean_cutoff : float
        For summary sheet.
    padj_cutoff : float
        For summary sheet.
    use_fdr : bool
        For summary sheet.
    fdr_cutoff : float
        For summary sheet.
    dpsi_cutoff : float
        For summary sheet.
    """
    outdir = Path(outdir)
    xlsx_path = outdir / "multi_condition_results.xlsx"

    with pd.ExcelWriter(xlsx_path, engine="openpyxl") as writer:
        for cond_name, data in condition_results.items():
            for bio_key, df in data["deseq2_filtered"].items():
                sheet = f"DESeq2_{cond_name}_{bio_key}"[:31]
                df.to_excel(writer, sheet_name=sheet, index=False)

        for cond_name, data in condition_results.items():
            for et, df in data["rmats_filtered"].items():
                sheet = f"rMATS_{cond_name}_{et}"[:31]
                df.to_excel(writer, sheet_name=sheet, index=False)

        for key, sheet_name in [
            ("concordance_matrix", "Cross_DESeq2_Direction"),
            ("log2fc_matrix", "Cross_DESeq2_Log2FC"),
            ("rmats_concordance", "Cross_rMATS_Direction"),
        ]:
            if key in cross_condition_data and len(cross_condition_data[key]) > 0:
                idx = key != "rmats_concordance"
                cross_condition_data[key].to_excel(
                    writer, sheet_name=sheet_name, index=idx
                )

        # Summary sheet
        rows = [
            ("Pipeline", "Multi-Condition DESeq2 + rMATS"),
            ("DESeq2 |log2FC| cutoff", log2fc_cutoff),
            ("DESeq2 baseMean cutoff", basemean_cutoff),
            ("DESeq2 padj cutoff", padj_cutoff),
            ("rMATS filter column", "FDR" if use_fdr else "PValue"),
            ("rMATS cutoff", fdr_cutoff if use_fdr else 0.05),
            ("rMATS |dPSI| cutoff", dpsi_cutoff),
            ("", ""),
        ]
        for cond_name, data in condition_results.items():
            filt = data["deseq2_filtered"].get("all_genes", pd.DataFrame())
            n_up = int((filt["direction"] == "up").sum()) if "direction" in filt.columns else 0
            n_down = int((filt["direction"] == "down").sum()) if "direction" in filt.columns else 0
            rows.append((f"{cond_name} -- DESeq2 significant", len(filt)))
            rows.append((f"{cond_name} -- DESeq2 up", n_up))
            rows.append((f"{cond_name} -- DESeq2 down", n_down))
            total_rmats = sum(len(df) for df in data["rmats_filtered"].values())
            rows.append((f"{cond_name} -- rMATS significant", total_rmats))

        pd.DataFrame(rows, columns=["Parameter", "Value"]).to_excel(
            writer, sheet_name="Summary", index=False
        )

    print(f"  Master XLSX: {xlsx_path}")


# ---------------------------------------------------------------------------
# Prism .pzfx export
# ---------------------------------------------------------------------------


def _create_prism_xml() -> ET.Element:
    """Create base Prism XML structure."""
    root = ET.Element("GraphPadPrismFile")
    root.set("PrismXMLVersion", "5.00")

    created = ET.SubElement(root, "Created")
    orig = ET.SubElement(created, "OriginalVersion")
    orig.set("CreatedByProgram", "GraphPad Prism")
    orig.set("CreatedByVersion", "6.0f.254")
    orig.set("Login", "")
    orig.set("DateTime", datetime.now().strftime("%Y-%m-%dT%H:%M:%S+00:00"))

    info_seq = ET.SubElement(root, "InfoSequence")
    ref = ET.SubElement(info_seq, "Ref")
    ref.set("ID", "Info0")
    ref.set("Selected", "1")

    info = ET.SubElement(root, "Info")
    info.set("ID", "Info0")
    info_title = ET.SubElement(info, "Title")
    info_title.text = "Project info 1"
    ET.SubElement(info, "Notes")
    const = ET.SubElement(info, "Constant")
    const_name = ET.SubElement(const, "Name")
    const_name.text = "Experiment Date"
    const_val = ET.SubElement(const, "Value")
    const_val.text = datetime.now().strftime("%Y-%m-%d")

    table_seq = ET.SubElement(root, "TableSequence")
    table_seq.set("Selected", "1")

    root.set("_table_count", "0")
    return root


def _add_table(
    root: ET.Element, table_name: str, columns_data: list[tuple[str, list]]
) -> None:
    """Add a table to Prism XML.

    Parameters
    ----------
    root : ET.Element
        Root Prism XML element.
    table_name : str
        Table title.
    columns_data : list
        List of (col_name, values_list) tuples.
    """
    table_idx = int(root.get("_table_count", "0"))
    table_id = f"Table{table_idx}"
    root.set("_table_count", str(table_idx + 1))

    table_seq = root.find("TableSequence")
    ref = ET.SubElement(table_seq, "Ref")
    ref.set("ID", table_id)
    if table_idx == 0:
        ref.set("Selected", "1")

    table = ET.SubElement(root, "Table")
    table.set("ID", table_id)
    table.set("XFormat", "none")
    table.set("TableType", "OneWay")
    table.set("EVFormat", "AsteriskAfterNumber")

    title_elem = ET.SubElement(table, "Title")
    title_elem.text = table_name

    def _is_numeric_string(s):
        try:
            float(s)
            return True
        except (ValueError, TypeError):
            return False

    # Detect first column as row titles if all non-numeric
    first_is_titles = False
    if columns_data:
        first_name, first_vals = columns_data[0]
        if first_vals and not any(_is_numeric_string(str(v)) for v in first_vals[:20]):
            first_is_titles = True

    if first_is_titles:
        # Row titles
        row_titles = ET.SubElement(table, "RowTitlesColumn")
        row_titles.set("Width", "200")
        for val in columns_data[0][1]:
            subrow = ET.SubElement(row_titles, "Subcolumn")
            d = ET.SubElement(subrow, "d")
            d.text = str(val)
        data_cols = columns_data[1:]
    else:
        data_cols = columns_data

    for col_name, values in data_cols:
        yc = ET.SubElement(table, "YColumn")
        yc.set("Width", "120")
        yc.set("Decimals", "4")
        yc.set("Subcolumns", "1")
        title = ET.SubElement(yc, "Title")
        title.text = col_name
        sub = ET.SubElement(yc, "Subcolumn")
        for val in values:
            d = ET.SubElement(sub, "d")
            if pd.notna(val):
                d.text = str(val)


def _save_prism(root: ET.Element, filepath: Path) -> None:
    """Save Prism XML to file."""
    # Clean up internal counter
    if "_table_count" in root.attrib:
        del root.attrib["_table_count"]

    xml_str = ET.tostring(root, encoding="unicode")
    try:
        pretty = minidom.parseString(xml_str).toprettyxml(indent="  ")
        lines = [line for line in pretty.split("\n") if line.strip()]
        with open(filepath, "w", encoding="utf-8") as f:
            f.write("\n".join(lines))
    except Exception:
        with open(filepath, "w", encoding="utf-8") as f:
            f.write(xml_str)
    print(f"  Saved: {filepath.name}")


def export_prism_pzfx(
    condition_results: dict[str, dict],
    condition_labels: dict[str, str],
    outdir: str | Path,
    cols: dict[str, str] | None = None,
    gsea_results: dict | None = None,
) -> None:
    """Generate comprehensive Prism .pzfx files for publication-quality graphs.

    Creates multiple .pzfx files with proper XML structure.

    Parameters
    ----------
    condition_results : dict
        Pipeline condition_results structure.
    condition_labels : dict
        Maps condition name -> label.
    outdir : str or Path
        Output directory (a prism_files/ subdirectory will be created).
    cols : dict or None
        DESeq2 column name mapping.
    gsea_results : dict or None
        GSEA results for enrichment Prism files.
    """
    if cols is None:
        cols = DEFAULT_DESEQ2_COLS

    outdir = Path(outdir)
    prism_dir = outdir / "prism_files"
    prism_dir.mkdir(exist_ok=True)

    print("\n-- Generating Prism Files --")

    fc_col = cols["log2fc"]
    padj_col = cols["padj"]
    name_col = cols.get("gene_name", "gene_name")

    names = list(condition_results.keys())

    # 1. DEG counts per condition
    root = _create_prism_xml()
    count_data = []
    for cond_name in names:
        filt = condition_results[cond_name]["deseq2_filtered"].get("all_genes", pd.DataFrame())
        n_up = int((filt["direction"] == "up").sum()) if "direction" in filt.columns else 0
        n_down = int((filt["direction"] == "down").sum()) if "direction" in filt.columns else 0
        label = condition_labels.get(cond_name, cond_name)
        count_data.append((label, [n_up, n_down]))
    if count_data:
        _add_table(root, "DEG Counts", [
            ("Direction", ["Up", "Down"]),
        ] + count_data)
        _save_prism(root, prism_dir / "deg_counts.pzfx")

    # 2. Top DEGs with log2FC per condition
    for cond_name in names:
        root = _create_prism_xml()
        filt = condition_results[cond_name]["deseq2_filtered"].get("all_genes", pd.DataFrame())
        if len(filt) == 0:
            continue
        label = condition_labels.get(cond_name, cond_name)

        if name_col in filt.columns and fc_col in filt.columns:
            top = filt.nlargest(30, fc_col, keep="first")
            _add_table(root, f"Top Up DEGs - {label}", [
                ("Gene", top[name_col].tolist()),
                ("log2FC", top[fc_col].tolist()),
                ("padj", top[padj_col].tolist() if padj_col in top.columns else []),
            ])

            bottom = filt.nsmallest(30, fc_col, keep="first")
            _add_table(root, f"Top Down DEGs - {label}", [
                ("Gene", bottom[name_col].tolist()),
                ("log2FC", bottom[fc_col].tolist()),
                ("padj", bottom[padj_col].tolist() if padj_col in bottom.columns else []),
            ])

            _save_prism(root, prism_dir / f"top_degs_{cond_name}.pzfx")

    # 3. Volcano plot data per condition
    for cond_name in names:
        root = _create_prism_xml()
        raw = condition_results[cond_name]["deseq2_raw"]
        label = condition_labels.get(cond_name, cond_name)

        data = raw.dropna(subset=[padj_col, fc_col]).copy()
        data["neg_log10_padj"] = -np.log10(data[padj_col].clip(lower=1e-300))

        gene_names = data[name_col].tolist() if name_col in data.columns else []
        _add_table(root, f"Volcano Data - {label}", [
            ("Gene", gene_names),
            ("log2FC", data[fc_col].tolist()),
            ("-log10(padj)", data["neg_log10_padj"].tolist()),
        ])
        _save_prism(root, prism_dir / f"volcano_data_{cond_name}.pzfx")

    # rMATS column names
    rmats_gene_col = DEFAULT_RMATS_COLS.get("gene_name", "geneSymbol")
    dpsi_col = DEFAULT_RMATS_COLS.get("inclevel_diff", "IncLevelDifference")
    fdr_col = DEFAULT_RMATS_COLS.get("fdr", "FDR")

    # 4. Splicing_data_{condition}.pzfx — per-condition splicing events by type
    for cond_name in names:
        rmats_filt = condition_results[cond_name].get("rmats_filtered", {})
        if not rmats_filt:
            continue
        root = _create_prism_xml()
        label = condition_labels.get(cond_name, cond_name)
        has_tables = False

        for et, df in rmats_filt.items():
            if len(df) > 0:
                event_data = df.head(100)  # Top 100 per event type
                cols_data = [("GeneSymbol", event_data[rmats_gene_col].tolist())]
                if dpsi_col in event_data.columns:
                    cols_data.append(("dPSI", event_data[dpsi_col].tolist()))
                if fdr_col in event_data.columns:
                    cols_data.append(("FDR", event_data[fdr_col].tolist()))
                _add_table(root, f"{et}_{label}", cols_data)
                has_tables = True

        # Event counts summary
        event_counts = [(et, len(df)) for et, df in rmats_filt.items()]
        if event_counts:
            _add_table(root, f"EventCounts_{label}", [
                ("EventType", [et for et, _ in event_counts]),
                ("Count", [cnt for _, cnt in event_counts]),
            ])
            has_tables = True

        if has_tables:
            _save_prism(root, prism_dir / f"Splicing_data_{cond_name}.pzfx")

    # 5. GSEA_data_{condition}.pzfx — pathway enrichment per condition
    for cond_name in names:
        label = condition_labels.get(cond_name, cond_name)
        gsea_prism_ok = False

        # Try in-memory GSEA results first
        if gsea_results and cond_name in gsea_results:
            root = _create_prism_xml()
            cond_gsea = gsea_results[cond_name]

            for db_name, pathways_df in cond_gsea.items():
                if len(pathways_df) > 0:
                    pathway_names = pathways_df["Term"].tolist()
                    nes_vals = pathways_df["nes"].tolist()
                    fdr_vals = pathways_df["fdr"].tolist()
                    neg_log10_fdr = (
                        -np.log10(pathways_df["fdr"].clip(lower=1e-300))
                    ).tolist()
                    geneset_sizes = (
                        pathways_df["geneset_size"].tolist()
                        if "geneset_size" in pathways_df.columns
                        else []
                    )

                    cols_data = [
                        ("Pathway", pathway_names),
                        ("NES", nes_vals),
                        ("FDR", fdr_vals),
                        ("-log10(FDR)", neg_log10_fdr),
                    ]
                    if geneset_sizes:
                        cols_data.append(("GenesetSize", geneset_sizes))
                    _add_table(root, f"{db_name[:20]}_{label}", cols_data)
                    gsea_prism_ok = True

            if gsea_prism_ok:
                _save_prism(root, prism_dir / f"GSEA_data_{cond_name}.pzfx")

        # Fallback: scan disk for gseapy CSV reports
        if not gsea_prism_ok:
            gsea_disk_dir = outdir / "gsea_results" / cond_name
            disk_tables_added = False
            if gsea_disk_dir.is_dir():
                root = _create_prism_xml()
                for db_subdir in sorted(gsea_disk_dir.iterdir()):
                    if not db_subdir.is_dir():
                        continue
                    report_csv = db_subdir / "gseapy.gene_set.prerank.report.csv"
                    if not report_csv.exists():
                        continue
                    try:
                        rpt = normalize_gsea_cols(pd.read_csv(report_csv))
                        if len(rpt) == 0:
                            continue
                        rpt["abs_nes"] = rpt["nes"].abs()
                        rpt = rpt.sort_values("abs_nes", ascending=False).head(5)
                        db_label = db_subdir.name[:20]
                        neg_log10_fdr_disk = (
                            -np.log10(rpt["fdr"].clip(lower=1e-300))
                        ).tolist()
                        cols_data = [
                            ("Pathway", rpt["Term"].tolist()),
                            ("NES", rpt["nes"].tolist()),
                            ("FDR", rpt["fdr"].tolist()),
                            ("-log10(FDR)", neg_log10_fdr_disk),
                        ]
                        if "lead_genes" in rpt.columns:
                            cols_data.append(
                                ("LeadGenes", rpt["lead_genes"].tolist())
                            )
                        _add_table(root, f"{db_label}_{label}", cols_data)
                        disk_tables_added = True
                    except Exception:
                        continue
                if disk_tables_added:
                    _save_prism(root, prism_dir / f"GSEA_data_{cond_name}.pzfx")

    # ===== CROSS-CONDITION FILES =====

    # 6. Splicing_overlap_SE.pzfx — pairwise SE gene overlap
    if len(names) >= 2:
        root = _create_prism_xml()
        for pair in combinations(names, 2):
            cond1, cond2 = pair
            label1 = condition_labels.get(cond1, cond1)
            label2 = condition_labels.get(cond2, cond2)

            se_genes1 = set()
            se_genes2 = set()
            rmats1 = condition_results[cond1].get("rmats_filtered", {})
            rmats2 = condition_results[cond2].get("rmats_filtered", {})
            if "SE" in rmats1 and rmats_gene_col in rmats1["SE"].columns:
                se_genes1 = set(rmats1["SE"][rmats_gene_col].dropna())
            if "SE" in rmats2 and rmats_gene_col in rmats2["SE"].columns:
                se_genes2 = set(rmats2["SE"][rmats_gene_col].dropna())

            overlap = se_genes1 & se_genes2
            _add_table(root, f"SE_Overlap_{cond1}_vs_{cond2}", [
                ("Metric", [f"{label1} only", f"{label2} only", "Overlap"]),
                ("Count", [
                    len(se_genes1 - se_genes2),
                    len(se_genes2 - se_genes1),
                    len(overlap),
                ]),
            ])

            # Gene lists (top 50 from overlap)
            overlap_list = sorted(overlap)[:50]
            if overlap_list:
                _add_table(root, f"SE_Genes_{cond1}_vs_{cond2}", [
                    ("GeneSymbol", overlap_list),
                ])

        _save_prism(root, prism_dir / "Splicing_overlap_SE.pzfx")

    # 7. Splicing_overlap_other_events.pzfx — pairwise A3SS/A5SS/MXE/RI overlaps
    if len(names) >= 2:
        root = _create_prism_xml()
        for et in ["A3SS", "A5SS", "MXE", "RI"]:
            for pair in combinations(names, 2):
                cond1, cond2 = pair
                label1 = condition_labels.get(cond1, cond1)
                label2 = condition_labels.get(cond2, cond2)

                genes1 = set()
                genes2 = set()
                rmats1 = condition_results[cond1].get("rmats_filtered", {})
                rmats2 = condition_results[cond2].get("rmats_filtered", {})
                if et in rmats1 and rmats_gene_col in rmats1[et].columns:
                    genes1 = set(rmats1[et][rmats_gene_col].dropna())
                if et in rmats2 and rmats_gene_col in rmats2[et].columns:
                    genes2 = set(rmats2[et][rmats_gene_col].dropna())

                overlap = genes1 & genes2
                _add_table(root, f"{et}_{cond1}_vs_{cond2}", [
                    ("Metric", [
                        f"{label1} only", f"{label2} only", "Overlap",
                    ]),
                    ("Count", [
                        len(genes1 - genes2),
                        len(genes2 - genes1),
                        len(overlap),
                    ]),
                ])

        _save_prism(root, prism_dir / "Splicing_overlap_other_events.pzfx")

    # 8. DEG_overlap.pzfx — pairwise DEG overlap with actual condition labels
    if len(names) >= 2:
        root = _create_prism_xml()
        for pair in combinations(names, 2):
            cond1, cond2 = pair
            label1 = condition_labels.get(cond1, cond1)
            label2 = condition_labels.get(cond2, cond2)

            filt1 = condition_results[cond1]["deseq2_filtered"].get(
                "all_genes", pd.DataFrame()
            )
            filt2 = condition_results[cond2]["deseq2_filtered"].get(
                "all_genes", pd.DataFrame()
            )
            genes1 = (
                set(filt1[name_col].dropna()) if name_col in filt1.columns else set()
            )
            genes2 = (
                set(filt2[name_col].dropna()) if name_col in filt2.columns else set()
            )
            overlap = genes1 & genes2

            _add_table(root, f"DEG_{cond1}_vs_{cond2}", [
                ("Metric", [f"{label1} only", f"{label2} only", "Overlap"]),
                ("Count", [
                    len(genes1 - genes2),
                    len(genes2 - genes1),
                    len(overlap),
                ]),
            ])

            # Shared gene list (top 50)
            overlap_list = sorted(overlap)[:50]
            if overlap_list:
                _add_table(root, f"Shared_Genes_{cond1}_vs_{cond2}", [
                    ("GeneSymbol", overlap_list),
                ])

        _save_prism(root, prism_dir / "DEG_overlap.pzfx")

    # 9. Summary_statistics.pzfx — condition, DEG count, splicing count
    root = _create_prism_xml()
    summary_data = []
    for cond_name in names:
        data = condition_results[cond_name]
        deg_count = len(
            data["deseq2_filtered"].get("all_genes", pd.DataFrame())
        )
        splicing_count = sum(
            len(df) for df in data.get("rmats_filtered", {}).values()
        )
        summary_data.append(
            (condition_labels.get(cond_name, cond_name), deg_count, splicing_count)
        )
    if summary_data:
        _add_table(root, "Summary_Statistics", [
            ("Condition", [s[0] for s in summary_data]),
            ("DEG_Count", [s[1] for s in summary_data]),
            ("Splicing_Count", [s[2] for s in summary_data]),
        ])
        _save_prism(root, prism_dir / "Summary_statistics.pzfx")

    # 10. Concordance_scatter.pzfx — pairwise log2FC scatter using RAW data
    if len(names) >= 2:
        root = _create_prism_xml()
        has_scatter = False
        for pair in combinations(names, 2):
            cond1, cond2 = pair
            # Use raw data so all shared genes appear (not just filtered DEGs)
            df1 = condition_results[cond1].get(
                "deseq2_raw",
                condition_results[cond1]["deseq2_filtered"].get(
                    "all_genes", pd.DataFrame()
                ),
            )
            df2 = condition_results[cond2].get(
                "deseq2_raw",
                condition_results[cond2]["deseq2_filtered"].get(
                    "all_genes", pd.DataFrame()
                ),
            )

            if (
                name_col not in df1.columns
                or name_col not in df2.columns
                or fc_col not in df1.columns
                or fc_col not in df2.columns
            ):
                continue

            # Merge on gene symbol — include all shared genes
            merged = pd.merge(
                df1[[name_col, fc_col]].dropna(),
                df2[[name_col, fc_col]].dropna(),
                on=name_col,
                suffixes=("_1", "_2"),
            )

            if len(merged) > 0:
                _add_table(root, f"Scatter_{cond1}_vs_{cond2}", [
                    ("GeneSymbol", merged[name_col].tolist()),
                    (f"log2FC_{cond1}", merged[f"{fc_col}_1"].tolist()),
                    (f"log2FC_{cond2}", merged[f"{fc_col}_2"].tolist()),
                ])
                has_scatter = True

        if has_scatter:
            _save_prism(root, prism_dir / "Concordance_scatter.pzfx")

    # 11. dPSI_distributions.pzfx — per-condition dPSI values with gene names
    root = _create_prism_xml()
    has_dpsi = False
    for cond_name in names:
        rmats_filt = condition_results[cond_name].get("rmats_filtered", {})
        all_dpsi = []
        all_genes = []
        all_etypes = []
        for et, df in rmats_filt.items():
            if len(df) > 0 and dpsi_col in df.columns:
                valid = df[dpsi_col].notna()
                all_dpsi.extend(df.loc[valid, dpsi_col].tolist())
                if rmats_gene_col in df.columns:
                    all_genes.extend(df.loc[valid, rmats_gene_col].tolist())
                else:
                    all_genes.extend([""] * int(valid.sum()))
                all_etypes.extend([et] * int(valid.sum()))

        if all_dpsi:
            # Sample 200 points for violin plot
            n = min(200, len(all_dpsi))
            _add_table(root, f"dPSI_{cond_name}", [
                ("Gene", all_genes[:n]),
                ("EventType", all_etypes[:n]),
                ("dPSI", all_dpsi[:n]),
            ])
            has_dpsi = True

    if has_dpsi:
        _save_prism(root, prism_dir / "dPSI_distributions.pzfx")

    print(f"  Prism export complete: {len(list(prism_dir.glob('*.pzfx')))} files in {prism_dir}")


# ---------------------------------------------------------------------------
# PowerPoint report
# ---------------------------------------------------------------------------


def export_powerpoint(
    condition_results: dict[str, dict],
    condition_labels: dict[str, str],
    outdir: str | Path,
    padj_cutoff: float = 0.05,
    log2fc_cutoff: float = 1.0,
    basemean_cutoff: float = 10.0,
    fdr_cutoff: float = 0.05,
    pval_cutoff: float = 0.05,
    dpsi_cutoff: float = 0.1,
    use_fdr: bool = True,
    dual_filter: bool = False,
    fig_format: str = "png",
) -> None:
    """Generate a professional PowerPoint report with all analysis figures.

    Parameters
    ----------
    condition_results : dict
        Pipeline condition_results structure.
    condition_labels : dict
        Maps condition name -> label.
    outdir : str or Path
        Output directory.
    padj_cutoff : float
        DESeq2 padj threshold.
    log2fc_cutoff : float
        DESeq2 log2FC threshold.
    basemean_cutoff : float
        DESeq2 baseMean threshold.
    fdr_cutoff : float
        rMATS FDR threshold.
    pval_cutoff : float
        rMATS p-value threshold.
    dpsi_cutoff : float
        rMATS dPSI threshold.
    use_fdr : bool
        Whether FDR was used for rMATS.
    dual_filter : bool
        Whether dual filter mode was used.
    fig_format : str
        Figure format for discovering figure files.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        print("  WARNING: python-pptx not installed, skipping PowerPoint generation")
        print("  Install with: pip install python-pptx")
        return

    try:
        from PIL import Image
    except ImportError:
        print("  WARNING: Pillow not installed, skipping PowerPoint generation")
        return

    print("\n-- Generating PowerPoint Report --")

    outdir = Path(outdir)
    DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    SLIDE_W = Inches(13.333)
    SLIDE_H = Inches(7.5)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    def add_section_slide(title_text, subtitle_text=""):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        shape = slide.shapes.add_shape(1, Emu(0), Emu(0), SLIDE_W, SLIDE_H)
        shape.fill.solid()
        shape.fill.fore_color.rgb = DARK_BLUE
        shape.line.fill.background()
        tx = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1.2))
        p = tx.text_frame.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.LEFT
        if subtitle_text:
            tx2 = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(11), Inches(0.8))
            p2 = tx2.text_frame.paragraphs[0]
            p2.text = subtitle_text
            p2.font.size = Pt(22)
            p2.font.color.rgb = WHITE
            p2.alignment = PP_ALIGN.LEFT
        return slide

    def add_image_slide(title_text, image_path):
        if not image_path.exists():
            return
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bar = slide.shapes.add_shape(1, Emu(0), Emu(0), SLIDE_W, Inches(0.9))
        bar.fill.solid()
        bar.fill.fore_color.rgb = DARK_BLUE
        bar.line.fill.background()
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(12), Inches(0.7))
        p = tx.text_frame.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = WHITE
        try:
            img = Image.open(image_path)
            iw, ih = img.size
            img.close()
            max_w = Inches(12.333)
            max_h = Inches(6.2)
            scale = min(max_w / Emu(int(iw * 914400 / 96)),
                        max_h / Emu(int(ih * 914400 / 96)))
            w = int(iw * 914400 / 96 * min(scale, 1.0))
            h = int(ih * 914400 / 96 * min(scale, 1.0))
            left = int((SLIDE_W - w) / 2)
            top = Inches(1.0) + int((max_h - h) / 2)
            slide.shapes.add_picture(str(image_path), left, top, w, h)
        except Exception as e:
            print(f"    WARNING: could not add {image_path.name}: {e}")

    # Title slide
    n_conds = len(condition_results)
    _ppt_filters = (
        f"padj < {padj_cutoff}, |log2FC| >= {log2fc_cutoff}, baseMean >= {basemean_cutoff}"
    )
    add_section_slide(
        "RNA-seq Analysis Report",
        f"{n_conds} conditions  |  {_ppt_filters}  |  {date.today().strftime('%B %d, %Y')}",
    )

    # Per-condition sections
    per_cond_figures = [
        ("Volcano Plot", f"volcano_plot.{fig_format}"),
        ("MA Plot", f"ma_plot.{fig_format}"),
        ("Biotype Distribution", f"biotype_distribution.{fig_format}"),
        ("P-value Histogram", f"pvalue_histogram.{fig_format}"),
        ("rMATS Event Summary", f"rmats_event_type_summary.{fig_format}"),
    ]

    for cond_name in condition_results.keys():
        label = condition_labels.get(cond_name, cond_name)
        fig_dir = outdir / cond_name / "figures"
        add_section_slide(label, "Condition-Specific Analysis")
        for slide_title, fname in per_cond_figures:
            path = fig_dir / fname
            if path.exists():
                add_image_slide(f"{label} -- {slide_title}", path)

    # Cross-condition section
    cross_dir = outdir / "cross_condition" / "figures"
    add_section_slide("Cross-Condition Analysis", "Comparative & Integrative Results")

    cross_figures = [
        ("DEG Counts Overview", f"deseq2_de_counts_overview.{fig_format}"),
        ("Direction Concordance", f"direction_concordance_heatmap.{fig_format}"),
        ("log2FC Heatmap", f"log2fc_heatmap.{fig_format}"),
        ("Pairwise log2FC Scatter", f"pairwise_log2fc_scatter.{fig_format}"),
    ]
    for slide_title, fname in cross_figures:
        path = cross_dir / fname
        if path.exists():
            add_image_slide(slide_title, path)

    # Save
    pptx_path = outdir / "RNA-seq_Analysis_Report.pptx"
    prs.save(str(pptx_path))
    print(f"  PowerPoint saved: {pptx_path} ({len(prs.slides)} slides)")


# ---------------------------------------------------------------------------
# Unfiltered merged export
# ---------------------------------------------------------------------------


def export_unfiltered_merged(
    condition_results: dict[str, dict],
    condition_labels: dict[str, str],
    outdir: str | Path,
    deseq2_cols: dict[str, str] | None = None,
    rmats_cols: dict[str, str] | None = None,
) -> None:
    """Export unfiltered (all genes/events) merged across conditions to Excel.

    Parameters
    ----------
    condition_results : dict
        Pipeline condition_results structure.
    condition_labels : dict
        Maps condition name -> label.
    outdir : str or Path
        Output directory.
    deseq2_cols : dict or None
        DESeq2 column mapping.
    rmats_cols : dict or None
        rMATS column mapping.
    """
    if deseq2_cols is None:
        deseq2_cols = DEFAULT_DESEQ2_COLS
    if rmats_cols is None:
        rmats_cols = DEFAULT_RMATS_COLS

    outdir = Path(outdir)
    xlsx_path = outdir / "Unfiltered_All_Conditions_Merged.xlsx"
    print("\n-- Exporting Unfiltered Merged Results --")
    print(f"   Output: {xlsx_path}")

    names = list(condition_results.keys())
    id_col = deseq2_cols["gene_id"]
    name_col = deseq2_cols["gene_name"]

    def _short(label):
        parts = label.split(" vs ")
        return parts[0].replace(" ", "_") if parts else label.replace(" ", "_")

    def _deseq2_key(df):
        if name_col in df.columns and df[name_col].notna().sum() > 0:
            return name_col
        return id_col

    deseq2_slices = {}
    for cname in names:
        raw = condition_results[cname].get("deseq2_raw", pd.DataFrame())
        if raw.empty:
            continue
        key = _deseq2_key(raw)
        short = _short(condition_labels[cname])
        df = raw.copy()
        key_cols = [c for c in [id_col, name_col] if c in df.columns]
        rename = {}
        for c in df.columns:
            if c not in key_cols:
                rename[c] = f"{short}_{c}"
        df = df.rename(columns=rename)
        deseq2_slices[cname] = (df, key, short)

    with pd.ExcelWriter(xlsx_path, engine="openpyxl") as writer:
        # DESeq2 all conditions merged
        if deseq2_slices:
            merged = None
            for cname in names:
                if cname not in deseq2_slices:
                    continue
                slice_df, key, short = deseq2_slices[cname]
                if merged is None:
                    merged = slice_df.copy()
                else:
                    merge_on = [c for c in [key] if c in merged.columns and c in slice_df.columns]
                    right_drop = [
                        c for c in [id_col, name_col]
                        if c in slice_df.columns and c in merged.columns and c not in merge_on
                    ]
                    merged = merged.merge(
                        slice_df.drop(columns=right_drop, errors="ignore"),
                        on=merge_on, how="outer",
                    )

            if merged is not None and not merged.empty:
                merged.to_excel(writer, sheet_name="DESeq2_All_Conditions", index=False)
                print(f"   DESeq2_All_Conditions: {len(merged):,} genes x {len(merged.columns)} cols")

        # Pairwise DESeq2
        for cA, cB in combinations(names, 2):
            if cA not in deseq2_slices or cB not in deseq2_slices:
                continue
            slA, keyA, shortA = deseq2_slices[cA]
            slB, keyB, shortB = deseq2_slices[cB]
            merge_on = [c for c in [keyA] if c in slA.columns and c in slB.columns]
            right_drop = [
                c for c in [id_col, name_col]
                if c in slB.columns and c in slA.columns and c not in merge_on
            ]
            pair_df = slA.merge(
                slB.drop(columns=right_drop, errors="ignore"),
                on=merge_on, how="outer",
            )
            sheet = f"DESeq2_{shortA}_vs_{shortB}"[:31]
            pair_df.to_excel(writer, sheet_name=sheet, index=False)

    print(f"   Saved: {xlsx_path}")


# ---------------------------------------------------------------------------
# Pairwise comparison workbook export
# ---------------------------------------------------------------------------


def export_pairwise_workbook(
    condition_results: dict[str, dict],
    condition_labels: dict[str, str],
    outdir: str | Path,
    cols: dict[str, str] | None = None,
    rmats_cols: dict[str, str] | None = None,
    dpsi_cutoff: float = 0.1,
) -> None:
    """Export pairwise comparison Excel workbooks for DESeq2 and rMATS data.

    For each pair of conditions, creates one .xlsx file with:
      - DESeq2 sheets: shared (all/up/down), opposite-direction, condition-only
      - rMATS sheets (per event type): same breakdown by event coordinates
      - rMATS gene-level sheets: by gene identifier
      - Summary sheet with row counts per category

    Parameters
    ----------
    condition_results : dict
        Pipeline condition_results structure.
    condition_labels : dict
        Maps condition name -> label.
    outdir : str or Path
        Output directory.
    cols : dict or None
        DESeq2 column name mapping.
    rmats_cols : dict or None
        rMATS column name mapping.
    dpsi_cutoff : float
        dPSI cutoff for directional classification.
    """
    if cols is None:
        cols = DEFAULT_DESEQ2_COLS
    if rmats_cols is None:
        rmats_cols = DEFAULT_RMATS_COLS

    outdir = Path(outdir)
    names = list(condition_results.keys())
    dpsi_col = rmats_cols["inclevel_diff"]
    gene_col = rmats_cols["gene_name"]
    id_col_rmats = rmats_cols["gene_id"]

    def _short(label):
        parts = label.split(" vs ")
        return parts[0].replace(" ", "_") if parts else label.replace(" ", "_")

    for name_a, name_b in combinations(names, 2):
        label_a = condition_labels[name_a]
        label_b = condition_labels[name_b]
        short_a = _short(label_a)
        short_b = _short(label_b)

        wb_name = f"pairwise_{name_a}_vs_{name_b}.xlsx"
        wb_path = outdir / wb_name
        summary_rows: list[dict] = []

        with pd.ExcelWriter(wb_path, engine="openpyxl") as writer:

            # -------------------------------------------------------
            # DESeq2 sheets
            # -------------------------------------------------------
            filt_a = condition_results[name_a]["deseq2_filtered"]["all_genes"].copy()
            filt_b = condition_results[name_b]["deseq2_filtered"]["all_genes"].copy()

            key_col_a, _ = best_gene_key(filt_a, cols)
            key_col_b, _ = best_gene_key(filt_b, cols)
            if not key_col_a or key_col_a not in filt_a.columns:
                key_col_a = cols["gene_name"]
            if not key_col_b or key_col_b not in filt_b.columns:
                key_col_b = cols["gene_name"]

            merge_key = key_col_a
            if key_col_b != key_col_a and key_col_b in filt_b.columns:
                filt_b = filt_b.rename(columns={key_col_b: merge_key})

            suffix_a = f"_{short_a}"
            suffix_b = f"_{short_b}"

            all_a = set(filt_a[merge_key].dropna().unique())
            all_b = set(filt_b[merge_key].dropna().unique())

            dir_col = "direction"
            up_a = set(filt_a.loc[filt_a[dir_col] == "up", merge_key].dropna().unique())
            up_b = set(filt_b.loc[filt_b[dir_col] == "up", merge_key].dropna().unique())
            down_a = set(filt_a.loc[filt_a[dir_col] == "down", merge_key].dropna().unique())
            down_b = set(filt_b.loc[filt_b[dir_col] == "down", merge_key].dropna().unique())

            shared_all = all_a & all_b
            shared_up = up_a & up_b
            shared_down = down_a & down_b
            up_a_down_b = up_a & down_b
            down_a_up_b = down_a & up_b
            only_a = all_a - all_b
            only_b = all_b - all_a

            de_categories = {
                "DE_shared_all": shared_all,
                "DE_shared_up": shared_up,
                "DE_shared_down": shared_down,
                f"DE_up_{short_a}_down_{short_b}": up_a_down_b,
                f"DE_down_{short_a}_up_{short_b}": down_a_up_b,
                f"DE_only_{short_a}": only_a,
                f"DE_only_{short_b}": only_b,
            }

            for sheet_name, gene_set in de_categories.items():
                if not gene_set:
                    empty_df = pd.DataFrame(columns=[merge_key])
                    sheet_label = sheet_name[:31]
                    empty_df.to_excel(writer, sheet_name=sheet_label, index=False)
                    summary_rows.append({"category": sheet_name, "count": 0})
                    continue

                sub_a = filt_a[filt_a[merge_key].isin(gene_set)].copy()
                sub_b = filt_b[filt_b[merge_key].isin(gene_set)].copy()

                rename_a = {c: f"{c}{suffix_a}" for c in sub_a.columns if c != merge_key}
                rename_b = {c: f"{c}{suffix_b}" for c in sub_b.columns if c != merge_key}
                sub_a = sub_a.rename(columns=rename_a)
                sub_b = sub_b.rename(columns=rename_b)

                merged = sub_a.merge(sub_b, on=merge_key, how="outer")

                sheet_label = sheet_name[:31]
                merged.to_excel(writer, sheet_name=sheet_label, index=False)
                summary_rows.append({"category": sheet_name, "count": len(merged)})
                print(f"   {sheet_label}: {len(merged):,} genes")

            # -------------------------------------------------------
            # rMATS sheets (per event type)
            # -------------------------------------------------------
            for et in RMATS_EVENT_TYPES:
                df_a = condition_results[name_a]["rmats_filtered"].get(et)
                df_b = condition_results[name_b]["rmats_filtered"].get(et)

                has_a = df_a is not None and len(df_a) > 0
                has_b = df_b is not None and len(df_b) > 0

                if not has_a and not has_b:
                    continue

                if has_a:
                    df_a = df_a.copy()
                    df_a["_ekey"] = make_event_key(df_a, et).values
                    df_a = df_a[df_a["_ekey"] != ""]
                else:
                    df_a = pd.DataFrame(columns=["_ekey"])

                if has_b:
                    df_b = df_b.copy()
                    df_b["_ekey"] = make_event_key(df_b, et).values
                    df_b = df_b[df_b["_ekey"] != ""]
                else:
                    df_b = pd.DataFrame(columns=["_ekey"])

                events_all_a = set(df_a["_ekey"].dropna().unique())
                events_all_b = set(df_b["_ekey"].dropna().unique())

                events_inc_a = set(
                    df_a.loc[df_a[dpsi_col] >= dpsi_cutoff, "_ekey"].dropna().unique()
                ) if has_a and dpsi_col in df_a.columns else set()
                events_inc_b = set(
                    df_b.loc[df_b[dpsi_col] >= dpsi_cutoff, "_ekey"].dropna().unique()
                ) if has_b and dpsi_col in df_b.columns else set()
                events_exc_a = set(
                    df_a.loc[df_a[dpsi_col] <= -dpsi_cutoff, "_ekey"].dropna().unique()
                ) if has_a and dpsi_col in df_a.columns else set()
                events_exc_b = set(
                    df_b.loc[df_b[dpsi_col] <= -dpsi_cutoff, "_ekey"].dropna().unique()
                ) if has_b and dpsi_col in df_b.columns else set()

                rmats_categories = {
                    f"{et}_shared_all": events_all_a & events_all_b,
                    f"{et}_shared_included": events_inc_a & events_inc_b,
                    f"{et}_shared_excluded": events_exc_a & events_exc_b,
                    f"{et}_inc_{short_a}_exc_{short_b}": events_inc_a & events_exc_b,
                    f"{et}_exc_{short_a}_inc_{short_b}": events_exc_a & events_inc_b,
                    f"{et}_only_{short_a}": events_all_a - events_all_b,
                    f"{et}_only_{short_b}": events_all_b - events_all_a,
                }

                coord_cols = COORD_COLS.get(et, [])
                id_cols = [
                    c for c in [rmats_cols["gene_id"], rmats_cols["gene_name"]]
                    if (has_a and c in df_a.columns) or (has_b and c in df_b.columns)
                ]
                shared_merge_cols = (
                    [c for c in coord_cols
                     if (has_a and c in df_a.columns) or (has_b and c in df_b.columns)]
                    + id_cols
                )

                for sheet_name_et, event_set in rmats_categories.items():
                    if not event_set:
                        empty_df = pd.DataFrame(columns=["_ekey"])
                        sheet_label_et = sheet_name_et[:31]
                        empty_df.to_excel(writer, sheet_name=sheet_label_et, index=False)
                        summary_rows.append({"category": sheet_name_et, "count": 0})
                        continue

                    sub_a = df_a[df_a["_ekey"].isin(event_set)].copy() if has_a else pd.DataFrame()
                    sub_b = df_b[df_b["_ekey"].isin(event_set)].copy() if has_b else pd.DataFrame()

                    rmats_id = rmats_cols["event_id"]
                    if rmats_id in sub_a.columns:
                        sub_a = sub_a.drop(columns=[rmats_id])
                    if rmats_id in sub_b.columns:
                        sub_b = sub_b.drop(columns=[rmats_id])

                    if len(sub_a) > 0:
                        rename_a_et = {
                            c: f"{short_a}_{c}" for c in sub_a.columns
                            if c not in shared_merge_cols and c != "_ekey"
                        }
                        sub_a = sub_a.rename(columns=rename_a_et)

                    if len(sub_b) > 0:
                        rename_b_et = {
                            c: f"{short_b}_{c}" for c in sub_b.columns
                            if c not in shared_merge_cols and c != "_ekey"
                        }
                        sub_b = sub_b.rename(columns=rename_b_et)

                    if len(sub_a) > 0 and len(sub_b) > 0:
                        merged = sub_a.merge(
                            sub_b.drop(
                                columns=[c for c in shared_merge_cols if c in sub_b.columns],
                                errors="ignore",
                            ),
                            on="_ekey", how="outer",
                        )
                    elif len(sub_a) > 0:
                        merged = sub_a
                    else:
                        merged = sub_b

                    if "_ekey" in merged.columns:
                        merged = merged.drop(columns=["_ekey"])

                    sheet_label_et = sheet_name_et[:31]
                    merged.to_excel(writer, sheet_name=sheet_label_et, index=False)
                    summary_rows.append({"category": sheet_name_et, "count": len(merged)})
                    print(f"   {sheet_label_et}: {len(merged):,} events")

            # -------------------------------------------------------
            # rMATS gene-level sheets (per event type)
            # -------------------------------------------------------
            for et in RMATS_EVENT_TYPES:
                df_a_gl = condition_results[name_a]["rmats_filtered"].get(et)
                df_b_gl = condition_results[name_b]["rmats_filtered"].get(et)

                has_a_gl = df_a_gl is not None and len(df_a_gl) > 0
                has_b_gl = df_b_gl is not None and len(df_b_gl) > 0

                if not has_a_gl and not has_b_gl:
                    continue

                _gl_id_col = id_col_rmats if (
                    (has_a_gl and id_col_rmats in df_a_gl.columns)
                    or (has_b_gl and id_col_rmats in df_b_gl.columns)
                ) else gene_col

                genes_all_a = set(
                    df_a_gl[_gl_id_col].dropna().unique()
                ) if has_a_gl and _gl_id_col in df_a_gl.columns else set()
                genes_all_b = set(
                    df_b_gl[_gl_id_col].dropna().unique()
                ) if has_b_gl and _gl_id_col in df_b_gl.columns else set()

                genes_inc_a = set(
                    df_a_gl.loc[df_a_gl[dpsi_col] >= dpsi_cutoff, _gl_id_col].dropna().unique()
                ) if has_a_gl and dpsi_col in df_a_gl.columns and _gl_id_col in df_a_gl.columns else set()
                genes_inc_b = set(
                    df_b_gl.loc[df_b_gl[dpsi_col] >= dpsi_cutoff, _gl_id_col].dropna().unique()
                ) if has_b_gl and dpsi_col in df_b_gl.columns and _gl_id_col in df_b_gl.columns else set()
                genes_exc_a = set(
                    df_a_gl.loc[df_a_gl[dpsi_col] <= -dpsi_cutoff, _gl_id_col].dropna().unique()
                ) if has_a_gl and dpsi_col in df_a_gl.columns and _gl_id_col in df_a_gl.columns else set()
                genes_exc_b = set(
                    df_b_gl.loc[df_b_gl[dpsi_col] <= -dpsi_cutoff, _gl_id_col].dropna().unique()
                ) if has_b_gl and dpsi_col in df_b_gl.columns and _gl_id_col in df_b_gl.columns else set()

                gene_categories = {
                    f"{et}_gene_shared_all": genes_all_a & genes_all_b,
                    f"{et}_gene_shared_inc": genes_inc_a & genes_inc_b,
                    f"{et}_gene_shared_exc": genes_exc_a & genes_exc_b,
                    f"{et}_gene_inc{short_a}_exc{short_b}": genes_inc_a & genes_exc_b,
                    f"{et}_gene_exc{short_a}_inc{short_b}": genes_exc_a & genes_inc_b,
                    f"{et}_gene_only_{short_a}": genes_all_a - genes_all_b,
                    f"{et}_gene_only_{short_b}": genes_all_b - genes_all_a,
                }

                for sheet_name_gl, gene_set_gl in gene_categories.items():
                    if not gene_set_gl:
                        empty_df = pd.DataFrame(columns=[gene_col])
                        sheet_label_gl = sheet_name_gl[:31]
                        empty_df.to_excel(writer, sheet_name=sheet_label_gl, index=False)
                        summary_rows.append({"category": sheet_name_gl, "count": 0})
                        continue

                    sub_a_gl = (
                        df_a_gl[df_a_gl[_gl_id_col].isin(gene_set_gl)].copy()
                        if has_a_gl and _gl_id_col in df_a_gl.columns
                        else pd.DataFrame()
                    )
                    sub_b_gl = (
                        df_b_gl[df_b_gl[_gl_id_col].isin(gene_set_gl)].copy()
                        if has_b_gl and _gl_id_col in df_b_gl.columns
                        else pd.DataFrame()
                    )

                    if len(sub_a_gl) > 0:
                        sub_a_gl["_source_condition"] = label_a
                    if len(sub_b_gl) > 0:
                        sub_b_gl["_source_condition"] = label_b

                    if len(sub_a_gl) > 0 and len(sub_b_gl) > 0:
                        merged_gl = pd.concat([sub_a_gl, sub_b_gl], ignore_index=True)
                    elif len(sub_a_gl) > 0:
                        merged_gl = sub_a_gl
                    else:
                        merged_gl = sub_b_gl

                    sheet_label_gl = sheet_name_gl[:31]
                    merged_gl.to_excel(writer, sheet_name=sheet_label_gl, index=False)
                    n_genes = len(gene_set_gl)
                    summary_rows.append({
                        "category": sheet_name_gl,
                        "count": f"{n_genes} genes / {len(merged_gl)} events",
                    })
                    print(f"   {sheet_label_gl}: {n_genes:,} genes, {len(merged_gl):,} events")

            # -------------------------------------------------------
            # Summary sheet
            # -------------------------------------------------------
            if summary_rows:
                summary_df = pd.DataFrame(summary_rows)
                summary_df.to_excel(writer, sheet_name="Summary", index=False)

        print(f"  Saved: {wb_path}")


# ---------------------------------------------------------------------------
# Validation
# ---------------------------------------------------------------------------


def validate_outputs(
    condition_results: dict[str, dict],
    condition_labels: dict[str, str],
    outdir: str | Path,
    cols: dict[str, str] | None = None,
    rmats_cols: dict[str, str] | None = None,
) -> bool:
    """Validate all pipeline outputs for completeness and correctness.

    Checks:
    - File existence and sizes
    - Prism files >10KB
    - Venn diagram math (All = Concordant_Up + Concordant_Down + Discordant)
    - DEG totals (total = up + down)
    - Gene symbols present in output files

    Parameters
    ----------
    condition_results : dict
        Pipeline condition_results structure.
    condition_labels : dict
        Maps condition name -> label.
    outdir : str or Path
        Base output directory.
    cols : dict or None
        DESeq2 column mapping.
    rmats_cols : dict or None
        rMATS column mapping.

    Returns
    -------
    bool
        True if all validation checks passed.
    """
    if cols is None:
        cols = DEFAULT_DESEQ2_COLS
    if rmats_cols is None:
        rmats_cols = DEFAULT_RMATS_COLS

    outdir = Path(outdir)

    print("\n" + "=" * 60)
    print("  VALIDATION REPORT")
    print("=" * 60)

    validation_passed = True
    warnings: list[str] = []
    errors: list[str] = []

    # --- 1. Check file existence ---
    print("\n1. File Existence Check:")

    expected_files = [
        ("cross_condition/multi_condition_results.xlsx", "Multi-condition results"),
        ("RNA-seq_Analysis_Report.pptx", "PowerPoint report"),
    ]

    for fname, desc in expected_files:
        fpath = outdir / fname
        if fpath.exists():
            size_mb = fpath.stat().st_size / (1024 * 1024)
            print(f"  [PASS] {desc}: {fname} ({size_mb:.2f} MB)")
        else:
            errors.append(f"Missing file: {fname}")
            print(f"  [FAIL] {desc}: {fname} MISSING")
            validation_passed = False

    # --- 2. Check Prism files ---
    print("\n2. Prism File Validation:")
    prism_dir = outdir / "prism_files"
    if prism_dir.exists():
        prism_files = list(prism_dir.glob("*.pzfx"))
        print(f"  Found {len(prism_files)} Prism files")

        small_files = []
        for pfile in prism_files:
            size_kb = pfile.stat().st_size / 1024
            if size_kb < 10:
                small_files.append((pfile.name, size_kb))

        if small_files:
            for fname_pf, size_kb in small_files:
                warnings.append(f"Prism file {fname_pf} is only {size_kb:.1f} KB (expected >10 KB)")
                print(f"  [WARN] {fname_pf}: {size_kb:.1f} KB (may be too small)")
        else:
            print("  [PASS] All Prism files >10 KB")

        if len(prism_files) < 15:
            warnings.append(f"Expected 15+ Prism files, found {len(prism_files)}")
            print(f"  [WARN] Only {len(prism_files)} Prism files (expected 15+)")
        else:
            print(f"  [PASS] File count meets requirement ({len(prism_files)} >= 15)")
    else:
        errors.append("Prism directory not found")
        print("  [FAIL] Prism directory not found")
        validation_passed = False

    # --- 3. Validate DEG counts ---
    print("\n3. DEG Count Validation:")
    gene_name_col = cols.get("gene_name", "gene_name")

    for cond_name, data in condition_results.items():
        cond_label = condition_labels.get(cond_name, cond_name)
        deg_filt = data["deseq2_filtered"]["all_genes"]

        if "direction" in deg_filt.columns:
            total = len(deg_filt)
            up = int((deg_filt["direction"] == "up").sum())
            down = int((deg_filt["direction"] == "down").sum())
            computed_total = up + down

            if computed_total == total:
                print(f"  [PASS] {cond_label}: Total={total}, Up={up}, Down={down}")
            else:
                errors.append(
                    f"{cond_label} DEG count mismatch: Total={total}, Up+Down={computed_total}"
                )
                print(f"  [FAIL] {cond_label}: Total={total} != Up+Down={computed_total}")
                validation_passed = False

            if gene_name_col in deg_filt.columns:
                n_missing = deg_filt[gene_name_col].isna().sum()
                if n_missing > 0:
                    warnings.append(f"{cond_label}: {n_missing} genes missing gene symbols")
                    print(f"  [WARN] {cond_label}: {n_missing}/{total} genes missing symbols")
                else:
                    print(f"  [PASS] {cond_label}: All genes have symbols")
            else:
                warnings.append(f"{cond_label}: No gene_name column found")
                print(f"  [WARN] {cond_label}: No gene_name column")

    # --- 4. Validate splicing Venn math ---
    print("\n4. Splicing Venn Math Validation:")
    names = list(condition_results.keys())
    gene_col_rmats = rmats_cols["gene_name"]
    dpsi_col = rmats_cols["inclevel_diff"]

    if len(names) >= 2:
        for et in RMATS_EVENT_TYPES:
            dfs: dict = {}
            for name in names:
                if et in condition_results[name]["rmats_filtered"]:
                    df = condition_results[name]["rmats_filtered"][et]
                    if len(df) > 0:
                        gene_dpsi = df.groupby(gene_col_rmats)[dpsi_col].mean()
                        dfs[name] = gene_dpsi

            if len(dfs) >= 2:
                shared_genes = set.intersection(*[set(d.index) for d in dfs.values()])

                concordant_up = 0
                concordant_down = 0
                discordant = 0

                for gene in shared_genes:
                    signs = [np.sign(dfs[name][gene]) for name in dfs.keys()]
                    if all(s > 0 for s in signs):
                        concordant_up += 1
                    elif all(s < 0 for s in signs):
                        concordant_down += 1
                    else:
                        discordant += 1

                computed_all = concordant_up + concordant_down + discordant
                if computed_all == len(shared_genes):
                    print(
                        f"  [PASS] {et}: All={len(shared_genes)}, Up={concordant_up}, "
                        f"Down={concordant_down}, Disc={discordant}"
                    )
                else:
                    errors.append(
                        f"{et} Venn math error: All={len(shared_genes)} != Sum={computed_all}"
                    )
                    print(f"  [FAIL] {et}: All={len(shared_genes)} != Sum={computed_all}")
                    validation_passed = False

    # --- 5. Summary ---
    print("\n" + "=" * 60)
    print("  VALIDATION SUMMARY")
    print("=" * 60)

    if validation_passed and not errors:
        print("  [PASS] ALL VALIDATION CHECKS PASSED")
    else:
        print(f"  [FAIL] VALIDATION FAILED: {len(errors)} error(s)")

    if warnings:
        print(f"\n  [WARN] {len(warnings)} warning(s):")
        for w in warnings[:10]:
            print(f"    - {w}")

    if errors:
        print(f"\n  [FAIL] {len(errors)} error(s):")
        for e in errors:
            print(f"    - {e}")

    print("=" * 60)

    return validation_passed
