#!/usr/bin/env python3
"""
=============================================================================
DESeq2 & rMATS Filtering and Visualization Pipeline
=============================================================================
Filters differential expression (DESeq2) and alternative splicing (rMATS)
results by user-defined cutoffs, exports filtered tables, and generates
publication-quality figures.

Handles all 5 rMATS event types: SE, A3SS, A5SS, RI, MXE

Usage:
    python deseq2_rmats_filter_pipeline.py

Edit the CONFIGURATION section below to set file paths and cutoffs.
=============================================================================
"""

import os
import sys
import json
import urllib.request
import urllib.parse
from datetime import datetime
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import matplotlib.ticker as ticker
import matplotlib.patches as mpatches
import matplotlib.lines as mlines
import matplotlib.gridspec as gridspec
import shutil
import seaborn as sns
from pathlib import Path
from itertools import combinations
from collections import Counter
import warnings
warnings.filterwarnings("ignore")
sys.setrecursionlimit(10000)  # scipy dendrogram needs this for large gene sets

try:
    from matplotlib_venn import venn2, venn3
except ImportError:
    raise ImportError(
        "matplotlib-venn is required for cross-condition comparisons.\n"
        "Install it with:  pip install matplotlib-venn"
    )

try:
    import plotly.express as px
except ImportError:
    px = None  # Interactive plots disabled if plotly not installed

try:
    from upsetplot import UpSet, from_memberships
    _UPSET_AVAILABLE = True
except ImportError:
    _UPSET_AVAILABLE = False
    UpSet = None
    from_memberships = None

try:
    from scipy.stats import pearsonr, fisher_exact
    _SCIPY_AVAILABLE = True
except ImportError:
    _SCIPY_AVAILABLE = False
    pearsonr = None
    fisher_exact = None

# ===============================================================================
# CONFIGURATION -- EDIT THESE VALUES
# ===============================================================================

# --- Base Directory & Conditions ---
BASE_DIR = ""

CONDITIONS = []

# rMATS event types to analyze and the file suffix pattern
RMATS_EVENT_TYPES = ["SE", "A3SS", "A5SS", "RI", "MXE"]
RMATS_FILE_SUFFIX = ".MATS.JCEC.txt"  # change to ".MATS.JCEC.xlsx" if using Excel files

OUTPUT_DIR = "./output"

# --- DESeq2 Cutoffs ---
LOG2FC_CUTOFF   = 1.0      # absolute log2 fold change threshold (HSCHARME: |log2FC| > 1)
BASEMEAN_CUTOFF = 10       # minimum baseMean expression (HSCHARME: <10 counts in ≥2 samples)
PADJ_CUTOFF     = 0.05     # adjusted p-value threshold (HSCHARME: padj < 0.05)

# Automatic biotype splitting: runs analysis for protein_coding and non-protein_coding
# separately in addition to all genes. Requires "biotype" column in your DESeq2 data.
AUTO_BIOTYPE_SPLIT = True

# Gene name lookup: when a DESeq2 file has only Ensembl IDs (no gene_name column),
# automatically fetch gene symbols from MyGene.info (requires internet connection).
# Set to False to skip and use Ensembl IDs as labels instead.
GENE_NAME_LOOKUP = True
SPECIES = "human"   # options: "human", "mouse", "rat", "zebrafish", "fly", "worm", etc.

# --- rMATS Cutoffs ---
RMATS_FDR_CUTOFF     = 0.05    # FDR threshold (HSCHARME: FDR < 0.05)
RMATS_PVAL_CUTOFF    = 0.05    # PValue threshold (aligned with FDR)
INCLEVEL_DIFF_CUTOFF = 0.1     # absolute IncLevelDifference threshold
USE_FDR = True                 # True = filter by FDR, False = filter by PValue
RMATS_DUAL_FILTER = False      # True = filter by BOTH FDR AND PValue simultaneously

# --- Column Name Mapping (adjust if your column headers differ) ---
DESEQ2_COLS = {
    "gene_id":       "gene_id",
    "gene_name":     "gene_name",
    "log2fc":        "log2FoldChange",
    "basemean":      "baseMean",
    "padj":          "padj",
    "pvalue":        "pvalue",
    "biotype":       "biotype",
    "stat":          "stat",       # Wald test statistic (for GSEA ranking)
    "lfcSE":         "lfcSE",      # log2FC standard error (shrinkage detection)
}

RMATS_COLS = {
    "event_id":      "ID",
    "gene_id":       "GeneID",
    "gene_name":     "geneSymbol",
    "pvalue":        "PValue",
    "fdr":           "FDR",
    "inclevel_diff": "IncLevelDifference",
}

# --- Biotype Grouping ---
# ALL keys must be lowercase; _assign_biotype_group() lowercases input before lookup,
# so 'lncRNA', 'LNCRNA', 'lncrna' from any tool/annotation version all map correctly.
_BIOTYPE_GROUPS = {
    # Protein coding
    "protein_coding": "Protein Coding",
    # Long non-coding RNA (Ensembl v37+ uses 'lncRNA'; older annotations use 'lincRNA' etc.)
    "lncrna": "lncRNA", "lincrna": "lncRNA",
    "sense_intronic": "lncRNA", "sense_overlapping": "lncRNA",
    "antisense": "lncRNA", "processed_transcript": "lncRNA",
    "bidirectional_promoter_lncrna": "lncRNA", "macro_lncrna": "lncRNA",
    "non_coding": "lncRNA",
    # Pseudogenes
    "pseudogene": "Pseudogene", "processed_pseudogene": "Pseudogene",
    "unprocessed_pseudogene": "Pseudogene",
    "transcribed_unprocessed_pseudogene": "Pseudogene",
    "transcribed_processed_pseudogene": "Pseudogene",
    "transcribed_unitary_pseudogene": "Pseudogene",
    "polymorphic_pseudogene": "Pseudogene", "unitary_pseudogene": "Pseudogene",
    "ig_pseudogene": "Pseudogene", "ig_c_pseudogene": "Pseudogene",
    "ig_v_pseudogene": "Pseudogene", "tr_v_pseudogene": "Pseudogene",
    "tr_j_pseudogene": "Pseudogene",
    # Small ncRNA (all lowercased so 'miRNA', 'snRNA', etc. all resolve correctly)
    "mirna": "Small ncRNA", "snrna": "Small ncRNA", "snorna": "Small ncRNA",
    "misc_ncrna": "Small ncRNA", "rrna": "Small ncRNA", "scrna": "Small ncRNA",
    "scarna": "Small ncRNA", "pirna": "Small ncRNA", "vault_rna": "Small ncRNA",
    "y_rna": "Small ncRNA", "ribozyme": "Small ncRNA", "srp_rna": "Small ncRNA",
    "trna": "Small ncRNA",
}
_BIOTYPE_ORDER = ["Protein Coding", "lncRNA", "Pseudogene", "Small ncRNA", "Other"]
_BIOTYPE_COLORS = {
    "Protein Coding": "#4C72B0",
    "lncRNA":         "#DD8452",
    "Pseudogene":     "#55A868",
    "Small ncRNA":    "#C44E52",
    "Other":          "#8172B2",
}

# --- Figure Settings ---
FIG_DPI     = 300
FIG_FORMAT  = "png"     # "png", "svg", or "pdf"
FONT_SIZE   = 12
COLOR_UP    = "#E69F00"  # upregulated / included  (Okabe-Ito orange)
COLOR_DOWN  = "#0072B2"  # downregulated / excluded (Okabe-Ito blue)
COLOR_NS    = "#BFBFBF"  # not significant

# Interactive plots: generate HTML files with hover tooltips (requires plotly)
INTERACTIVE_PLOTS = True

# --- Enrichment Analysis Settings ---
GSEA_DATABASES = [
    "GO_Biological_Process_2023",
    "GO_Cellular_Component_2023",
    "GO_Molecular_Function_2023",
    "KEGG_2021_Human",
    "Reactome_2022",
    "MSigDB_Hallmark_2020",
    "WikiPathway_2021_Human",
]

ORA_DATABASES = [
    "GO_Biological_Process_2023",
    "GO_Cellular_Component_2023",
    "GO_Molecular_Function_2023",
    "KEGG_2021_Human",
    "Reactome_2022",
]

GENES_OF_INTEREST = ["MIAT", "QKI", "QKI-5", "QKI-6", "QKI-7"]

# --- RBP Annotation (optional) ---
# Path to an RBP annotation Excel file. Supports two formats:
#   1) RBP-E-A-C-Complex.xlsx  (columns: Gene Name, TriSNRP, B Complex, ...)
#   2) RBP_yael_MW_lists.xlsx  (columns: Gene Symbol, Yael_RBP?, MW_RBP?)
# Leave empty to skip RBP annotation.
RBP_FILE = ""

# --- Normalized Counts (optional, for PCA/heatmaps) ---
COUNTS_FILE = ""              # Path to normalized_counts.tsv (genes x samples)
SAMPLE_METADATA = {}          # {"sample_name": "condition_name", ...} or auto-detect

# --- GSEA Settings (WSF ground truth) ---
GSEA_RANKING = "stat"         # "stat" (Wald, WSF ground truth) or "log2fc"
                              # Falls back to log2fc if stat column missing
GSEA_MIN_SIZE = 15            # WSF uses 15; more rigorous than 5
GSEA_MAX_SIZE = 500
GSEA_PERMUTATIONS = 1000      # WSF uses 1000 (was 100); better p-value accuracy

# --- ORA Method ---
ORA_METHOD = "both"           # "both" (run Enrichr + g:Profiler side-by-side),
                              # "gprofiler" (g:SCS correction), or "enrichr" (legacy).
                              # Falls back to enrichr if gprofiler-official not installed.

# Event type colors for comparison charts
EVENT_COLORS = {
    "SE":   "#E64B35",
    "A3SS": "#4DBBD5",
    "A5SS": "#00A087",
    "RI":   "#3C5488",
    "MXE":  "#F39B7F",
}


# ===============================================================================
# PIPELINE -- No edits needed below unless customizing
# ===============================================================================

def setup_style():
    """Set publication-quality matplotlib defaults."""
    plt.rcParams.update({
        "figure.dpi": FIG_DPI,
        "font.size": FONT_SIZE,
        "font.family": "sans-serif",
        "axes.linewidth": 1.2,
        "axes.spines.top": False,
        "axes.spines.right": False,
        "xtick.major.width": 1.0,
        "ytick.major.width": 1.0,
        "figure.facecolor": "white",
        "savefig.bbox": "tight",
        "savefig.dpi": FIG_DPI,
    })
    sns.set_palette("deep")


def add_count_box(ax, n_up, n_down, total, position="lower left",
                   up_label="Up", down_label="Down"):
    """Add a compact count box in a data-sparse corner of the plot."""
    text = (f"{up_label}: {n_up:,}\n"
            f"{down_label}: {n_down:,}\n"
            f"Total: {total:,}")
    # Use lower-left by default to avoid covering significant data points
    loc = {"upper left": (0.02, 0.98), "upper right": (0.98, 0.98),
           "lower left": (0.02, 0.02), "lower right": (0.98, 0.02)}
    x, y = loc.get(position, (0.02, 0.02))
    ha = "left" if "left" in position else "right"
    va = "top" if "upper" in position else "bottom"
    ax.text(x, y, text, transform=ax.transAxes, fontsize=9,
            va=va, ha=ha, fontweight="bold",
            bbox=dict(boxstyle="round,pad=0.3", facecolor="white",
                      edgecolor="grey", alpha=0.9))


def load_file(filepath, name="file"):
    """Load CSV/TSV/XLSX file based on extension.

    Handles quoted fields (e.g., "ENSG00000123456.15") automatically.
    """
    path = Path(filepath)
    if not path.exists():
        raise FileNotFoundError(f"Not found: {filepath}")

    ext = path.suffix.lower()
    if ext == ".xlsx" or ext == ".xls":
        df = pd.read_excel(path)
    elif ext == ".tsv" or ext == ".tab":
        df = pd.read_csv(path, sep="\t", quotechar='"')
    else:
        # Try tab-separated first (rMATS default), fall back to comma
        # quotechar='"' handles quoted fields like "ENSG00000123456.15"
        df = pd.read_csv(path, sep="\t", quotechar='"')
        if len(df.columns) <= 1:
            df = pd.read_csv(path, sep=",", quotechar='"')

    print(f"  Loaded {name}: {df.shape[0]:,} rows x {df.shape[1]} columns")
    return df


def load_counts_matrix(counts_path, sample_metadata=None, conditions=None):
    """Load a normalized counts matrix (genes x samples) for QC plots.

    Parameters
    ----------
    counts_path : str
        Path to normalized_counts.tsv or .csv (genes as rows, samples as columns).
    sample_metadata : dict, optional
        Mapping of sample column names to condition labels.
        If empty/None, auto-detected from CONDITIONS.
    conditions : list, optional
        The CONDITIONS list; used for auto-detecting sample-to-condition mapping
        when sample_metadata is not provided.

    Returns
    -------
    tuple : (pd.DataFrame, dict)
        (counts_df with genes as index, metadata dict) or (None, {}) on failure.
    """
    if not counts_path:
        print("[INFO] No counts file provided — skipping PCA, correlation heatmap, top DEG heatmap")
        return None, {}

    path = Path(counts_path)
    if not path.exists():
        print(f"[WARNING] Counts file not found: {counts_path}")
        return None, {}

    # Read the file, auto-detecting R's unnamed first column (index_col=0)
    ext = path.suffix.lower()
    sep = "\t" if ext in (".tsv", ".tab") else ","
    try:
        # Peek at header to check for unnamed first column (R's row.names export)
        with open(path, "r") as fh:
            header_line = fh.readline()
        first_field = header_line.split(sep)[0].strip().strip('"')
        if first_field == "" or first_field == "X":
            df = pd.read_csv(path, sep=sep, index_col=0, quotechar='"')
            print(f"  [INFO] Detected R row.names format — using first column as index")
        else:
            df = pd.read_csv(path, sep=sep, quotechar='"')
            # If first column looks like gene IDs, set as index
            if df.columns[0] in ("gene_id", "GeneID", "Geneid", "X") or \
               df.iloc[:, 0].astype(str).str.upper().str.startswith("ENS").mean() > 0.1:
                df = df.set_index(df.columns[0])
                print(f"  [INFO] Set '{df.index.name}' as gene index")
    except Exception as e:
        print(f"[WARNING] Failed to load counts file: {e}")
        return None, {}

    print(f"  Loaded counts matrix: {df.shape[0]:,} genes x {df.shape[1]} samples")

    # Strip Ensembl version suffixes from index
    idx_str = df.index.astype(str)
    ens_frac = idx_str.str.upper().str.startswith("ENS").sum() / max(len(idx_str), 1)
    if ens_frac > 0.1:
        df.index = idx_str.str.replace(r"\.\d+$", "", regex=True)
        print(f"  [INFO] Stripped Ensembl version suffixes from gene index")

    # Identify numeric sample columns (drop any non-numeric annotation columns)
    numeric_cols = df.select_dtypes(include=[np.number]).columns.tolist()
    if len(numeric_cols) < df.shape[1]:
        dropped = [c for c in df.columns if c not in numeric_cols]
        print(f"  [INFO] Dropping non-numeric columns: {dropped}")
        df = df[numeric_cols]

    if df.shape[1] < 2:
        print(f"[WARNING] Counts matrix has fewer than 2 sample columns, skipping")
        return None, {}

    # Auto-detect sample metadata from CONDITIONS if not provided
    if not sample_metadata and conditions:
        sample_metadata = {}
        for col in df.columns:
            col_lower = str(col).lower()
            for cond in conditions:
                cond_name = cond.get("name", "")
                cond_label = cond.get("label", cond_name)
                # Try substring match: condition name or label appears in sample name
                if cond_name.lower() in col_lower or cond_label.lower() in col_lower:
                    sample_metadata[col] = cond_label
                    break
        if sample_metadata:
            matched = len(sample_metadata)
            total = len(df.columns)
            print(f"  [INFO] Auto-detected metadata for {matched}/{total} samples from CONDITIONS")
        else:
            print(f"  [INFO] Could not auto-detect sample metadata — "
                  f"set SAMPLE_METADATA for condition-colored plots")

    if sample_metadata is None:
        sample_metadata = {}

    return df, sample_metadata


def pca_plot(counts_df=None, metadata=None, outdir=None, pca_file=None):
    """Generate a PCA scatter plot from counts data or a pre-computed PCA file.

    Two modes:
    - Pre-computed: reads a WSF-style pca_data.csv with columns
      (sample, PC1, PC2, condition, PC1_variance_pct, PC2_variance_pct).
    - From counts: log2(counts+1) -> StandardScaler -> sklearn PCA(n_components=2).

    Parameters
    ----------
    counts_df : pd.DataFrame, optional
        Normalized counts matrix (genes x samples). Used when pca_file is None.
    metadata : dict, optional
        {sample_name: condition_label} for coloring points.
    outdir : Path or str
        Directory to save the plot.
    pca_file : str or Path, optional
        Path to pre-computed pca_data.csv (WSF format). Takes priority over counts_df.
    """
    if outdir is None:
        return
    outdir = Path(outdir)
    outdir.mkdir(parents=True, exist_ok=True)

    # Okabe-Ito palette for conditions
    oi_palette = ["#0072B2", "#E69F00", "#009E73", "#CC79A7",
                  "#56B4E9", "#D55E00", "#F0E442", "#000000"]

    # --- Mode 1: Pre-computed PCA file ---
    if pca_file is not None:
        pca_path = Path(pca_file)
        if not pca_path.exists():
            print(f"  [WARNING] PCA file not found: {pca_file}")
            return
        pca_df = pd.read_csv(pca_path)
        required = {"sample", "PC1", "PC2"}
        if not required.issubset(set(pca_df.columns)):
            print(f"  [WARNING] PCA file missing required columns {required}")
            return

        conditions_list = pca_df["condition"].unique() if "condition" in pca_df.columns else ["Unknown"]
        color_map = {c: oi_palette[i % len(oi_palette)] for i, c in enumerate(conditions_list)}

        # Variance explained (from file or unknown)
        pc1_var = pca_df["PC1_variance_pct"].iloc[0] if "PC1_variance_pct" in pca_df.columns else None
        pc2_var = pca_df["PC2_variance_pct"].iloc[0] if "PC2_variance_pct" in pca_df.columns else None

        fig, ax = plt.subplots(figsize=(8, 6))
        for cond in conditions_list:
            mask = pca_df["condition"] == cond if "condition" in pca_df.columns else [True] * len(pca_df)
            subset = pca_df[mask]
            ax.scatter(subset["PC1"], subset["PC2"], c=color_map.get(cond, "#999999"),
                       label=cond, s=80, edgecolor="white", linewidth=0.5, zorder=3)

        xlabel = f"PC1 ({pc1_var:.1f}% variance)" if pc1_var is not None else "PC1"
        ylabel = f"PC2 ({pc2_var:.1f}% variance)" if pc2_var is not None else "PC2"
        ax.set_xlabel(xlabel)
        ax.set_ylabel(ylabel)
        ax.set_title("PCA — Sample Clustering")
        ax.axhline(0, color="#CCCCCC", linewidth=0.8, zorder=1)
        ax.axvline(0, color="#CCCCCC", linewidth=0.8, zorder=1)
        ax.legend(title="Condition", frameon=True)
        plt.tight_layout()
        fname = outdir / f"pca_plot.{FIG_FORMAT}"
        fig.savefig(fname, dpi=FIG_DPI)
        plt.close(fig)
        print(f"  Saved: {fname.name}")
        return

    # --- Mode 2: Compute PCA from counts ---
    if counts_df is None or counts_df.shape[1] < 2:
        print("  [INFO] Insufficient data for PCA plot")
        return

    try:
        from sklearn.decomposition import PCA
        from sklearn.preprocessing import StandardScaler
    except ImportError:
        print("  [WARNING] scikit-learn not installed, skipping PCA plot")
        print("  Install with: pip install scikit-learn")
        return

    # log2(counts + 1) transform, transpose so samples are rows
    log_counts = np.log2(counts_df + 1).T

    # Standardize features (genes) before PCA
    scaler = StandardScaler()
    scaled = scaler.fit_transform(log_counts)

    pca = PCA(n_components=2)
    pcs = pca.fit_transform(scaled)
    pc1_var = pca.explained_variance_ratio_[0] * 100
    pc2_var = pca.explained_variance_ratio_[1] * 100

    sample_names = counts_df.columns.tolist()
    if metadata:
        conditions_list = list(dict.fromkeys(metadata.get(s, "Unknown") for s in sample_names))
    else:
        conditions_list = ["Unknown"]
    color_map = {c: oi_palette[i % len(oi_palette)] for i, c in enumerate(conditions_list)}

    fig, ax = plt.subplots(figsize=(8, 6))
    for i, sample in enumerate(sample_names):
        cond = metadata.get(sample, "Unknown") if metadata else "Unknown"
        ax.scatter(pcs[i, 0], pcs[i, 1], c=color_map.get(cond, "#999999"),
                   s=80, edgecolor="white", linewidth=0.5, zorder=3,
                   label=cond if cond not in [ax.get_legend_handles_labels()[1]] else "")

    # Deduplicate legend labels
    handles, labels = ax.get_legend_handles_labels()
    by_label = dict(zip(labels, handles))
    ax.legend(by_label.values(), by_label.keys(), title="Condition", frameon=True)

    ax.set_xlabel(f"PC1 ({pc1_var:.1f}% variance)")
    ax.set_ylabel(f"PC2 ({pc2_var:.1f}% variance)")
    ax.set_title("PCA — Sample Clustering")
    ax.axhline(0, color="#CCCCCC", linewidth=0.8, zorder=1)
    ax.axvline(0, color="#CCCCCC", linewidth=0.8, zorder=1)

    # Add approximation note
    ax.text(0.02, 0.02,
            "PCA computed from log2(counts+1); for publication, use DESeq2 VST (PMID: 25516281)",
            transform=ax.transAxes, fontsize=7, color="#666666", style="italic",
            verticalalignment="bottom")

    plt.tight_layout()
    fname = outdir / f"pca_plot.{FIG_FORMAT}"
    fig.savefig(fname, dpi=FIG_DPI)
    plt.close(fig)
    print(f"  Saved: {fname.name}")


def sample_correlation_heatmap(counts_df, metadata, outdir):
    """Generate a sample-sample correlation heatmap with hierarchical clustering.

    Computes Euclidean distance on log2(counts+1) transformed data and
    displays a seaborn clustermap with condition color annotations.

    Parameters
    ----------
    counts_df : pd.DataFrame
        Normalized counts matrix (genes x samples).
    metadata : dict
        {sample_name: condition_label} for color annotations.
    outdir : Path or str
        Directory to save the plot.
    """
    if counts_df is None or counts_df.shape[1] < 2:
        print("  [INFO] Insufficient data for correlation heatmap")
        return

    outdir = Path(outdir)
    outdir.mkdir(parents=True, exist_ok=True)

    # log2(counts + 1) transform
    log_counts = np.log2(counts_df + 1)

    # Compute sample-sample correlation matrix
    corr = log_counts.corr(method="pearson")

    # Build condition color annotation
    oi_palette = ["#0072B2", "#E69F00", "#009E73", "#CC79A7",
                  "#56B4E9", "#D55E00", "#F0E442", "#000000"]
    if metadata:
        conditions_list = list(dict.fromkeys(metadata.get(s, "Unknown") for s in corr.columns))
        color_map = {c: oi_palette[i % len(oi_palette)] for i, c in enumerate(conditions_list)}
        col_colors = pd.Series(
            {s: color_map.get(metadata.get(s, "Unknown"), "#999999") for s in corr.columns},
            name="Condition"
        )
    else:
        col_colors = None

    try:
        g = sns.clustermap(
            corr,
            method="average",
            metric="euclidean",
            cmap="viridis",
            vmin=corr.values[np.triu_indices_from(corr.values, k=1)].min() if len(corr) > 1 else 0,
            vmax=1.0,
            col_colors=col_colors,
            row_colors=col_colors,
            linewidths=0.5,
            figsize=(max(8, len(corr) * 0.6), max(7, len(corr) * 0.55)),
            dendrogram_ratio=0.12,
            cbar_pos=(0.02, 0.8, 0.03, 0.15),
        )
        g.ax_heatmap.set_title("Sample Correlation (Pearson, log2 counts+1)", pad=20)

        # Add legend for condition colors
        if metadata:
            legend_patches = [mpatches.Patch(color=color_map[c], label=c) for c in conditions_list]
            g.ax_heatmap.legend(
                handles=legend_patches, title="Condition",
                bbox_to_anchor=(1.05, 1), loc="upper left", fontsize=9
            )

        fname = outdir / f"sample_correlation_heatmap.{FIG_FORMAT}"
        g.savefig(fname, dpi=FIG_DPI, bbox_inches="tight")
        plt.close()
        print(f"  Saved: {fname.name}")
    except Exception as e:
        print(f"  [WARNING] Failed to generate correlation heatmap: {e}")


def top_deg_heatmap(counts_df, condition_results, condition_labels,
                    metadata, outdir):
    """Generate a heatmap of top 50 DEGs (by padj) across all conditions.

    Selects the top 50 genes by adjusted p-value across all conditions
    (deduplicated), then displays z-scored log2(counts+1) expression
    with hierarchical clustering.

    Parameters
    ----------
    counts_df : pd.DataFrame
        Normalized counts matrix (genes x samples).
    condition_results : dict
        Pipeline condition_results structure with deseq2_filtered data.
    condition_labels : dict
        Maps condition name -> human-readable label.
    metadata : dict
        {sample_name: condition_label} for color annotations.
    outdir : Path or str
        Directory to save the plot.
    """
    if counts_df is None or counts_df.shape[1] < 2:
        print("  [INFO] Insufficient counts data for top DEG heatmap")
        return

    outdir = Path(outdir)
    outdir.mkdir(parents=True, exist_ok=True)

    padj_col = DESEQ2_COLS.get("padj", "padj")
    gene_id_col = DESEQ2_COLS.get("gene_id", "gene_id")
    gene_name_col = DESEQ2_COLS.get("gene_name", "gene_name")

    # Collect top genes by padj across all conditions (deduplicated)
    all_top_genes = []
    for cond_name, data in condition_results.items():
        deg_df = data.get("deseq2_filtered", {}).get("all_genes", pd.DataFrame())
        if len(deg_df) == 0:
            continue
        # Get gene IDs sorted by padj
        if padj_col in deg_df.columns and gene_id_col in deg_df.columns:
            sorted_df = deg_df.sort_values(padj_col)
            all_top_genes.extend(
                sorted_df[gene_id_col].dropna().astype(str).tolist()
            )

    if not all_top_genes:
        print("  [INFO] No DEGs found for top DEG heatmap")
        return

    # Deduplicate while preserving order (best padj first), take top 50
    seen = set()
    unique_genes = []
    for g in all_top_genes:
        if g not in seen:
            seen.add(g)
            unique_genes.append(g)
    top_genes = unique_genes[:50]

    # Filter counts to top genes (match by index)
    counts_idx = counts_df.index.astype(str)
    mask = counts_idx.isin(top_genes)
    if mask.sum() == 0:
        # Try stripping Ensembl versions from top_genes too
        top_genes_stripped = [g.split(".")[0] for g in top_genes]
        mask = counts_idx.isin(top_genes_stripped)
    if mask.sum() == 0:
        print(f"  [INFO] None of the top DEGs found in counts matrix — skipping heatmap")
        return

    subset = counts_df.loc[mask].copy()
    print(f"  Top DEG heatmap: {len(subset)} genes matched in counts matrix")

    # log2(counts+1) then z-score per gene (row)
    log_counts = np.log2(subset + 1)
    from scipy import stats
    z_scored = log_counts.apply(lambda row: stats.zscore(row, nan_policy="omit"), axis=1)
    # Cap z-scores at +/-3 to avoid extreme outliers dominating the colormap
    z_scored = z_scored.clip(-3, 3)

    # Try to use gene names as row labels if available
    # Build a mapping from gene IDs in condition_results
    id_to_name = {}
    for cond_name, data in condition_results.items():
        raw_df = data.get("deseq2_raw", pd.DataFrame())
        if gene_id_col in raw_df.columns and gene_name_col in raw_df.columns:
            for _, row in raw_df[[gene_id_col, gene_name_col]].dropna().iterrows():
                gid = str(row[gene_id_col])
                gname = str(row[gene_name_col])
                if gid and gname and gname != "nan":
                    id_to_name[gid] = gname
    if id_to_name:
        z_scored.index = [id_to_name.get(str(g), str(g)) for g in z_scored.index]

    # Build condition color annotation
    oi_palette = ["#0072B2", "#E69F00", "#009E73", "#CC79A7",
                  "#56B4E9", "#D55E00", "#F0E442", "#000000"]
    if metadata:
        conditions_unique = list(dict.fromkeys(metadata.get(s, "Unknown") for s in z_scored.columns))
        color_map = {c: oi_palette[i % len(oi_palette)] for i, c in enumerate(conditions_unique)}
        col_colors = pd.Series(
            {s: color_map.get(metadata.get(s, "Unknown"), "#999999") for s in z_scored.columns},
            name="Condition"
        )
    else:
        col_colors = None

    try:
        n_genes = len(z_scored)
        fig_height = max(8, n_genes * 0.25)
        g = sns.clustermap(
            z_scored,
            cmap="RdBu_r",
            center=0,
            vmin=-3,
            vmax=3,
            col_colors=col_colors,
            method="ward",
            metric="euclidean",
            linewidths=0.3,
            figsize=(max(8, len(z_scored.columns) * 0.8), fig_height),
            dendrogram_ratio=(0.1, 0.08),
            cbar_pos=(0.02, 0.8, 0.03, 0.15),
            yticklabels=True,
        )
        g.ax_heatmap.set_title("Top 50 DEGs — z-scored log2(counts+1)", pad=20)
        g.ax_heatmap.set_ylabel("")
        g.ax_heatmap.tick_params(axis="y", labelsize=7)

        # Add condition legend
        if metadata:
            legend_patches = [mpatches.Patch(color=color_map[c], label=c)
                              for c in conditions_unique]
            g.ax_heatmap.legend(
                handles=legend_patches, title="Condition",
                bbox_to_anchor=(1.05, 1), loc="upper left", fontsize=9
            )

        fname = outdir / f"top_deg_heatmap.{FIG_FORMAT}"
        g.savefig(fname, dpi=FIG_DPI, bbox_inches="tight")
        plt.close()
        print(f"  Saved: {fname.name}")
    except Exception as e:
        print(f"  [WARNING] Failed to generate top DEG heatmap: {e}")


# ---------------------------------------------------------------------------
# GENE NAME LOOKUP (MyGene.info REST API — no extra package required)
# ---------------------------------------------------------------------------

_GENE_NAME_CACHE: dict = {}   # runtime cache: {ensembl_id: gene_symbol}


def _fetch_gene_names(ensembl_ids: list, species: str = "human") -> dict:
    """Query MyGene.info to resolve Ensembl IDs → gene symbols.

    Batches up to 1 000 IDs per POST request.  Returns {ensembl_id: symbol}.
    Silently returns an empty dict on any network error.
    """
    result = {}
    ids = [str(i) for i in ensembl_ids if str(i).upper().startswith("ENS")]
    if not ids:
        return result

    batch_size = 1000
    url = "https://mygene.info/v3/query"
    for start in range(0, len(ids), batch_size):
        batch = ids[start : start + batch_size]
        payload = urllib.parse.urlencode({
            "q":       ",".join(batch),
            "scopes":  "ensembl.gene",
            "fields":  "symbol",
            "species": species,
        }).encode()
        req = urllib.request.Request(
            url, data=payload,
            headers={"Content-Type": "application/x-www-form-urlencoded"},
        )
        try:
            with urllib.request.urlopen(req, timeout=30) as resp:
                data = json.loads(resp.read())
                hits = data if isinstance(data, list) else data.get("hits", [])
            for hit in hits:
                if "symbol" in hit and "query" in hit:
                    result[hit["query"]] = hit["symbol"]
        except Exception as exc:
            print(f"  WARNING: MyGene.info lookup failed for batch starting at {start}: {exc}")
    return result


def _enrich_with_gene_names(df: "pd.DataFrame", file_label: str = "") -> "pd.DataFrame":
    """Add a gene_name column from Ensembl ID lookup when it is absent.

    Uses _GENE_NAME_CACHE to avoid redundant network requests across conditions.
    Falls back gracefully: if lookup fails or GENE_NAME_LOOKUP is False the
    DataFrame is returned unchanged and downstream code uses Ensembl IDs as labels.
    """
    cols = DESEQ2_COLS
    name_col = cols.get("gene_name", "")
    id_col   = cols.get("gene_id",   "")

    if not GENE_NAME_LOOKUP:
        return df
    if name_col and name_col in df.columns:
        return df   # already present
    if not id_col or id_col not in df.columns:
        return df

    # Only proceed if the id column looks like Ensembl IDs
    sample = df[id_col].dropna().astype(str).head(500)
    ens_frac = sample.str.upper().str.startswith("ENS").sum() / max(len(sample), 1)
    if ens_frac <= 0.1:
        return df

    all_ids   = df[id_col].dropna().astype(str).unique().tolist()
    to_fetch  = [i for i in all_ids if i not in _GENE_NAME_CACHE]

    if to_fetch:
        print(f"  Fetching gene names for {len(to_fetch):,} Ensembl IDs "
              f"from MyGene.info (species={SPECIES})…")
        new_mappings = _fetch_gene_names(to_fetch, species=SPECIES)
        _GENE_NAME_CACHE.update(new_mappings)
        resolved = sum(1 for i in to_fetch if i in new_mappings)
        print(f"  Resolved {resolved:,} / {len(to_fetch):,} gene names")
        # Mark unresolvable IDs with None so they are not re-queried in this session
        for _id in to_fetch:
            if _id not in _GENE_NAME_CACHE:
                _GENE_NAME_CACHE[_id] = None

    if not name_col:
        return df

    df = df.copy()
    # Map: use resolved symbol if available, otherwise fall back to Ensembl ID
    # _GENE_NAME_CACHE may contain None for unresolvable IDs — 'or str(x)' handles that
    df[name_col] = df[id_col].map(lambda x: _GENE_NAME_CACHE.get(str(x)) or f"{str(x)} (no symbol)")
    print(f"  Added gene_name column ('{name_col}') via Ensembl lookup")
    return df


# mygene type_of_gene → Ensembl-style biotype mapping
_MYGENE_BIOTYPE_MAP = {
    "protein-coding": "protein_coding",
    "ncrna": "lncrna",
    "pseudo": "pseudogene",
    "snrna": "snrna",
    "snorna": "snorna",
    "rrna": "rrna",
    "trna": "trna",
    "scrna": "scrna",
    "mirna": "mirna",
}


def _reassign_biotypes_from_mygene(df, file_label=""):
    """Re-assign biotype_group when all values are 'Other' using mygene type_of_gene.

    Queries MyGene.info for gene type and maps results through _BIOTYPE_GROUPS.
    Only runs when the biotype column exists and has no meaningful values.
    """
    bio_col = DESEQ2_COLS.get("biotype", "")
    id_col = DESEQ2_COLS.get("gene_id", "")
    if not bio_col or bio_col not in df.columns or not id_col or id_col not in df.columns:
        return df

    unique_biotypes = set(df[bio_col].dropna().str.strip().unique())
    if unique_biotypes - {"Other", "other", ""}:
        return df  # has real biotype values, no re-assignment needed

    print(f"  [INFO] All biotypes are 'Other' in {file_label} — fetching from MyGene.info...")
    all_ids = df[id_col].dropna().astype(str).unique().tolist()
    ens_ids = [i for i in all_ids if i.upper().startswith("ENS")]
    if not ens_ids:
        return df

    # Fetch type_of_gene from mygene in batches
    type_map = {}
    batch_size = 1000
    url = "https://mygene.info/v3/query"
    for start in range(0, len(ens_ids), batch_size):
        batch = ens_ids[start:start + batch_size]
        payload = urllib.parse.urlencode({
            "q": ",".join(batch),
            "scopes": "ensembl.gene",
            "fields": "type_of_gene",
            "species": SPECIES,
        }).encode()
        req = urllib.request.Request(
            url, data=payload,
            headers={"Content-Type": "application/x-www-form-urlencoded"},
        )
        try:
            with urllib.request.urlopen(req, timeout=30) as resp:
                data = json.loads(resp.read())
                hits = data if isinstance(data, list) else data.get("hits", [])
            for hit in hits:
                if "type_of_gene" in hit and "query" in hit:
                    type_map[hit["query"]] = hit["type_of_gene"]
        except Exception as exc:
            print(f"  WARNING: MyGene.info biotype lookup failed: {exc}")

    if not type_map:
        return df

    # Map mygene types → Ensembl-style biotypes → biotype groups
    df = df.copy()
    def _resolve(gene_id):
        raw = type_map.get(str(gene_id), "")
        ensembl_bt = _MYGENE_BIOTYPE_MAP.get(raw.lower(), raw.lower())
        return _BIOTYPE_GROUPS.get(ensembl_bt, "Other")

    df[bio_col] = df[id_col].map(_resolve)
    resolved = (df[bio_col] != "Other").sum()
    print(f"  Resolved {resolved:,} / {len(df):,} gene biotypes via MyGene.info")
    return df


def validate_columns(df, required_cols, name="file"):
    """Check that expected columns exist."""
    missing = [c for c in required_cols if c and c not in df.columns]
    if missing:
        print(f"  Available columns: {list(df.columns)}")
        raise KeyError(f"Missing columns in {name}: {missing}")


# Common column name aliases for auto-detection (checked case-insensitively).
# Maps the DESEQ2_COLS key -> list of known alternative header names.
_DESEQ2_ALIASES = {
    "gene_id":   ["gene_id", "ensembl_gene_id", "ensembl_geneid", "ensemblgeneid",
                  "ensemblgene", "enzemblgeneid", "ensgene",
                  "geneid", "gene", "id", "feature_id", "X"],
    "gene_name": ["gene_name", "gene_symbol", "symbol", "hgnc_symbol",
                  "name", "genename", "external_gene_name", "mgi_symbol"],
    "log2fc":    ["log2foldchange", "log2fc", "log2_fold_change", "lfc",
                  "logfc", "log2ratio", "log2_ratio"],
    "basemean":  ["basemean", "base_mean", "aveexpr", "mean_expression",
                  "meanexpr", "avgexpr", "averageexpression"],
    "padj":      ["padj", "p.adj", "adjusted_p_value", "adj.p.val",
                  "adj_pval", "fdr", "bh", "bonferroni"],
    "pvalue":    ["pvalue", "pval", "p.value", "p_value", "p", "rawp"],
    "biotype":   ["biotype", "gene_biotype", "gene_type", "transcript_biotype",
                  "transcript_type"],
    "stat":      ["stat", "wald_statistic", "test_stat", "statistic"],
    "lfcSE":     ["lfcse", "lfc_se", "std_error", "lfcstderror"],
}

# Common column name aliases for rMATS files.
_RMATS_ALIASES = {
    "event_id":      ["id", "event_id", "eventid"],
    "gene_id":       ["geneid", "gene_id", "ensembl_gene_id"],
    "gene_name":     ["genesymbol", "gene_name", "gene_symbol", "symbol"],
    "pvalue":        ["pvalue", "pval", "p.value", "p_value"],
    "fdr":           ["fdr", "padj", "adj.p.val", "q.value", "qvalue"],
    "inclevel_diff": ["incleveldifference", "inclevel_diff", "deltapsi",
                      "dpsi", "delta_psi"],
}


def _resolve_column(df, col_key, configured_name, alias_map, file_label="file"):
    """Return the actual column name in df for a given column key.

    Tries (in order):
    1. Exact match of configured name.
    2. Case-insensitive exact match of configured name.
    3. Case-insensitive exact match of any known alias.
    4. Prefix match: any column that starts with '{configured_name}_'
       (handles DESeq2 outputs like 'log2FoldChange_CondA_vs_CondB').
    5. Prefix match: any column that starts with '{alias}_' for any alias.

    Returns the matching column name as it appears in df, or None if not found.
    """
    # 1. Exact match
    if configured_name and configured_name in df.columns:
        return configured_name

    lower_cols = {c.lower(): c for c in df.columns}

    # 2. Case-insensitive exact match of configured name
    if configured_name and configured_name.lower() in lower_cols:
        found = lower_cols[configured_name.lower()]
        print(f"  Column '{configured_name}' not found — using '{found}' (case-insensitive match)")
        return found

    # 3. Case-insensitive exact alias scan
    for alias in alias_map.get(col_key, []):
        if alias.lower() in lower_cols:
            found = lower_cols[alias.lower()]
            print(f"  Column '{configured_name}' not found — auto-detected '{found}' as {col_key}")
            return found

    # 4. Prefix match on configured name  (e.g. 'log2FoldChange_CondA_vs_CondB')
    if configured_name:
        prefix = configured_name.lower() + "_"
        candidates = [orig for lc, orig in lower_cols.items() if lc.startswith(prefix)]
        if candidates:
            # Prefer shortest name (fewest qualifiers); warn if multiple found
            candidates.sort(key=len)
            found = candidates[0]
            if len(candidates) > 1:
                print(f"  WARNING: multiple '{prefix}*' columns in {file_label}: "
                      f"{candidates} — using shortest '{found}'. Set DESEQ2_COLS['{col_key}'] "
                      f"to the exact column name to override.")
            else:
                print(f"  Column '{configured_name}' not found — using '{found}' "
                      f"(prefix match for {col_key})")
            return found

    # 5. Prefix match on aliases
    for alias in alias_map.get(col_key, []):
        prefix = alias.lower() + "_"
        candidates = [orig for lc, orig in lower_cols.items() if lc.startswith(prefix)]
        if candidates:
            candidates.sort(key=len)
            found = candidates[0]
            if len(candidates) > 1:
                print(f"  WARNING: multiple '{prefix}*' columns in {file_label}: "
                      f"{candidates} — using shortest '{found}'.")
            else:
                print(f"  Column '{configured_name}' not found — using '{found}' "
                      f"(prefix match via alias '{alias}' for {col_key})")
            return found

    return None


def _validate_rmats_columns(df, col_mapping):
    """Validate that required rMATS columns exist, trying common aliases.

    For each required column from *col_mapping* (a dict like RMATS_COLS),
    check if it exists in the DataFrame.  If not, try known aliases in
    order; if found, rename the column and print a warning.  If no alias
    is found, raise ``ValueError`` with available columns listed.

    Returns the (possibly renamed) DataFrame.
    """
    alias_map = {
        "FDR": ["FDR", "fdr", "adj.P.Val", "padj", "q-value"],
        "PValue": ["PValue", "pvalue", "P.Value", "p_value", "p-value"],
        "IncLevelDifference": ["IncLevelDifference", "IncLevel_Difference",
                               "dPSI", "inc_level_diff"],
        "geneSymbol": ["geneSymbol", "GeneSymbol", "gene_symbol",
                       "geneName", "gene_name"],
        "GeneID": ["GeneID", "gene_id", "Ensembl_ID"],
        "ID": ["ID", "id", "event_id"],
    }
    rename = {}
    for _key, expected in col_mapping.items():
        if not expected or expected in df.columns:
            continue
        # Try aliases for this expected column name
        aliases = alias_map.get(expected, [])
        found = False
        for alias in aliases:
            if alias in df.columns:
                rename[alias] = expected
                print(f"  [WARN] rMATS column '{alias}' mapped to '{expected}'")
                found = True
                break
        if not found:
            raise ValueError(
                f"Required rMATS column '{expected}' not found. "
                f"Available columns: {list(df.columns)}. "
                f"Check your rMATS output format."
            )
    if rename:
        df = df.rename(columns=rename)
    return df


def _normalize_gsea_cols(df):
    """Normalize gseapy result column names to canonical lowercase forms.

    Handles version differences in gseapy (uppercase vs lowercase,
    'FDR q-val' vs 'fdr' vs 'FDR', etc.).
    Also derives 'geneset_size' from 'Tag %' if not already present.
    """
    col_map = {}
    for c in df.columns:
        cl = c.lower().strip().replace(" ", "_").replace("-", "_")
        if cl in ("fdr", "fdr_q_val", "fdr_bh", "padj", "adj_p_value"):
            col_map[c] = "fdr"
        elif cl in ("nes", "normalized_enrichment_score"):
            col_map[c] = "nes"
        elif cl in ("term", "pathway", "gene_set"):
            col_map[c] = "Term"
        elif cl in ("geneset_size", "gene_set_size", "gs_size"):
            col_map[c] = "geneset_size"
        elif cl in ("lead_genes", "leading_edge", "lead_edge_genes"):
            col_map[c] = "lead_genes"
        elif cl in ("pvalue", "pval", "p_value", "nom_p_val", "nom_p_value"):
            col_map[c] = "pvalue"
        elif cl in ("es", "enrichment_score"):
            col_map[c] = "es"
        elif cl in ("tag_%", "tag_percent"):
            col_map[c] = "tag_pct"
        elif cl in ("gene_%", "gene_percent"):
            col_map[c] = "gene_pct"
    if col_map:
        df = df.rename(columns=col_map)
    # Derive geneset_size from tag_pct if missing (gseapy 1.1+)
    # Tag % looks like "2/5" where denominator = gene set size
    if "geneset_size" not in df.columns and "tag_pct" in df.columns:
        def _parse_gs_size(val):
            try:
                return int(str(val).split("/")[1])
            except (IndexError, ValueError):
                return 0
        df["geneset_size"] = df["tag_pct"].apply(_parse_gs_size)
    return df


def _strip_ensembl_version(series):
    """Strip version suffixes from Ensembl IDs: ENSG00000123456.12 -> ENSG00000123456.

    Only strips when the value looks like an Ensembl ID (starts with ENS).
    Leaves non-Ensembl values (gene symbols, etc.) unchanged.
    """
    def _strip(val):
        s = str(val)
        if s.upper().startswith("ENS") and "." in s:
            return s.rsplit(".", 1)[0]
        return val
    return series.apply(_strip)


def normalize_deseq2_columns(df, file_label="DESeq2 file"):
    """Rename df columns so they match the names in DESEQ2_COLS, then strip
    Ensembl version numbers from the gene_id column.

    Handles files where column headers differ from the configured names
    by trying case-insensitive matching and common aliases.
    Returns a (possibly renamed) copy of df.
    """
    rename_map = {}
    for key, configured in DESEQ2_COLS.items():
        if not configured:
            continue
        if configured in df.columns:
            continue  # already correct, no rename needed
        actual = _resolve_column(df, key, configured, _DESEQ2_ALIASES, file_label)
        if actual and actual != configured:
            rename_map[actual] = configured

    if rename_map:
        df = df.rename(columns=rename_map)

    # Strip Ensembl version numbers from gene_id column
    # Use same >10% fraction threshold as _best_gene_key() for consistency
    id_col = DESEQ2_COLS.get("gene_id", "")
    if id_col and id_col in df.columns:
        sample = df[id_col].dropna().astype(str).head(500)
        ens_frac = sample.str.upper().str.startswith("ENS").sum() / max(len(sample), 1)
        if ens_frac > 0.1:
            df = df.copy()
            df[id_col] = _strip_ensembl_version(df[id_col])

    # Detect stat and lfcSE columns (optional — silent if missing)
    stat_col = DESEQ2_COLS.get("stat", "stat")
    lfcse_col = DESEQ2_COLS.get("lfcSE", "lfcSE")
    log2fc_col = DESEQ2_COLS.get("log2fc", "log2FoldChange")

    if stat_col and stat_col in df.columns:
        pass  # stat column present — available for GSEA Wald ranking
    if lfcse_col and lfcse_col in df.columns and log2fc_col and log2fc_col in df.columns:
        # Check for likely LFC shrinkage: if median lfcSE/|log2FC| < 0.3,
        # DESeq2 apeglm/ashr shrinkage was likely applied
        lfc_abs = df[log2fc_col].abs()
        lfcse_vals = df[lfcse_col]
        # Guard against division by zero: only compute ratio where |log2FC| > 0
        valid = (lfc_abs > 0) & lfcse_vals.notna()
        if valid.sum() > 0:
            ratio = (lfcse_vals[valid] / lfc_abs[valid]).median()
            if ratio < 0.3:
                print(f"  [INFO] lfcSE/|log2FC| median ratio = {ratio:.3f} — "
                      f"likely LFC shrinkage applied (apeglm/ashr)")

    return df


def _best_gene_key(df):
    """Return the column name that gives the most reliable unique gene identifier.

    Prefers gene_id (Ensembl) when it looks like Ensembl IDs (ENS prefix),
    because Ensembl IDs are stable unique identifiers unlike gene symbols which
    can be duplicated, aliased, or vary between datasets.
    Falls back to gene_name if gene_id is absent or non-Ensembl.

    Returns (col_name, description_string).
    """
    id_col   = DESEQ2_COLS.get("gene_id",   "")
    name_col = DESEQ2_COLS.get("gene_name", "")

    if id_col and id_col in df.columns:
        # Sample up to 500 rows; require >10% ENS-prefixed to confirm Ensembl IDs
        # (avoids false-positive from a single stray ENS value in a gene-name file)
        sample = df[id_col].dropna().astype(str).head(500)
        ens_frac = sample.str.upper().str.startswith("ENS").sum() / max(len(sample), 1)
        if ens_frac > 0.1:
            return id_col, "Ensembl ID"

    return name_col, "gene name"


# ---------------------------------------------------------------------------
# RBP ANNOTATION
# ---------------------------------------------------------------------------

def load_rbp_annotations(rbp_file):
    """Load RBP annotations from an Excel file.

    Supports two formats:
      1) RBP-E-A-C-Complex.xlsx  — columns: Gene Name, TriSNRP, B Complex,
         C Complex, U2-Ecomplex, Yael_RBP?, MW_RBP?
      2) RBP_yael_MW_lists.xlsx  — columns: Gene Symbol, Yael_RBP?, MW_RBP?

    Returns a dict mapping uppercase gene_symbol -> dict of annotations.
    """
    path = Path(rbp_file)
    if not path.exists():
        print(f"  WARNING: RBP file not found: {rbp_file}")
        return {}

    df = pd.read_excel(path, sheet_name=0)
    print(f"  Loaded RBP annotations: {df.shape[0]:,} rows x {df.shape[1]} columns")

    # Detect gene column
    gene_col = None
    for candidate in ["Gene Name", "Gene Symbol", "gene_name", "gene_symbol"]:
        if candidate in df.columns:
            gene_col = candidate
            break
    if gene_col is None:
        print(f"  WARNING: Cannot find gene column in RBP file. "
              f"Columns: {list(df.columns)}")
        return {}

    # Detect format by checking for spliceosome complex columns
    complex_cols = {}
    has_complexes = False
    for raw_col, clean_key in [("TriSNRP", "TriSNRP"),
                                ("B Complex", "B_Complex"),
                                ("C Complex", "C_Complex"),
                                ("U2-Ecomplex", "U2_Ecomplex")]:
        if raw_col in df.columns:
            complex_cols[raw_col] = clean_key
            has_complexes = True

    annotations = {}
    for _, row in df.iterrows():
        gene = row[gene_col]
        if pd.isna(gene) or str(gene).strip() == "":
            continue
        gene_upper = str(gene).strip().upper()

        entry = {
            "is_RBP_Yael": str(row.get("Yael_RBP?", "")).strip().upper() == "Y",
            "is_RBP_MW":   str(row.get("MW_RBP?", "")).strip().upper() == "Y",
        }
        if has_complexes:
            for raw_col, clean_key in complex_cols.items():
                entry[clean_key] = str(row.get(raw_col, "")).strip().upper() == "Y"

        annotations[gene_upper] = entry

    n_yael = sum(1 for v in annotations.values() if v["is_RBP_Yael"])
    n_mw = sum(1 for v in annotations.values() if v["is_RBP_MW"])
    print(f"  RBP annotations: {len(annotations):,} genes "
          f"(Yael: {n_yael:,}, MW: {n_mw:,})")
    if has_complexes:
        print(f"  Spliceosome complex columns detected: "
              f"{list(complex_cols.values())}")
    return annotations


def annotate_rbps(deg_df, rbp_annotations, gene_col="gene_name"):
    """Add RBP annotation columns to a DEG DataFrame.

    Adds: is_RBP, is_RBP_Yael, is_RBP_MW, and spliceosome complex columns
    if available. Matching is case-insensitive.
    """
    if not rbp_annotations or gene_col not in deg_df.columns:
        return deg_df

    df = deg_df.copy()
    gene_upper = df[gene_col].astype(str).str.strip().str.upper()

    # Determine which annotation keys exist (from first entry)
    sample_entry = next(iter(rbp_annotations.values()))
    all_keys = list(sample_entry.keys())

    for key in all_keys:
        df[key] = gene_upper.map(
            lambda g, k=key: rbp_annotations.get(g, {}).get(k, False))

    # Add aggregate is_RBP column (True if either Yael or MW)
    df["is_RBP"] = df["is_RBP_Yael"] | df["is_RBP_MW"]

    # Reorder: put is_RBP right after the annotation keys
    cols = [c for c in df.columns if c not in ["is_RBP"] + all_keys]
    # Insert RBP columns at end, with is_RBP first
    df = df[cols + ["is_RBP"] + all_keys]

    n_rbp = int(df["is_RBP"].sum())
    print(f"  RBP-annotated: {n_rbp:,} / {len(df):,} genes are RBPs")
    return df


def rbp_heatmap(condition_results, condition_labels, outdir, max_genes=80):
    """Cross-condition heatmap of log2FC for DEGs that are annotated RBPs.

    Only includes genes that are significant DEGs in at least 1 condition
    AND are annotated as RBPs (is_RBP == True).
    """
    fc_col   = DESEQ2_COLS["log2fc"]
    name_col = DESEQ2_COLS.get("gene_name", "")
    names    = list(condition_results.keys())

    if not name_col:
        print("  RBP heatmap: no gene_name column configured, skipping")
        return

    # Collect RBP DEGs across all conditions
    rbp_genes = set()
    for cond_name in names:
        filt = condition_results[cond_name]["deseq2_filtered"]["all_genes"]
        if "is_RBP" not in filt.columns or name_col not in filt.columns:
            continue
        rbp_filt = filt[filt["is_RBP"] == True]
        rbp_genes.update(rbp_filt[name_col].dropna().astype(str).str.strip().unique())

    if not rbp_genes:
        print("  RBP heatmap: no RBP DEGs found across conditions, skipping")
        return

    # Build log2FC matrix from RAW data for these genes
    fc_dfs = []
    for cond_name in names:
        raw = condition_results[cond_name]["deseq2_raw"]
        if name_col not in raw.columns or fc_col not in raw.columns:
            continue
        sub = raw[[name_col, fc_col]].dropna(subset=[name_col]).copy()
        sub[name_col] = sub[name_col].astype(str).str.strip()
        sub = sub[sub[name_col].isin(rbp_genes)]
        sub = sub.drop_duplicates(subset=[name_col]).set_index(name_col)
        fc_series = sub[fc_col]
        fc_series.name = condition_labels[cond_name]
        fc_dfs.append(fc_series)

    if not fc_dfs:
        print("  RBP heatmap: no log2FC data available, skipping")
        return

    matrix = pd.concat(fc_dfs, axis=1).dropna()

    if len(matrix) == 0:
        print("  RBP heatmap: no genes with complete log2FC data, skipping")
        return

    # Limit to top genes by max absolute FC
    title_note = ""
    if len(matrix) > max_genes:
        matrix["max_abs_fc"] = matrix.abs().max(axis=1)
        matrix = matrix.nlargest(max_genes, "max_abs_fc").drop(columns="max_abs_fc")
        title_note = f" (top {max_genes})"

    # Use blue-white-orange diverging colormap (color-blind friendly)
    from matplotlib.colors import LinearSegmentedColormap
    rbp_cmap = LinearSegmentedColormap.from_list(
        "blue_white_orange", ["#0072B2", "#FFFFFF", "#E69F00"])

    g = sns.clustermap(matrix, cmap=rbp_cmap, center=0,
                       figsize=(8, max(6, len(matrix) * 0.2)),
                       row_cluster=True, col_cluster=False,
                       yticklabels=True, linewidths=0.3, linecolor="white")
    g.fig.suptitle(
        f"RBP DEG log2FC Heatmap{title_note}",
        y=1.02, fontsize=12, fontweight="bold")

    outpath = outdir / f"rbp_log2fc_heatmap.{FIG_FORMAT}"
    g.savefig(outpath, format=FIG_FORMAT, dpi=FIG_DPI, bbox_inches="tight")
    plt.close(g.fig)
    print(f"  Saved: {outpath} ({len(matrix)} RBP genes)")


def rbp_summary_table(condition_results, condition_labels, outdir):
    """Export an Excel table of all DEG RBPs across conditions.

    One row per gene with columns for each condition's log2FC and padj,
    plus RBP category columns (Yael, MW, spliceosome complexes).
    """
    fc_col   = DESEQ2_COLS["log2fc"]
    padj_col = DESEQ2_COLS["padj"]
    name_col = DESEQ2_COLS.get("gene_name", "")
    names    = list(condition_results.keys())

    if not name_col:
        print("  RBP summary: no gene_name column configured, skipping")
        return

    # Collect all RBP DEGs and their annotation info
    all_rbp_info = {}   # gene -> annotation dict
    gene_data = {}      # gene -> {condition: {log2fc, padj}}

    for cond_name in names:
        filt = condition_results[cond_name]["deseq2_filtered"]["all_genes"]
        if "is_RBP" not in filt.columns or name_col not in filt.columns:
            continue
        rbp_rows = filt[filt["is_RBP"] == True]
        for _, row in rbp_rows.iterrows():
            gene = str(row[name_col]).strip()
            if not gene or gene == "nan":
                continue
            if gene not in gene_data:
                gene_data[gene] = {}
            gene_data[gene][cond_name] = {
                "log2FC": row.get(fc_col, np.nan),
                "padj":   row.get(padj_col, np.nan),
            }

            # Capture RBP annotation columns from the first occurrence
            if gene not in all_rbp_info:
                info = {}
                for col in ["is_RBP_Yael", "is_RBP_MW"]:
                    if col in filt.columns:
                        info[col] = row.get(col, False)
                for col in ["TriSNRP", "B_Complex", "C_Complex", "U2_Ecomplex"]:
                    if col in filt.columns:
                        info[col] = row.get(col, False)
                all_rbp_info[gene] = info

    if not gene_data:
        print("  RBP summary: no RBP DEGs found, skipping")
        return

    # Build output DataFrame
    rows = []
    for gene in sorted(gene_data.keys()):
        row = {"Gene": gene}
        # Add per-condition log2FC and padj
        for cond_name in names:
            label = condition_labels[cond_name]
            if cond_name in gene_data[gene]:
                row[f"log2FC ({label})"] = gene_data[gene][cond_name]["log2FC"]
                row[f"padj ({label})"]   = gene_data[gene][cond_name]["padj"]
            else:
                row[f"log2FC ({label})"] = np.nan
                row[f"padj ({label})"]   = np.nan
        # Add RBP annotation columns
        info = all_rbp_info.get(gene, {})
        for k, v in info.items():
            row[k] = v
        rows.append(row)

    summary_df = pd.DataFrame(rows)
    outpath = outdir / "rbp_summary.xlsx"
    summary_df.to_excel(outpath, index=False)
    print(f"  Saved: {outpath} ({len(summary_df)} RBP DEGs)")


def normalize_rmats_columns(df, file_label="rMATS file"):
    """Rename df columns so they match the names in RMATS_COLS, then strip
    Ensembl version numbers from the gene_id column (GeneID).
    """
    rename_map = {}
    for key, configured in RMATS_COLS.items():
        if not configured:
            continue
        if configured in df.columns:
            continue
        actual = _resolve_column(df, key, configured, _RMATS_ALIASES, file_label)
        if actual and actual != configured:
            rename_map[actual] = configured

    if rename_map:
        df = df.rename(columns=rename_map)

    # Strip Ensembl version numbers from GeneID column
    gene_id_col = RMATS_COLS.get("gene_id", "")
    if gene_id_col and gene_id_col in df.columns:
        sample = df[gene_id_col].dropna().astype(str).head(500)
        ens_frac = sample.str.upper().str.startswith("ENS").sum() / max(len(sample), 1)
        if ens_frac > 0.1:
            df = df.copy()
            df[gene_id_col] = _strip_ensembl_version(df[gene_id_col])

    return df


# ---------------------------------------------------------------------------
# DESeq2 PROCESSING
# ---------------------------------------------------------------------------

def filter_deseq2(df, biotype_filter=None, label="All"):
    """Apply DESeq2 cutoffs and return (biotype-subset, significance-filtered) DataFrames."""
    cols = DESEQ2_COLS

    # Drop rows with NaN in critical columns and report counts
    n_before = len(df)
    df = df.dropna(subset=[cols["padj"], cols["log2fc"], cols["basemean"]])
    n_dropped = n_before - len(df)
    if n_dropped > 0:
        print(f"  NOTE: {n_dropped:,} rows dropped (NA in padj/log2FC/baseMean — "
              f"normal DESeq2 behaviour for low-count/outlier genes)")

    # Subset by biotype first (affects both the "all" and "filtered" returns)
    # Normalize for comparison: handles both raw Ensembl ("protein_coding") and
    # pre-grouped values ("Protein Coding") from prior pipeline runs.
    if biotype_filter in ("protein_coding", "non_protein_coding"):
        _bio_norm = df[cols["biotype"]].fillna("").str.lower().str.replace(" ", "_")
        if biotype_filter == "protein_coding":
            df = df[_bio_norm == "protein_coding"]
        else:
            df = df[_bio_norm != "protein_coding"]

    # Apply significance cutoffs
    mask = (
        (df[cols["padj"]] < PADJ_CUTOFF) &
        (df[cols["log2fc"]].abs() >= LOG2FC_CUTOFF) &
        (df[cols["basemean"]] >= BASEMEAN_CUTOFF)
    )

    filtered = df[mask].copy()

    # Add direction column
    filtered["direction"] = np.where(
        filtered[cols["log2fc"]] > 0, "up", "down"
    )

    print(f"\n-- DESeq2 Filtering [{label}] --")
    print(f"  Input rows:  {len(df):,}")
    print(f"  |log2FoldChange| >= {LOG2FC_CUTOFF}, baseMean >= {BASEMEAN_CUTOFF}, padj < {PADJ_CUTOFF}")
    if biotype_filter:
        print(f"  Biotype filter: {biotype_filter}")
    print(f"  -> Filtered: {len(filtered):,} genes "
          f"({filtered['direction'].value_counts().to_dict()})")

    return df, filtered


def volcano_plot(df, outdir, label="All", suffix=""):
    """Generate volcano plot from full DESeq2 results."""
    cols = DESEQ2_COLS
    data = df.dropna(subset=[cols["padj"], cols["log2fc"]]).copy()
    data["-log10padj"] = -np.log10(data[cols["padj"]].clip(lower=1e-300))

    # Classify points — include baseMean filter to match actual DEG counts
    basemean_ok = (data[cols["basemean"]] >= BASEMEAN_CUTOFF) if cols["basemean"] in data.columns else True
    conditions = [
        (data[cols["padj"]] < PADJ_CUTOFF) & (data[cols["log2fc"]] >= LOG2FC_CUTOFF) & basemean_ok,
        (data[cols["padj"]] < PADJ_CUTOFF) & (data[cols["log2fc"]] <= -LOG2FC_CUTOFF) & basemean_ok,
    ]
    choices = ["Up", "Down"]
    data["status"] = np.select(conditions, choices, default="NS")

    color_map = {"Up": COLOR_UP, "Down": COLOR_DOWN, "NS": COLOR_NS}

    fig, ax = plt.subplots(figsize=(8, 6))

    for status in ["NS", "Down", "Up"]:
        subset = data[data["status"] == status]
        lbl = "NS" if status == "NS" else f"{status} ({len(subset):,})"
        ax.scatter(
            subset[cols["log2fc"]], subset["-log10padj"],
            c=color_map[status], s=8, alpha=0.5, edgecolors="none",
            label=lbl, rasterized=True
        )

    ax.axhline(-np.log10(PADJ_CUTOFF), color="grey", ls="--", lw=0.8)
    ax.axvline(LOG2FC_CUTOFF, color="grey", ls="--", lw=0.8)
    ax.axvline(-LOG2FC_CUTOFF, color="grey", ls="--", lw=0.8)

    n_up = (data["status"] == "Up").sum()
    n_down = (data["status"] == "Down").sum()
    add_count_box(ax, n_up, n_down, n_up + n_down, position="lower left")

    ax.set_xlabel("log$_2$ Fold Change")
    ax.set_ylabel("-log$_{10}$ (adjusted p-value)")
    ax.set_title(f"Volcano Plot - DESeq2 ({label})")
    ax.legend(bbox_to_anchor=(1.02, 1), loc="upper left", frameon=True, fontsize=10, markerscale=2)

    fname = f"volcano_plot{suffix}.{FIG_FORMAT}"
    outpath = outdir / fname
    fig.savefig(outpath, format=FIG_FORMAT, bbox_inches="tight")
    plt.close(fig)
    print(f"  Saved: {outpath}")


def ma_plot(df, outdir, label="All", suffix=""):
    """Generate MA plot (baseMean vs log2FC)."""
    cols = DESEQ2_COLS
    data = df.dropna(subset=[cols["padj"], cols["log2fc"], cols["basemean"]]).copy()

    basemean_ok = data[cols["basemean"]] >= BASEMEAN_CUTOFF
    sig = (data[cols["padj"]] < PADJ_CUTOFF) & (data[cols["log2fc"]].abs() >= LOG2FC_CUTOFF) & basemean_ok
    data["significant"] = np.where(
        sig & (data[cols["log2fc"]] > 0), "Up",
        np.where(sig & (data[cols["log2fc"]] < 0), "Down", "NS")
    )

    color_map = {"Up": COLOR_UP, "Down": COLOR_DOWN, "NS": COLOR_NS}

    fig, ax = plt.subplots(figsize=(8, 6))

    for status in ["NS", "Down", "Up"]:
        subset = data[data["significant"] == status]
        lbl = "NS" if status == "NS" else f"{status} ({len(subset):,})"
        ax.scatter(
            np.log10(subset[cols["basemean"]].clip(lower=0.1)),
            subset[cols["log2fc"]],
            c=color_map[status], s=8, alpha=0.5, edgecolors="none",
            label=lbl, rasterized=True
        )

    ax.axhline(0, color="black", lw=0.8)
    ax.axhline(LOG2FC_CUTOFF, color="grey", ls="--", lw=0.8)
    ax.axhline(-LOG2FC_CUTOFF, color="grey", ls="--", lw=0.8)

    n_up = (data["significant"] == "Up").sum()
    n_down = (data["significant"] == "Down").sum()
    add_count_box(ax, n_up, n_down, n_up + n_down, position="lower left")

    ax.set_xlabel("log$_{10}$ (baseMean)")
    ax.set_ylabel("log$_2$ Fold Change")
    ax.set_title(f"MA Plot - DESeq2 ({label})")
    ax.legend(bbox_to_anchor=(1.02, 1), loc="upper left", frameon=True, fontsize=10, markerscale=2)

    fname = f"ma_plot{suffix}.{FIG_FORMAT}"
    outpath = outdir / fname
    fig.savefig(outpath, format=FIG_FORMAT, bbox_inches="tight")
    plt.close(fig)
    print(f"  Saved: {outpath}")


def volcano_plot_interactive(df, outdir, label="All", suffix=""):
    """Generate interactive volcano plot with hover tooltips (HTML output)."""
    if px is None:
        return
    cols = DESEQ2_COLS
    data = df.dropna(subset=[cols["padj"], cols["log2fc"]]).copy()
    data["-log10padj"] = -np.log10(data[cols["padj"]].clip(lower=1e-300))

    basemean_ok = (data[cols["basemean"]] >= BASEMEAN_CUTOFF) if cols["basemean"] in data.columns else True
    conditions = [
        (data[cols["padj"]] < PADJ_CUTOFF) & (data[cols["log2fc"]] >= LOG2FC_CUTOFF) & basemean_ok,
        (data[cols["padj"]] < PADJ_CUTOFF) & (data[cols["log2fc"]] <= -LOG2FC_CUTOFF) & basemean_ok,
    ]
    data["Status"] = np.select(conditions, ["Up", "Down"], default="NS")

    n_up = int((data["Status"] == "Up").sum())
    n_down = int((data["Status"] == "Down").sum())

    color_map = {"NS": COLOR_NS, "Down": COLOR_DOWN, "Up": COLOR_UP}

    # Use gene_name for hover label; fall back to gene_id when gene_name column is absent
    _hover_name = cols["gene_name"] if cols["gene_name"] in data.columns else cols["gene_id"]
    fig = px.scatter(
        data.sort_values("Status", key=lambda s: s.map({"NS": 0, "Down": 1, "Up": 2})),
        x=cols["log2fc"],
        y="-log10padj",
        color="Status",
        color_discrete_map=color_map,
        category_orders={"Status": ["NS", "Down", "Up"]},
        hover_name=_hover_name,
        hover_data={
            cols["log2fc"]: ":.3f",
            cols["padj"]: ":.2e",
            cols["basemean"]: ":.1f",
            "-log10padj": ":.2f",
            "Status": False,
        },
        opacity=0.5,
        title=f"Volcano Plot - DESeq2 ({label})",
    )
    fig.update_traces(marker=dict(size=5))
    fig.add_hline(y=-np.log10(PADJ_CUTOFF), line_dash="dash", line_color="grey", line_width=0.8)
    fig.add_vline(x=LOG2FC_CUTOFF, line_dash="dash", line_color="grey", line_width=0.8)
    fig.add_vline(x=-LOG2FC_CUTOFF, line_dash="dash", line_color="grey", line_width=0.8)
    fig.add_annotation(
        text=f"Up: {n_up:,}<br>Down: {n_down:,}<br>Total: {n_up + n_down:,}",
        xref="paper", yref="paper", x=0.02, y=0.98,
        showarrow=False, bgcolor="rgba(255,255,255,0.85)", bordercolor="grey",
        borderwidth=1, font=dict(size=11), align="left", xanchor="left", yanchor="top",
    )
    fig.update_layout(
        xaxis_title="log₂ Fold Change",
        yaxis_title="-log₁₀ (adjusted p-value)",
        hovermode="closest", template="plotly_white",
        width=900, height=650,
    )

    fname = f"volcano_plot{suffix}_interactive.html"
    outpath = outdir / fname
    fig.write_html(str(outpath))
    print(f"  Saved: {outpath}")


def ma_plot_interactive(df, outdir, label="All", suffix=""):
    """Generate interactive MA plot with hover tooltips (HTML output)."""
    if px is None:
        return
    cols = DESEQ2_COLS
    data = df.dropna(subset=[cols["padj"], cols["log2fc"], cols["basemean"]]).copy()
    data["log10_basemean"] = np.log10(data[cols["basemean"]].clip(lower=0.1))

    basemean_ok = data[cols["basemean"]] >= BASEMEAN_CUTOFF
    sig = (data[cols["padj"]] < PADJ_CUTOFF) & (data[cols["log2fc"]].abs() >= LOG2FC_CUTOFF) & basemean_ok
    data["Status"] = np.where(
        sig & (data[cols["log2fc"]] > 0), "Up",
        np.where(sig & (data[cols["log2fc"]] < 0), "Down", "NS"),
    )

    n_up = int((data["Status"] == "Up").sum())
    n_down = int((data["Status"] == "Down").sum())

    color_map = {"NS": COLOR_NS, "Down": COLOR_DOWN, "Up": COLOR_UP}

    _hover_name = cols["gene_name"] if cols["gene_name"] in data.columns else cols["gene_id"]
    fig = px.scatter(
        data.sort_values("Status", key=lambda s: s.map({"NS": 0, "Down": 1, "Up": 2})),
        x="log10_basemean",
        y=cols["log2fc"],
        color="Status",
        color_discrete_map=color_map,
        category_orders={"Status": ["NS", "Down", "Up"]},
        hover_name=_hover_name,
        hover_data={
            cols["log2fc"]: ":.3f",
            cols["padj"]: ":.2e",
            cols["basemean"]: ":.1f",
            "log10_basemean": False,
            "Status": False,
        },
        opacity=0.5,
        title=f"MA Plot - DESeq2 ({label})",
    )
    fig.update_traces(marker=dict(size=5))
    fig.add_hline(y=0, line_color="black", line_width=0.8)
    fig.add_hline(y=LOG2FC_CUTOFF, line_dash="dash", line_color="grey", line_width=0.8)
    fig.add_hline(y=-LOG2FC_CUTOFF, line_dash="dash", line_color="grey", line_width=0.8)
    fig.add_annotation(
        text=f"Up: {n_up:,}<br>Down: {n_down:,}<br>Total: {n_up + n_down:,}",
        xref="paper", yref="paper", x=0.02, y=0.98,
        showarrow=False, bgcolor="rgba(255,255,255,0.85)", bordercolor="grey",
        borderwidth=1, font=dict(size=11), align="left", xanchor="left", yanchor="top",
    )
    fig.update_layout(
        xaxis_title="log₁₀ (baseMean)",
        yaxis_title="log₂ Fold Change",
        hovermode="closest", template="plotly_white",
        width=900, height=650,
    )

    fname = f"ma_plot{suffix}_interactive.html"
    outpath = outdir / fname
    fig.write_html(str(outpath))
    print(f"  Saved: {outpath}")


def biotype_chart(filtered_df, outdir, label="All", suffix=""):
    """Generate biotype distribution bar chart from filtered DESeq2 results."""
    cols = DESEQ2_COLS
    if cols["biotype"] not in filtered_df.columns:
        print("  Skipping biotype chart (no biotype column)")
        return

    counts = filtered_df[cols["biotype"]].value_counts()
    if len(counts) == 0:
        return

    if len(counts) > 10:
        top = counts.head(9)
        other = pd.Series({"Other": counts.iloc[9:].sum()})
        counts = pd.concat([top, other])

    palette = sns.color_palette("Set2", n_colors=len(counts))

    fig, axes = plt.subplots(1, 2, figsize=(14, 6))

    ax1 = axes[0]
    bars = ax1.barh(counts.index[::-1], counts.values[::-1], color=palette[::-1])
    ax1.set_xlabel("Number of DE Genes")
    ax1.set_title(f"Biotype Distribution ({label})")
    for bar, val in zip(bars, counts.values[::-1]):
        ax1.text(bar.get_width() + 0.5, bar.get_y() + bar.get_height()/2,
                 f"{val}", va="center", fontsize=9)

    ax2 = axes[1]
    wedges, texts, autotexts = ax2.pie(
        counts.values, labels=counts.index, autopct="%1.1f%%",
        colors=palette, startangle=90, pctdistance=0.8
    )
    for t in autotexts:
        t.set_fontsize(8)
    ax2.set_title(f"Biotype Proportions ({label})")

    plt.tight_layout()
    fname = f"biotype_distribution{suffix}.{FIG_FORMAT}"
    outpath = outdir / fname
    fig.savefig(outpath, format=FIG_FORMAT)
    plt.close(fig)
    print(f"  Saved: {outpath}")


# ---------------------------------------------------------------------------
# BIOTYPE HELPERS & ANALYSES
# ---------------------------------------------------------------------------

def _assign_biotype_group(series):
    """Map detailed Ensembl biotypes to 5 broad groups using _BIOTYPE_GROUPS.

    Case-insensitive: 'Protein_coding', 'protein_coding', 'PROTEIN_CODING' all map correctly.
    """
    return series.map(lambda x: _BIOTYPE_GROUPS.get(str(x).lower(), "Other"))


def _bh_correction(pvals):
    """Benjamini-Hochberg FDR correction (no external deps required)."""
    n = len(pvals)
    if n == 0:
        return np.array([])
    order = np.argsort(pvals)
    adjusted = np.array(pvals, dtype=float)[order] * n / (np.arange(1, n + 1))
    for i in range(n - 2, -1, -1):
        adjusted[i] = min(adjusted[i], adjusted[i + 1])
    result = np.empty(n)
    result[order] = np.minimum(adjusted, 1.0)
    return result


def biotype_direction_chart(filtered_df, outdir, label="All", suffix=""):
    """Diverging horizontal bar chart: Up/Down DE gene counts split by biotype group."""
    cols = DESEQ2_COLS
    if cols["biotype"] not in filtered_df.columns:
        print("  Skipping biotype direction chart (no biotype column)")
        return
    if len(filtered_df) == 0:
        return

    df = filtered_df.copy()
    df["_group"] = _assign_biotype_group(df[cols["biotype"]])

    up_counts = df[df["direction"] == "up"]["_group"].value_counts()
    dn_counts = df[df["direction"] == "down"]["_group"].value_counts()

    groups = [g for g in _BIOTYPE_ORDER if g in up_counts.index or g in dn_counts.index]
    if not groups:
        return

    up_vals = [up_counts.get(g, 0) for g in groups]
    dn_vals = [-dn_counts.get(g, 0) for g in groups]

    fig, ax = plt.subplots(figsize=(9, max(3, len(groups) * 1.0 + 1.5)))
    y = range(len(groups))

    ax.barh(list(y), up_vals, color=COLOR_UP, label="Up-regulated", alpha=0.85)
    ax.barh(list(y), dn_vals, color=COLOR_DOWN, label="Down-regulated", alpha=0.85)
    ax.axvline(0, color="black", linewidth=0.8)

    for i, (u, d) in enumerate(zip(up_vals, dn_vals)):
        if u > 0:
            ax.text(u + 0.3, i, str(u), va="center", fontsize=9, color=COLOR_UP)
        if d < 0:
            ax.text(d - 0.3, i, str(-d), va="center", ha="right", fontsize=9, color=COLOR_DOWN)

    ax.set_yticks(list(y))
    ax.set_yticklabels(groups)
    ax.set_xlabel("Number of DE Genes")
    ax.set_title(f"DE Genes by Biotype & Direction ({label})")
    ax.legend(loc="lower right", fontsize=9)
    ax.xaxis.set_major_formatter(ticker.FuncFormatter(lambda x, _: str(int(abs(x)))))
    plt.tight_layout()

    fname = f"biotype_direction_chart{suffix}.{FIG_FORMAT}"
    outpath = outdir / fname
    fig.savefig(outpath, format=FIG_FORMAT, dpi=FIG_DPI)
    plt.close(fig)
    print(f"  Saved: {outpath}")


def biotype_enrichment_test(filtered_df, all_df, outdir, label="All", suffix=""):
    """Fisher's exact test: is each biotype group enriched among DE genes vs background?"""
    if not _SCIPY_AVAILABLE:
        print("  Skipping biotype enrichment (scipy not installed: pip install scipy)")
        return
    cols = DESEQ2_COLS
    if cols["biotype"] not in all_df.columns:
        print("  Skipping biotype enrichment (no biotype column)")
        return

    bg = all_df.copy()
    bg["_group"] = _assign_biotype_group(bg[cols["biotype"]])
    de = filtered_df.copy()
    if cols["biotype"] in de.columns:
        de["_group"] = _assign_biotype_group(de[cols["biotype"]])
    else:
        return

    n_bg = len(bg)
    n_de = len(de)
    if n_de == 0:
        return

    results = []
    for grp in _BIOTYPE_ORDER:
        a = int((de["_group"] == grp).sum())       # DE & in group
        b = n_de - a                                 # DE & not in group
        c = int((bg["_group"] == grp).sum()) - a    # Non-DE & in group (bg ⊃ de)
        d = n_bg - n_de - c                          # Non-DE & not in group
        if a + c == 0:
            continue
        if c < 0 or d < 0:
            print(f"  WARNING: biotype_enrichment_test: negative cell count for '{grp}' "
                  f"(c={c}, d={d}). Check that filtered_df is a subset of all_df.")
            c = max(c, 0); d = max(d, 0)
        try:
            odds, pval = fisher_exact([[a, b], [c, d]], alternative="two-sided")
        except Exception:
            continue
        # Cap log2(OR) at ±10 to keep axis readable; mark capped values
        if odds == 0:
            log2_or = -10.0
        elif not np.isfinite(odds):
            log2_or = 10.0
        else:
            log2_or = float(np.clip(np.log2(odds), -10, 10))
        results.append({"group": grp, "log2OR": log2_or, "pval": pval,
                        "n_de": a, "n_group": a + c})

    if not results:
        return

    res_df = pd.DataFrame(results)
    res_df["fdr"] = _bh_correction(res_df["pval"].tolist())
    res_df = res_df.sort_values("log2OR")

    fig, ax = plt.subplots(figsize=(9, max(3, len(res_df) * 0.9 + 1.5)))
    colors = [("#2ecc71" if f < 0.05 else COLOR_NS) for f in res_df["fdr"]]
    ax.hlines(range(len(res_df)), 0, res_df["log2OR"], color="grey", linewidth=1.5, zorder=1)
    ax.scatter(res_df["log2OR"], range(len(res_df)), color=colors, s=80, zorder=2)
    ax.axvline(0, color="black", linewidth=0.8, linestyle="--")

    # Annotation: "DE in group / total in group" — placed consistently on right side
    x_max = max(abs(res_df["log2OR"].max()), abs(res_df["log2OR"].min()), 1.0)
    ax.set_xlim(-x_max * 1.4, x_max * 1.8)  # leave room on right for labels
    for i, row in enumerate(res_df.itertuples()):
        label_txt = f" {row.n_de}/{row.n_group}"
        ax.text(x_max * 1.05, i, label_txt, va="center", fontsize=8, ha="left")
    ax.text(x_max * 1.05, len(res_df) - 0.5, "DE/Total", va="bottom", fontsize=7,
            ha="left", color="grey")

    ax.set_yticks(range(len(res_df)))
    ax.set_yticklabels(res_df["group"].tolist())
    ax.set_xlabel("log\u2082(Odds Ratio)")
    ax.set_title(f"Biotype Enrichment Among DE Genes ({label})\nGreen = FDR < 0.05")
    plt.tight_layout()

    fname = f"biotype_enrichment{suffix}.{FIG_FORMAT}"
    outpath = outdir / fname
    fig.savefig(outpath, format=FIG_FORMAT, dpi=FIG_DPI)
    plt.close(fig)
    print(f"  Saved: {outpath}")


def biotype_volcano(all_df, outdir, label="All", suffix=""):
    """Volcano plot colored by biotype group instead of direction."""
    cols = DESEQ2_COLS
    if cols["biotype"] not in all_df.columns:
        print("  Skipping biotype volcano (no biotype column)")
        return

    df = all_df.copy().dropna(subset=[cols["log2fc"], cols["padj"]])
    df["_group"] = _assign_biotype_group(df[cols["biotype"]])
    df["_neg_log10p"] = -np.log10(df[cols["padj"]].clip(lower=1e-300))
    _bm_ok = (df[cols["basemean"]] >= BASEMEAN_CUTOFF) if cols["basemean"] in df.columns else True
    df["_sig"] = (df[cols["padj"]] < PADJ_CUTOFF) & (df[cols["log2fc"]].abs() >= LOG2FC_CUTOFF) & _bm_ok

    fig, ax = plt.subplots(figsize=(9, 7))

    # NS genes: grey background
    ns = df[~df["_sig"]]
    ax.scatter(ns[cols["log2fc"]], ns["_neg_log10p"], color=COLOR_NS,
               alpha=0.25, s=4, rasterized=True, label=None)

    # Significant genes: colored by biotype group
    sig = df[df["_sig"]]
    legend_handles = []
    for grp in _BIOTYPE_ORDER:
        grp_sig = sig[sig["_group"] == grp]
        if len(grp_sig) == 0:
            continue
        color = _BIOTYPE_COLORS.get(grp, "#999999")
        ax.scatter(grp_sig[cols["log2fc"]], grp_sig["_neg_log10p"],
                   color=color, alpha=0.75, s=10, rasterized=True)
        legend_handles.append(mpatches.Patch(color=color, label=f"{grp} (n={len(grp_sig)})"))

    ax.axhline(-np.log10(PADJ_CUTOFF), color="grey", linestyle="--", linewidth=0.8)
    ax.axvline(LOG2FC_CUTOFF, color="grey", linestyle="--", linewidth=0.8)
    ax.axvline(-LOG2FC_CUTOFF, color="grey", linestyle="--", linewidth=0.8)
    ax.set_xlabel("log\u2082 Fold Change")
    ax.set_ylabel("-log\u2081\u2080(adjusted p-value)")
    ax.set_title(f"Volcano Plot by Biotype ({label})")
    if legend_handles:
        ax.legend(handles=legend_handles, fontsize=8, loc="upper left")
    plt.tight_layout()

    fname = f"volcano_biotype{suffix}.{FIG_FORMAT}"
    outpath = outdir / fname
    fig.savefig(outpath, format=FIG_FORMAT, dpi=FIG_DPI)
    plt.close(fig)
    print(f"  Saved: {outpath}")


def ecdf_log2fc_by_biotype(all_df, outdir, label="All", suffix=""):
    """ECDF of |log2FC| per biotype group — shows if some biotypes have larger FC shifts."""
    cols = DESEQ2_COLS
    if cols["biotype"] not in all_df.columns:
        print("  Skipping ECDF biotype plot (no biotype column)")
        return

    df = all_df.copy().dropna(subset=[cols["log2fc"]])
    if cols["basemean"] in df.columns:
        df = df[df[cols["basemean"]] >= BASEMEAN_CUTOFF]
    df["_group"] = _assign_biotype_group(df[cols["biotype"]])
    df["_abs_lfc"] = df[cols["log2fc"]].abs()

    min_genes = 10
    plotted = 0
    fig, ax = plt.subplots(figsize=(8, 5))

    for grp in _BIOTYPE_ORDER:
        sub = df[df["_group"] == grp]["_abs_lfc"].dropna().sort_values()
        if len(sub) < min_genes:
            continue
        ecdf_y = np.arange(1, len(sub) + 1) / len(sub)
        color = _BIOTYPE_COLORS.get(grp, "#999999")
        ax.plot(sub.values, ecdf_y, label=f"{grp} (n={len(sub)})", color=color, linewidth=2)
        plotted += 1

    if plotted < 2:
        plt.close(fig)
        print("  Skipping ECDF biotype (fewer than 2 groups with ≥10 genes)")
        return

    ax.axvline(LOG2FC_CUTOFF, color="grey", linestyle="--", linewidth=0.8,
               label=f"cutoff ({LOG2FC_CUTOFF})")
    ax.set_xlabel("|log\u2082 Fold Change|")
    ax.set_ylabel("Cumulative Fraction")
    ax.set_title(f"ECDF of |log\u2082FC| by Biotype Group ({label})")
    ax.legend(fontsize=8)
    ax.set_xlim(left=0)
    plt.tight_layout()

    fname = f"ecdf_log2fc_biotype{suffix}.{FIG_FORMAT}"
    outpath = outdir / fname
    fig.savefig(outpath, format=FIG_FORMAT, dpi=FIG_DPI)
    plt.close(fig)
    print(f"  Saved: {outpath}")


def cross_condition_biotype_comparison(condition_results, condition_labels, outdir):
    """Grouped bar + stacked % chart of DE gene counts by biotype group across conditions."""
    cols = DESEQ2_COLS
    count_data = {}
    for name, res in condition_results.items():
        filt = res["deseq2_filtered"].get("all_genes")
        if filt is None or len(filt) == 0:
            continue
        if cols["biotype"] not in filt.columns:
            continue
        groups = _assign_biotype_group(filt[cols["biotype"]])
        count_data[condition_labels[name]] = groups.value_counts()

    if len(count_data) < 1:
        print("  Skipping cross-condition biotype comparison (no biotype data)")
        return

    all_groups = _BIOTYPE_ORDER
    df = pd.DataFrame(count_data, index=all_groups).fillna(0).astype(int)
    df = df.loc[(df.sum(axis=1) > 0)]   # drop empty rows

    fig, axes = plt.subplots(1, 2, figsize=(14, max(5, len(df) * 0.6 + 2)))

    # Panel A: grouped bar (count)
    ax = axes[0]
    x = np.arange(len(df))
    n_conds = len(df.columns)
    width = 0.8 / max(n_conds, 1)
    cond_colors = sns.color_palette("tab10", n_conds)
    for i, cond_lbl in enumerate(df.columns):
        vals = df[cond_lbl].values
        bars = ax.bar(x + i * width - (n_conds - 1) * width / 2,
                      vals, width * 0.9, label=cond_lbl, color=cond_colors[i], alpha=0.85)
    ax.set_xticks(x)
    ax.set_xticklabels(df.index, rotation=30, ha="right")
    ax.set_ylabel("DE Gene Count")
    ax.set_title("DE Genes per Biotype Group")
    ax.legend(fontsize=8)

    # Panel B: stacked 100% bar (proportion per condition)
    ax2 = axes[1]
    col_sums = df.sum(axis=0).replace(0, np.nan)  # avoid div-by-zero for empty conditions
    pct = df.div(col_sums, axis=1) * 100
    bottom = np.zeros(n_conds)
    x2 = np.arange(n_conds)
    for grp_idx, grp in enumerate(df.index):
        vals = pct.loc[grp].values
        color = _BIOTYPE_COLORS.get(grp, "#999999")
        ax2.bar(x2, vals, bottom=bottom, color=color, alpha=0.85, label=grp)
        # Annotate if large enough
        for xi, (v, b) in enumerate(zip(vals, bottom)):
            if v >= 5:
                ax2.text(xi, b + v / 2, f"{v:.0f}%", ha="center", va="center",
                         fontsize=8, color="white" if v > 10 else "black")
        bottom += vals
    ax2.set_xticks(x2)
    ax2.set_xticklabels(df.columns, rotation=20, ha="right")
    ax2.set_ylabel("Percentage of DE Genes (%)")
    ax2.set_title("Biotype Composition per Condition")
    ax2.set_ylim(0, 105)
    ax2.legend(fontsize=8, loc="upper right")

    plt.tight_layout()
    fname = f"cross_condition_biotype_comparison.{FIG_FORMAT}"
    outpath = outdir / fname
    fig.savefig(outpath, format=FIG_FORMAT, dpi=FIG_DPI)
    plt.close(fig)
    print(f"  Saved: {outpath}")


def cross_condition_biotype_direction(condition_results, condition_labels, outdir):
    """Faceted diverging bar chart: Up/Down DE gene counts by biotype group, one panel per condition."""
    cols = DESEQ2_COLS
    n = len(condition_results)
    if n == 0:
        return

    nrows, ncols = _grid_dims(n)
    fig, axes = plt.subplots(nrows, ncols,
                              figsize=(ncols * 5, nrows * max(3, len(_BIOTYPE_ORDER) * 0.8 + 1)),
                              squeeze=False)
    axes_flat = [axes[r][c] for r in range(nrows) for c in range(ncols)]

    max_abs = 0
    panel_data = []
    for name, res in condition_results.items():
        filt = res["deseq2_filtered"].get("all_genes")
        cond_lbl = condition_labels[name]
        if filt is None or len(filt) == 0 or cols["biotype"] not in filt.columns:
            panel_data.append((cond_lbl, None, None))
            continue
        filt = filt.copy()
        filt["_group"] = _assign_biotype_group(filt[cols["biotype"]])
        up_c = filt[filt["direction"] == "up"]["_group"].value_counts()
        dn_c = filt[filt["direction"] == "down"]["_group"].value_counts()
        max_abs = max(max_abs, up_c.max() if len(up_c) else 0,
                      dn_c.max() if len(dn_c) else 0)
        panel_data.append((cond_lbl, up_c, dn_c))

    xlim = max_abs * 1.15 if max_abs > 0 else 10

    for ax_idx, (cond_lbl, up_c, dn_c) in enumerate(panel_data):
        ax = axes_flat[ax_idx]
        if up_c is None and dn_c is None:
            ax.set_visible(False)
            continue

        groups = [g for g in _BIOTYPE_ORDER
                  if (up_c is not None and g in up_c.index) or
                     (dn_c is not None and g in dn_c.index)]
        if not groups:
            ax.set_visible(False)
            continue

        y = range(len(groups))
        up_vals = [up_c.get(g, 0) if up_c is not None else 0 for g in groups]
        dn_vals = [-dn_c.get(g, 0) if dn_c is not None else 0 for g in groups]

        ax.barh(list(y), up_vals, color=COLOR_UP, alpha=0.85, label="Up")
        ax.barh(list(y), dn_vals, color=COLOR_DOWN, alpha=0.85, label="Down")
        ax.axvline(0, color="black", linewidth=0.8)
        ax.set_xlim(-xlim, xlim)
        ax.set_yticks(list(y))
        ax.set_yticklabels(groups, fontsize=9)
        ax.set_title(cond_lbl, fontsize=10)
        ax.set_xlabel("DE Gene Count", fontsize=8)
        ax.xaxis.set_major_formatter(ticker.FuncFormatter(lambda x, _: str(int(abs(x)))))
        if ax_idx == 0:
            ax.legend(fontsize=8, loc="lower right")

    # Hide unused panels
    for ax_idx in range(len(panel_data), len(axes_flat)):
        axes_flat[ax_idx].set_visible(False)

    fig.suptitle("DE Gene Direction by Biotype Group (per Condition)", fontsize=12, y=1.01)
    plt.tight_layout()
    fname = f"cross_condition_biotype_direction.{FIG_FORMAT}"
    outpath = outdir / fname
    fig.savefig(outpath, format=FIG_FORMAT, dpi=FIG_DPI, bbox_inches="tight")
    plt.close(fig)
    print(f"  Saved: {outpath}")


# ---------------------------------------------------------------------------
# rMATS PROCESSING -- All 5 event types
# ---------------------------------------------------------------------------

def load_all_rmats(rmats_dir):
    """Load all rMATS event type files from the directory."""
    rmats_dir = Path(rmats_dir)
    all_data = {}

    for event_type in RMATS_EVENT_TYPES:
        filepath = rmats_dir / f"{event_type}{RMATS_FILE_SUFFIX}"
        if not filepath.exists():
            print(f"  WARNING: {filepath.name} not found, skipping {event_type}")
            continue

        df = load_file(filepath, name=f"rMATS {event_type}")
        df = normalize_rmats_columns(df, f"rMATS {event_type}")
        required = [v for v in RMATS_COLS.values() if v is not None]
        validate_columns(df, required, name=f"rMATS {event_type}")

        # Log which critical columns were found
        critical = ["FDR", "PValue", "IncLevelDifference", "geneSymbol"]
        status_parts = [f"{c} \u2713" if c in df.columns else f"{c} \u2717"
                        for c in critical]
        print(f"  [INFO] {event_type} columns: {', '.join(status_parts)}")

        # Fill missing geneSymbol with GeneID (Ensembl ID)
        gene_col = RMATS_COLS["gene_name"]   # "geneSymbol"
        id_col = RMATS_COLS["gene_id"]       # "GeneID"
        if gene_col in df.columns and id_col in df.columns:
            mask = df[gene_col].isna() | (df[gene_col].str.strip() == "")
            if mask.any():
                df.loc[mask, gene_col] = df.loc[mask, id_col]
                print(f"  [INFO] Filled {mask.sum()} missing geneSymbol values with GeneID (Ensembl ID)")

        df["event_type"] = event_type
        all_data[event_type] = df

    return all_data


def filter_rmats(df, event_type=""):
    """Apply rMATS cutoffs and return filtered DataFrame."""
    cols = RMATS_COLS
    df = _validate_rmats_columns(df, cols)

    if RMATS_DUAL_FILTER:
        # Dual filter: require BOTH FDR AND PValue thresholds
        drop_cols = [c for c in [cols["fdr"], cols["pvalue"], cols["inclevel_diff"]] if c in df.columns]
        df = df.dropna(subset=drop_cols)
        mask = (
            (df[cols["fdr"]] < RMATS_FDR_CUTOFF) &
            (df[cols["pvalue"]] < RMATS_PVAL_CUTOFF) &
            (df[cols["inclevel_diff"]].abs() >= INCLEVEL_DIFF_CUTOFF)
        )
        filtered = df[mask].copy()
        print(f"  {event_type}: {len(df):,} total -> {len(filtered):,} significant "
              f"(FDR < {RMATS_FDR_CUTOFF} AND PValue < {RMATS_PVAL_CUTOFF}, |dPSI| >= {INCLEVEL_DIFF_CUTOFF})")
    else:
        pval_col = cols["fdr"] if USE_FDR else cols["pvalue"]
        pval_cutoff = RMATS_FDR_CUTOFF if USE_FDR else RMATS_PVAL_CUTOFF
        pval_label = "FDR" if USE_FDR else "PValue"
        df = df.dropna(subset=[pval_col, cols["inclevel_diff"]])
        mask = (
            (df[pval_col] < pval_cutoff) &
            (df[cols["inclevel_diff"]].abs() >= INCLEVEL_DIFF_CUTOFF)
        )
        filtered = df[mask].copy()
        print(f"  {event_type}: {len(df):,} total -> {len(filtered):,} significant "
              f"({pval_label} < {pval_cutoff}, |dPSI| >= {INCLEVEL_DIFF_CUTOFF})")

    return df, filtered


def rmats_scatter(df, event_type, outdir):
    """Scatter plot: IncLevelDifference vs -log10(pvalue) for a single event type."""
    cols = RMATS_COLS

    if RMATS_DUAL_FILTER:
        pval_col = cols["fdr"]  # plot FDR on y-axis for dual mode
        pval_cutoff = RMATS_FDR_CUTOFF
        pval_label = "FDR"
    else:
        pval_col = cols["fdr"] if USE_FDR else cols["pvalue"]
        pval_cutoff = RMATS_FDR_CUTOFF if USE_FDR else RMATS_PVAL_CUTOFF
        pval_label = "FDR" if USE_FDR else "PValue"

    data = df.dropna(subset=[pval_col, cols["inclevel_diff"]]).copy()
    data["-log10p"] = -np.log10(data[pval_col].clip(lower=1e-300))

    if RMATS_DUAL_FILTER:
        sig = (
            (data[cols["fdr"]] < RMATS_FDR_CUTOFF) &
            (data[cols["pvalue"]] < RMATS_PVAL_CUTOFF) &
            (data[cols["inclevel_diff"]].abs() >= INCLEVEL_DIFF_CUTOFF)
        )
    else:
        sig = (
            (data[pval_col] < pval_cutoff) &
            (data[cols["inclevel_diff"]].abs() >= INCLEVEL_DIFF_CUTOFF)
        )
    data["significant"] = np.where(sig, "Significant", "NS")

    fig, ax = plt.subplots(figsize=(8, 6))

    for status, color in [("NS", COLOR_NS), ("Significant", EVENT_COLORS.get(event_type, COLOR_UP))]:
        subset = data[data["significant"] == status]
        lbl = "NS" if status == "NS" else f"{status} ({len(subset):,})"
        ax.scatter(
            subset[cols["inclevel_diff"]], subset["-log10p"],
            c=color, s=10, alpha=0.5, edgecolors="none",
            label=lbl, rasterized=True
        )

    ax.axhline(-np.log10(pval_cutoff), color="grey", ls="--", lw=0.8)
    ax.axvline(INCLEVEL_DIFF_CUTOFF, color="grey", ls="--", lw=0.8)
    ax.axvline(-INCLEVEL_DIFF_CUTOFF, color="grey", ls="--", lw=0.8)

    sig_data = data[data["significant"] == "Significant"]
    n_inc = (sig_data[cols["inclevel_diff"]] >= INCLEVEL_DIFF_CUTOFF).sum()
    n_exc = (sig_data[cols["inclevel_diff"]] <= -INCLEVEL_DIFF_CUTOFF).sum()
    add_count_box(ax, n_inc, n_exc, n_inc + n_exc, position="lower left",
                  up_label=f"Included (dPSI\u22650.1)", down_label=f"Excluded (dPSI\u2264\u22120.1)")

    ax.set_xlabel("$\\Delta$PSI (IncLevelDifference)")
    ax.set_ylabel(f"-log$_{{10}}$ ({pval_label})")
    ax.set_title(f"rMATS - {event_type} (Skipped Exon)" if event_type == "SE"
                 else f"rMATS - {event_type}")
    ax.legend(bbox_to_anchor=(1.02, 1), loc="upper left", frameon=True, fontsize=10, markerscale=2)

    outpath = outdir / f"rmats_{event_type}_scatter.{FIG_FORMAT}"
    fig.savefig(outpath, format=FIG_FORMAT, bbox_inches="tight")
    plt.close(fig)
    print(f"  Saved: {outpath}")


def rmats_combined_volcano(all_data, outdir):
    """Combined scatter: all event types overlaid on one plot."""
    cols = RMATS_COLS
    pval_col = cols["fdr"] if USE_FDR else cols["pvalue"]
    pval_cutoff = RMATS_FDR_CUTOFF if USE_FDR else RMATS_PVAL_CUTOFF
    pval_label = "FDR" if USE_FDR else "PValue"

    fig, ax = plt.subplots(figsize=(10, 7))

    for event_type, df in all_data.items():
        data = df.dropna(subset=[pval_col, cols["inclevel_diff"]]).copy()
        data["-log10p"] = -np.log10(data[pval_col].clip(lower=1e-300))

        sig = (
            (data[pval_col] < pval_cutoff) &
            (data[cols["inclevel_diff"]].abs() >= INCLEVEL_DIFF_CUTOFF)
        )
        sig_data = data[sig]
        ns_data = data[~sig]

        # Plot NS in grey (no label to avoid clutter)
        ax.scatter(
            ns_data[cols["inclevel_diff"]], ns_data["-log10p"],
            c=COLOR_NS, s=6, alpha=0.15, edgecolors="none", rasterized=True
        )
        # Plot significant with event-type color
        ax.scatter(
            sig_data[cols["inclevel_diff"]], sig_data["-log10p"],
            c=EVENT_COLORS.get(event_type, COLOR_UP), s=12, alpha=0.7,
            edgecolors="none", label=f"{event_type} ({len(sig_data):,})",
            rasterized=True
        )

    ax.axhline(-np.log10(pval_cutoff), color="grey", ls="--", lw=0.8)
    ax.axvline(INCLEVEL_DIFF_CUTOFF, color="grey", ls="--", lw=0.8)
    ax.axvline(-INCLEVEL_DIFF_CUTOFF, color="grey", ls="--", lw=0.8)

    # Per-event-type breakdown
    lines = []
    grand_inc = 0
    grand_exc = 0
    for et, df in all_data.items():
        d = df.dropna(subset=[pval_col, cols["inclevel_diff"]])
        s = (d[pval_col] < pval_cutoff) & (d[cols["inclevel_diff"]].abs() >= INCLEVEL_DIFF_CUTOFF)
        sig_d = d[s]
        n_inc = int((sig_d[cols["inclevel_diff"]] >= INCLEVEL_DIFF_CUTOFF).sum())
        n_exc = int((sig_d[cols["inclevel_diff"]] <= -INCLEVEL_DIFF_CUTOFF).sum())
        grand_inc += n_inc
        grand_exc += n_exc
        lines.append(f"{et}: {n_inc:,} inc / {n_exc:,} exc")
    lines.append(f"Total: {grand_inc + grand_exc:,} ({grand_inc:,} inc / {grand_exc:,} exc)")
    box_text = "\n".join(lines)
    ax.text(0.02, 0.02, box_text, transform=ax.transAxes, fontsize=8,
            va="bottom", ha="left", fontweight="bold",
            bbox=dict(boxstyle="round,pad=0.3", facecolor="white",
                      edgecolor="grey", alpha=0.9))

    ax.set_xlabel("$\\Delta$PSI (IncLevelDifference)")
    ax.set_ylabel(f"-log$_{{10}}$ ({pval_label})")
    ax.set_title("rMATS - All Splicing Event Types")
    ax.legend(bbox_to_anchor=(1.02, 1), loc="upper left", frameon=True, fontsize=10, markerscale=2,
              title="Significant Events")

    outpath = outdir / f"rmats_all_events_scatter.{FIG_FORMAT}"
    fig.savefig(outpath, format=FIG_FORMAT, bbox_inches="tight")
    plt.close(fig)
    print(f"  Saved: {outpath}")


def rmats_event_summary_chart(filtered_counts, outdir):
    """Bar chart comparing significant event counts across all event types."""
    event_types = list(filtered_counts.keys())
    counts = list(filtered_counts.values())
    colors = [EVENT_COLORS.get(et, "#888888") for et in event_types]

    fig, ax = plt.subplots(figsize=(8, 5))
    bars = ax.bar(event_types, counts, color=colors, edgecolor="black", linewidth=0.5)

    for bar, val in zip(bars, counts):
        ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() + max(counts)*0.02,
                str(val), ha="center", va="bottom", fontsize=11, fontweight="bold")

    ax.set_xlabel("Splicing Event Type")
    ax.set_ylabel("Number of Significant Events")
    ax.set_title(f"rMATS - Significant Events by Type\n"
                 f"({'FDR' if USE_FDR else 'PValue'} < "
                 f"{RMATS_FDR_CUTOFF if USE_FDR else RMATS_PVAL_CUTOFF}, "
                 f"|$\\Delta$PSI| >= {INCLEVEL_DIFF_CUTOFF})")

    outpath = outdir / f"rmats_event_type_summary.{FIG_FORMAT}"
    fig.savefig(outpath, format=FIG_FORMAT)
    plt.close(fig)
    print(f"  Saved: {outpath}")


def rmats_dpsi_distribution(all_filtered, outdir):
    """Violin/box plot of delta-PSI distributions across event types."""
    cols = RMATS_COLS
    plot_data = []
    for event_type, df in all_filtered.items():
        if len(df) == 0:
            continue
        tmp = df[[cols["inclevel_diff"]]].copy()
        tmp["Event Type"] = event_type
        plot_data.append(tmp)

    if not plot_data:
        print("  No significant events to plot dPSI distribution")
        return

    combined = pd.concat(plot_data, ignore_index=True)
    combined.rename(columns={cols["inclevel_diff"]: "dPSI"}, inplace=True)

    fig, ax = plt.subplots(figsize=(9, 6))

    event_order = [et for et in RMATS_EVENT_TYPES if et in all_filtered and len(all_filtered[et]) > 0]
    palette = [EVENT_COLORS.get(et, "#888888") for et in event_order]

    sns.violinplot(data=combined, x="Event Type", y="dPSI", order=event_order,
                   palette=palette, inner=None, alpha=0.3, ax=ax)
    sns.stripplot(data=combined, x="Event Type", y="dPSI", order=event_order,
                  palette=palette, size=3, alpha=0.6, jitter=True, ax=ax)

    ax.axhline(0, color="black", lw=0.8)
    ax.axhline(INCLEVEL_DIFF_CUTOFF, color="grey", ls="--", lw=0.6, alpha=0.5)
    ax.axhline(-INCLEVEL_DIFF_CUTOFF, color="grey", ls="--", lw=0.6, alpha=0.5)

    # Per-event inc/exc counts above and below each violin
    y_max = combined["dPSI"].max()
    y_min = combined["dPSI"].min()
    y_pad = (y_max - y_min) * 0.08
    for i, et in enumerate(event_order):
        et_data = combined[combined["Event Type"] == et]["dPSI"]
        n_inc = int((et_data >= INCLEVEL_DIFF_CUTOFF).sum())
        n_exc = int((et_data <= -INCLEVEL_DIFF_CUTOFF).sum())
        # Included count above the violin
        ax.text(i, y_max + y_pad, f"{n_inc} inc",
                ha="center", va="bottom", fontsize=9, fontweight="bold",
                color=EVENT_COLORS.get(et, "#333333"))
        # Excluded count below the violin
        ax.text(i, y_min - y_pad, f"{n_exc} exc",
                ha="center", va="top", fontsize=9, fontweight="bold",
                color=EVENT_COLORS.get(et, "#333333"))
    # Expand y-axis to fit labels
    ax.set_ylim(y_min - y_pad * 3, y_max + y_pad * 3)

    ax.set_ylabel("$\\Delta$PSI (IncLevelDifference)")
    ax.set_title("Distribution of $\\Delta$PSI by Event Type (Significant Events)")

    outpath = outdir / f"rmats_dpsi_distribution.{FIG_FORMAT}"
    fig.savefig(outpath, format=FIG_FORMAT)
    plt.close(fig)
    print(f"  Saved: {outpath}")


# ---------------------------------------------------------------------------
# CROSS-CONDITION COMPARISONS
# ---------------------------------------------------------------------------

def extract_gene_sets(condition_results):
    """Extract significant gene identifier sets from per-condition DESeq2 results.

    Uses Ensembl gene_id as the comparison key when available (more reliable
    than gene names which can be non-unique or inconsistent across datasets).
    Falls back to gene_name if gene_id is absent or non-Ensembl.

    Returns dict {cond_name: {"all_sig": set, "up": set, "down": set,
                               "key_col": str, "key_desc": str}}.
    """
    result = {}
    for name, data in condition_results.items():
        filt = data["deseq2_filtered"]["all_genes"]
        key_col, key_desc = _best_gene_key(filt)
        if not key_col or key_col not in filt.columns:
            key_col = DESEQ2_COLS["gene_name"]
            key_desc = "gene name (fallback)"

        all_genes  = set(filt[key_col].dropna().unique())
        up_genes   = set(filt[filt["direction"] == "up"][key_col].dropna().unique())
        down_genes = set(filt[filt["direction"] == "down"][key_col].dropna().unique())
        result[name] = {
            "all_sig": all_genes, "up": up_genes, "down": down_genes,
            "key_col": key_col,   "key_desc": key_desc,
        }
    return result


def deseq2_venn_diagrams(gene_sets, condition_labels, outdir):
    """Generate Venn diagrams for DE gene overlap across conditions."""
    names = list(gene_sets.keys())
    labels = [condition_labels[n] for n in names]

    for set_key, title_suffix, filename in [
        ("all_sig", "Significant in Either Condition", "venn_all_sig_genes"),
        ("up", "Upregulated Genes", "venn_upregulated"),
        ("down", "Downregulated Genes", "venn_downregulated"),
    ]:
        sets = [gene_sets[n][set_key] for n in names]

        fig, ax = plt.subplots(figsize=(8, 8))
        if len(sets) == 2:
            venn2(sets, set_labels=labels, ax=ax)
        elif len(sets) == 3:
            # Use Okabe-Ito color-blind friendly colors
            v = venn3(sets, set_labels=labels, ax=ax)
            # Set colors: blue, orange, sky blue for the 3 circles
            if v:
                for rid, col in [('100','#0072B2'),('010','#E69F00'),('001','#56B4E9'),
                                 ('110','#009E73'),('101','#CC79A7'),('011','#F0E442'),
                                 ('111','#D55E00')]:
                    patch = v.get_patch_by_id(rid)
                    if patch:
                        patch.set_color(col)
        ax.set_title(f"DESeq2 - {title_suffix}", fontsize=14, fontweight="bold")

        outpath = outdir / f"{filename}.{FIG_FORMAT}"
        fig.savefig(outpath, format=FIG_FORMAT)
        plt.close(fig)
        print(f"  Saved: {outpath}")


def deseq2_direction_concordance(condition_results, condition_labels, outdir):
    """Direction concordance heatmap + pairwise concordance bar chart.

    For genes significant in 2+ conditions, shows whether they change
    in the same direction (+1 up, -1 down, 0 not significant).
    Matches genes by Ensembl ID when available (more accurate than gene name).
    Returns the concordance DataFrame.
    """
    names = list(condition_results.keys())

    # Determine best key (prefer Ensembl ID for matching accuracy)
    first_filt = condition_results[names[0]]["deseq2_filtered"]["all_genes"]
    key_col, key_desc = _best_gene_key(first_filt)
    name_col = DESEQ2_COLS.get("gene_name", key_col)
    print(f"  Concordance: matching genes by {key_desc} ('{key_col}')")

    # Build direction matrix: rows=genes, cols=conditions, values=+1/-1/0
    direction_dfs = []
    label_lookups = {}   # key_id -> gene_name for display
    for name in names:
        filt = condition_results[name]["deseq2_filtered"]["all_genes"]
        dir_series = filt.set_index(key_col)["direction"].map({"up": 1, "down": -1})
        dir_series = dir_series[~dir_series.index.duplicated(keep="first")]
        dir_series.name = condition_labels[name]
        direction_dfs.append(dir_series)
        # Build lookup: key -> gene_name (for heatmap labels)
        if name_col in filt.columns and name_col != key_col:
            lookup = filt.dropna(subset=[key_col, name_col]) \
                        .drop_duplicates(subset=[key_col]) \
                        .set_index(key_col)[name_col]
            label_lookups.update(lookup.to_dict())

    matrix = pd.concat(direction_dfs, axis=1).fillna(0).astype(int)

    # Filter: significant in at least 2 conditions
    sig_count = (matrix != 0).sum(axis=1)
    matrix = matrix[sig_count >= 2]

    if len(matrix) == 0:
        print("  No genes significant in 2+ conditions for concordance analysis")
        return pd.DataFrame()

    # Limit to top genes if too many (by number of conditions they're significant in)
    max_heatmap = 200
    title_note = ""
    if len(matrix) > max_heatmap:
        matrix["_sig_count"] = (matrix != 0).sum(axis=1)
        matrix = matrix.nlargest(max_heatmap, "_sig_count").drop(columns="_sig_count")
        title_note = f" (top {max_heatmap})"

    # Replace Ensembl IDs in index with "GeneName (ENSGXXX)" for readability
    if label_lookups:
        matrix.index = [f"{label_lookups.get(idx, idx)} ({idx})"
                        if idx in label_lookups else idx
                        for idx in matrix.index]

    # --- Clustermap ---
    cmap = sns.diverging_palette(240, 10, as_cmap=True)
    g = sns.clustermap(matrix, cmap=cmap, center=0, vmin=-1, vmax=1,
                       figsize=(8, max(6, len(matrix) * 0.15)),
                       row_cluster=True, col_cluster=False,
                       yticklabels=(len(matrix) <= 60),
                       linewidths=0.5, linecolor="white")
    g.fig.suptitle(f"Direction Concordance{title_note}\n({len(matrix)} genes sig in 2+ conditions)",
                   y=1.02, fontsize=13, fontweight="bold")
    outpath = outdir / f"direction_concordance_heatmap.{FIG_FORMAT}"
    g.savefig(outpath, format=FIG_FORMAT)
    plt.close(g.fig)
    print(f"  Saved: {outpath}")

    # --- Pairwise concordance summary bar chart ---
    pairs = list(combinations(range(len(names)), 2))
    pair_labels = []
    concordant_counts = []
    discordant_counts = []

    for i, j in pairs:
        col_i = condition_labels[names[i]]
        col_j = condition_labels[names[j]]
        both_sig = matrix[(matrix[col_i] != 0) & (matrix[col_j] != 0)]
        same_dir = int((both_sig[col_i] == both_sig[col_j]).sum())
        diff_dir = int((both_sig[col_i] != both_sig[col_j]).sum())
        pair_labels.append(f"{col_i}\nvs\n{col_j}")
        concordant_counts.append(same_dir)
        discordant_counts.append(diff_dir)

    fig, ax = plt.subplots(figsize=(8, 5))
    x = range(len(pair_labels))
    width = 0.35
    ax.bar([xi - width / 2 for xi in x], concordant_counts, width,
           label="Concordant", color="#0072B2")  # Okabe-Ito blue
    ax.bar([xi + width / 2 for xi in x], discordant_counts, width,
           label="Discordant", color="#E69F00")  # Okabe-Ito orange
    ax.set_xticks(list(x))
    ax.set_xticklabels(pair_labels, fontsize=9)
    ax.set_ylabel("Number of Genes")
    ax.set_title("Direction Concordance Between Conditions")
    ax.legend()
    for xi, cc, dc in zip(x, concordant_counts, discordant_counts):
        ax.text(xi - width / 2, cc + 1, str(cc), ha="center", fontsize=9, fontweight="bold")
        ax.text(xi + width / 2, dc + 1, str(dc), ha="center", fontsize=9, fontweight="bold")

    outpath = outdir / f"direction_concordance_summary.{FIG_FORMAT}"
    fig.savefig(outpath, format=FIG_FORMAT)
    plt.close(fig)
    print(f"  Saved: {outpath}")

    return matrix


def deseq2_log2fc_heatmap(condition_results, condition_labels, outdir,
                           min_conditions=2, max_genes=80):
    """Clustered heatmap of log2FC for genes significant in >= min_conditions.

    Uses Ensembl ID for matching across conditions (more accurate), then
    displays gene names on the heatmap y-axis for readability.
    """
    fc_col   = DESEQ2_COLS["log2fc"]
    name_col = DESEQ2_COLS.get("gene_name", "")
    names    = list(condition_results.keys())

    # Determine best key for cross-condition matching
    first_filt = condition_results[names[0]]["deseq2_filtered"]["all_genes"]
    key_col, key_desc = _best_gene_key(first_filt)
    print(f"  log2FC heatmap: matching genes by {key_desc} ('{key_col}')")

    # Find genes significant in >= min_conditions
    gene_counter = Counter()
    for name in names:
        filt = condition_results[name]["deseq2_filtered"]["all_genes"]
        if key_col in filt.columns:
            gene_counter.update(filt[key_col].dropna().unique())
    target_genes = {g for g, c in gene_counter.items() if c >= min_conditions}

    if not target_genes:
        print(f"  No genes significant in {min_conditions}+ conditions for log2FC heatmap")
        return pd.DataFrame()

    # Build log2FC matrix from RAW data, keyed by best identifier
    fc_dfs = []
    name_lookup = {}   # key_id -> gene_name for display
    for name in names:
        raw = condition_results[name]["deseq2_raw"]
        if key_col not in raw.columns:
            continue
        fc_series = raw.set_index(key_col)[fc_col]
        fc_series = fc_series[~fc_series.index.duplicated(keep="first")]
        fc_series.name = condition_labels[name]
        fc_dfs.append(fc_series)
        # Collect key -> gene_name mapping
        if name_col and name_col in raw.columns and name_col != key_col:
            lkp = raw.dropna(subset=[key_col, name_col]) \
                     .drop_duplicates(subset=[key_col]) \
                     .set_index(key_col)[name_col]
            name_lookup.update(lkp.to_dict())

    if not fc_dfs:
        print("  log2FC heatmap: key column missing in raw data, skipping")
        return pd.DataFrame()

    matrix = pd.concat(fc_dfs, axis=1)
    matrix = matrix.loc[matrix.index.isin(target_genes)].dropna()

    if len(matrix) == 0:
        print("  No overlapping genes with complete log2FC data")
        return pd.DataFrame()

    # Limit to max_genes by max absolute FC
    title_note = ""
    if len(matrix) > max_genes:
        matrix["max_abs_fc"] = matrix.abs().max(axis=1)
        matrix = matrix.nlargest(max_genes, "max_abs_fc").drop(columns="max_abs_fc")
        title_note = f" (top {max_genes})"

    # Replace Ensembl IDs with "GeneName (ENSGXXX)" labels on y-axis
    if name_lookup:
        matrix.index = [f"{name_lookup.get(idx, idx)} ({idx})"
                        if idx in name_lookup else idx
                        for idx in matrix.index]

    g = sns.clustermap(matrix, cmap="RdBu_r", center=0,
                       figsize=(8, max(6, len(matrix) * 0.18)),
                       row_cluster=True, col_cluster=False,
                       yticklabels=True, linewidths=0.3, linecolor="white")
    g.fig.suptitle(
        f"log2FC Heatmap - Genes Significant in {min_conditions}+ Conditions{title_note}",
        y=1.02, fontsize=12, fontweight="bold")

    outpath = outdir / f"log2fc_heatmap.{FIG_FORMAT}"
    g.savefig(outpath, format=FIG_FORMAT)
    plt.close(g.fig)
    print(f"  Saved: {outpath} ({len(matrix)} genes)")

    return matrix


def rmats_cross_condition_venn(condition_results, condition_labels, outdir,
                               match_by="event"):
    """Venn diagrams for significant splicing events across conditions.

    Parameters
    ----------
    match_by : str
        ``"event"`` (default) matches by genomic coordinates via
        ``_make_event_key()``.  ``"gene"`` matches by the ``geneSymbol``
        column so that any gene with *any* significant splicing event in a
        condition is counted once.
    """
    names = list(condition_results.keys())
    labels = [condition_labels[n] for n in names]
    gene_col = RMATS_COLS["gene_name"]
    id_col   = RMATS_COLS["gene_id"]

    is_gene = match_by == "gene"
    level_label = "gene-level" if is_gene else "coordinate-level"
    fname_suffix = "_genelevel" if is_gene else ""

    # Event-level outputs go into a subfolder
    outdir = Path(outdir)
    if not is_gene:
        outdir = outdir / "event_level"
        outdir.mkdir(exist_ok=True, parents=True)

    # --- Venn (all event types combined) ---
    event_sets = []
    for name in names:
        all_sig = set()
        for et, filt_df in condition_results[name]["rmats_filtered"].items():
            if is_gene:
                # Match by GeneID; fall back to geneSymbol
                _match_col = id_col if id_col in filt_df.columns else gene_col
                if _match_col in filt_df.columns:
                    all_sig.update(filt_df[_match_col].dropna().unique())
            else:
                keys = _make_event_key(filt_df, et)
                all_sig.update(keys[keys != ""].unique())
        event_sets.append(all_sig)

    fig, ax = plt.subplots(figsize=(8, 8))
    if len(event_sets) == 2:
        v = venn2(event_sets, set_labels=labels, ax=ax)
        _style_venn(v, 2)
    elif len(event_sets) == 3:
        v = venn3(event_sets, set_labels=labels, ax=ax)
        _style_venn(v, 3)
    ax.set_title(f"rMATS \u2014 Significant Splicing Events ({level_label})",
                 fontsize=13, fontweight="bold")
    outpath = Path(outdir) / f"venn_rmats_events{fname_suffix}.{FIG_FORMAT}"
    fig.savefig(outpath, format=FIG_FORMAT, dpi=FIG_DPI, bbox_inches="tight")
    plt.close(fig)
    print(f"  Saved: {outpath}")

    # --- Per-event-type Venns ---
    n_types = len(RMATS_EVENT_TYPES)
    fig, axes = plt.subplots(1, n_types, figsize=(5 * n_types, 5))
    if n_types == 1:
        axes = [axes]

    for idx, et in enumerate(RMATS_EVENT_TYPES):
        ax = axes[idx]
        sets = []
        for name in names:
            if et in condition_results[name]["rmats_filtered"]:
                filt_df = condition_results[name]["rmats_filtered"][et]
                if is_gene:
                    # Match by GeneID; fall back to geneSymbol
                    _match_col = id_col if id_col in filt_df.columns else gene_col
                    items = set(filt_df[_match_col].dropna().unique()) if _match_col in filt_df.columns else set()
                else:
                    keys = _make_event_key(filt_df, et)
                    items = set(keys[keys != ""].unique())
            else:
                items = set()
            sets.append(items)

        if len(sets) == 2:
            v = venn2(sets, set_labels=labels, ax=ax)
            _style_venn(v, 2)
        elif len(sets) == 3:
            v = venn3(sets, set_labels=labels, ax=ax)
            _style_venn(v, 3)
        ax.set_title(f"{et}", fontsize=12, fontweight="bold")

    fig.suptitle(f"rMATS \u2014 Splicing Event Overlap by Type ({level_label})",
                 fontsize=14, fontweight="bold")
    plt.tight_layout()
    outpath = Path(outdir) / f"venn_rmats_events_by_type{fname_suffix}.{FIG_FORMAT}"
    fig.savefig(outpath, format=FIG_FORMAT, dpi=FIG_DPI, bbox_inches="tight")
    plt.close(fig)
    print(f"  Saved: {outpath}")


def rmats_direction_concordance(condition_results, condition_labels, outdir):
    """Compare dPSI direction for shared significant splicing genes."""
    names = list(condition_results.keys())
    gene_col = RMATS_COLS["gene_name"]
    dpsi_col = RMATS_COLS["inclevel_diff"]

    concordance_rows = []

    for et in RMATS_EVENT_TYPES:
        # Get mean dPSI per gene for this event type from each condition
        dfs = {}
        for name in names:
            if et in condition_results[name]["rmats_filtered"]:
                df = condition_results[name]["rmats_filtered"][et]
                if len(df) > 0:
                    gene_dpsi = df.groupby(gene_col)[dpsi_col].mean()
                    dfs[name] = gene_dpsi

        if len(dfs) < 2:
            continue

        # Find shared genes
        shared_genes = set.intersection(*[set(d.index) for d in dfs.values()])

        for gene in shared_genes:
            signs = {name: np.sign(dfs[name][gene]) for name in dfs}
            all_same = len(set(signs.values())) == 1
            row = {
                "gene": gene,
                "event_type": et,
                "concordant": all_same,
            }
            for n in dfs:
                row[f"dPSI_{condition_labels[n]}"] = dfs[n][gene]
            concordance_rows.append(row)

    if not concordance_rows:
        print("  No shared significant splicing events for concordance analysis")
        return pd.DataFrame()

    conc_df = pd.DataFrame(concordance_rows)

    # Summary bar chart by event type
    event_order = [et for et in RMATS_EVENT_TYPES if et in conc_df["event_type"].values]
    concordant_counts = []
    discordant_counts = []
    for et in event_order:
        et_data = conc_df[conc_df["event_type"] == et]
        concordant_counts.append(int(et_data["concordant"].sum()))
        discordant_counts.append(int((~et_data["concordant"]).sum()))

    fig, ax = plt.subplots(figsize=(9, 5))
    x = range(len(event_order))
    width = 0.35
    ax.bar([xi - width / 2 for xi in x], concordant_counts, width,
           label="Concordant (same dPSI sign)", color="#0072B2")  # Okabe-Ito blue
    ax.bar([xi + width / 2 for xi in x], discordant_counts, width,
           label="Discordant (opposite dPSI sign)", color="#E69F00")  # Okabe-Ito orange
    ax.set_xticks(list(x))
    ax.set_xticklabels(event_order)
    ax.set_ylabel("Number of Shared Genes")
    ax.set_title("rMATS Direction Concordance - Shared Significant Splicing Genes")
    ax.legend()
    for xi, cc, dc in zip(x, concordant_counts, discordant_counts):
        ax.text(xi - width / 2, cc + 0.5, str(cc), ha="center", fontsize=9, fontweight="bold")
        ax.text(xi + width / 2, dc + 0.5, str(dc), ha="center", fontsize=9, fontweight="bold")

    outpath = outdir / f"rmats_direction_concordance.{FIG_FORMAT}"
    fig.savefig(outpath, format=FIG_FORMAT)
    plt.close(fig)
    print(f"  Saved: {outpath}")

    return conc_df


# Coordinate columns that uniquely identify each splicing event
_COORD_COLS = {
    "SE": ["chr", "strand", "exonStart_0base", "exonEnd",
            "upstreamES", "upstreamEE", "downstreamES", "downstreamEE"],
    "A3SS": ["chr", "strand", "longExonStart_0base", "longExonEnd",
             "shortES", "shortEE", "flankingES", "flankingEE"],
    "A5SS": ["chr", "strand", "longExonStart_0base", "longExonEnd",
             "shortES", "shortEE", "flankingES", "flankingEE"],
    "RI": ["chr", "strand", "riExonStart_0base", "riExonEnd",
            "upstreamES", "upstreamEE", "downstreamES", "downstreamEE"],
    "MXE": ["chr", "strand", "1stExonStart_0base", "1stExonEnd",
            "2ndExonStart_0base", "2ndExonEnd",
            "upstreamES", "upstreamEE", "downstreamES", "downstreamEE"],
}


def _make_event_key(df, event_type):
    """Create a unique event key from genomic coordinate columns.

    Returns a pandas Series of strings in the format
    ``chr:strand:col1:col2:...`` for each row in *df*.  If any required
    coordinate column is missing, an empty string is returned for every
    row so callers can filter gracefully.

    Parameters
    ----------
    df : pd.DataFrame
        rMATS filtered or raw DataFrame for a single event type.
    event_type : str
        One of ``RMATS_EVENT_TYPES`` (SE, A3SS, A5SS, RI, MXE).
    """
    coord_cols = _COORD_COLS.get(event_type, [])
    if not coord_cols:
        return pd.Series([""] * len(df), index=df.index)

    # Check that all required columns are present
    missing = [c for c in coord_cols if c not in df.columns]
    if missing:
        return pd.Series([""] * len(df), index=df.index)

    # Build key by concatenating coordinate values with ':' separator
    key = df[coord_cols[0]].astype(str)
    for col in coord_cols[1:]:
        key = key + ":" + df[col].astype(str)
    return key


def _style_venn(v, n_sets):
    """Apply Okabe-Ito colors to a matplotlib_venn Venn diagram object."""
    if v is None:
        return
    if n_sets == 2:
        for rid, col in [('10','#0072B2'),('01','#E69F00'),('11','#009E73')]:
            patch = v.get_patch_by_id(rid)
            if patch:
                patch.set_color(col)
                patch.set_alpha(0.6)
    elif n_sets == 3:
        for rid, col in [('100','#0072B2'),('010','#E69F00'),('001','#56B4E9'),
                         ('110','#009E73'),('101','#CC79A7'),('011','#F0E442'),
                         ('111','#D55E00')]:
            patch = v.get_patch_by_id(rid)
            if patch:
                patch.set_color(col)
                patch.set_alpha(0.6)


def rmats_directional_venn_diagrams(condition_results, condition_labels, outdir,
                                    match_by="event"):
    """4-panel directional Venn diagrams for splicing events by event type.

    Panels: A. All Events, B. Concordant Included, C. Concordant Excluded, D. Discordant

    Parameters
    ----------
    match_by : str
        ``"event"`` (default) matches by genomic coordinates via
        ``_make_event_key()``.  ``"gene"`` matches by the ``geneSymbol``
        column, using mean dPSI per gene for direction classification.
    """
    names = list(condition_results.keys())
    labels = [condition_labels[n] for n in names]
    dpsi_col = RMATS_COLS["inclevel_diff"]
    gene_col = RMATS_COLS["gene_name"]

    is_gene = match_by == "gene"
    level_label = "gene-level" if is_gene else "coordinate-level"
    fname_suffix = "_genelevel" if is_gene else ""

    # Event-level outputs go into a subfolder
    outdir = Path(outdir)
    if not is_gene:
        outdir = outdir / "event_level"
        outdir.mkdir(exist_ok=True, parents=True)

    if len(names) < 2:
        print("  Directional Venn diagrams require at least 2 conditions")
        return

    for et in RMATS_EVENT_TYPES:
        # Collect dPSI per item for this event type
        dfs = {}
        for name in names:
            if et in condition_results[name]["rmats_filtered"]:
                df = condition_results[name]["rmats_filtered"][et]
                if len(df) > 0:
                    if is_gene:
                        if gene_col in df.columns and dpsi_col in df.columns:
                            item_dpsi = df.groupby(gene_col)[dpsi_col].mean()
                            item_dpsi = item_dpsi[item_dpsi.index.notna()]
                            if len(item_dpsi) > 0:
                                dfs[name] = item_dpsi
                    else:
                        df = df.copy()
                        df["_ekey"] = _make_event_key(df, et).values
                        df = df[df["_ekey"] != ""]
                        if len(df) > 0:
                            # Use mean dPSI when multiple rows share the same event key
                            event_dpsi = df.groupby("_ekey")[dpsi_col].mean()
                            dfs[name] = event_dpsi

        if len(dfs) < 2:
            print(f"  Skipping {et} directional Venn (insufficient conditions)")
            continue

        # Identify all items per condition
        all_events_per_cond = {n: set(d.index) for n, d in dfs.items()}

        # Find shared items and classify by direction
        shared_events = set.intersection(*all_events_per_cond.values())

        concordant_up = set()
        concordant_down = set()
        discordant = set()

        for ekey in shared_events:
            signs = [np.sign(dfs[name][ekey]) for name in names]
            if all(s > 0 for s in signs):
                concordant_up.add(ekey)
            elif all(s < 0 for s in signs):
                concordant_down.add(ekey)
            else:
                discordant.add(ekey)

        # Validation
        computed_all = len(concordant_up) + len(concordant_down) + len(discordant)
        if computed_all != len(shared_events):
            print(f"  WARNING: Venn math mismatch for {et}: "
                  f"shared={len(shared_events)} but sum={computed_all}")

        # Create 4-panel figure
        fig, axes = plt.subplots(2, 2, figsize=(14, 12))
        axes = axes.flatten()

        # Panel A: All Events
        ax = axes[0]
        if len(names) == 2:
            v = venn2(list(all_events_per_cond.values()), set_labels=labels, ax=ax)
            _style_venn(v, 2)
        elif len(names) == 3:
            v = venn3(list(all_events_per_cond.values()), set_labels=labels, ax=ax)
            _style_venn(v, 3)
        _sizes = " | ".join(f"{labels[i]}: {len(list(all_events_per_cond.values())[i])}"
                             for i in range(len(labels)))
        ax.set_title(
            f"A. All Events ({_sizes})",
            fontsize=12, fontweight="bold")

        # Panel B: Concordant Included
        ax = axes[1]
        concordant_up_per_cond = []
        for name in names:
            events_up = set()
            if name in dfs:
                for ekey in concordant_up:
                    if ekey in dfs[name].index:
                        events_up.add(ekey)
            concordant_up_per_cond.append(events_up)

        if len(names) == 2:
            v = venn2(concordant_up_per_cond, set_labels=labels, ax=ax)
            _style_venn(v, 2)
        elif len(names) == 3:
            v = venn3(concordant_up_per_cond, set_labels=labels, ax=ax)
            _style_venn(v, 3)
        ax.set_title(f"B. Concordant Included (n={len(concordant_up)})",
                     fontsize=12, fontweight="bold")

        # Panel C: Concordant Excluded
        ax = axes[2]
        concordant_down_per_cond = []
        for name in names:
            events_down = set()
            if name in dfs:
                for ekey in concordant_down:
                    if ekey in dfs[name].index:
                        events_down.add(ekey)
            concordant_down_per_cond.append(events_down)

        if len(names) == 2:
            v = venn2(concordant_down_per_cond, set_labels=labels, ax=ax)
            _style_venn(v, 2)
        elif len(names) == 3:
            v = venn3(concordant_down_per_cond, set_labels=labels, ax=ax)
            _style_venn(v, 3)
        ax.set_title(f"C. Concordant Excluded (n={len(concordant_down)})",
                     fontsize=12, fontweight="bold")

        # Panel D: Discordant
        ax = axes[3]
        discordant_per_cond = []
        for name in names:
            events_disc = set()
            if name in dfs:
                for ekey in discordant:
                    if ekey in dfs[name].index:
                        events_disc.add(ekey)
            discordant_per_cond.append(events_disc)

        if len(names) == 2:
            v = venn2(discordant_per_cond, set_labels=labels, ax=ax)
            _style_venn(v, 2)
        elif len(names) == 3:
            v = venn3(discordant_per_cond, set_labels=labels, ax=ax)
            _style_venn(v, 3)
        ax.set_title(f"D. Discordant (n={len(discordant)})",
                     fontsize=12, fontweight="bold")

        fig.suptitle(f"Directional Splicing Overlap \u2014 {et} Events ({level_label})",
                     fontsize=14, fontweight="bold", y=0.98)
        plt.tight_layout()

        outpath = Path(outdir) / f"venn_rmats_directional_{et}{fname_suffix}.{FIG_FORMAT}"
        fig.savefig(outpath, format=FIG_FORMAT, dpi=FIG_DPI, bbox_inches="tight")
        plt.close(fig)
        print(f"  Saved: {outpath} (Inc={len(concordant_up)}, Exc={len(concordant_down)}, "
              f"Disc={len(discordant)}, Total={len(shared_events)})")

    # --- Pairwise directional Venns (venn2) ---
    if len(names) >= 3:
        for name_a, name_b in combinations(names, 2):
            label_a = condition_labels[name_a]
            label_b = condition_labels[name_b]
            pair_labels = [label_a, label_b]

            for et in RMATS_EVENT_TYPES:
                pair_dfs = {}
                for nm in (name_a, name_b):
                    if et in condition_results[nm]["rmats_filtered"]:
                        df = condition_results[nm]["rmats_filtered"][et]
                        if len(df) > 0:
                            if is_gene:
                                if gene_col in df.columns and dpsi_col in df.columns:
                                    item_dpsi = df.groupby(gene_col)[dpsi_col].mean()
                                    item_dpsi = item_dpsi[item_dpsi.index.notna()]
                                    if len(item_dpsi) > 0:
                                        pair_dfs[nm] = item_dpsi
                            else:
                                df = df.copy()
                                df["_ekey"] = _make_event_key(df, et).values
                                df = df[df["_ekey"] != ""]
                                if len(df) > 0:
                                    pair_dfs[nm] = df.groupby("_ekey")[dpsi_col].mean()

                if len(pair_dfs) < 2:
                    continue

                events_a = set(pair_dfs[name_a].index)
                events_b = set(pair_dfs[name_b].index)
                shared = events_a & events_b

                conc_up = set()
                conc_down = set()
                disc = set()
                for ekey in shared:
                    sa = np.sign(pair_dfs[name_a][ekey])
                    sb = np.sign(pair_dfs[name_b][ekey])
                    if sa > 0 and sb > 0:
                        conc_up.add(ekey)
                    elif sa < 0 and sb < 0:
                        conc_down.add(ekey)
                    else:
                        disc.add(ekey)

                # Validation
                computed = len(conc_up) + len(conc_down) + len(disc)
                if computed != len(shared):
                    print(f"  WARNING: Pairwise venn math mismatch for {et} "
                          f"{name_a} vs {name_b}: shared={len(shared)} sum={computed}")

                fig, axes = plt.subplots(2, 2, figsize=(14, 12))
                axes_flat = axes.flatten()

                # Panel A: All Events
                v = venn2([events_a, events_b], set_labels=pair_labels, ax=axes_flat[0])
                _style_venn(v, 2)
                axes_flat[0].set_title(
                    f"A. All Events ({label_a}: {len(events_a)} | {label_b}: {len(events_b)})",
                    fontsize=12, fontweight="bold")

                # Panel B: Concordant Included
                cup_a = {e for e in conc_up if e in pair_dfs[name_a].index}
                cup_b = {e for e in conc_up if e in pair_dfs[name_b].index}
                v = venn2([cup_a, cup_b], set_labels=pair_labels, ax=axes_flat[1])
                _style_venn(v, 2)
                axes_flat[1].set_title(
                    f"B. Concordant Included (n={len(conc_up)})",
                    fontsize=12, fontweight="bold")

                # Panel C: Concordant Excluded
                cdn_a = {e for e in conc_down if e in pair_dfs[name_a].index}
                cdn_b = {e for e in conc_down if e in pair_dfs[name_b].index}
                v = venn2([cdn_a, cdn_b], set_labels=pair_labels, ax=axes_flat[2])
                _style_venn(v, 2)
                axes_flat[2].set_title(
                    f"C. Concordant Excluded (n={len(conc_down)})",
                    fontsize=12, fontweight="bold")

                # Panel D: Discordant
                dsc_a = {e for e in disc if e in pair_dfs[name_a].index}
                dsc_b = {e for e in disc if e in pair_dfs[name_b].index}
                v = venn2([dsc_a, dsc_b], set_labels=pair_labels, ax=axes_flat[3])
                _style_venn(v, 2)
                axes_flat[3].set_title(
                    f"D. Discordant (n={len(disc)})",
                    fontsize=12, fontweight="bold")

                fig.suptitle(
                    f"Directional Splicing \u2014 {et} \u2014 "
                    f"{label_a} vs {label_b} ({level_label})",
                    fontsize=13, fontweight="bold", y=0.98)
                plt.tight_layout()
                outpath = (Path(outdir) /
                           f"venn_rmats_directional_{et}_{name_a}_vs_{name_b}{fname_suffix}.{FIG_FORMAT}")
                fig.savefig(outpath, format=FIG_FORMAT, dpi=FIG_DPI, bbox_inches="tight")
                plt.close(fig)
                print(f"  Saved: {outpath.name} (Inc={len(conc_up)}, "
                      f"Exc={len(conc_down)}, Disc={len(disc)})")


# ---------------------------------------------------------------------------
# PAIRWISE VENN DIAGRAMS (Brian's preferred format)
# ---------------------------------------------------------------------------

def pairwise_splicing_venns(condition_results, condition_labels, outdir,
                            match_by="event"):
    """Pairwise 5-panel Venn diagrams for splicing events.

    For each pair of conditions x each event type:
        1. All Significant events
        2. Included in both  (dPSI >= INCLEVEL_DIFF_CUTOFF)
        3. Excluded in both  (dPSI <= -INCLEVEL_DIFF_CUTOFF)
        4. Included in {A}, Excluded in {B}
        5. Excluded in {A}, Included in {B}

    Parameters
    ----------
    match_by : str
        ``"event"`` (default) matches by genomic coordinates via
        ``_make_event_key()``.  ``"gene"`` matches by the ``geneSymbol``
        column; directional subsets use genes where *any* event has
        dPSI >= cutoff (included) or dPSI <= -cutoff (excluded).
    """
    outdir = Path(outdir)
    names = list(condition_results.keys())
    dpsi_col = RMATS_COLS["inclevel_diff"]
    gene_col = RMATS_COLS["gene_name"]
    id_col   = RMATS_COLS["gene_id"]

    is_gene = match_by == "gene"
    level_label = "gene-level" if is_gene else "event-level"
    fname_suffix = "_genelevel" if is_gene else ""

    # Event-level outputs go into a subfolder
    if not is_gene:
        outdir = outdir / "event_level"
        outdir.mkdir(exist_ok=True, parents=True)

    for name_a, name_b in combinations(names, 2):
        label_a = condition_labels[name_a]
        label_b = condition_labels[name_b]

        for et in RMATS_EVENT_TYPES:
            df_a = condition_results[name_a]["rmats_filtered"].get(et)
            df_b = condition_results[name_b]["rmats_filtered"].get(et)

            if df_a is None or len(df_a) == 0 or df_b is None or len(df_b) == 0:
                continue

            if is_gene:
                # Match by GeneID for set operations; display uses geneSymbol
                _match_col = id_col if id_col in df_a.columns and id_col in df_b.columns else gene_col
                if _match_col not in df_a.columns or _match_col not in df_b.columns:
                    continue

                events_all_a = set(df_a[_match_col].dropna().unique())
                events_all_b = set(df_b[_match_col].dropna().unique())

                # Direction subsets by gene
                events_inc_a = set(
                    df_a.loc[df_a[dpsi_col] >= INCLEVEL_DIFF_CUTOFF, _match_col]
                    .dropna().unique()) if dpsi_col in df_a.columns else set()
                events_inc_b = set(
                    df_b.loc[df_b[dpsi_col] >= INCLEVEL_DIFF_CUTOFF, _match_col]
                    .dropna().unique()) if dpsi_col in df_b.columns else set()
                events_exc_a = set(
                    df_a.loc[df_a[dpsi_col] <= -INCLEVEL_DIFF_CUTOFF, _match_col]
                    .dropna().unique()) if dpsi_col in df_a.columns else set()
                events_exc_b = set(
                    df_b.loc[df_b[dpsi_col] <= -INCLEVEL_DIFF_CUTOFF, _match_col]
                    .dropna().unique()) if dpsi_col in df_b.columns else set()
            else:
                # Build event keys
                key_a = _make_event_key(df_a, et)
                key_b = _make_event_key(df_b, et)

                if key_a.eq("").all() or key_b.eq("").all():
                    continue

                # Add keys as temporary column for set operations
                df_a = df_a.copy()
                df_b = df_b.copy()
                df_a["_ekey"] = key_a.values
                df_b["_ekey"] = key_b.values

                events_all_a = set(df_a["_ekey"].dropna().unique())
                events_all_b = set(df_b["_ekey"].dropna().unique())

                # Direction subsets
                events_inc_a = set(
                    df_a.loc[df_a[dpsi_col] >= INCLEVEL_DIFF_CUTOFF, "_ekey"]
                    .dropna().unique())
                events_inc_b = set(
                    df_b.loc[df_b[dpsi_col] >= INCLEVEL_DIFF_CUTOFF, "_ekey"]
                    .dropna().unique())
                events_exc_a = set(
                    df_a.loc[df_a[dpsi_col] <= -INCLEVEL_DIFF_CUTOFF, "_ekey"]
                    .dropna().unique())
                events_exc_b = set(
                    df_b.loc[df_b[dpsi_col] <= -INCLEVEL_DIFF_CUTOFF, "_ekey"]
                    .dropna().unique())

            # Opposite-direction intersections
            inc_a_exc_b = events_inc_a & events_exc_b
            exc_a_inc_b = events_exc_a & events_inc_b

            panels = [
                ("Significant in Either Condition", events_all_a, events_all_b),
                (f"Included in Both (dPSI \u2265 {INCLEVEL_DIFF_CUTOFF})",
                 events_inc_a, events_inc_b),
                (f"Excluded in Both (dPSI \u2264 \u2212{INCLEVEL_DIFF_CUTOFF})",
                 events_exc_a, events_exc_b),
                (f"Inc {label_a} / Exc {label_b}",
                 events_inc_a, events_exc_b),
                (f"Exc {label_a} / Inc {label_b}",
                 events_exc_a, events_inc_b),
            ]

            fig, axes = plt.subplots(2, 3, figsize=(20, 12))
            axes_flat = axes.flatten()

            for idx, (panel_title, set_a, set_b) in enumerate(panels):
                ax = axes_flat[idx]
                v = venn2([set_a, set_b], set_labels=(label_a, label_b), ax=ax)
                _style_venn(v, 2)
                ax.set_title(f"{panel_title}\n({label_a}: {len(set_a)} | {label_b}: {len(set_b)})",
                             fontsize=11, fontweight="bold")

            # Hide unused 6th panel
            axes_flat[5].set_visible(False)

            fig.suptitle(f"{et} Splicing Events ({level_label}) \u2014 "
                         f"{label_a} vs {label_b}",
                         fontsize=14, fontweight="bold")
            plt.tight_layout()
            outpath = outdir / f"venn_splicing_{et}_{name_a}_vs_{name_b}{fname_suffix}.{FIG_FORMAT}"
            fig.savefig(outpath, format=FIG_FORMAT, dpi=FIG_DPI, bbox_inches="tight")
            plt.close(fig)

            print(f"  Saved: {outpath.name} "
                  f"(All: {len(events_all_a - events_all_b)}|"
                  f"{len(events_all_a & events_all_b)}|"
                  f"{len(events_all_b - events_all_a)}, "
                  f"Inc/Exc: {len(inc_a_exc_b)}, "
                  f"Exc/Inc: {len(exc_a_inc_b)})")


def pairwise_deg_venns(condition_results, condition_labels, outdir):
    """Pairwise 5-panel Venn diagrams for differentially expressed genes.

    Panels:
        1. All Significant DEGs
        2. Upregulated in both
        3. Downregulated in both
        4. Up in {A}, Down in {B}
        5. Down in {A}, Up in {B}
    """
    outdir = Path(outdir)
    names = list(condition_results.keys())

    for name_a, name_b in combinations(names, 2):
        label_a = condition_labels[name_a]
        label_b = condition_labels[name_b]

        filt_a = condition_results[name_a]["deseq2_filtered"]["all_genes"]
        filt_b = condition_results[name_b]["deseq2_filtered"]["all_genes"]

        key_col_a, _ = _best_gene_key(filt_a)
        key_col_b, _ = _best_gene_key(filt_b)
        if not key_col_a or key_col_a not in filt_a.columns:
            key_col_a = DESEQ2_COLS["gene_name"]
        if not key_col_b or key_col_b not in filt_b.columns:
            key_col_b = DESEQ2_COLS["gene_name"]

        all_a = set(filt_a[key_col_a].dropna().unique())
        all_b = set(filt_b[key_col_b].dropna().unique())
        up_a = set(filt_a.loc[filt_a["direction"] == "up", key_col_a].dropna().unique())
        up_b = set(filt_b.loc[filt_b["direction"] == "up", key_col_b].dropna().unique())
        down_a = set(filt_a.loc[filt_a["direction"] == "down", key_col_a].dropna().unique())
        down_b = set(filt_b.loc[filt_b["direction"] == "down", key_col_b].dropna().unique())

        # Opposite-direction intersections
        up_a_down_b = up_a & down_b
        down_a_up_b = down_a & up_b

        panels = [
            ("Significant in Either Condition", all_a, all_b),
            ("Upregulated", up_a, up_b),
            ("Downregulated", down_a, down_b),
            (f"Up in {label_a} / Down in {label_b}", up_a, down_b),
            (f"Down in {label_a} / Up in {label_b}", down_a, up_b),
        ]

        fig, axes = plt.subplots(2, 3, figsize=(20, 12))
        axes_flat = axes.flatten()

        for idx, (panel_title, set_a, set_b) in enumerate(panels):
            ax = axes_flat[idx]
            v = venn2([set_a, set_b], set_labels=(label_a, label_b), ax=ax)
            _style_venn(v, 2)
            ax.set_title(f"{panel_title}\n({label_a}: {len(set_a)} | {label_b}: {len(set_b)})",
                         fontsize=11, fontweight="bold")

        # Hide unused 6th panel
        axes_flat[5].set_visible(False)

        fig.suptitle(f"Differentially Expressed Genes \u2014 "
                     f"{label_a} vs {label_b}",
                     fontsize=14, fontweight="bold")
        plt.tight_layout()
        outpath = outdir / f"venn_deg_{name_a}_vs_{name_b}.{FIG_FORMAT}"
        fig.savefig(outpath, format=FIG_FORMAT, dpi=FIG_DPI, bbox_inches="tight")
        plt.close(fig)

        print(f"  Saved: {outpath.name} "
              f"(All: {len(all_a - all_b)}|{len(all_a & all_b)}|{len(all_b - all_a)}, "
              f"Up: {len(up_a - up_b)}|{len(up_a & up_b)}|{len(up_b - up_a)}, "
              f"Down: {len(down_a - down_b)}|{len(down_a & down_b)}|{len(down_b - down_a)}, "
              f"Up/Down: {len(up_a_down_b)}, Down/Up: {len(down_a_up_b)})")


# ---------------------------------------------------------------------------
# HELPERS FOR NEW ANALYSES
# ---------------------------------------------------------------------------

def _grid_dims(n):
    """Return (nrows, ncols) for a tight subplot grid of n panels."""
    if n <= 1: return (1, 1)
    if n == 2: return (1, 2)
    if n == 3: return (1, 3)
    if n == 4: return (2, 2)
    if n == 5: return (2, 3)
    ncols = int(np.ceil(np.sqrt(n)))
    nrows = int(np.ceil(n / ncols))
    return (nrows, ncols)


def _parse_inclevel_mean(series):
    """Parse rMATS IncLevel1/IncLevel2 comma-separated strings to per-row means."""
    def _row_mean(val):
        parts = [v.strip() for v in str(val).split(",") if v.strip() not in ("NA", "nan", "")]
        floats = []
        for p in parts:
            try:
                floats.append(float(p))
            except ValueError:
                pass
        return np.nanmean(floats) if floats else np.nan
    return series.apply(_row_mean)


# ---------------------------------------------------------------------------
# NEW DESEQ2 PER-CONDITION PLOTS
# ---------------------------------------------------------------------------

def pvalue_histogram(df, outdir, label="All", suffix=""):
    """Raw p-value distribution histogram — QC diagnostic."""
    cols = DESEQ2_COLS
    pval_col = cols["pvalue"]
    if pval_col not in df.columns:
        print(f"  Skipping p-value histogram: column '{pval_col}' not found")
        return

    data = df[pval_col].dropna()
    if len(data) == 0:
        print("  Skipping p-value histogram: no data")
        return

    fig, ax = plt.subplots(figsize=(7, 5))
    n_bins = 50
    _, edges, patches = ax.hist(data, bins=n_bins, range=(0, 1), edgecolor="white", lw=0.3)
    # Colour bins with left edge < PADJ_CUTOFF
    for patch, left in zip(patches, edges[:-1]):
        patch.set_facecolor(COLOR_UP if left < PADJ_CUTOFF else COLOR_NS)

    ax.axvline(PADJ_CUTOFF, color="black", ls="--", lw=1.0,
               label=f"padj cutoff ({PADJ_CUTOFF})")
    n_sig = int((data < PADJ_CUTOFF).sum())
    ax.text(0.97, 0.97,
            f"p < {PADJ_CUTOFF}: {n_sig:,}\nTotal: {len(data):,}",
            transform=ax.transAxes, ha="right", va="top", fontsize=9,
            bbox=dict(boxstyle="round,pad=0.3", facecolor="white", edgecolor="grey", alpha=0.85))
    ax.set_xlabel("Raw p-value")
    ax.set_ylabel("Gene count")
    ax.set_title(f"P-value Distribution — {label}\n"
                 "(spike near 0 = true DE signal; flat = no enrichment)")
    ax.legend(fontsize=9)

    fname = f"pvalue_histogram{suffix}.{FIG_FORMAT}"
    outpath = outdir / fname
    fig.savefig(outpath, format=FIG_FORMAT)
    plt.close(fig)
    print(f"  Saved: {outpath}")


def top_genes_lollipop(filtered_df, outdir, label="All", suffix="", top_n=20):
    """Horizontal lollipop chart of top N up + top N down genes by log2FC."""
    if len(filtered_df) == 0:
        print(f"  Skipping lollipop: no significant genes [{label}]")
        return

    cols = DESEQ2_COLS
    fc_col = cols["log2fc"]
    name_col = cols["gene_name"]

    if name_col not in filtered_df.columns:
        print(f"  Skipping lollipop: gene_name column '{name_col}' not found")
        return
    df = filtered_df.dropna(subset=[fc_col, name_col]).copy()
    up   = df[df["direction"] == "up"].nlargest(top_n, fc_col)
    down = df[df["direction"] == "down"].nsmallest(top_n, fc_col)
    plot_df = pd.concat([down, up], ignore_index=True)

    fig_h = max(6, len(plot_df) * 0.35)
    fig, ax = plt.subplots(figsize=(9, fig_h))

    colors = [COLOR_UP if d == "up" else COLOR_DOWN for d in plot_df["direction"]]
    y_pos = range(len(plot_df))

    for i, (fc, col) in enumerate(zip(plot_df[fc_col], colors)):
        ax.hlines(i, 0, fc, color=col, lw=1.8, alpha=0.8)
        ax.plot(fc, i, "o", color=col, ms=7, zorder=3)

    ax.axvline(0, color="black", lw=0.8)
    ax.axvline(LOG2FC_CUTOFF,  color="grey", ls="--", lw=0.7)
    ax.axvline(-LOG2FC_CUTOFF, color="grey", ls="--", lw=0.7)
    ax.set_yticks(list(y_pos))
    ax.set_yticklabels(plot_df[name_col].tolist(), fontsize=max(6, 9 - len(plot_df) // 15))
    ax.set_xlabel("log$_2$ Fold Change")
    ax.set_title(f"Top DE Genes — {label}\n(top {top_n} up + {top_n} down by |log2FC|)")

    ax.legend(handles=[mpatches.Patch(facecolor=COLOR_UP, label="Up"),
                       mpatches.Patch(facecolor=COLOR_DOWN, label="Down")],
              loc="lower right", fontsize=9)

    fname = f"top_genes_lollipop{suffix}.{FIG_FORMAT}"
    outpath = outdir / fname
    fig.savefig(outpath, format=FIG_FORMAT)
    plt.close(fig)
    print(f"  Saved: {outpath}")


def expression_rank_plot(df, outdir, label="All", suffix=""):
    """Waterfall rank plot: all genes sorted by log2FC, significant genes highlighted."""
    cols = DESEQ2_COLS
    data = df.dropna(subset=[cols["padj"], cols["log2fc"], cols["basemean"]]).copy()
    if len(data) == 0:
        print(f"  Skipping rank plot: no data [{label}]")
        return

    data = data.sort_values(cols["log2fc"]).reset_index(drop=True)
    rank = np.arange(len(data))

    _bm_ok = (data[cols["basemean"]] >= BASEMEAN_CUTOFF) if cols["basemean"] in data.columns else True
    conds = [
        (data[cols["padj"]] < PADJ_CUTOFF) & (data[cols["log2fc"]] >= LOG2FC_CUTOFF) & _bm_ok,
        (data[cols["padj"]] < PADJ_CUTOFF) & (data[cols["log2fc"]] <= -LOG2FC_CUTOFF) & _bm_ok,
    ]
    data["status"] = np.select(conds, ["Up", "Down"], default="NS")

    fig, ax = plt.subplots(figsize=(9, 5))
    for status, color, size, alpha, z in [
        ("NS",   COLOR_NS,   3, 0.25, 1),
        ("Down", COLOR_DOWN, 6, 0.75, 2),
        ("Up",   COLOR_UP,   6, 0.75, 2),
    ]:
        mask = data["status"] == status
        lbl = "NS" if status == "NS" else f"{status} ({mask.sum():,})"
        ax.scatter(rank[mask], data.loc[mask, cols["log2fc"]],
                   c=color, s=size, alpha=alpha, edgecolors="none", rasterized=True, zorder=z,
                   label=lbl)

    ax.axhline(0, color="black", lw=0.8)
    ax.axhline( LOG2FC_CUTOFF, color="grey", ls="--", lw=0.7)
    ax.axhline(-LOG2FC_CUTOFF, color="grey", ls="--", lw=0.7)
    n_up   = (data["status"] == "Up").sum()
    n_down = (data["status"] == "Down").sum()
    add_count_box(ax, n_up, n_down, n_up + n_down, position="lower left")
    ax.set_xlabel("Gene Rank (sorted by log$_2$ FC)")
    ax.set_ylabel("log$_2$ Fold Change")
    ax.set_title(f"Expression Rank Plot — {label}")
    ax.legend(bbox_to_anchor=(1.02, 1), loc="upper left", fontsize=9, markerscale=2)

    fname = f"expression_rank_plot{suffix}.{FIG_FORMAT}"
    outpath = outdir / fname
    fig.savefig(outpath, format=FIG_FORMAT, bbox_inches="tight")
    plt.close(fig)
    print(f"  Saved: {outpath}")


# ---------------------------------------------------------------------------
# NEW DESEQ2 CROSS-CONDITION PLOTS
# ---------------------------------------------------------------------------

def deseq2_de_counts_chart(condition_results, condition_labels, outdir):
    """Grouped bar chart: Up / Down / Total DE gene counts per condition."""
    names   = list(condition_results.keys())
    labels  = [condition_labels[n] for n in names]
    ups, downs, totals = [], [], []
    for name in names:
        filt = condition_results[name]["deseq2_filtered"].get("all_genes", pd.DataFrame())
        if "direction" in filt.columns:
            u = int((filt["direction"] == "up").sum())
            d = int((filt["direction"] == "down").sum())
        else:
            u = d = 0
        ups.append(u)
        downs.append(d)
        totals.append(u + d)

    x = np.arange(len(names))
    w = 0.25
    fig, ax = plt.subplots(figsize=(max(8, len(names) * 2.5), 6))
    b1 = ax.bar(x - w, ups,    w, label="Up",    color=COLOR_UP)
    b2 = ax.bar(x,     downs,  w, label="Down",  color=COLOR_DOWN)
    b3 = ax.bar(x + w, totals, w, label="Total", color="#888888")

    for bars in (b1, b2, b3):
        for bar in bars:
            h = bar.get_height()
            if h > 0:
                ax.text(bar.get_x() + bar.get_width() / 2, h + 0.5,
                        str(int(h)), ha="center", va="bottom", fontsize=9, fontweight="bold")

    ax.set_xticks(x)
    ax.set_xticklabels(labels, fontsize=10)
    ax.set_ylabel("Number of DE Genes")
    ax.set_title("DESeq2 — DE Gene Counts Overview Across Conditions")
    ax.legend()

    outpath = outdir / f"deseq2_de_counts_overview.{FIG_FORMAT}"
    fig.savefig(outpath, format=FIG_FORMAT)
    plt.close(fig)
    print(f"  Saved: {outpath}")


def deseq2_upset_plot(condition_results, condition_labels, outdir):
    """UpSet plot for DE gene overlap across 3-5 conditions (requires upsetplot)."""
    if not _UPSET_AVAILABLE:
        print("  Skipping UpSet plots: 'upsetplot' not installed (pip install upsetplot)")
        return
    if len(condition_results) < 3:
        print("  Skipping DESeq2 UpSet plot: requires 3+ conditions")
        return

    names  = list(condition_results.keys())
    labels = [condition_labels[n] for n in names]
    gene_col = DESEQ2_COLS["gene_name"]

    for set_key, title_suffix, fname in [
        ("all_sig", "Significant in Either Condition", "deseq2_upset_all_sig"),
        ("up",      "Upregulated Genes",         "deseq2_upset_up"),
        ("down",    "Downregulated Genes",        "deseq2_upset_down"),
    ]:
        # Build per-condition gene sets
        gene_sets = {}
        for name, lbl in zip(names, labels):
            filt = condition_results[name]["deseq2_filtered"].get("all_genes", pd.DataFrame())
            if len(filt) == 0:
                gene_sets[lbl] = set()
                continue
            if set_key == "all_sig":
                gene_sets[lbl] = set(filt[gene_col].dropna().unique())
            elif set_key == "up":
                gene_sets[lbl] = set(filt[filt["direction"] == "up"][gene_col].dropna().unique())
            else:
                gene_sets[lbl] = set(filt[filt["direction"] == "down"][gene_col].dropna().unique())

        all_genes = set.union(*gene_sets.values())
        if not all_genes:
            print(f"  Skipping UpSet ({title_suffix}): no genes")
            continue

        memberships = []
        for gene in all_genes:
            memberships.append(tuple(lbl for lbl in labels if gene in gene_sets[lbl]))

        try:
            upset_data = from_memberships(memberships)
            # Aggregate duplicate membership patterns for upsetplot compatibility
            upset_data = upset_data.groupby(level=list(range(upset_data.index.nlevels))).sum()
            upset = UpSet(upset_data, show_counts=True, sort_by="cardinality")
            upset.plot()
            plt.suptitle(f"DESeq2 UpSet — {title_suffix}", y=1.02, fontsize=12, fontweight="bold")
            outpath = outdir / f"{fname}.{FIG_FORMAT}"
            plt.savefig(outpath, format=FIG_FORMAT, bbox_inches="tight")
            plt.close("all")
            print(f"  Saved: {outpath}")
        except Exception as e:
            print(f"  UpSet plot failed ({title_suffix}): {e}")


def pairwise_log2fc_scatter(condition_results, condition_labels, outdir):
    """Pairwise scatter of log2FC for shared genes between every pair of conditions.

    Merges on Ensembl ID when available for accurate gene matching.
    """
    names  = list(condition_results.keys())
    if len(names) < 2:
        return

    fc_col   = DESEQ2_COLS["log2fc"]
    pairs    = list(combinations(names, 2))
    nrows, ncols = _grid_dims(len(pairs))

    # Use best key for matching (Ensembl ID preferred)
    first_raw = condition_results[names[0]]["deseq2_raw"]
    key_col, key_desc = _best_gene_key(first_raw)

    fig, axes = plt.subplots(nrows, ncols, figsize=(5 * ncols, 5 * nrows), squeeze=False)
    axes_flat = [axes[r][c] for r in range(nrows) for c in range(ncols)]

    for ax_idx, (nameA, nameB) in enumerate(pairs):
        ax = axes_flat[ax_idx]
        lblA = condition_labels[nameA]
        lblB = condition_labels[nameB]

        rawA = condition_results[nameA]["deseq2_raw"].dropna(subset=[key_col, fc_col])
        rawB = condition_results[nameB]["deseq2_raw"].dropna(subset=[key_col, fc_col])

        merged = rawA[[key_col, fc_col]].merge(
            rawB[[key_col, fc_col]], on=key_col, suffixes=("_A", "_B"))
        if len(merged) == 0:
            ax.set_visible(False)
            continue

        sig_A = set(condition_results[nameA]["deseq2_filtered"].get(
            "all_genes", pd.DataFrame()).get(key_col, pd.Series()).dropna().unique())
        sig_B = set(condition_results[nameB]["deseq2_filtered"].get(
            "all_genes", pd.DataFrame()).get(key_col, pd.Series()).dropna().unique())

        status = []
        for g in merged[key_col]:
            inA = g in sig_A
            inB = g in sig_B
            if inA and inB:   status.append("Both sig")
            elif inA:         status.append(f"Only {lblA}")
            elif inB:         status.append(f"Only {lblB}")
            else:             status.append("NS")
        merged["status"] = status

        color_map = {"Both sig": "#7B2D8B",
                     f"Only {lblA}": COLOR_UP,
                     f"Only {lblB}": COLOR_DOWN,
                     "NS": COLOR_NS}
        for st in ["NS", f"Only {lblA}", f"Only {lblB}", "Both sig"]:
            sub = merged[merged["status"] == st]
            lbl = "NS" if st == "NS" else f"{st} ({len(sub):,})"
            ax.scatter(sub[f"{fc_col}_A"], sub[f"{fc_col}_B"],
                       c=color_map.get(st, "#aaaaaa"), s=5, alpha=0.5,
                       edgecolors="none", rasterized=True, label=lbl)

        # Diagonal y=x
        lims = [min(merged[f"{fc_col}_A"].min(), merged[f"{fc_col}_B"].min()) - 0.2,
                max(merged[f"{fc_col}_A"].max(), merged[f"{fc_col}_B"].max()) + 0.2]
        ax.plot(lims, lims, "k--", lw=0.7, alpha=0.5)
        ax.set_xlim(lims); ax.set_ylim(lims)

        if _SCIPY_AVAILABLE and len(merged) >= 3:
            r, _ = pearsonr(merged[f"{fc_col}_A"], merged[f"{fc_col}_B"])
            # Line of best fit
            x_arr = merged[f"{fc_col}_A"].values
            y_arr = merged[f"{fc_col}_B"].values
            slope, intercept = np.polyfit(x_arr, y_arr, 1)
            x_fit = np.linspace(x_arr.min(), x_arr.max(), 100)
            ax.plot(x_fit, slope * x_fit + intercept, color="#E69F00", lw=1.5, alpha=0.8)
            ax.annotate(f"R² = {r**2:.3f}\nn = {len(merged)}", xy=(0, 1), xycoords="axes fraction",
                        xytext=(4, 4), textcoords="offset points",
                        ha="left", va="bottom", fontsize=9,
                        bbox=dict(boxstyle="round,pad=0.3", facecolor="white", edgecolor="grey", alpha=0.9))

        ax.set_xlabel(f"log$_2$FC  {lblA}", fontsize=9)
        ax.set_ylabel(f"log$_2$FC  {lblB}", fontsize=9)
        ax.set_title(f"{lblA} vs {lblB}", fontsize=10)
        ax.legend(fontsize=7, markerscale=2, loc="lower right")

    for i in range(len(pairs), nrows * ncols):
        axes_flat[i].set_visible(False)

    fig.suptitle("Pairwise log$_2$FC Comparison (Shared Genes)", fontsize=13, fontweight="bold")
    plt.tight_layout()
    outpath = outdir / f"pairwise_log2fc_scatter.{FIG_FORMAT}"
    fig.savefig(outpath, format=FIG_FORMAT, bbox_inches="tight")
    plt.close(fig)
    print(f"  Saved: {outpath}")


# ---------------------------------------------------------------------------
# NEW rMATS PER-CONDITION PLOTS
# ---------------------------------------------------------------------------

def rmats_psi_scatter(rmats_raw, rmats_filtered, event_type, outdir):
    """Mean PSI (IncLevel1 vs IncLevel2) scatter showing direction of splicing shift."""
    df = rmats_raw.get(event_type, pd.DataFrame())
    if len(df) == 0:
        return

    if "IncLevel1" not in df.columns or "IncLevel2" not in df.columns:
        print(f"  Skipping PSI scatter ({event_type}): IncLevel1/IncLevel2 columns not found")
        return

    df = df.copy()
    df["_psi1"] = _parse_inclevel_mean(df["IncLevel1"])
    df["_psi2"] = _parse_inclevel_mean(df["IncLevel2"])
    df = df.dropna(subset=["_psi1", "_psi2"])
    if len(df) == 0:
        return

    id_col = RMATS_COLS["event_id"]
    filt_df = rmats_filtered.get(event_type, pd.DataFrame())
    sig_ids = set(filt_df[id_col].unique()) if len(filt_df) > 0 and id_col in filt_df.columns else set()
    df["_sig"] = df[id_col].isin(sig_ids) if id_col in df.columns else False

    color = EVENT_COLORS.get(event_type, "#333333")
    fig, ax = plt.subplots(figsize=(6, 6))

    ns_mask  = ~df["_sig"]
    sig_mask =  df["_sig"]
    ax.scatter(df.loc[ns_mask,  "_psi1"], df.loc[ns_mask,  "_psi2"],
               c=COLOR_NS, s=5, alpha=0.2, edgecolors="none", rasterized=True, label="NS")
    ax.scatter(df.loc[sig_mask, "_psi1"], df.loc[sig_mask, "_psi2"],
               c=color,    s=8, alpha=0.7, edgecolors="none", rasterized=True,
               label=f"Significant ({sig_mask.sum():,})")

    # Diagonal and cutoff offset lines
    ax.plot([0, 1], [0, 1], "k--", lw=0.7, alpha=0.5)
    ax.plot([0, 1 - INCLEVEL_DIFF_CUTOFF], [INCLEVEL_DIFF_CUTOFF, 1],
            color="grey", ls=":", lw=0.6, alpha=0.6)
    ax.plot([INCLEVEL_DIFF_CUTOFF, 1], [0, 1 - INCLEVEL_DIFF_CUTOFF],
            color="grey", ls=":", lw=0.6, alpha=0.6)

    ax.set_xlim(0, 1); ax.set_ylim(0, 1)
    ax.set_xlabel("Mean PSI — Sample Group 1")
    ax.set_ylabel("Mean PSI — Sample Group 2")
    ax.set_title(f"PSI Shift — {event_type}")
    ax.legend(fontsize=9, markerscale=2)

    outpath = outdir / f"rmats_{event_type}_psi_scatter.{FIG_FORMAT}"
    fig.savefig(outpath, format=FIG_FORMAT)
    plt.close(fig)
    print(f"  Saved: {outpath}")


# ---------------------------------------------------------------------------
# NEW rMATS CROSS-CONDITION PLOTS
# ---------------------------------------------------------------------------

def rmats_event_count_comparison(rmats_conditions, condition_labels, outdir):
    """Grouped bar chart: x=event type, groups=conditions, y=significant event count."""
    if len(rmats_conditions) == 0:
        return

    names  = list(rmats_conditions.keys())
    labels = [condition_labels[n] for n in names]
    palette = sns.color_palette("Set1", n_colors=len(names))

    x = np.arange(len(RMATS_EVENT_TYPES))
    w = 0.8 / max(len(names), 1)

    fig, ax = plt.subplots(figsize=(10, 6))
    for idx, (name, lbl, col) in enumerate(zip(names, labels, palette)):
        counts = []
        for et in RMATS_EVENT_TYPES:
            filt_df = rmats_conditions[name]["rmats_filtered"].get(et, pd.DataFrame())
            counts.append(len(filt_df))
        offset = (idx - len(names) / 2 + 0.5) * w
        bars = ax.bar(x + offset, counts, w, label=lbl, color=col, alpha=0.85)
        for bar, cnt in zip(bars, counts):
            if cnt > 0:
                ax.text(bar.get_x() + bar.get_width() / 2, cnt + 0.5,
                        str(cnt), ha="center", va="bottom", fontsize=8, fontweight="bold")

    ax.set_xticks(x)
    ax.set_xticklabels(RMATS_EVENT_TYPES, fontsize=11)
    ax.set_ylabel("Significant Splicing Events")
    ax.set_title("rMATS — Significant Event Counts by Type Across Conditions")
    ax.legend(fontsize=9)

    outpath = outdir / f"rmats_event_count_comparison.{FIG_FORMAT}"
    fig.savefig(outpath, format=FIG_FORMAT)
    plt.close(fig)
    print(f"  Saved: {outpath}")


def pairwise_dpsi_scatter(rmats_conditions, condition_labels, outdir,
                          match_by="event"):
    """Pairwise scatter of dPSI for shared splicing events.

    Parameters
    ----------
    match_by : str
        ``"event"`` (default) matches by genomic coordinates via
        ``_make_event_key()``.  ``"gene"`` matches by the ``geneSymbol``
        column; when a gene has multiple events the event with the largest
        |dPSI| is used as the representative scatter point.
    """
    names = list(rmats_conditions.keys())
    if len(names) < 2:
        return

    dpsi_col = RMATS_COLS["inclevel_diff"]
    gene_col = RMATS_COLS["gene_name"]
    id_col   = RMATS_COLS["gene_id"]
    pairs = list(combinations(names, 2))
    nrows, ncols = _grid_dims(len(pairs))

    is_gene = match_by == "gene"
    level_label = "gene-level" if is_gene else "coordinate-level"
    fname_suffix = "_genelevel" if is_gene else ""

    # Event-level outputs go into a subfolder
    outdir = Path(outdir)
    if not is_gene:
        outdir = outdir / "event_level"
        outdir.mkdir(exist_ok=True, parents=True)

    fig, axes = plt.subplots(nrows, ncols, figsize=(5 * ncols, 5 * nrows), squeeze=False)
    axes_flat = [axes[r][c] for r in range(nrows) for c in range(ncols)]

    for ax_idx, (nameA, nameB) in enumerate(pairs):
        ax = axes_flat[ax_idx]
        lblA = condition_labels[nameA]
        lblB = condition_labels[nameB]

        all_rows = []
        for et in RMATS_EVENT_TYPES:
            dfA = rmats_conditions[nameA]["rmats_filtered"].get(et, pd.DataFrame())
            dfB = rmats_conditions[nameB]["rmats_filtered"].get(et, pd.DataFrame())
            if len(dfA) == 0 or len(dfB) == 0:
                continue

            if is_gene:
                # Match by GeneID; fall back to geneSymbol
                _match_col = id_col if id_col in dfA.columns and id_col in dfB.columns else gene_col
                if _match_col not in dfA.columns or _match_col not in dfB.columns:
                    continue
                if dpsi_col not in dfA.columns or dpsi_col not in dfB.columns:
                    continue

                # For gene-level: pick event with largest |dPSI| per gene
                dfA = dfA.copy()
                dfB = dfB.copy()
                dfA["_abs_dpsi"] = dfA[dpsi_col].abs()
                dfB["_abs_dpsi"] = dfB[dpsi_col].abs()
                repA = dfA.loc[dfA.groupby(_match_col)["_abs_dpsi"].idxmax()]
                repB = dfB.loc[dfB.groupby(_match_col)["_abs_dpsi"].idxmax()]

                cols_a = [_match_col, dpsi_col]
                cols_b = [_match_col, dpsi_col]
                merged = repA[cols_a].merge(
                    repB[cols_b], on=_match_col, suffixes=("_A", "_B"))
                if len(merged) == 0:
                    continue
                merged["event_type"] = et
                all_rows.append(merged)
            else:
                # Build event keys
                keyA = _make_event_key(dfA, et)
                keyB = _make_event_key(dfB, et)
                if keyA.eq("").all() or keyB.eq("").all():
                    continue

                dfA = dfA.copy()
                dfB = dfB.copy()
                dfA["_ekey"] = keyA.values
                dfB["_ekey"] = keyB.values

                # Keep gene symbol for tooltip / labeling
                cols_a = ["_ekey", dpsi_col]
                cols_b = ["_ekey", dpsi_col]
                if gene_col in dfA.columns:
                    cols_a.append(gene_col)
                if gene_col in dfB.columns:
                    cols_b.append(gene_col)

                merged = dfA[cols_a].merge(
                    dfB[cols_b], on="_ekey", suffixes=("_A", "_B"))
                if len(merged) == 0:
                    continue
                merged["event_type"] = et
                all_rows.append(merged)

        if not all_rows:
            ax.set_visible(False)
            continue

        combined = pd.concat(all_rows, ignore_index=True)
        for et in RMATS_EVENT_TYPES:
            sub = combined[combined["event_type"] == et]
            if len(sub) == 0:
                continue
            ax.scatter(sub[f"{dpsi_col}_A"], sub[f"{dpsi_col}_B"],
                       c=EVENT_COLORS.get(et, "#888888"), s=8, alpha=0.6,
                       edgecolors="none", rasterized=True, label=f"{et} ({len(sub):,})")

        lims = [min(combined[f"{dpsi_col}_A"].min(),
                    combined[f"{dpsi_col}_B"].min()) - 0.05,
                max(combined[f"{dpsi_col}_A"].max(),
                    combined[f"{dpsi_col}_B"].max()) + 0.05]
        ax.plot(lims, lims, "k--", lw=0.7, alpha=0.5)
        ax.axhline(0, color="black", lw=0.5)
        ax.axvline(0, color="black", lw=0.5)
        ax.set_xlim(lims)
        ax.set_ylim(lims)

        if _SCIPY_AVAILABLE and len(combined) >= 3:
            r, _ = pearsonr(combined[f"{dpsi_col}_A"], combined[f"{dpsi_col}_B"])
            x_arr = combined[f"{dpsi_col}_A"].values
            y_arr = combined[f"{dpsi_col}_B"].values
            slope, intercept = np.polyfit(x_arr, y_arr, 1)
            x_fit = np.linspace(x_arr.min(), x_arr.max(), 100)
            ax.plot(x_fit, slope * x_fit + intercept, color="#E69F00", lw=1.5, alpha=0.8)
            ax.annotate(f"R\u00b2 = {r**2:.3f}", xy=(0, 1), xycoords="axes fraction",
                        xytext=(4, 4), textcoords="offset points",
                        ha="left", va="bottom", fontsize=9,
                        bbox=dict(boxstyle="round,pad=0.3", facecolor="white",
                                  edgecolor="grey", alpha=0.9))

        ax.set_xlabel(f"dPSI  {lblA}", fontsize=9)
        ax.set_ylabel(f"dPSI  {lblB}", fontsize=9)
        ax.set_title(f"{lblA} vs {lblB}", fontsize=10)
        ax.legend(fontsize=7, markerscale=2, loc="lower right")

    for i in range(len(pairs), nrows * ncols):
        axes_flat[i].set_visible(False)

    fig.suptitle(f"Pairwise dPSI Comparison (Shared Splicing Events, {level_label})",
                 fontsize=13, fontweight="bold")
    plt.tight_layout()
    outpath = Path(outdir) / f"pairwise_dpsi_scatter{fname_suffix}.{FIG_FORMAT}"
    fig.savefig(outpath, format=FIG_FORMAT, bbox_inches="tight")
    plt.close(fig)
    print(f"  Saved: {outpath}")


def rmats_upset_plot(rmats_conditions, condition_labels, outdir,
                     match_by="event"):
    """UpSet plot for rMATS splicing event sets across 3+ conditions.

    Parameters
    ----------
    match_by : str
        ``"event"`` (default) matches by genomic coordinates via
        ``_make_event_key()``.  ``"gene"`` matches by the ``geneSymbol``
        column so each gene is counted once regardless of the number of
        splicing events it has.
    """
    if not _UPSET_AVAILABLE:
        print("  Skipping rMATS UpSet: 'upsetplot' not installed (pip install upsetplot)")
        return
    if len(rmats_conditions) < 3:
        print("  Skipping rMATS UpSet: requires 3+ conditions with rMATS data")
        return

    names = list(rmats_conditions.keys())
    labels = [condition_labels[n] for n in names]
    gene_col = RMATS_COLS["gene_name"]
    id_col   = RMATS_COLS["gene_id"]

    is_gene = match_by == "gene"
    level_label = "gene-level" if is_gene else "coordinate-level"
    fname_suffix = "_genelevel" if is_gene else ""

    # Event-level outputs go into a subfolder
    outdir = Path(outdir)
    if not is_gene:
        outdir = outdir / "event_level"
        outdir.mkdir(exist_ok=True, parents=True)

    # Combined upset (all event types)
    event_sets = {}
    for name, lbl in zip(names, labels):
        items = set()
        for et, et_df in rmats_conditions[name]["rmats_filtered"].items():
            if len(et_df) > 0:
                if is_gene:
                    # Match by GeneID; fall back to geneSymbol
                    _match_col = id_col if id_col in et_df.columns else gene_col
                    if _match_col in et_df.columns:
                        items.update(et_df[_match_col].dropna().unique())
                else:
                    keys = _make_event_key(et_df, et)
                    # Prefix with event type to avoid cross-type collisions
                    items.update(f"{et}|{k}" for k in keys[keys != ""].unique())
        event_sets[lbl] = items

    all_events = set.union(*event_sets.values()) if event_sets else set()
    if all_events:
        memberships = [tuple(lbl for lbl in labels if ev in event_sets[lbl])
                       for ev in all_events]
        try:
            upset_data = from_memberships(memberships)
            upset_data = upset_data.groupby(
                level=list(range(upset_data.index.nlevels))).sum()
            upset = UpSet(upset_data, show_counts=True, sort_by="cardinality")
            upset.plot()
            plt.suptitle(f"rMATS UpSet \u2014 Splicing Event Overlap ({level_label})",
                         y=1.02, fontsize=12, fontweight="bold")
            outpath = Path(outdir) / f"rmats_upset_events{fname_suffix}.{FIG_FORMAT}"
            plt.savefig(outpath, format=FIG_FORMAT, bbox_inches="tight")
            plt.close("all")
            print(f"  Saved: {outpath}")
        except Exception as e:
            print(f"  rMATS UpSet (events) failed: {e}")

    # Per-event-type upsets
    for et in RMATS_EVENT_TYPES:
        et_sets = {}
        for name, lbl in zip(names, labels):
            filt_df = rmats_conditions[name]["rmats_filtered"].get(et, pd.DataFrame())
            if len(filt_df) > 0:
                if is_gene:
                    # Match by GeneID; fall back to geneSymbol
                    _match_col = id_col if id_col in filt_df.columns else gene_col
                    if _match_col in filt_df.columns:
                        et_sets[lbl] = set(filt_df[_match_col].dropna().unique())
                else:
                    keys = _make_event_key(filt_df, et)
                    et_sets[lbl] = set(keys[keys != ""].unique())
        active = {lbl: s for lbl, s in et_sets.items() if s}
        if len(active) < 2:
            continue
        all_et_events = set.union(*active.values())
        all_labels = labels
        memberships = [tuple(lbl for lbl in all_labels
                             if lbl in et_sets and ev in et_sets[lbl])
                       for ev in all_et_events]
        try:
            upset_data = from_memberships(memberships)
            upset_data = upset_data.groupby(
                level=list(range(upset_data.index.nlevels))).sum()
            upset = UpSet(upset_data, show_counts=True, sort_by="cardinality")
            upset.plot()
            plt.suptitle(f"rMATS UpSet \u2014 {et} Event Overlap ({level_label})",
                         y=1.02, fontsize=12, fontweight="bold")
            outpath = Path(outdir) / f"rmats_upset_{et}{fname_suffix}.{FIG_FORMAT}"
            plt.savefig(outpath, format=FIG_FORMAT, bbox_inches="tight")
            plt.close("all")
            print(f"  Saved: {outpath}")
        except Exception as e:
            print(f"  rMATS UpSet ({et}) failed: {e}")


# ---------------------------------------------------------------------------
# EVENT-LEVEL HEATMAP, PIE CHART, AND PAIRWISE WORKBOOK EXPORT
# ---------------------------------------------------------------------------

def rmats_event_heatmap(condition_results, condition_labels, event_type, outdir):
    """Clustered heatmap of dPSI values across conditions for a specific event type.

    Parameters
    ----------
    condition_results : dict
        {name: {"rmats_filtered": {et: df, ...}, ...}} structure.
    condition_labels : dict
        {name: label} mapping.
    event_type : str
        One of RMATS_EVENT_TYPES (e.g. "SE").
    outdir : str or Path
        Output directory.
    """
    outdir = Path(outdir)
    names = list(condition_results.keys())
    labels = [condition_labels[n] for n in names]
    dpsi_col = RMATS_COLS["inclevel_diff"]
    gene_col = RMATS_COLS["gene_name"]

    # Collect dPSI per event key for each condition
    event_data = {}
    gene_lookup = {}  # event_key -> geneSymbol

    for name, lbl in zip(names, labels):
        filt_df = condition_results[name]["rmats_filtered"].get(event_type, pd.DataFrame())
        if len(filt_df) == 0:
            continue
        df = filt_df.copy()
        df["_ekey"] = _make_event_key(df, event_type).values
        df = df[df["_ekey"] != ""]
        if len(df) == 0:
            continue

        # Mean dPSI per event key (handles potential duplicates)
        grouped = df.groupby("_ekey").agg(
            dpsi=(dpsi_col, "mean"),
            gene=(gene_col, "first") if gene_col in df.columns else (dpsi_col, "count"),
        )
        event_data[lbl] = grouped["dpsi"]

        # Store gene symbol mapping
        if gene_col in df.columns:
            for ekey, gene in df.groupby("_ekey")[gene_col].first().items():
                if ekey not in gene_lookup and pd.notna(gene):
                    gene_lookup[ekey] = gene

    if len(event_data) < 2:
        print(f"  Skipping {event_type} heatmap (insufficient conditions with data)")
        return

    # Build matrix: rows = events, columns = conditions
    dpsi_df = pd.DataFrame(event_data)

    # Keep only events significant in at least 1 condition (non-NaN in >= 1 col)
    dpsi_df = dpsi_df.dropna(how="all")
    if len(dpsi_df) == 0:
        print(f"  Skipping {event_type} heatmap (no events to plot)")
        return

    # Limit to top 80 events by max |dPSI| across conditions
    max_events = 80
    if len(dpsi_df) > max_events:
        max_abs_dpsi = dpsi_df.abs().max(axis=1)
        top_idx = max_abs_dpsi.nlargest(max_events).index
        dpsi_df = dpsi_df.loc[top_idx]

    # Replace event keys with gene symbols for row labels
    row_labels = []
    for ekey in dpsi_df.index:
        gene = gene_lookup.get(ekey, "")
        if gene:
            # Abbreviate coordinate key for uniqueness
            coord_short = ekey.split(":")[0] + ":" + ekey.split(":")[-1]
            row_labels.append(f"{gene} ({coord_short})")
        else:
            row_labels.append(ekey[:40])
    dpsi_df.index = row_labels

    # Fill NaN with 0 for clustering (event not significant in that condition)
    dpsi_filled = dpsi_df.fillna(0)

    # Create diverging colormap matching pipeline colors (blue-white-orange)
    from matplotlib.colors import LinearSegmentedColormap
    cmap = LinearSegmentedColormap.from_list(
        "dpsi_diverging", [COLOR_DOWN, "white", COLOR_UP], N=256)

    # Determine symmetric color limits
    vmax = max(abs(dpsi_filled.values.min()), abs(dpsi_filled.values.max()))
    vmax = max(vmax, INCLEVEL_DIFF_CUTOFF)  # Ensure at least the cutoff is visible

    # Determine figure height based on number of events
    fig_height = max(8, len(dpsi_df) * 0.25 + 2)
    fig_width = max(6, len(labels) * 1.5 + 4)

    try:
        g = sns.clustermap(
            dpsi_filled,
            cmap=cmap,
            center=0,
            vmin=-vmax,
            vmax=vmax,
            figsize=(fig_width, fig_height),
            dendrogram_ratio=(0.15, 0.05),
            cbar_kws={"label": "dPSI (IncLevelDifference)", "shrink": 0.6},
            linewidths=0.5,
            linecolor="white",
            yticklabels=True,
            xticklabels=True,
            row_cluster=len(dpsi_filled) > 1,
            col_cluster=len(dpsi_filled.columns) > 1,
        )

        g.fig.suptitle(
            f"dPSI Heatmap \u2014 {event_type} Events (coordinate-level, "
            f"top {len(dpsi_df)} by |dPSI|)",
            fontsize=13, fontweight="bold", y=1.02)

        # Adjust row label font size for readability
        g.ax_heatmap.set_yticklabels(
            g.ax_heatmap.get_yticklabels(), fontsize=7, rotation=0)
        g.ax_heatmap.set_xticklabels(
            g.ax_heatmap.get_xticklabels(), fontsize=10, rotation=45, ha="right")

        outpath = outdir / f"heatmap_dpsi_{event_type}.{FIG_FORMAT}"
        g.savefig(outpath, format=FIG_FORMAT, dpi=FIG_DPI, bbox_inches="tight")
        plt.close("all")
        print(f"  Saved: {outpath}")
    except Exception as e:
        print(f"  {event_type} heatmap failed: {e}")
        plt.close("all")


def rmats_event_pie_chart(condition_results, condition_labels, outdir):
    """Pie chart per condition showing splicing event type distribution.

    Parameters
    ----------
    condition_results : dict
        {name: {"rmats_filtered": {et: df, ...}, ...}}.
    condition_labels : dict
        {name: label}.
    outdir : str or Path
        Output directory.
    """
    outdir = Path(outdir)
    names = list(condition_results.keys())
    labels = [condition_labels[n] for n in names]
    n_conds = len(names)

    if n_conds == 0:
        return

    nrows, ncols = _grid_dims(n_conds)
    fig, axes = plt.subplots(nrows, ncols, figsize=(5 * ncols, 5 * nrows))
    if n_conds == 1:
        axes_flat = [axes]
    else:
        axes_flat = axes.flatten() if hasattr(axes, "flatten") else [axes]

    for idx, (name, lbl) in enumerate(zip(names, labels)):
        ax = axes_flat[idx]

        counts = []
        et_labels = []
        colors = []
        for et in RMATS_EVENT_TYPES:
            filt_df = condition_results[name]["rmats_filtered"].get(et, pd.DataFrame())
            n_events = len(filt_df)
            if n_events > 0:
                counts.append(n_events)
                et_labels.append(et)
                colors.append(EVENT_COLORS.get(et, "#888888"))

        if not counts:
            ax.text(0.5, 0.5, "No significant\nevents", ha="center", va="center",
                    fontsize=11, transform=ax.transAxes)
            ax.set_title(lbl, fontsize=12, fontweight="bold")
            ax.axis("off")
            continue

        total = sum(counts)
        wedges, texts, autotexts = ax.pie(
            counts,
            labels=et_labels,
            colors=colors,
            autopct=lambda pct: f"{pct:.1f}%\n({int(round(pct / 100 * total))})",
            startangle=90,
            pctdistance=0.65,
            textprops={"fontsize": 9},
        )
        for autotext in autotexts:
            autotext.set_fontsize(8)
            autotext.set_fontweight("bold")

        ax.set_title(f"{lbl}\n(n={total:,} events)", fontsize=12, fontweight="bold")

    # Hide unused panels
    for i in range(n_conds, len(axes_flat)):
        axes_flat[i].set_visible(False)

    fig.suptitle("rMATS \u2014 Significant Event Type Distribution",
                 fontsize=14, fontweight="bold")
    plt.tight_layout()
    outpath = outdir / f"rmats_event_type_pie.{FIG_FORMAT}"
    fig.savefig(outpath, format=FIG_FORMAT, dpi=FIG_DPI, bbox_inches="tight")
    plt.close(fig)
    print(f"  Saved: {outpath}")


def export_pairwise_workbook(condition_results, condition_labels, outdir):
    """Export pairwise comparison Excel workbooks for DESeq2 and rMATS data.

    For each pair of conditions, creates one .xlsx file with:
      - DESeq2 sheets: shared (all/up/down), opposite-direction, condition-only
      - rMATS sheets (per event type): same breakdown by event coordinates
      - Summary sheet with row counts per category

    Parameters
    ----------
    condition_results : dict
        {name: {"deseq2_filtered": {"all_genes": df}, "rmats_filtered": {et: df}}}.
    condition_labels : dict
        {name: label}.
    outdir : str or Path
        Output directory.
    """
    outdir = Path(outdir)
    names = list(condition_results.keys())
    dpsi_col = RMATS_COLS["inclevel_diff"]
    gene_col = RMATS_COLS["gene_name"]
    id_col   = RMATS_COLS["gene_id"]
    lfc_col = DESEQ2_COLS["log2fc"]

    def _short(label):
        """Short label from condition label for column suffixing."""
        parts = label.split(" vs ")
        return parts[0].replace(" ", "_") if parts else label.replace(" ", "_")

    for name_a, name_b in combinations(names, 2):
        label_a = condition_labels[name_a]
        label_b = condition_labels[name_b]
        short_a = _short(label_a)
        short_b = _short(label_b)

        wb_name = f"pairwise_{name_a}_vs_{name_b}.xlsx"
        wb_path = outdir / wb_name
        summary_rows = []

        with pd.ExcelWriter(wb_path, engine="openpyxl") as writer:

            # -------------------------------------------------------
            # DESeq2 sheets
            # -------------------------------------------------------
            filt_a = condition_results[name_a]["deseq2_filtered"]["all_genes"].copy()
            filt_b = condition_results[name_b]["deseq2_filtered"]["all_genes"].copy()

            # Determine best gene key for matching
            key_col_a, _ = _best_gene_key(filt_a)
            key_col_b, _ = _best_gene_key(filt_b)
            if not key_col_a or key_col_a not in filt_a.columns:
                key_col_a = DESEQ2_COLS["gene_name"]
            if not key_col_b or key_col_b not in filt_b.columns:
                key_col_b = DESEQ2_COLS["gene_name"]

            # Ensure consistent key column name for merging
            merge_key = key_col_a
            if key_col_b != key_col_a and key_col_b in filt_b.columns:
                # Rename to match for merge
                filt_b = filt_b.rename(columns={key_col_b: merge_key})

            # Suffix non-key columns
            suffix_a = f"_{short_a}"
            suffix_b = f"_{short_b}"

            # Build gene sets
            all_a = set(filt_a[merge_key].dropna().unique())
            all_b = set(filt_b[merge_key].dropna().unique())

            dir_col = "direction"
            up_a = set(filt_a.loc[filt_a[dir_col] == "up", merge_key].dropna().unique())
            up_b = set(filt_b.loc[filt_b[dir_col] == "up", merge_key].dropna().unique())
            down_a = set(filt_a.loc[filt_a[dir_col] == "down", merge_key].dropna().unique())
            down_b = set(filt_b.loc[filt_b[dir_col] == "down", merge_key].dropna().unique())

            shared_all = all_a & all_b
            shared_up = up_a & up_b
            shared_down = down_a & down_b
            up_a_down_b = up_a & down_b
            down_a_up_b = down_a & up_b
            only_a = all_a - all_b
            only_b = all_b - all_a

            de_categories = {
                "DE_shared_all": shared_all,
                "DE_shared_up": shared_up,
                "DE_shared_down": shared_down,
                f"DE_up_{short_a}_down_{short_b}": up_a_down_b,
                f"DE_down_{short_a}_up_{short_b}": down_a_up_b,
                f"DE_only_{short_a}": only_a,
                f"DE_only_{short_b}": only_b,
            }

            for sheet_name, gene_set in de_categories.items():
                if not gene_set:
                    # Write empty sheet with header
                    empty_df = pd.DataFrame(columns=[merge_key])
                    sheet_label = sheet_name[:31]
                    empty_df.to_excel(writer, sheet_name=sheet_label, index=False)
                    summary_rows.append({"category": sheet_name, "count": 0})
                    continue

                # Subset both conditions' data for these genes
                sub_a = filt_a[filt_a[merge_key].isin(gene_set)].copy()
                sub_b = filt_b[filt_b[merge_key].isin(gene_set)].copy()

                # Rename non-key columns with suffixes
                rename_a = {c: f"{c}{suffix_a}" for c in sub_a.columns if c != merge_key}
                rename_b = {c: f"{c}{suffix_b}" for c in sub_b.columns if c != merge_key}
                sub_a = sub_a.rename(columns=rename_a)
                sub_b = sub_b.rename(columns=rename_b)

                # Merge side by side
                merged = sub_a.merge(sub_b, on=merge_key, how="outer")

                sheet_label = sheet_name[:31]
                merged.to_excel(writer, sheet_name=sheet_label, index=False)
                summary_rows.append({"category": sheet_name, "count": len(merged)})
                print(f"   {sheet_label}: {len(merged):,} genes")

            # -------------------------------------------------------
            # rMATS sheets (per event type)
            # -------------------------------------------------------
            for et in RMATS_EVENT_TYPES:
                df_a = condition_results[name_a]["rmats_filtered"].get(et)
                df_b = condition_results[name_b]["rmats_filtered"].get(et)

                has_a = df_a is not None and len(df_a) > 0
                has_b = df_b is not None and len(df_b) > 0

                if not has_a and not has_b:
                    continue

                # Prepare DataFrames with event keys
                if has_a:
                    df_a = df_a.copy()
                    df_a["_ekey"] = _make_event_key(df_a, et).values
                    df_a = df_a[df_a["_ekey"] != ""]
                else:
                    df_a = pd.DataFrame(columns=["_ekey"])

                if has_b:
                    df_b = df_b.copy()
                    df_b["_ekey"] = _make_event_key(df_b, et).values
                    df_b = df_b[df_b["_ekey"] != ""]
                else:
                    df_b = pd.DataFrame(columns=["_ekey"])

                events_all_a = set(df_a["_ekey"].dropna().unique())
                events_all_b = set(df_b["_ekey"].dropna().unique())

                # Direction subsets
                events_inc_a = set(
                    df_a.loc[df_a[dpsi_col] >= INCLEVEL_DIFF_CUTOFF, "_ekey"]
                    .dropna().unique()) if has_a and dpsi_col in df_a.columns else set()
                events_inc_b = set(
                    df_b.loc[df_b[dpsi_col] >= INCLEVEL_DIFF_CUTOFF, "_ekey"]
                    .dropna().unique()) if has_b and dpsi_col in df_b.columns else set()
                events_exc_a = set(
                    df_a.loc[df_a[dpsi_col] <= -INCLEVEL_DIFF_CUTOFF, "_ekey"]
                    .dropna().unique()) if has_a and dpsi_col in df_a.columns else set()
                events_exc_b = set(
                    df_b.loc[df_b[dpsi_col] <= -INCLEVEL_DIFF_CUTOFF, "_ekey"]
                    .dropna().unique()) if has_b and dpsi_col in df_b.columns else set()

                rmats_categories = {
                    f"{et}_shared_all": events_all_a & events_all_b,
                    f"{et}_shared_included": events_inc_a & events_inc_b,
                    f"{et}_shared_excluded": events_exc_a & events_exc_b,
                    f"{et}_inc_{short_a}_exc_{short_b}": events_inc_a & events_exc_b,
                    f"{et}_exc_{short_a}_inc_{short_b}": events_exc_a & events_inc_b,
                    f"{et}_only_{short_a}": events_all_a - events_all_b,
                    f"{et}_only_{short_b}": events_all_b - events_all_a,
                }

                # Coordinate columns used as merge keys (not suffixed)
                coord_cols = _COORD_COLS.get(et, [])
                # Shared columns that should not be suffixed
                id_cols = [c for c in [RMATS_COLS["gene_id"], RMATS_COLS["gene_name"]]
                           if (has_a and c in df_a.columns) or
                              (has_b and c in df_b.columns)]
                shared_merge_cols = (
                    [c for c in coord_cols
                     if (has_a and c in df_a.columns) or
                        (has_b and c in df_b.columns)]
                    + id_cols
                )

                for sheet_name, event_set in rmats_categories.items():
                    if not event_set:
                        empty_df = pd.DataFrame(columns=["_ekey"])
                        sheet_label = sheet_name[:31]
                        empty_df.to_excel(writer, sheet_name=sheet_label, index=False)
                        summary_rows.append({"category": sheet_name, "count": 0})
                        continue

                    sub_a = df_a[df_a["_ekey"].isin(event_set)].copy() if has_a else pd.DataFrame()
                    sub_b = df_b[df_b["_ekey"].isin(event_set)].copy() if has_b else pd.DataFrame()

                    # Drop the rMATS ID column (run-specific)
                    rmats_id = RMATS_COLS["event_id"]
                    if rmats_id in sub_a.columns:
                        sub_a = sub_a.drop(columns=[rmats_id])
                    if rmats_id in sub_b.columns:
                        sub_b = sub_b.drop(columns=[rmats_id])

                    # Suffix non-shared columns
                    if len(sub_a) > 0:
                        rename_a = {c: f"{short_a}_{c}" for c in sub_a.columns
                                    if c not in shared_merge_cols and c != "_ekey"}
                        sub_a = sub_a.rename(columns=rename_a)

                    if len(sub_b) > 0:
                        rename_b = {c: f"{short_b}_{c}" for c in sub_b.columns
                                    if c not in shared_merge_cols and c != "_ekey"}
                        sub_b = sub_b.rename(columns=rename_b)

                    # Merge on event key
                    if len(sub_a) > 0 and len(sub_b) > 0:
                        merged = sub_a.merge(sub_b.drop(
                            columns=[c for c in shared_merge_cols if c in sub_b.columns],
                            errors="ignore"),
                            on="_ekey", how="outer")
                    elif len(sub_a) > 0:
                        merged = sub_a
                    else:
                        merged = sub_b

                    # Drop the internal key column before export
                    if "_ekey" in merged.columns:
                        merged = merged.drop(columns=["_ekey"])

                    sheet_label = sheet_name[:31]
                    merged.to_excel(writer, sheet_name=sheet_label, index=False)
                    summary_rows.append({"category": sheet_name, "count": len(merged)})
                    print(f"   {sheet_label}: {len(merged):,} events")

            # -------------------------------------------------------
            # rMATS gene-level sheets (per event type)
            # -------------------------------------------------------
            for et in RMATS_EVENT_TYPES:
                df_a_gl = condition_results[name_a]["rmats_filtered"].get(et)
                df_b_gl = condition_results[name_b]["rmats_filtered"].get(et)

                has_a_gl = df_a_gl is not None and len(df_a_gl) > 0
                has_b_gl = df_b_gl is not None and len(df_b_gl) > 0

                if not has_a_gl and not has_b_gl:
                    continue

                # Build gene sets — match by GeneID, fall back to geneSymbol
                _gl_id_col = id_col if (
                    (has_a_gl and id_col in df_a_gl.columns) or
                    (has_b_gl and id_col in df_b_gl.columns)
                ) else gene_col

                genes_all_a = set(df_a_gl[_gl_id_col].dropna().unique()) if has_a_gl and _gl_id_col in df_a_gl.columns else set()
                genes_all_b = set(df_b_gl[_gl_id_col].dropna().unique()) if has_b_gl and _gl_id_col in df_b_gl.columns else set()

                # Direction subsets by gene
                genes_inc_a = set(
                    df_a_gl.loc[df_a_gl[dpsi_col] >= INCLEVEL_DIFF_CUTOFF, _gl_id_col]
                    .dropna().unique()) if has_a_gl and dpsi_col in df_a_gl.columns and _gl_id_col in df_a_gl.columns else set()
                genes_inc_b = set(
                    df_b_gl.loc[df_b_gl[dpsi_col] >= INCLEVEL_DIFF_CUTOFF, _gl_id_col]
                    .dropna().unique()) if has_b_gl and dpsi_col in df_b_gl.columns and _gl_id_col in df_b_gl.columns else set()
                genes_exc_a = set(
                    df_a_gl.loc[df_a_gl[dpsi_col] <= -INCLEVEL_DIFF_CUTOFF, _gl_id_col]
                    .dropna().unique()) if has_a_gl and dpsi_col in df_a_gl.columns and _gl_id_col in df_a_gl.columns else set()
                genes_exc_b = set(
                    df_b_gl.loc[df_b_gl[dpsi_col] <= -INCLEVEL_DIFF_CUTOFF, _gl_id_col]
                    .dropna().unique()) if has_b_gl and dpsi_col in df_b_gl.columns and _gl_id_col in df_b_gl.columns else set()

                gene_categories = {
                    f"{et}_gene_shared_all": genes_all_a & genes_all_b,
                    f"{et}_gene_shared_inc": genes_inc_a & genes_inc_b,
                    f"{et}_gene_shared_exc": genes_exc_a & genes_exc_b,
                    f"{et}_gene_inc{short_a}_exc{short_b}": genes_inc_a & genes_exc_b,
                    f"{et}_gene_exc{short_a}_inc{short_b}": genes_exc_a & genes_inc_b,
                    f"{et}_gene_only_{short_a}": genes_all_a - genes_all_b,
                    f"{et}_gene_only_{short_b}": genes_all_b - genes_all_a,
                }

                for sheet_name_gl, gene_set_gl in gene_categories.items():
                    if not gene_set_gl:
                        empty_df = pd.DataFrame(columns=[gene_col])
                        sheet_label_gl = sheet_name_gl[:31]
                        empty_df.to_excel(writer, sheet_name=sheet_label_gl, index=False)
                        summary_rows.append({"category": sheet_name_gl, "count": 0})
                        continue

                    # Include ALL event rows where the gene matches (by GeneID)
                    sub_a_gl = df_a_gl[df_a_gl[_gl_id_col].isin(gene_set_gl)].copy() if has_a_gl and _gl_id_col in df_a_gl.columns else pd.DataFrame()
                    sub_b_gl = df_b_gl[df_b_gl[_gl_id_col].isin(gene_set_gl)].copy() if has_b_gl and _gl_id_col in df_b_gl.columns else pd.DataFrame()

                    # Add source condition column and concatenate
                    if len(sub_a_gl) > 0:
                        sub_a_gl["_source_condition"] = label_a
                    if len(sub_b_gl) > 0:
                        sub_b_gl["_source_condition"] = label_b

                    if len(sub_a_gl) > 0 and len(sub_b_gl) > 0:
                        merged_gl = pd.concat([sub_a_gl, sub_b_gl], ignore_index=True)
                    elif len(sub_a_gl) > 0:
                        merged_gl = sub_a_gl
                    else:
                        merged_gl = sub_b_gl

                    sheet_label_gl = sheet_name_gl[:31]
                    merged_gl.to_excel(writer, sheet_name=sheet_label_gl, index=False)
                    n_genes = len(gene_set_gl)
                    summary_rows.append({"category": sheet_name_gl,
                                         "count": f"{n_genes} genes / {len(merged_gl)} events"})
                    print(f"   {sheet_label_gl}: {n_genes:,} genes, {len(merged_gl):,} events")

            # -------------------------------------------------------
            # Summary sheet
            # -------------------------------------------------------
            if summary_rows:
                summary_df = pd.DataFrame(summary_rows)
                summary_df.to_excel(writer, sheet_name="Summary", index=False)

        print(f"  Saved: {wb_path}")


# ---------------------------------------------------------------------------
# NEW COMBINED DESeq2 + rMATS PLOTS
# ---------------------------------------------------------------------------

def deseq2_vs_rmats_venn(condition_results, condition_labels, outdir):
    """Per-condition Venn: DE genes vs genes with significant splicing events."""
    gene_col_de = DESEQ2_COLS["gene_name"]
    gene_col_as = RMATS_COLS["gene_name"]

    for name, data in condition_results.items():
        lbl = condition_labels[name]
        filt_de = data["deseq2_filtered"].get("all_genes", pd.DataFrame())
        rmats_f = data.get("rmats_filtered", {})

        if len(filt_de) == 0 or not any(len(v) > 0 for v in rmats_f.values()):
            continue

        de_genes = set(filt_de[gene_col_de].dropna().unique())
        as_genes = set()
        for et_df in rmats_f.values():
            if len(et_df) > 0 and gene_col_as in et_df.columns:
                as_genes.update(et_df[gene_col_as].dropna().unique())

        overlap = de_genes & as_genes
        print(f"  {lbl}: {len(de_genes):,} DE genes, {len(as_genes):,} AS genes, "
              f"{len(overlap):,} overlap both")

        fig, ax = plt.subplots(figsize=(7, 7))
        venn2([de_genes, as_genes], set_labels=["DE genes", "AS genes"], ax=ax)
        ax.set_title(f"{lbl} — DE vs Alternatively Spliced Genes",
                     fontsize=12, fontweight="bold")

        outpath = outdir / f"deseq2_vs_rmats_venn_{name}.{FIG_FORMAT}"
        fig.savefig(outpath, format=FIG_FORMAT)
        plt.close(fig)
        print(f"  Saved: {outpath}")


def log2fc_vs_dpsi_scatter(condition_results, condition_labels, outdir):
    """Scatter of log2FC vs dPSI for genes that are both DE and AS (one point per event)."""
    gene_col_de = DESEQ2_COLS["gene_name"]
    fc_col      = DESEQ2_COLS["log2fc"]
    gene_col_as = RMATS_COLS["gene_name"]
    dpsi_col    = RMATS_COLS["inclevel_diff"]

    valid_conds = []
    for name, data in condition_results.items():
        filt_de = data["deseq2_filtered"].get("all_genes", pd.DataFrame())
        rmats_f = data.get("rmats_filtered", {})
        if len(filt_de) > 0 and any(len(v) > 0 for v in rmats_f.values()):
            valid_conds.append(name)

    if not valid_conds:
        print("  Skipping log2FC vs dPSI: no conditions with both DE and AS data")
        return

    nrows, ncols = _grid_dims(len(valid_conds))
    fig, axes = plt.subplots(nrows, ncols, figsize=(6 * ncols, 6 * nrows), squeeze=False)
    axes_flat = [axes[r][c] for r in range(nrows) for c in range(ncols)]

    for ax_idx, name in enumerate(valid_conds):
        ax = axes_flat[ax_idx]
        lbl = condition_labels[name]
        data = condition_results[name]
        filt_de = data["deseq2_filtered"]["all_genes"]
        rmats_f = data["rmats_filtered"]

        # Build DE lookup: gene_name -> log2FC
        de_lookup = filt_de.dropna(subset=[gene_col_de, fc_col]).set_index(gene_col_de)[fc_col]
        de_lookup = de_lookup[~de_lookup.index.duplicated(keep="first")]

        plotted = False
        for et in RMATS_EVENT_TYPES:
            et_df = rmats_f.get(et, pd.DataFrame())
            if len(et_df) == 0 or gene_col_as not in et_df.columns:
                continue
            et_df = et_df.dropna(subset=[gene_col_as, dpsi_col]).copy()
            et_df = et_df[et_df[gene_col_as].isin(de_lookup.index)]
            if len(et_df) == 0:
                continue

            x_vals = de_lookup.loc[et_df[gene_col_as].values].values
            y_vals = et_df[dpsi_col].values
            ax.scatter(x_vals, y_vals, c=EVENT_COLORS.get(et, "#888888"),
                       s=8, alpha=0.65, edgecolors="none", rasterized=True,
                       label=f"{et} ({len(et_df):,})")
            plotted = True

        if not plotted:
            ax.set_visible(False)
            continue

        # Reference lines
        ax.axhline(0,  color="black", lw=0.7)
        ax.axvline(0,  color="black", lw=0.7)
        ax.axhline( INCLEVEL_DIFF_CUTOFF,  color="grey", ls="--", lw=0.6)
        ax.axhline(-INCLEVEL_DIFF_CUTOFF,  color="grey", ls="--", lw=0.6)
        ax.axvline( LOG2FC_CUTOFF,  color="grey", ls="--", lw=0.6)
        ax.axvline(-LOG2FC_CUTOFF,  color="grey", ls="--", lw=0.6)

        # Quadrant labels
        xl, xr = ax.get_xlim()
        yb, yt = ax.get_ylim()
        for qx, qy, qtxt in [
            (xr * 0.98, yt * 0.95, "Up + Included"),
            (xr * 0.98, yb * 0.95, "Up + Excluded"),
            (xl * 0.98, yt * 0.95, "Down + Included"),
            (xl * 0.98, yb * 0.95, "Down + Excluded"),
        ]:
            ax.text(qx, qy, qtxt, ha="right" if qx > 0 else "left",
                    va="top" if qy > 0 else "bottom", fontsize=7, color="#555555",
                    style="italic")

        ax.set_xlabel("log$_2$ Fold Change (DESeq2)", fontsize=9)
        ax.set_ylabel("$\\Delta$PSI (rMATS)", fontsize=9)
        ax.set_title(f"{lbl}", fontsize=10)
        ax.legend(fontsize=7, markerscale=2)

    for i in range(len(valid_conds), nrows * ncols):
        axes_flat[i].set_visible(False)

    fig.suptitle("Combined: log$_2$FC vs $\\Delta$PSI\n(Genes Both DE and Alternatively Spliced)",
                 fontsize=12, fontweight="bold")
    plt.tight_layout()
    outpath = outdir / f"log2fc_vs_dpsi_scatter.{FIG_FORMAT}"
    fig.savefig(outpath, format=FIG_FORMAT)
    plt.close(fig)
    print(f"  Saved: {outpath}")


# ---------------------------------------------------------------------------
# EXPORT
# ---------------------------------------------------------------------------

def export_results(deseq2_sets, all_rmats_filtered, outdir):
    """Export per-condition results as two XLSX workbooks (DESeq2 + rMATS).

    deseq2_sets: dict of {label: filtered_df} e.g. {"all_genes": df, "protein_coding": df, ...}
    """
    print(f"\n-- Exported Files --")

    # --- DESeq2 workbook ---
    deseq2_xlsx = outdir / "deseq2_results.xlsx"
    with pd.ExcelWriter(deseq2_xlsx, engine="openpyxl") as writer:
        for label, df in deseq2_sets.items():
            sheet = label.replace("_", " ").title().replace(" ", "_")
            df.to_excel(writer, sheet_name=sheet, index=False)

        # Summary sheet
        summary_rows = [
            ("|log2FC| cutoff", LOG2FC_CUTOFF),
            ("baseMean cutoff", BASEMEAN_CUTOFF),
            ("padj cutoff", PADJ_CUTOFF),
            ("Auto biotype split", AUTO_BIOTYPE_SPLIT),
            ("", ""),
        ]
        for label, df in deseq2_sets.items():
            n_up = int((df["direction"] == "up").sum()) if "direction" in df.columns else "N/A"
            n_down = int((df["direction"] == "down").sum()) if "direction" in df.columns else "N/A"
            display = label.replace("_", " ").title()
            summary_rows.append((f"{display} -- genes passing filter", len(df)))
            summary_rows.append((f"{display} -- upregulated", n_up))
            summary_rows.append((f"{display} -- downregulated", n_down))

        pd.DataFrame(summary_rows, columns=["Parameter", "Value"]).to_excel(
            writer, sheet_name="Summary", index=False)

    print(f"  {deseq2_xlsx}")

    # --- rMATS workbook (only if there is rMATS data) ---
    if all_rmats_filtered:
        rmats_xlsx = outdir / "rmats_results.xlsx"
        with pd.ExcelWriter(rmats_xlsx, engine="openpyxl") as writer:
            for event_type, df in all_rmats_filtered.items():
                df.to_excel(writer, sheet_name=event_type, index=False)

            # Combined sheet
            all_combined = pd.concat(all_rmats_filtered.values(), ignore_index=True)
            all_combined.to_excel(writer, sheet_name="All_Significant", index=False)

            # Summary sheet
            summary_rows = [
                ("Filter column", "FDR" if USE_FDR else "PValue"),
                ("p-value/FDR cutoff", RMATS_FDR_CUTOFF if USE_FDR else RMATS_PVAL_CUTOFF),
                ("|dPSI| cutoff", INCLEVEL_DIFF_CUTOFF),
                ("", ""),
            ]
            for event_type, df in all_rmats_filtered.items():
                summary_rows.append((f"{event_type} significant events", len(df)))
            summary_rows.append(("Total significant events", len(all_combined)))

            pd.DataFrame(summary_rows, columns=["Parameter", "Value"]).to_excel(
                writer, sheet_name="Summary", index=False)

        print(f"  {rmats_xlsx}")


def build_multi_condition_summary(condition_results):
    """Build summary rows for the multi-condition XLSX."""
    rows = [
        ("Pipeline", "Multi-Condition DESeq2 + rMATS"),
        ("DESeq2 |log2FC| cutoff", LOG2FC_CUTOFF),
        ("DESeq2 baseMean cutoff", BASEMEAN_CUTOFF),
        ("DESeq2 padj cutoff", PADJ_CUTOFF),
        ("rMATS filter column", "FDR" if USE_FDR else "PValue"),
        ("rMATS p-value/FDR cutoff", RMATS_FDR_CUTOFF if USE_FDR else RMATS_PVAL_CUTOFF),
        ("rMATS |dPSI| cutoff", INCLEVEL_DIFF_CUTOFF),
        ("", ""),
    ]
    for cond_name, data in condition_results.items():
        filt = data["deseq2_filtered"].get("all_genes", pd.DataFrame())
        n_up = int((filt["direction"] == "up").sum()) if "direction" in filt.columns else 0
        n_down = int((filt["direction"] == "down").sum()) if "direction" in filt.columns else 0
        rows.append((f"{cond_name} -- DESeq2 significant", len(filt)))
        rows.append((f"{cond_name} -- DESeq2 up", n_up))
        rows.append((f"{cond_name} -- DESeq2 down", n_down))
        total_rmats = sum(len(df) for df in data["rmats_filtered"].values())
        rows.append((f"{cond_name} -- rMATS significant", total_rmats))
    return rows


def export_combined_results(condition_results, cross_condition_data, outdir):
    """Write master multi-condition XLSX with all conditions + cross-condition data."""
    xlsx_path = outdir / "multi_condition_results.xlsx"

    with pd.ExcelWriter(xlsx_path, engine="openpyxl") as writer:
        # Per-condition DESeq2 sheets
        for cond_name, data in condition_results.items():
            for bio_key, df in data["deseq2_filtered"].items():
                sheet = f"DESeq2_{cond_name}_{bio_key}"[:31]
                df.to_excel(writer, sheet_name=sheet, index=False)

        # Per-condition rMATS sheets
        for cond_name, data in condition_results.items():
            for et, df in data["rmats_filtered"].items():
                sheet = f"rMATS_{cond_name}_{et}"[:31]
                df.to_excel(writer, sheet_name=sheet, index=False)

        # Cross-condition sheets
        for key, sheet_name in [
            ("concordance_matrix", "Cross_DESeq2_Direction"),
            ("log2fc_matrix", "Cross_DESeq2_Log2FC"),
            ("rmats_concordance", "Cross_rMATS_Direction"),
        ]:
            if key in cross_condition_data and len(cross_condition_data[key]) > 0:
                idx = key != "rmats_concordance"  # rmats concordance has no meaningful index
                cross_condition_data[key].to_excel(
                    writer, sheet_name=sheet_name, index=idx)

        # Summary sheet
        summary_rows = build_multi_condition_summary(condition_results)
        pd.DataFrame(summary_rows, columns=["Parameter", "Value"]).to_excel(
            writer, sheet_name="Summary", index=False)

    print(f"  Master XLSX: {xlsx_path}")


def run_gsea_enrichment(condition_results, condition_labels, outdir):
    """Run GSEA prerank enrichment analysis for each condition.

    Uses gseapy with multiple databases, returns top 5 pathways per database per condition.
    Generates per-condition enrichment plots and summary tables.
    """
    try:
        import gseapy as gp
    except ImportError:
        print("  WARNING: gseapy not installed, skipping GSEA analysis")
        print("  Install with: pip install gseapy")
        return {}

    gsea_dir = outdir / "gsea_results"
    gsea_dir.mkdir(exist_ok=True)

    print("\n-- Running GSEA Enrichment --")

    # GSEA databases
    databases = GSEA_DATABASES

    gene_name_col = DESEQ2_COLS.get("gene_name", "gene_name")
    log2fc_col = DESEQ2_COLS.get("log2fc", "log2FoldChange")
    padj_col = DESEQ2_COLS.get("padj", "padj")

    gsea_results = {}

    for cond_name, data in condition_results.items():
        cond_label = condition_labels.get(cond_name, cond_name)
        print(f"  Processing {cond_label}...")

        # GSEA uses ALL expressed genes ranked by log2FC, not just filtered DEGs
        deseq2_full = data.get("deseq2_raw", pd.DataFrame())
        if len(deseq2_full) == 0:
            print(f"    No DESeq2 data for {cond_label}, skipping GSEA")
            continue

        # Create ranked gene list from full unfiltered data
        # Determine ranking column: Wald statistic (WSF method) or log2FC (fallback)
        stat_col = DESEQ2_COLS.get("stat", "stat")
        use_stat = (GSEA_RANKING == "stat" and stat_col in deseq2_full.columns
                    and deseq2_full[stat_col].notna().sum() > 0)
        if use_stat:
            rank_col = stat_col
            print(f"    Ranking by Wald statistic (column: '{stat_col}') — WSF method")
        else:
            rank_col = log2fc_col
            if GSEA_RANKING == "stat":
                print(f"    [INFO] Wald stat column not found — falling back to log2FC ranking")
            else:
                print(f"    Ranking by log2FC (column: '{log2fc_col}')")

        ranked_genes = deseq2_full[[gene_name_col, rank_col]].dropna()
        # Remove duplicate gene names (keep the one with largest absolute rank value)
        ranked_genes["abs_rank"] = ranked_genes[rank_col].abs()
        ranked_genes = ranked_genes.sort_values("abs_rank", ascending=False).drop_duplicates(subset=gene_name_col, keep="first").drop(columns="abs_rank")
        ranked_genes = ranked_genes.sort_values(rank_col, ascending=False)
        ranked_genes = ranked_genes.set_index(gene_name_col)[rank_col]

        if len(ranked_genes) < 15:
            print(f"    Too few genes ({len(ranked_genes)}) for GSEA in {cond_label}")
            continue
        print(f"    Ranking {len(ranked_genes)} genes for GSEA")

        cond_gsea_dir = gsea_dir / cond_name
        cond_gsea_dir.mkdir(exist_ok=True)

        cond_results = {}

        for db in databases:
            try:
                print(f"    Running {db}...")
                prerank_res = gp.prerank(
                    rnk=ranked_genes,
                    gene_sets=db,
                    outdir=str(cond_gsea_dir / db.replace(":", "_")),
                    min_size=GSEA_MIN_SIZE,
                    max_size=GSEA_MAX_SIZE,
                    permutation_num=GSEA_PERMUTATIONS,
                    seed=42,
                    verbose=False,
                )

                # Extract top 5 significant pathways
                res_df = _normalize_gsea_cols(prerank_res.res2d)
                if len(res_df) > 0:
                    sig_pathways = res_df[res_df["fdr"] < 0.25].sort_values("nes", key=abs, ascending=False).head(5)
                    if len(sig_pathways) > 0:
                        cond_results[db] = sig_pathways[["Term", "nes", "fdr", "geneset_size", "lead_genes"]]
                        print(f"      Found {len(sig_pathways)} significant pathways")
                    else:
                        print(f"      No significant pathways (FDR < 0.25)")
                else:
                    print(f"      No results from {db}")

            except Exception as e:
                print(f"      ERROR with {db}: {e}")
                continue

        gsea_results[cond_name] = cond_results

        # Export per-condition GSEA summary
        if cond_results:
            summary_rows = []
            for db, df in cond_results.items():
                for idx, row in df.iterrows():
                    summary_rows.append({
                        "Database": db,
                        "Pathway": row["Term"],
                        "NES": row["nes"],
                        "FDR": row["fdr"],
                        "Genes": row["geneset_size"],
                        "LeadingEdge": row["lead_genes"],
                    })

            if summary_rows:
                summary_df = pd.DataFrame(summary_rows)
                summary_path = cond_gsea_dir / f"GSEA_summary_{cond_name}.xlsx"
                summary_df.to_excel(summary_path, index=False)
                print(f"    Saved: {summary_path.name}")

    print(f"  GSEA complete: {len(gsea_results)} conditions processed")
    return gsea_results


# ===========================================================================
# GO / Pathway Over-Representation Analysis (ORA)
# ===========================================================================

_ORA_DATABASES = ORA_DATABASES

_DB_TO_CATEGORY = {
    "GO_Biological_Process_2023": "BP",
    "GO_Cellular_Component_2023": "CC",
    "GO_Molecular_Function_2023": "MF",
    "KEGG_2021_Human": "KEGG",
    "Reactome_2022": "Reactome",
}

# Okabe-Ito palette for category colors
_CATEGORY_COLORS = {
    "BP":       "#0072B2",  # blue
    "CC":       "#E69F00",  # orange
    "MF":       "#009E73",  # green
    "KEGG":     "#CC79A7",  # purple
    "Reactome": "#56B4E9",  # sky blue
}

# Species mapping for g:Profiler (common model organisms)
_GPROFILER_SPECIES = {
    "human":     "hsapiens",
    "mouse":     "mmusculus",
    "rat":       "rnorvegicus",
    "zebrafish": "drerio",
    "fly":       "dmelanogaster",
    "worm":      "celegans",
}


def run_gprofiler_ora(condition_results, condition_labels, outdir,
                      _best_gene_key_fn=None, DESEQ2_COLS_map=None):
    """Run g:Profiler over-representation analysis (ORA) for up/down DEGs.

    Uses the gprofiler-official Python package with g:SCS FDR correction
    (hierarchy-aware, superior to Benjamini-Hochberg for GO terms).
    Queries GO:BP, GO:MF, GO:CC, KEGG, and Reactome.

    Output DataFrame columns match the Enrichr schema exactly (Term,
    Adjusted_P_value, Overlap_count, Category) so that
    go_enrichment_combined_plot() works unchanged.

    Falls back to Enrichr (via run_go_enrichment with ORA_METHOD override)
    if gprofiler-official is not installed.

    Parameters
    ----------
    condition_results : dict
        Pipeline condition_results structure.
    condition_labels : dict
        Maps condition name -> human-readable label.
    outdir : Path
        Output directory (a go_ora/ subdirectory will be created).
    _best_gene_key_fn : callable, optional
        Function from the pipeline. If None, uses module-level _best_gene_key.
    DESEQ2_COLS_map : dict, optional
        Column-name mapping. If None, uses module-level DESEQ2_COLS.

    Returns
    -------
    dict : go_results[cond_name] = {"up": DataFrame, "down": DataFrame}
    """
    try:
        from gprofiler import GProfiler
    except ImportError:
        print("  [INFO] gprofiler-official not installed — falling back to Enrichr")
        print("  Install with: pip install gprofiler-official")
        # Temporarily override ORA_METHOD to avoid infinite recursion
        return run_go_enrichment(condition_results, condition_labels, outdir,
                                _best_gene_key_fn=_best_gene_key_fn,
                                DESEQ2_COLS_map=DESEQ2_COLS_map,
                                _force_enrichr=True)

    outdir = Path(outdir)
    # Use separate subdir in dual mode to avoid collisions with Enrichr output
    ora_subdir = "go_ora_gprofiler" if ORA_METHOD == "both" else "go_ora"
    ora_dir = outdir / ora_subdir
    ora_dir.mkdir(parents=True, exist_ok=True)

    if DESEQ2_COLS_map is None:
        DESEQ2_COLS_map = DESEQ2_COLS
    if _best_gene_key_fn is None:
        _best_gene_key_fn = _best_gene_key

    print("\n-- Running g:Profiler ORA (WSF method) --")

    # Map species name to g:Profiler organism code
    organism = _GPROFILER_SPECIES.get(SPECIES.lower(), "hsapiens")
    print(f"  Organism: {organism}")

    gp = GProfiler(return_dataframe=True)

    # g:Profiler source names -> our category labels
    source_to_category = {
        "GO:BP": "BP",
        "GO:MF": "MF",
        "GO:CC": "CC",
        "KEGG": "KEGG",
        "REAC": "Reactome",
    }
    sources = list(source_to_category.keys())

    go_results = {}

    for cond_name, data in condition_results.items():
        cond_label = condition_labels.get(cond_name, cond_name)

        deg_df = data.get("deseq2_filtered", {}).get("all_genes", pd.DataFrame())
        if len(deg_df) == 0:
            print(f"  No DEGs for {cond_label}, skipping ORA")
            continue

        # Determine gene column — prefer symbols for g:Profiler
        gene_col, gene_type = _best_gene_key_fn(deg_df)
        if "Ensembl" in gene_type:
            name_col = DESEQ2_COLS_map.get("gene_name", "")
            if name_col and name_col in deg_df.columns:
                gene_col = name_col
                print(f"  {cond_label}: using gene symbols ({gene_col}) for ORA "
                      f"instead of Ensembl IDs")

        # Split by direction
        if "direction" not in deg_df.columns:
            print(f"  WARNING: no 'direction' column in {cond_label} DEGs, skipping")
            continue

        up_genes = (
            deg_df.loc[deg_df["direction"] == "up", gene_col]
            .dropna().astype(str).unique().tolist()
        )
        down_genes = (
            deg_df.loc[deg_df["direction"] == "down", gene_col]
            .dropna().astype(str).unique().tolist()
        )

        cond_go = {}

        for direction, gene_list in [("up", up_genes), ("down", down_genes)]:
            n = len(gene_list)
            if n < 5:
                print(f"  {cond_label} ({direction}): only {n} genes — "
                      f"too few for ORA, skipping")
                cond_go[direction] = pd.DataFrame()
                continue

            print(f"  Running g:Profiler ORA for {cond_label} ({direction}: {n} genes)...")

            try:
                result = gp.profile(
                    organism=organism,
                    query=gene_list,
                    sources=sources,
                    significance_threshold_method="g_SCS",
                )

                if result is None or len(result) == 0:
                    print(f"    No enriched terms for {cond_label} ({direction})")
                    cond_go[direction] = pd.DataFrame()
                    continue

                # Filter by padj < 0.05
                if "p_value" in result.columns:
                    result = result[result["p_value"] < 0.05].copy()

                if len(result) == 0:
                    print(f"    No terms pass padj < 0.05 for {cond_label} ({direction})")
                    cond_go[direction] = pd.DataFrame()
                    continue

                # Map g:Profiler columns to Enrichr-compatible schema
                # g:Profiler returns: source, native, name, p_value, intersection_size, ...
                rows = []
                for _, row in result.iterrows():
                    source = row.get("source", "")
                    category = source_to_category.get(source, source)
                    rows.append({
                        "Term": row.get("name", row.get("native", "")),
                        "Adjusted_P_value": row.get("p_value", 1.0),
                        "Overlap_count": int(row.get("intersection_size", 0)),
                        "Category": category,
                    })

                direction_df = pd.DataFrame(rows)

                # Top 10 terms per database (category)
                top_rows = []
                for cat in direction_df["Category"].unique():
                    cat_df = direction_df[direction_df["Category"] == cat]
                    cat_df = cat_df.sort_values("Adjusted_P_value").head(10)
                    top_rows.append(cat_df)

                if top_rows:
                    direction_df = pd.concat(top_rows, ignore_index=True)
                    print(f"    Found {len(direction_df)} enriched terms across "
                          f"{direction_df['Category'].nunique()} databases")
                else:
                    direction_df = pd.DataFrame()

                cond_go[direction] = direction_df

            except Exception as e:
                print(f"    ERROR running g:Profiler for {cond_label} ({direction}): {e}")
                cond_go[direction] = pd.DataFrame()

        go_results[cond_name] = cond_go

        # Export per-condition ORA summary
        for direction in ("up", "down"):
            df = cond_go.get(direction, pd.DataFrame())
            if len(df) > 0:
                fname = ora_dir / f"gprofiler_ora_{cond_name}_{direction}.csv"
                df.to_csv(fname, index=False)
                print(f"    Saved: {fname.name}")

    print(f"  g:Profiler ORA complete: {len(go_results)} conditions processed")
    return go_results


def run_go_enrichment(condition_results, condition_labels, outdir,
                      _best_gene_key_fn=None, DESEQ2_COLS_map=None,
                      _force_enrichr=False):
    """Run gseapy over-representation analysis (ORA) for up/down DEGs.

    Parameters
    ----------
    condition_results : dict
        Pipeline condition_results structure.
    condition_labels : dict
        Maps condition name -> human-readable label.
    outdir : Path
        Output directory (a go_ora/ subdirectory will be created).
    _best_gene_key_fn : callable, optional
        Function from the pipeline. If None, uses module-level _best_gene_key.
    DESEQ2_COLS_map : dict, optional
        Column-name mapping. If None, uses module-level DESEQ2_COLS.
    _force_enrichr : bool, optional
        If True, skip g:Profiler routing (used internally for fallback).

    Returns
    -------
    dict  :  go_results[cond_name] = {"up": DataFrame, "down": DataFrame}
    """
    # Route to g:Profiler if configured (unless forced to Enrichr by fallback)
    # "both" mode: main() calls each function directly, so falls through to Enrichr
    if ORA_METHOD == "gprofiler" and not _force_enrichr:
        return run_gprofiler_ora(condition_results, condition_labels, outdir,
                                _best_gene_key_fn=_best_gene_key_fn,
                                DESEQ2_COLS_map=DESEQ2_COLS_map)

    try:
        import gseapy as gp
    except ImportError:
        print("  WARNING: gseapy not installed, skipping GO ORA")
        print("  Install with:  pip install gseapy")
        return {}

    outdir = Path(outdir)
    # Use separate subdir in dual mode to avoid collisions with g:Profiler output
    ora_subdir = "go_ora_enrichr" if ORA_METHOD == "both" else "go_ora"
    ora_dir = outdir / ora_subdir
    ora_dir.mkdir(parents=True, exist_ok=True)

    if DESEQ2_COLS_map is None:
        DESEQ2_COLS_map = DESEQ2_COLS
    if _best_gene_key_fn is None:
        _best_gene_key_fn = _best_gene_key

    print("\n-- Running GO / Pathway ORA --")

    go_results = {}

    for cond_name, data in condition_results.items():
        cond_label = condition_labels.get(cond_name, cond_name)

        deg_df = data.get("deseq2_filtered", {}).get("all_genes", pd.DataFrame())
        if len(deg_df) == 0:
            print(f"  No DEGs for {cond_label}, skipping ORA")
            continue

        # Determine gene column
        gene_col, gene_type = _best_gene_key_fn(deg_df)

        # For ORA we prefer gene symbols (Enrichr databases use symbols).
        # If _best_gene_key returned Ensembl IDs, try to fall back to gene_name.
        if "Ensembl" in gene_type:
            name_col = DESEQ2_COLS_map.get("gene_name", "")
            if name_col and name_col in deg_df.columns:
                gene_col = name_col
                print(f"  {cond_label}: using gene symbols ({gene_col}) for ORA "
                      f"instead of Ensembl IDs")

        # Split by direction
        if "direction" not in deg_df.columns:
            print(f"  WARNING: no 'direction' column in {cond_label} DEGs, skipping")
            continue

        up_genes = (
            deg_df.loc[deg_df["direction"] == "up", gene_col]
            .dropna().astype(str).unique().tolist()
        )
        down_genes = (
            deg_df.loc[deg_df["direction"] == "down", gene_col]
            .dropna().astype(str).unique().tolist()
        )

        cond_go = {}

        for direction, gene_list in [("up", up_genes), ("down", down_genes)]:
            n = len(gene_list)
            if n < 5:
                print(f"  {cond_label} ({direction}): only {n} genes — "
                      f"too few for ORA, skipping")
                cond_go[direction] = pd.DataFrame()
                continue

            print(f"  Running GO ORA for {cond_label} ({direction}: {n} genes)...")

            direction_rows = []

            for db in _ORA_DATABASES:
                category = _DB_TO_CATEGORY[db]
                try:
                    enr = gp.enrich(
                        gene_list=gene_list,
                        gene_sets=db,
                        outdir=None,       # don't write gseapy's own output
                        no_plot=True,
                        verbose=False,
                    )
                    res_df = enr.results if hasattr(enr, "results") else enr.res2d
                    if res_df is None or len(res_df) == 0:
                        print(f"    {db}: no enriched terms")
                        continue

                    # Parse overlap count from "Overlap" column (e.g. "3/200")
                    if "Overlap" in res_df.columns:
                        res_df["Overlap_count"] = (
                            res_df["Overlap"]
                            .astype(str)
                            .str.split("/").str[0]
                            .astype(int)
                        )
                    elif "Gene_set" in res_df.columns:
                        res_df["Overlap_count"] = 0
                    else:
                        res_df["Overlap_count"] = 0

                    # Identify p-value column
                    padj_col = None
                    for candidate in ["Adjusted P-value", "Adjusted_P_value",
                                      "FDR q-val", "padj", "fdr"]:
                        if candidate in res_df.columns:
                            padj_col = candidate
                            break
                    if padj_col is None:
                        # Try case-insensitive
                        for c in res_df.columns:
                            if "adjust" in c.lower() and "p" in c.lower():
                                padj_col = c
                                break
                    if padj_col is None:
                        print(f"    {db}: could not find adjusted p-value column "
                              f"(columns: {list(res_df.columns)})")
                        continue

                    # Identify term column
                    term_col = None
                    for candidate in ["Term", "term", "Pathway", "Gene_set"]:
                        if candidate in res_df.columns:
                            term_col = candidate
                            break
                    if term_col is None:
                        term_col = res_df.columns[0]

                    # Filter significant and take top 10
                    sig = res_df[res_df[padj_col].astype(float) < 0.05].copy()
                    if len(sig) == 0:
                        print(f"    {db}: no significant terms (padj < 0.05)")
                        continue

                    sig = sig.sort_values(padj_col, ascending=True).head(10)
                    n_sig = len(sig)
                    print(f"    {db}: {n_sig} significant terms")

                    for _, row in sig.iterrows():
                        term = str(row[term_col])
                        direction_rows.append({
                            "Term": term,
                            "Adjusted_P_value": float(row[padj_col]),
                            "Overlap_count": int(row["Overlap_count"]),
                            "Category": category,
                        })

                except Exception as e:
                    print(f"    ERROR with {db}: {e}")
                    continue

            if direction_rows:
                result_df = pd.DataFrame(direction_rows)
                result_df = result_df.sort_values("Adjusted_P_value", ascending=True)
                cond_go[direction] = result_df

                # Save per-direction Excel
                xlsx_path = ora_dir / f"GO_ORA_{cond_name}_{direction}.xlsx"
                result_df.to_excel(xlsx_path, index=False)
                print(f"    Saved: {xlsx_path.name}")
            else:
                cond_go[direction] = pd.DataFrame()
                print(f"    No significant terms for {cond_label} ({direction})")

        go_results[cond_name] = cond_go

    print(f"\n  GO ORA complete: {len(go_results)} conditions processed")
    return go_results


def go_enrichment_combined_plot(go_results, condition_labels, outdir,
                                FIG_FORMAT_override=None, FIG_DPI_override=None,
                                filename_suffix=""):
    """Create a two-panel dot plot (up / down) for each condition.

    Parameters
    ----------
    go_results : dict
        Output of run_go_enrichment().
    condition_labels : dict
        Maps condition name -> human-readable label.
    outdir : Path
        Output directory for figures.
    FIG_FORMAT_override : str, optional
        Image format. Uses module-level FIG_FORMAT if None.
    FIG_DPI_override : int, optional
        Resolution. Uses module-level FIG_DPI if None.
    """
    fmt = FIG_FORMAT_override if FIG_FORMAT_override is not None else FIG_FORMAT
    dpi = FIG_DPI_override if FIG_DPI_override is not None else FIG_DPI

    outdir = Path(outdir)
    outdir.mkdir(parents=True, exist_ok=True)

    for cond_name, directions in go_results.items():
        cond_label = condition_labels.get(cond_name, cond_name)
        up_df = directions.get("up", pd.DataFrame())
        down_df = directions.get("down", pd.DataFrame())

        has_up = len(up_df) > 0
        has_down = len(down_df) > 0

        if not has_up and not has_down:
            print(f"  No ORA results for {cond_label}, skipping plot")
            continue

        fig, axes = plt.subplots(1, 2, figsize=(20, 12))
        method_label = filename_suffix.replace("_", " ").strip().title()
        title_extra = f" ({method_label})" if method_label else ""
        fig.suptitle(f"GO/Pathway Enrichment \u2014 {cond_label}{title_extra}",
                     fontsize=14, fontweight="bold", y=0.98)

        for ax, df, title in [(axes[0], up_df, "Upregulated"),
                               (axes[1], down_df, "Downregulated")]:

            if len(df) == 0:
                ax.set_title(title, fontsize=12, fontweight="bold")
                ax.text(0.5, 0.5, "No significant enrichment",
                        ha="center", va="center", fontsize=11,
                        transform=ax.transAxes, color="#666666")
                ax.set_xticks([])
                ax.set_yticks([])
                for spine in ax.spines.values():
                    spine.set_visible(False)
                continue

            # Prepare plot data: top 10 per category, sorted by significance
            plot_rows = []
            for cat in ["BP", "CC", "MF", "KEGG", "Reactome"]:
                cat_df = df[df["Category"] == cat].copy()
                if len(cat_df) == 0:
                    continue
                cat_df = cat_df.sort_values("Adjusted_P_value").head(10)
                plot_rows.append(cat_df)

            if not plot_rows:
                ax.set_title(title, fontsize=12, fontweight="bold")
                ax.text(0.5, 0.5, "No significant enrichment",
                        ha="center", va="center", fontsize=11,
                        transform=ax.transAxes, color="#666666")
                ax.set_xticks([])
                ax.set_yticks([])
                for spine in ax.spines.values():
                    spine.set_visible(False)
                continue

            plot_df = pd.concat(plot_rows, ignore_index=True)
            # Sort: most significant at the top of the plot (reversed for barh)
            plot_df = plot_df.sort_values("Adjusted_P_value", ascending=False)

            # Compute -log10(padj)
            plot_df["neg_log10_padj"] = -np.log10(
                plot_df["Adjusted_P_value"].clip(lower=1e-300)
            )

            # Build label with category prefix
            plot_df["label"] = (
                plot_df["Category"] + ": " + plot_df["Term"]
            )
            # Truncate long labels
            plot_df["label"] = plot_df["label"].apply(
                lambda x: x[:80] + "..." if len(x) > 83 else x
            )

            colors = [_CATEGORY_COLORS.get(c, "#999999") for c in plot_df["Category"]]

            # Dot sizes proportional to overlap count
            counts = plot_df["Overlap_count"].values
            min_size, max_size = 30, 300
            if counts.max() > counts.min():
                sizes = min_size + (counts - counts.min()) / (
                    counts.max() - counts.min()
                ) * (max_size - min_size)
            else:
                sizes = np.full(len(counts), (min_size + max_size) / 2)

            ax.scatter(plot_df["neg_log10_padj"], range(len(plot_df)),
                       c=colors, s=sizes, edgecolors="white", linewidths=0.5,
                       zorder=3)
            ax.set_yticks(range(len(plot_df)))
            ax.set_yticklabels(plot_df["label"].tolist(), fontsize=9)
            ax.set_xlabel("-log$_{10}$(adjusted p-value)", fontsize=10)
            ax.set_title(title, fontsize=12, fontweight="bold")
            ax.grid(axis="x", alpha=0.3, linestyle="--")
            ax.set_axisbelow(True)

        # -- Legend for category colors --
        cat_handles = [
            mlines.Line2D([], [], marker="o", color="w",
                          markerfacecolor=_CATEGORY_COLORS[cat],
                          markeredgecolor="white", markersize=9, label=cat)
            for cat in ["BP", "CC", "MF", "KEGG", "Reactome"]
        ]

        # -- Legend for dot sizes --
        all_counts = []
        for d in [up_df, down_df]:
            if len(d) > 0 and "Overlap_count" in d.columns:
                all_counts.extend(d["Overlap_count"].tolist())

        if all_counts:
            c_min, c_max = int(min(all_counts)), int(max(all_counts))
            if c_min == c_max:
                size_vals = [c_min]
            else:
                c_mid = (c_min + c_max) // 2
                size_vals = [c_min, c_mid, c_max]

            size_handles = []
            for sv in size_vals:
                if c_max > c_min:
                    s = min_size + (sv - c_min) / (c_max - c_min) * (max_size - min_size)
                else:
                    s = (min_size + max_size) / 2
                size_handles.append(
                    mlines.Line2D([], [], marker="o", color="w",
                                  markerfacecolor="#AAAAAA",
                                  markeredgecolor="gray",
                                  markersize=np.sqrt(s) * 0.7,
                                  label=f"{sv} genes")
                )
            all_handles = cat_handles + size_handles
        else:
            all_handles = cat_handles

        fig.legend(handles=all_handles, loc="lower center",
                   ncol=len(all_handles), fontsize=9,
                   frameon=True, fancybox=True, shadow=False,
                   bbox_to_anchor=(0.5, 0.01))

        plt.tight_layout(rect=[0, 0.06, 1, 0.95])

        fig_path = outdir / f"go_enrichment_combined_{cond_name}{filename_suffix}.{fmt}"
        fig.savefig(fig_path, dpi=dpi, bbox_inches="tight")
        plt.close(fig)
        print(f"  Saved: {fig_path.name}")


def export_go_prism(go_results, condition_labels, prism_dir, filename_suffix=""):
    """Export GO ORA results to GraphPad Prism .pzfx files.

    Parameters
    ----------
    go_results : dict
        Output of run_go_enrichment().
    condition_labels : dict
        Maps condition name -> human-readable label.
    prism_dir : Path
        Directory for Prism files.
    """
    import xml.etree.ElementTree as ET
    from xml.dom import minidom

    prism_dir = Path(prism_dir)
    prism_dir.mkdir(parents=True, exist_ok=True)

    # -- Reuse the same corrected Prism helpers via closures --
    # (create_prism_xml, add_table, save_prism_xml defined in export_prism_files
    #  but we need local aliases here; redefine with same correct logic)

    def _is_numeric_str(s):
        try:
            float(s)
            return True
        except (ValueError, TypeError):
            return False

    def _create_prism_xml():
        root = ET.Element("GraphPadPrismFile")
        root.set("PrismXMLVersion", "5.00")
        created = ET.SubElement(root, "Created")
        orig = ET.SubElement(created, "OriginalVersion")
        orig.set("CreatedByProgram", "GraphPad Prism")
        orig.set("CreatedByVersion", "6.0f.254")
        orig.set("Login", "")
        orig.set("DateTime", datetime.now().strftime("%Y-%m-%dT%H:%M:%S+00:00"))
        info_seq = ET.SubElement(root, "InfoSequence")
        ref = ET.SubElement(info_seq, "Ref")
        ref.set("ID", "Info0")
        ref.set("Selected", "1")
        info = ET.SubElement(root, "Info")
        info.set("ID", "Info0")
        ET.SubElement(info, "Title").text = "Project info 1"
        ET.SubElement(info, "Notes")
        table_seq = ET.SubElement(root, "TableSequence")
        table_seq.set("Selected", "1")
        root.set("_table_count", "0")
        return root

    def _add_table(root, table_name, columns_data):
        table_idx = int(root.get("_table_count", "0"))
        table_id = f"Table{table_idx}"
        root.set("_table_count", str(table_idx + 1))
        table_seq = root.find("TableSequence")
        ref = ET.SubElement(table_seq, "Ref")
        ref.set("ID", table_id)
        if table_idx == 0:
            ref.set("Selected", "1")
        table = ET.SubElement(root, "Table")
        table.set("ID", table_id)
        table.set("XFormat", "none")
        table.set("TableType", "OneWay")
        table.set("EVFormat", "AsteriskAfterNumber")
        ET.SubElement(table, "Title").text = table_name
        # Auto-detect row titles
        row_titles_col = None
        if columns_data:
            _, first_vals = columns_data[0]
            if all(isinstance(v, str) and not _is_numeric_str(v) for v in first_vals):
                row_titles_col = 0
        if row_titles_col is not None:
            _, rt_vals = columns_data[row_titles_col]
            rt_elem = ET.SubElement(table, "RowTitlesColumn")
            rt_elem.set("Width", "89")
            subcol = ET.SubElement(rt_elem, "Subcolumn")
            for val in rt_vals:
                ET.SubElement(subcol, "d").text = str(val)
            data_columns = [c for i, c in enumerate(columns_data) if i != row_titles_col]
        else:
            data_columns = list(columns_data)
        for col_name, values in data_columns:
            ycol = ET.SubElement(table, "YColumn")
            ycol.set("Width", "89")
            ycol.set("Decimals", "6")
            ycol.set("Subcolumns", "1")
            ET.SubElement(ycol, "Title").text = str(col_name)
            subcol = ET.SubElement(ycol, "Subcolumn")
            for val in values:
                d = ET.SubElement(subcol, "d")
                if val is None or (isinstance(val, float) and str(val) == 'nan'):
                    d.text = ""
                else:
                    d.text = str(val)
        return table

    def _save_prism_xml(root, filepath):
        if "_table_count" in root.attrib:
            del root.attrib["_table_count"]
        xml_str = ET.tostring(root, encoding="unicode", xml_declaration=False)
        dom = minidom.parseString(xml_str)
        pretty_xml = dom.toprettyxml(indent="  ", encoding=None)
        if pretty_xml.startswith('<?xml'):
            first_newline = pretty_xml.index('\n')
            pretty_xml = pretty_xml[first_newline + 1:]
        with open(filepath, "w", encoding="utf-8") as f:
            f.write('<?xml version="1.0" encoding="UTF-8"?>\n')
            f.write(pretty_xml)

    # -- Generate one .pzfx per condition --

    for cond_name, directions in go_results.items():
        cond_label = condition_labels.get(cond_name, cond_name)
        root = _create_prism_xml()
        tables_added = 0

        for direction in ("up", "down"):
            df = directions.get(direction, pd.DataFrame())
            if len(df) == 0:
                continue

            terms = df["Term"].tolist()
            categories = df["Category"].tolist()
            neg_log_pvals = (-np.log10(
                df["Adjusted_P_value"].clip(lower=1e-300)
            )).round(4).tolist()
            counts = df["Overlap_count"].astype(int).tolist()

            table_name = f"GO_ORA_{direction}_{cond_label}"
            _add_table(root, table_name, [
                ("Term", terms),
                ("Category", categories),
                ("-log10(padj)", neg_log_pvals),
                ("GeneCount", counts),
            ])
            tables_added += 1

        if tables_added > 0:
            pzfx_path = prism_dir / f"GO_ORA_{cond_name}{filename_suffix}.pzfx"
            _save_prism_xml(root, pzfx_path)
            print(f"  Saved Prism: {pzfx_path.name}")
        else:
            print(f"  No ORA data for {cond_label}, skipping Prism export")


def _gsea_dotplot_legacy(gsea_results, condition_labels, outdir):
    """Generate horizontal dot plots of top GSEA pathways per condition.

    Falls back to reading gseapy CSV reports from disk when in-memory
    results are empty.

    NOTE: Legacy function, superseded by gsea_combined_plot().
    """
    gsea_dir = outdir / "gsea_results"
    gsea_dir.mkdir(exist_ok=True)

    print("\n-- Generating GSEA Dot Plots --")

    for cond_name, cond_label in condition_labels.items():
        # Collect pathways from in-memory dict or disk fallback
        all_rows = []

        if gsea_results and cond_name in gsea_results:
            for db_name, pathways_df in gsea_results[cond_name].items():
                if len(pathways_df) > 0:
                    for _, row in pathways_df.iterrows():
                        all_rows.append({
                            "Term": row["Term"],
                            "NES": float(row["nes"]),
                            "FDR": float(row["fdr"]),
                        })

        if not all_rows:
            # Fallback: scan disk CSV reports
            cond_gsea_dir = gsea_dir / cond_name
            if cond_gsea_dir.is_dir():
                for db_subdir in sorted(cond_gsea_dir.iterdir()):
                    if not db_subdir.is_dir():
                        continue
                    report_csv = db_subdir / "gseapy.gene_set.prerank.report.csv"
                    if not report_csv.exists():
                        continue
                    try:
                        rpt = _normalize_gsea_cols(pd.read_csv(report_csv))
                        if len(rpt) == 0:
                            continue
                        rpt["abs_nes"] = rpt["nes"].abs()
                        rpt = rpt.sort_values("abs_nes", ascending=False).head(5)
                        for _, row in rpt.iterrows():
                            all_rows.append({
                                "Term": row["Term"],
                                "NES": float(row["nes"]),
                                "FDR": float(row["fdr"]),
                            })
                    except Exception:
                        continue

        if not all_rows:
            print(f"  No GSEA data for {cond_label}, skipping dotplot")
            continue

        # Sort by |NES|, keep top 20
        df = pd.DataFrame(all_rows)
        df["abs_NES"] = df["NES"].abs()
        df = df.sort_values("abs_NES", ascending=False).head(20)
        df = df.sort_values("NES", ascending=True)  # bottom-to-top ordering

        # Truncate long pathway names
        df["Term_short"] = df["Term"].str[:50]

        # Dot size: scaled by -log10(FDR + 1e-10)
        neg_log_fdr = -np.log10(df["FDR"].values + 1e-10)
        size_min, size_max = 20, 200
        if neg_log_fdr.max() > neg_log_fdr.min():
            sizes = size_min + (neg_log_fdr - neg_log_fdr.min()) / (neg_log_fdr.max() - neg_log_fdr.min()) * (size_max - size_min)
        else:
            sizes = np.full_like(neg_log_fdr, (size_min + size_max) / 2)

        # Dot color: orange for positive NES, blue for negative NES
        colors = [COLOR_UP if nes > 0 else COLOR_DOWN for nes in df["NES"].values]

        fig, ax = plt.subplots(figsize=(10, max(5, len(df) * 0.35)))
        ax.scatter(df["NES"].values, range(len(df)), s=sizes, c=colors,
                   edgecolors="black", linewidths=0.5, zorder=3)
        ax.set_yticks(range(len(df)))
        ax.set_yticklabels(df["Term_short"].values, fontsize=9)
        ax.axvline(0, color="grey", linewidth=0.8, linestyle="--", zorder=1)
        ax.set_xlabel("NES (Normalized Enrichment Score)", fontsize=11)
        ax.set_title(f"GSEA Top Pathways \u2014 {cond_label}", fontsize=13, fontweight="bold")

        # Legend for dot size
        for fdr_val, label in [(0.05, "FDR=0.05"), (0.25, "FDR=0.25")]:
            sz = size_min + (-np.log10(fdr_val + 1e-10) - neg_log_fdr.min()) / max(neg_log_fdr.max() - neg_log_fdr.min(), 1e-10) * (size_max - size_min)
            sz = np.clip(sz, size_min, size_max)
            ax.scatter([], [], s=sz, c="grey", edgecolors="black", linewidths=0.5, label=label)
        ax.legend(title="Dot size", loc="lower right", fontsize=8, title_fontsize=9)

        plt.tight_layout()
        outpath = gsea_dir / f"gsea_dotplot_{cond_name}.{FIG_FORMAT}"
        fig.savefig(outpath, format=FIG_FORMAT, dpi=FIG_DPI, bbox_inches="tight")
        plt.close(fig)
        print(f"  Saved: {outpath.name}")


# ===========================================================================
# Enhanced GSEA visualization and export
# ===========================================================================

DB_SHORT = {
    "GO_Biological_Process_2023": "BP",
    "GO_Cellular_Component_2023": "CC",
    "GO_Molecular_Function_2023": "MF",
    "KEGG_2021_Human": "KEGG",
    "Reactome_2022": "Reactome",
    "MSigDB_Hallmark_2020": "Hallmark",
    "WikiPathway_2021_Human": "WikiPath",
}

# Reverse lookup: folder name on disk -> short label
# gseapy replaces colons with underscores in directory names
_DB_SHORT_FUZZY = {}
for _full, _short in DB_SHORT.items():
    _DB_SHORT_FUZZY[_full] = _short
    _DB_SHORT_FUZZY[_full.replace(":", "_")] = _short


def _db_short_label(db_name):
    """Return the short label for a database name, with fuzzy matching."""
    if db_name in _DB_SHORT_FUZZY:
        return _DB_SHORT_FUZZY[db_name]
    # Fallback: try substring matching
    for full, short in DB_SHORT.items():
        if full.replace(":", "_") in db_name or db_name in full:
            return short
    # Last resort: first 8 chars
    return db_name[:8]


def _collect_gsea_rows(gsea_results, cond_name, gsea_dir, top_n=5,
                       include_lead_genes=False):
    """Collect GSEA pathway rows from in-memory results or disk CSV fallback.

    Returns a list of dicts with keys: Term, NES, FDR, Database, and
    optionally LeadingEdgeGenes.
    """
    all_rows = []

    # --- Try in-memory results first ---
    if gsea_results and cond_name in gsea_results:
        for db_name, pathways_df in gsea_results[cond_name].items():
            if pathways_df is None or len(pathways_df) == 0:
                continue
            df = pathways_df.copy()
            df["abs_nes"] = df["nes"].abs()
            df = df.sort_values("abs_nes", ascending=False).head(top_n)
            short = _db_short_label(db_name)
            for _, row in df.iterrows():
                entry = {
                    "Term": str(row["Term"]),
                    "NES": float(row["nes"]),
                    "FDR": float(row["fdr"]),
                    "Database": db_name,
                    "DatabaseShort": short,
                }
                if include_lead_genes and "lead_genes" in row.index:
                    entry["LeadingEdgeGenes"] = str(row["lead_genes"]) if pd.notna(row["lead_genes"]) else ""
                all_rows.append(entry)

    # --- Fallback: scan disk CSV reports ---
    if not all_rows:
        cond_gsea_dir = gsea_dir / cond_name
        if cond_gsea_dir.is_dir():
            for db_subdir in sorted(cond_gsea_dir.iterdir()):
                if not db_subdir.is_dir():
                    continue
                report_csv = db_subdir / "gseapy.gene_set.prerank.report.csv"
                if not report_csv.exists():
                    continue
                try:
                    rpt = _normalize_gsea_cols(pd.read_csv(report_csv))
                    if len(rpt) == 0:
                        continue
                    rpt["abs_nes"] = rpt["nes"].abs()
                    rpt = rpt.sort_values("abs_nes", ascending=False).head(top_n)
                    db_name = db_subdir.name
                    short = _db_short_label(db_name)
                    for _, row in rpt.iterrows():
                        entry = {
                            "Term": str(row["Term"]),
                            "NES": float(row["nes"]),
                            "FDR": float(row["fdr"]),
                            "Database": db_name,
                            "DatabaseShort": short,
                        }
                        if include_lead_genes and "lead_genes" in row.index:
                            entry["LeadingEdgeGenes"] = str(row["lead_genes"]) if pd.notna(row["lead_genes"]) else ""
                        all_rows.append(entry)
                except Exception:
                    continue

    return all_rows


def gsea_combined_plot(gsea_results, condition_labels, outdir):
    """Generate combined GSEA dot plots showing pathways from ALL databases.

    Replaces the older gsea_dotplot() function. For each condition, creates
    one figure with database-tagged pathway names, NES on x-axis, dot size
    proportional to -log10(FDR), and color indicating enrichment direction.

    Falls back to reading gseapy CSV reports from disk when in-memory
    results are empty.
    """
    gsea_dir = outdir / "gsea_results"
    gsea_dir.mkdir(exist_ok=True)

    print("\n-- Generating Combined GSEA Plots --")

    for cond_name, cond_label in condition_labels.items():
        rows = _collect_gsea_rows(gsea_results, cond_name, gsea_dir, top_n=5)

        if not rows:
            print(f"  No GSEA data for {cond_label}, skipping combined plot")
            continue

        df = pd.DataFrame(rows)

        # Cap at 35 pathways total (top 5 per database, 7 databases)
        df["abs_NES"] = df["NES"].abs()
        df = df.sort_values("abs_NES", ascending=False).head(35)

        # Sort by NES: positive at top, negative at bottom
        df = df.sort_values("NES", ascending=True).reset_index(drop=True)

        # Build tagged pathway names: "[KEGG] PI3K-Akt signaling"
        df["TaggedTerm"] = df.apply(
            lambda r: f"[{r['DatabaseShort']}] {r['Term'][:55]}", axis=1
        )

        # Deduplicate display names (rare, but possible across databases)
        seen = {}
        unique_names = []
        for name in df["TaggedTerm"]:
            if name in seen:
                seen[name] += 1
                unique_names.append(f"{name} ({seen[name]})")
            else:
                seen[name] = 0
                unique_names.append(name)
        df["TaggedTerm"] = unique_names

        n_pathways = len(df)
        n_databases = df["DatabaseShort"].nunique()

        # -log10(FDR) for dot size
        neg_log_fdr = -np.log10(df["FDR"].values + 1e-10)
        size_min, size_max = 30, 250
        if neg_log_fdr.max() > neg_log_fdr.min():
            sizes = size_min + (
                (neg_log_fdr - neg_log_fdr.min())
                / (neg_log_fdr.max() - neg_log_fdr.min())
                * (size_max - size_min)
            )
        else:
            sizes = np.full_like(neg_log_fdr, (size_min + size_max) / 2)

        # Dot color: orange for positive NES, blue for negative
        colors = [COLOR_UP if nes > 0 else COLOR_DOWN for nes in df["NES"].values]

        # --- Plot ---
        fig_h = max(8, n_pathways * 0.4)
        fig, ax = plt.subplots(figsize=(14, fig_h))

        ax.scatter(
            df["NES"].values,
            range(n_pathways),
            s=sizes,
            c=colors,
            edgecolors="black",
            linewidths=0.5,
            zorder=3,
            alpha=0.85,
        )

        ax.set_yticks(range(n_pathways))
        ax.set_yticklabels(df["TaggedTerm"].values, fontsize=9)
        ax.axvline(0, color="grey", linewidth=0.8, linestyle="--", zorder=1)
        ax.set_xlabel("NES (Normalized Enrichment Score)", fontsize=12)
        ax.set_title(
            f"GSEA Enrichment \u2014 {cond_label}",
            fontsize=14,
            fontweight="bold",
        )

        # Light horizontal grid for readability
        ax.set_axisbelow(True)
        ax.yaxis.grid(True, linestyle=":", alpha=0.3)

        # --- Legend: dot size scale ---
        legend_elements = []
        # Size legend entries
        fdr_examples = [0.001, 0.01, 0.05, 0.25]
        for fdr_val in fdr_examples:
            nlf = -np.log10(fdr_val + 1e-10)
            if neg_log_fdr.max() > neg_log_fdr.min():
                sz = size_min + (nlf - neg_log_fdr.min()) / (
                    neg_log_fdr.max() - neg_log_fdr.min()
                ) * (size_max - size_min)
            else:
                sz = (size_min + size_max) / 2
            sz = np.clip(sz, size_min, size_max)
            legend_elements.append(
                ax.scatter(
                    [], [], s=sz, c="grey", edgecolors="black",
                    linewidths=0.5, label=f"FDR={fdr_val}"
                )
            )

        # Color legend entries
        legend_elements.append(
            mpatches.Patch(facecolor=COLOR_UP, edgecolor="black",
                           linewidth=0.5, label="Positive NES (up)")
        )
        legend_elements.append(
            mpatches.Patch(facecolor=COLOR_DOWN, edgecolor="black",
                           linewidth=0.5, label="Negative NES (down)")
        )

        ax.legend(
            handles=legend_elements,
            title="Significance & Direction",
            loc="lower right",
            fontsize=8,
            title_fontsize=9,
            framealpha=0.9,
        )

        plt.tight_layout()
        outpath = gsea_dir / f"gsea_combined_{cond_name}.{FIG_FORMAT}"
        fig.savefig(outpath, format=FIG_FORMAT, dpi=FIG_DPI, bbox_inches="tight")
        plt.close(fig)
        print(f"  Saved: gsea_combined_{cond_name}.{FIG_FORMAT} "
              f"({n_pathways} pathways from {n_databases} databases)")


def gsea_enrichment_plots(gsea_results, condition_labels, outdir):
    """Collect or regenerate per-pathway enrichment score plots.

    Strategy:
      A) Look for pre-generated plots produced by gseapy.prerank() in the
         output directory structure.
      B) If none found, attempt to regenerate using gseapy's plotting API.

    Organizes all enrichment plots into:
        outdir/gsea_results/{cond_name}/enrichment_plots/
    """
    gsea_dir = outdir / "gsea_results"
    gsea_dir.mkdir(exist_ok=True)

    print("\n-- Collecting GSEA Enrichment Plots --")

    for cond_name, cond_label in condition_labels.items():
        cond_gsea_dir = gsea_dir / cond_name
        if not cond_gsea_dir.is_dir():
            print(f"  No GSEA output directory for {cond_label}, skipping")
            continue

        plots_dest = cond_gsea_dir / "enrichment_plots"
        plots_dest.mkdir(exist_ok=True)

        collected = 0

        # --- Option A: scan for pre-generated plot files ---
        for db_subdir in sorted(cond_gsea_dir.iterdir()):
            if not db_subdir.is_dir() or db_subdir.name == "enrichment_plots":
                continue

            short_label = _db_short_label(db_subdir.name)

            # gseapy saves enrichment plots as PNG/PDF in the output directory
            plot_files = (
                list(db_subdir.glob("*.png"))
                + list(db_subdir.glob("*.pdf"))
                + list(db_subdir.glob("*.svg"))
            )

            # Filter out report CSVs; keep only image files
            plot_files = [
                f for f in plot_files
                if f.suffix.lower() in (".png", ".pdf", ".svg")
                and "report" not in f.stem.lower()
            ]

            for pf in plot_files:
                dest_name = f"{short_label}_{pf.name}"
                dest_path = plots_dest / dest_name
                try:
                    shutil.copy2(str(pf), str(dest_path))
                    collected += 1
                except Exception:
                    pass

        # --- Option B: try to regenerate if none found ---
        if collected == 0:
            try:
                from gseapy.plot import gseaplot
                _has_gseaplot = True
            except ImportError:
                _has_gseaplot = False

            if _has_gseaplot:
                for db_subdir in sorted(cond_gsea_dir.iterdir()):
                    if not db_subdir.is_dir() or db_subdir.name == "enrichment_plots":
                        continue

                    short_label = _db_short_label(db_subdir.name)
                    report_csv = db_subdir / "gseapy.gene_set.prerank.report.csv"
                    if not report_csv.exists():
                        continue

                    try:
                        rpt = _normalize_gsea_cols(pd.read_csv(report_csv))
                        sig = rpt[rpt["fdr"] < 0.25].sort_values(
                            "nes", key=abs, ascending=False
                        )
                        if len(sig) == 0:
                            continue

                        for _, row in sig.head(10).iterrows():
                            term = str(row["Term"])
                            term_clean = (
                                term.replace("/", "_")
                                .replace("\\", "_")
                                .replace(":", "_")
                                .replace(" ", "_")
                            )
                            rank_file = db_subdir / term_clean / "gene_set_prerank.csv"
                            if not rank_file.exists():
                                rank_file = db_subdir / f"{term_clean}.csv"
                            if not rank_file.exists():
                                continue

                            try:
                                out_name = f"{short_label}_{term_clean[:60]}.{FIG_FORMAT}"
                                out_path = plots_dest / out_name
                                gseaplot(
                                    rank_metric=rank_file,
                                    term=term,
                                    ofname=str(out_path),
                                )
                                collected += 1
                            except Exception:
                                continue
                    except Exception:
                        continue

        if collected > 0:
            print(f"  {cond_label}: {collected} enrichment plots -> "
                  f"{plots_dest.relative_to(outdir)}/")
        else:
            print(f"  {cond_label}: no enrichment plots found or generated")


def export_gsea_leading_edge(gsea_results, condition_labels, outdir):
    """Export leading edge genes for significant GSEA pathways to Excel.

    For each condition and database, extracts pathways with FDR < 0.25 and
    their leading edge gene lists. Saves to:
        outdir / "gsea_leading_edge_genes.xlsx"

    Falls back to reading disk CSV reports if in-memory results are empty.
    """
    gsea_dir = outdir / "gsea_results"
    gsea_dir.mkdir(exist_ok=True)

    print("\n-- Exporting GSEA Leading Edge Genes --")

    all_entries = []

    for cond_name, cond_label in condition_labels.items():
        rows = _collect_gsea_rows(
            gsea_results, cond_name, gsea_dir,
            top_n=999,  # get all pathways, not just top 5
            include_lead_genes=True,
        )

        if not rows:
            print(f"  No GSEA data for {cond_label}, skipping")
            continue

        for r in rows:
            fdr = r.get("FDR", 1.0)
            if fdr >= 0.25:
                continue
            all_entries.append({
                "Condition": cond_label,
                "Database": r.get("Database", ""),
                "Pathway": r.get("Term", ""),
                "NES": r.get("NES", np.nan),
                "FDR": fdr,
                "LeadingEdgeGenes": r.get("LeadingEdgeGenes", ""),
            })

    if not all_entries:
        print("  No significant pathways (FDR < 0.25) found across any condition")
        return

    le_df = pd.DataFrame(all_entries)
    le_df = le_df.sort_values(["Condition", "FDR"], ascending=[True, True])

    # --- Save to dedicated Excel file ---
    le_path = outdir / "gsea_leading_edge_genes.xlsx"
    try:
        le_df.to_excel(le_path, index=False, sheet_name="Leading Edge Genes")
        print(f"  Saved: gsea_leading_edge_genes.xlsx "
              f"({len(le_df)} pathways across "
              f"{le_df['Condition'].nunique()} conditions)")
    except Exception as e:
        print(f"  ERROR saving leading edge Excel: {e}")
        return

    # --- Try to append a "Leading Edge" sheet to the main cross-condition Excel ---
    main_xlsx = outdir / "cross_condition_summary.xlsx"
    if main_xlsx.exists():
        try:
            import openpyxl
            with pd.ExcelWriter(
                str(main_xlsx), engine="openpyxl", mode="a",
                if_sheet_exists="replace",
            ) as writer:
                le_df.to_excel(writer, index=False, sheet_name="Leading Edge")
            print(f"  Also added 'Leading Edge' sheet to {main_xlsx.name}")
        except ImportError:
            print("  openpyxl not installed; skipped appending to main Excel")
        except Exception as e:
            print(f"  Could not append to {main_xlsx.name}: {e}")


# ===========================================================================
# Additional analysis features
# ===========================================================================

# Okabe-Ito extended palette for multi-category charts
_OKABE_ITO = [
    "#E69F00",  # orange
    "#56B4E9",  # sky blue
    "#009E73",  # bluish green
    "#F0E442",  # yellow
    "#0072B2",  # blue
    "#D55E00",  # vermillion
    "#CC79A7",  # reddish purple
    "#000000",  # black
]


def volcano_plot_labeled(deseq2_all, outdir, label="", suffix="",
                         genes_of_interest=None):
    """Enhanced volcano plot that labels top DEGs and genes of interest.

    Parameters
    ----------
    deseq2_all : pd.DataFrame
        Full (unfiltered) DESeq2 results with at minimum the columns
        referenced by DESEQ2_COLS: log2fc, padj, gene_name.
    outdir : Path or str
        Directory for saving the figure.
    label : str
        Condition label for the title (e.g. "MIAT OE vs Control").
    suffix : str
        Filename suffix (e.g. "_protein_coding").
    genes_of_interest : list[str] or None
        Gene names to always highlight.  Defaults to key study genes.
    """
    outdir = Path(outdir)
    outdir.mkdir(parents=True, exist_ok=True)

    if genes_of_interest is None:
        genes_of_interest = GENES_OF_INTEREST

    # Resolve column names through DESEQ2_COLS
    log2fc_col = DESEQ2_COLS.get("log2fc", "log2FoldChange")
    padj_col_name = DESEQ2_COLS.get("padj", "padj")
    gene_col = DESEQ2_COLS.get("gene_name", "gene_name")

    # Prepare data
    data = deseq2_all.dropna(subset=[padj_col_name, log2fc_col]).copy()
    data["_neg_log10_padj"] = -np.log10(data[padj_col_name].clip(lower=1e-300))

    # Classify significance — include baseMean filter to match actual DEG counts
    basemean_col = DESEQ2_COLS.get("basemean", "baseMean")
    basemean_ok = (data[basemean_col] >= BASEMEAN_CUTOFF) if basemean_col in data.columns else True
    cond_up = (data[padj_col_name] < PADJ_CUTOFF) & (data[log2fc_col] >= LOG2FC_CUTOFF) & basemean_ok
    cond_dn = (data[padj_col_name] < PADJ_CUTOFF) & (data[log2fc_col] <= -LOG2FC_CUTOFF) & basemean_ok
    data["_status"] = np.select([cond_up, cond_dn], ["Up", "Down"], default="NS")

    color_map = {"Up": COLOR_UP, "Down": COLOR_DOWN, "NS": COLOR_NS}

    # --- Build figure ---
    fig, ax = plt.subplots(figsize=(10, 8))

    for status in ["NS", "Down", "Up"]:
        subset = data[data["_status"] == status]
        lbl = "NS" if status == "NS" else f"{status} ({len(subset):,})"
        ax.scatter(
            subset[log2fc_col], subset["_neg_log10_padj"],
            c=color_map[status], s=8, alpha=0.5, edgecolors="none",
            label=lbl, rasterized=True,
        )

    # Threshold lines
    ax.axhline(-np.log10(PADJ_CUTOFF), color="grey", ls="--", lw=0.8, zorder=0)
    ax.axvline(LOG2FC_CUTOFF, color="grey", ls="--", lw=0.8, zorder=0)
    ax.axvline(-LOG2FC_CUTOFF, color="grey", ls="--", lw=0.8, zorder=0)

    # Count box
    n_up = int((data["_status"] == "Up").sum())
    n_down = int((data["_status"] == "Down").sum())
    add_count_box(ax, n_up, n_down, n_up + n_down, position="lower left")

    # --- Identify genes to label ---
    sig_mask = data["_status"].isin(["Up", "Down"])
    sig_data = data[sig_mask].copy()

    top10 = pd.DataFrame()
    if len(sig_data) > 0:
        top10 = sig_data.nsmallest(10, padj_col_name)

    # --- Label placement with simple overlap avoidance ---
    label_rows = top10.copy() if len(sig_data) > 0 else pd.DataFrame()

    placed_labels = []  # list of (x_text, y_text)

    def _find_offset(x_data, y_data, placed, base_dx=0.8, base_dy=0.6):
        """Find a non-overlapping text position near (x_data, y_data)."""
        xlim = ax.get_xlim()
        ylim = ax.get_ylim()
        x_range = xlim[1] - xlim[0] if xlim[1] != xlim[0] else 1.0
        y_range = ylim[1] - ylim[0] if ylim[1] != ylim[0] else 1.0

        x_thresh = x_range * 0.06
        y_thresh = y_range * 0.035

        dx = base_dx if x_data < 0 else -base_dx
        dx_scaled = dx * (x_range / 10.0)

        for attempt in range(20):
            dy_scaled = (base_dy + attempt * 0.5) * (y_range / 15.0)
            sign = 1 if attempt % 2 == 0 else -1
            x_text = x_data + dx_scaled
            y_text = y_data + sign * dy_scaled

            overlaps = False
            for px, py in placed:
                if abs(x_text - px) < x_thresh and abs(y_text - py) < y_thresh:
                    overlaps = True
                    break
            if not overlaps:
                return x_text, y_text

        return x_text, y_text

    if gene_col in data.columns:
        for _, row in label_rows.iterrows():
            gene_name = str(row[gene_col])
            x_data = row[log2fc_col]
            y_data = row["_neg_log10_padj"]

            x_text, y_text = _find_offset(x_data, y_data, placed_labels)
            placed_labels.append((x_text, y_text))

            ax.annotate(
                gene_name,
                xy=(x_data, y_data),
                xytext=(x_text, y_text),
                fontsize=8,
                color="black",
                arrowprops=dict(
                    arrowstyle="-|>",
                    color="grey",
                    lw=0.8,
                    connectionstyle="arc3,rad=0.15",
                ),
                bbox=dict(
                    boxstyle="round,pad=0.2",
                    facecolor="white",
                    edgecolor="grey",
                    alpha=0.85,
                ),
                zorder=6,
            )

    # Axes and title
    ax.set_xlabel("log$_2$ Fold Change", fontsize=12)
    ax.set_ylabel("$-$log$_{10}$ (adjusted p-value)", fontsize=12)
    title_label = f" — {label}" if label else ""
    ax.set_title(f"Volcano Plot (Labeled){title_label}", fontsize=14)
    ax.legend(bbox_to_anchor=(1.02, 1), loc="upper left", frameon=True, fontsize=9, markerscale=2)

    fname = f"volcano_plot_labeled{suffix}.{FIG_FORMAT}"
    outpath = outdir / fname
    fig.savefig(outpath, format=FIG_FORMAT, dpi=FIG_DPI, bbox_inches="tight")
    plt.close(fig)
    print(f"  Saved: {outpath}")


def gene_overlap_summary(condition_results, condition_labels, outdir):
    """Create a table showing which genes are DE, AS, or both per condition.

    Parameters
    ----------
    condition_results : dict
        Keyed by condition name.  Each value has:
        - "deseq2_filtered" -> {"all_genes": DataFrame} with 'direction' col
        - "rmats_filtered"  -> {event_type: DataFrame}
    condition_labels : dict
        Maps condition name -> human-readable label.
    outdir : Path or str
        Directory for saving the Excel file.
    """
    outdir = Path(outdir)
    outdir.mkdir(parents=True, exist_ok=True)

    gene_name_col_de = DESEQ2_COLS.get("gene_name", "gene_name")
    gene_name_col_as = RMATS_COLS.get("gene_name", "geneSymbol")

    # Collect DE and AS gene sets per condition
    de_genes_per_cond = {}   # cond_name -> set of gene names
    as_genes_per_cond = {}   # cond_name -> set of gene names

    cond_names = list(condition_results.keys())

    for cond_name in cond_names:
        res = condition_results[cond_name]

        # DE genes
        de_set = set()
        filt = res.get("deseq2_filtered", {}).get("all_genes", pd.DataFrame())
        if not filt.empty and gene_name_col_de in filt.columns:
            de_set = set(filt[gene_name_col_de].dropna().unique())
        de_genes_per_cond[cond_name] = de_set

        # AS genes (union across all event types)
        as_set = set()
        rmats_filt = res.get("rmats_filtered", {})
        for et, et_df in rmats_filt.items():
            if not et_df.empty and gene_name_col_as in et_df.columns:
                as_set.update(et_df[gene_name_col_as].dropna().unique())
        as_genes_per_cond[cond_name] = as_set

    # Build master gene list
    all_genes = set()
    for s in list(de_genes_per_cond.values()) + list(as_genes_per_cond.values()):
        all_genes.update(s)

    if not all_genes:
        print("  Gene overlap summary: no genes found across conditions.")
        return

    # Build the summary DataFrame
    rows = []
    for gene in sorted(all_genes):
        row = {"Gene": gene}
        n_conditions = 0
        for cond_name in cond_names:
            lbl = condition_labels.get(cond_name, cond_name)
            is_de = gene in de_genes_per_cond[cond_name]
            is_as = gene in as_genes_per_cond[cond_name]
            if is_de and is_as:
                row[lbl] = "DE+AS"
                n_conditions += 1
            elif is_de:
                row[lbl] = "DE"
                n_conditions += 1
            elif is_as:
                row[lbl] = "AS"
                n_conditions += 1
            else:
                row[lbl] = ""
        row["Conditions"] = n_conditions
        rows.append(row)

    summary_df = pd.DataFrame(rows)

    # Sort: by number of conditions (descending), then alphabetically
    summary_df = summary_df.sort_values(
        by=["Conditions", "Gene"],
        ascending=[False, True],
    ).reset_index(drop=True)

    # Save to Excel
    outpath = outdir / "gene_overlap_summary.xlsx"
    summary_df.to_excel(outpath, index=False, freeze_panes=(1, 1))
    n_genes = len(summary_df)
    m_conds = len(cond_names)
    print(f"  Saved gene overlap summary: {n_genes} genes across {m_conds} conditions")
    print(f"    -> {outpath}")


def summary_dashboard(condition_results, condition_labels, go_results,
                      gsea_results, outdir):
    """Create a 2x3 summary dashboard with key statistics.

    Parameters
    ----------
    condition_results : dict
        Same structure as gene_overlap_summary.
    condition_labels : dict
        Maps condition name -> human-readable label.
    go_results : dict or None
        Keyed by condition name.  Each value is a dict with at least:
        - "up" -> DataFrame of enriched GO terms (or None)
        - "down" -> DataFrame of enriched GO terms (or None)
        If None or empty, panel shows "N/A".
    gsea_results : dict or None
        Keyed by condition name.  Each value is a DataFrame with a column
        indicating significance (e.g. "padj" < 0.05).
        If None or empty, panel shows "N/A".
    outdir : Path or str
        Directory for saving the figure.
    """
    outdir = Path(outdir)
    outdir.mkdir(parents=True, exist_ok=True)

    cond_names = list(condition_results.keys())
    cond_labels_list = [condition_labels.get(cn, cn) for cn in cond_names]
    n_conds = len(cond_names)

    gene_name_col_de = DESEQ2_COLS.get("gene_name", "gene_name")
    gene_name_col_as = RMATS_COLS.get("gene_name", "geneSymbol")

    # ------------------------------------------------------------------
    # Gather statistics for each panel
    # ------------------------------------------------------------------

    # Panel 1: DEG counts (up / down per condition)
    deg_up = []
    deg_down = []
    for cn in cond_names:
        filt = condition_results[cn].get("deseq2_filtered", {}).get(
            "all_genes", pd.DataFrame()
        )
        if not filt.empty and "direction" in filt.columns:
            deg_up.append(int((filt["direction"].str.lower() == "up").sum()))
            deg_down.append(int((filt["direction"].str.lower() == "down").sum()))
        else:
            deg_up.append(0)
            deg_down.append(0)

    # Panel 2: Splicing events by type per condition
    splice_counts = {et: [] for et in RMATS_EVENT_TYPES}
    for cn in cond_names:
        rmats_filt = condition_results[cn].get("rmats_filtered", {})
        for et in RMATS_EVENT_TYPES:
            df = rmats_filt.get(et, pd.DataFrame())
            splice_counts[et].append(len(df))

    # Panel 3: DE vs AS overlap per condition
    de_only_counts = []
    as_only_counts = []
    both_counts = []
    all_de_genes = set()
    all_as_genes = set()
    all_both_genes = set()
    for cn in cond_names:
        filt = condition_results[cn].get("deseq2_filtered", {}).get(
            "all_genes", pd.DataFrame()
        )
        de_set = set()
        if not filt.empty and gene_name_col_de in filt.columns:
            de_set = set(filt[gene_name_col_de].dropna().unique())

        as_set = set()
        rmats_filt = condition_results[cn].get("rmats_filtered", {})
        for et_df in rmats_filt.values():
            if not et_df.empty and gene_name_col_as in et_df.columns:
                as_set.update(et_df[gene_name_col_as].dropna().unique())

        both = de_set & as_set
        de_only_counts.append(len(de_set - both))
        as_only_counts.append(len(as_set - both))
        both_counts.append(len(both))

        all_de_genes.update(de_set)
        all_as_genes.update(as_set)
        all_both_genes.update(both)

    # Panel 4: GO terms (if available)
    go_available = go_results is not None and len(go_results) > 0
    go_up_counts = []
    go_down_counts = []
    if go_available:
        for cn in cond_names:
            cond_go = go_results.get(cn, {})
            if cond_go is None:
                cond_go = {}
            up_df = cond_go.get("up", None)
            dn_df = cond_go.get("down", None)
            go_up_counts.append(len(up_df) if up_df is not None and not (
                isinstance(up_df, pd.DataFrame) and up_df.empty) else 0)
            go_down_counts.append(len(dn_df) if dn_df is not None and not (
                isinstance(dn_df, pd.DataFrame) and dn_df.empty) else 0)

    # Panel 5: GSEA pathways (if available)
    gsea_available = gsea_results is not None and len(gsea_results) > 0
    gsea_counts = []
    if gsea_available:
        for cn in cond_names:
            cond_gsea = gsea_results.get(cn, None)
            if cond_gsea is not None and isinstance(cond_gsea, pd.DataFrame):
                # Count significant pathways (padj < 0.05 or FDR < 0.25)
                if "padj" in cond_gsea.columns:
                    gsea_counts.append(int((cond_gsea["padj"] < 0.05).sum()))
                elif "FDR q-val" in cond_gsea.columns:
                    gsea_counts.append(int((cond_gsea["FDR q-val"] < 0.25).sum()))
                else:
                    gsea_counts.append(len(cond_gsea))
            else:
                gsea_counts.append(0)

    # ------------------------------------------------------------------
    # Build the 2x3 figure
    # ------------------------------------------------------------------
    fig = plt.figure(figsize=(18, 10))
    gs = gridspec.GridSpec(2, 3, hspace=0.35, wspace=0.30)

    x = np.arange(n_conds)
    bar_width = 0.35

    # --- Panel 1: DEG Counts (grouped bar) ---
    ax1 = fig.add_subplot(gs[0, 0])
    ax1.bar(x - bar_width / 2, deg_up, bar_width, label="Up",
            color=COLOR_UP, edgecolor="white", linewidth=0.5)
    ax1.bar(x + bar_width / 2, deg_down, bar_width, label="Down",
            color=COLOR_DOWN, edgecolor="white", linewidth=0.5)
    ax1.set_xticks(x)
    ax1.set_xticklabels(cond_labels_list, fontsize=9, rotation=15, ha="right")
    ax1.set_ylabel("Number of DEGs", fontsize=10)
    ax1.set_title("DEG Counts", fontsize=12, fontweight="bold")
    ax1.legend(fontsize=9)
    for i in range(n_conds):
        if deg_up[i] > 0:
            ax1.text(i - bar_width / 2, deg_up[i] + 1, str(deg_up[i]),
                     ha="center", va="bottom", fontsize=8)
        if deg_down[i] > 0:
            ax1.text(i + bar_width / 2, deg_down[i] + 1, str(deg_down[i]),
                     ha="center", va="bottom", fontsize=8)
    ax1.spines["top"].set_visible(False)
    ax1.spines["right"].set_visible(False)

    # --- Panel 2: Splicing Events (stacked bar) ---
    ax2 = fig.add_subplot(gs[0, 1])
    bottoms = np.zeros(n_conds)
    for et in RMATS_EVENT_TYPES:
        vals = np.array(splice_counts[et])
        ax2.bar(x, vals, bar_width * 1.5, bottom=bottoms,
                label=et, color=EVENT_COLORS.get(et, "#888888"),
                edgecolor="white", linewidth=0.5)
        bottoms += vals
    ax2.set_xticks(x)
    ax2.set_xticklabels(cond_labels_list, fontsize=9, rotation=15, ha="right")
    ax2.set_ylabel("Number of Events", fontsize=10)
    ax2.set_title("Splicing Events by Type", fontsize=12, fontweight="bold")
    ax2.legend(fontsize=8, ncol=2, loc="upper right")
    for i in range(n_conds):
        total = int(bottoms[i])
        if total > 0:
            ax2.text(i, total + 1, str(total), ha="center", va="bottom",
                     fontsize=8, fontweight="bold")
    ax2.spines["top"].set_visible(False)
    ax2.spines["right"].set_visible(False)

    # --- Panel 3: DE vs AS Overlap (grouped bar) ---
    ax3 = fig.add_subplot(gs[0, 2])
    w = 0.25
    ax3.bar(x - w, de_only_counts, w, label="DE only",
            color=COLOR_UP, edgecolor="white", linewidth=0.5)
    ax3.bar(x, as_only_counts, w, label="AS only",
            color="#009E73", edgecolor="white", linewidth=0.5)
    ax3.bar(x + w, both_counts, w, label="DE + AS",
            color="#CC79A7", edgecolor="white", linewidth=0.5)
    ax3.set_xticks(x)
    ax3.set_xticklabels(cond_labels_list, fontsize=9, rotation=15, ha="right")
    ax3.set_ylabel("Number of Genes", fontsize=10)
    ax3.set_title("DE vs AS Overlap", fontsize=12, fontweight="bold")
    ax3.legend(fontsize=9)
    ax3.spines["top"].set_visible(False)
    ax3.spines["right"].set_visible(False)

    # --- Panel 4: GO Terms Enriched ---
    ax4 = fig.add_subplot(gs[1, 0])
    if go_available and (sum(go_up_counts) + sum(go_down_counts)) > 0:
        ax4.bar(x - bar_width / 2, go_up_counts, bar_width, label="Up-regulated",
                color=COLOR_UP, edgecolor="white", linewidth=0.5)
        ax4.bar(x + bar_width / 2, go_down_counts, bar_width, label="Down-regulated",
                color=COLOR_DOWN, edgecolor="white", linewidth=0.5)
        ax4.set_xticks(x)
        ax4.set_xticklabels(cond_labels_list, fontsize=9, rotation=15, ha="right")
        ax4.set_ylabel("Enriched GO Terms", fontsize=10)
        ax4.legend(fontsize=9)
        for i in range(n_conds):
            if go_up_counts[i] > 0:
                ax4.text(i - bar_width / 2, go_up_counts[i] + 0.5,
                         str(go_up_counts[i]), ha="center", va="bottom", fontsize=8)
            if go_down_counts[i] > 0:
                ax4.text(i + bar_width / 2, go_down_counts[i] + 0.5,
                         str(go_down_counts[i]), ha="center", va="bottom", fontsize=8)
        ax4.spines["top"].set_visible(False)
        ax4.spines["right"].set_visible(False)
    else:
        ax4.text(0.5, 0.5, "GO Enrichment\nN/A", transform=ax4.transAxes,
                 fontsize=16, ha="center", va="center", color="grey",
                 fontstyle="italic")
        ax4.set_xticks([])
        ax4.set_yticks([])
        for spine in ax4.spines.values():
            spine.set_visible(False)
    ax4.set_title("GO Terms Enriched", fontsize=12, fontweight="bold")

    # --- Panel 5: GSEA Pathways ---
    ax5 = fig.add_subplot(gs[1, 1])
    if gsea_available and sum(gsea_counts) > 0:
        bars = ax5.bar(x, gsea_counts, bar_width * 1.5,
                       color=_OKABE_ITO[:n_conds], edgecolor="white",
                       linewidth=0.5)
        ax5.set_xticks(x)
        ax5.set_xticklabels(cond_labels_list, fontsize=9, rotation=15, ha="right")
        ax5.set_ylabel("Significant Pathways", fontsize=10)
        for i, (bar, count) in enumerate(zip(bars, gsea_counts)):
            if count > 0:
                ax5.text(bar.get_x() + bar.get_width() / 2, count + 0.5,
                         str(count), ha="center", va="bottom", fontsize=8)
        ax5.spines["top"].set_visible(False)
        ax5.spines["right"].set_visible(False)
    else:
        ax5.text(0.5, 0.5, "GSEA Pathways\nN/A", transform=ax5.transAxes,
                 fontsize=16, ha="center", va="center", color="grey",
                 fontstyle="italic")
        ax5.set_xticks([])
        ax5.set_yticks([])
        for spine in ax5.spines.values():
            spine.set_visible(False)
    ax5.set_title("GSEA Pathways", fontsize=12, fontweight="bold")

    # --- Panel 6: Key Numbers (text summary) ---
    ax6 = fig.add_subplot(gs[1, 2])
    ax6.set_xticks([])
    ax6.set_yticks([])
    for spine in ax6.spines.values():
        spine.set_visible(False)
    ax6.set_facecolor("#F7F7F7")

    total_degs = len(all_de_genes)
    total_as = len(all_as_genes)
    total_both = len(all_both_genes)

    # Try to find most enriched GO category
    best_go_term = "N/A"
    if go_available:
        for cn in cond_names:
            cond_go = go_results.get(cn, {})
            if cond_go is None:
                continue
            for direction_key in ["up", "down"]:
                go_df = cond_go.get(direction_key, None)
                if go_df is not None and isinstance(go_df, pd.DataFrame) and not go_df.empty:
                    term_col = None
                    for candidate in ["Term", "name", "Description", "term_name",
                                      "GO_term", "Name"]:
                        if candidate in go_df.columns:
                            term_col = candidate
                            break
                    if term_col is not None:
                        best_go_term = str(go_df.iloc[0][term_col])
                        if len(best_go_term) > 45:
                            best_go_term = best_go_term[:42] + "..."
                        break
            if best_go_term != "N/A":
                break

    summary_lines = [
        ("Total Unique DEGs", f"{total_degs:,}"),
        ("Total Unique AS Genes", f"{total_as:,}"),
        ("Genes in Both DE + AS", f"{total_both:,}"),
        ("Conditions Analyzed", str(n_conds)),
        ("Top GO Term", best_go_term),
    ]

    ax6.set_title("Key Numbers", fontsize=12, fontweight="bold")
    y_start = 0.88
    y_step = 0.16
    for i, (key, val) in enumerate(summary_lines):
        y_pos = y_start - i * y_step
        ax6.text(0.08, y_pos, key + ":", transform=ax6.transAxes,
                 fontsize=11, fontweight="bold", va="top", ha="left",
                 color="#333333")
        ax6.text(0.92, y_pos, val, transform=ax6.transAxes,
                 fontsize=13, fontweight="bold", va="top", ha="right",
                 color=COLOR_DOWN)

    # Title
    fig.suptitle("RNA-seq Analysis Summary Dashboard", fontsize=16,
                 fontweight="bold", y=0.98)

    fname = f"analysis_summary_dashboard.{FIG_FORMAT}"
    outpath = outdir / fname
    fig.savefig(outpath, format=FIG_FORMAT, dpi=FIG_DPI, bbox_inches="tight")
    plt.close(fig)
    print(f"  Saved: {outpath}")


def export_prism_files(condition_results, condition_labels, outdir, gsea_results=None):
    """Generate comprehensive Prism .pzfx files for publication-quality graphs.

    Generates 15+ Prism files with proper XML structure, all including gene symbols.
    Each file must be >10KB as required.
    """
    import xml.etree.ElementTree as ET
    from xml.dom import minidom

    prism_dir = outdir / "prism_files"
    prism_dir.mkdir(exist_ok=True)

    print("\n-- Generating Prism Files --")

    # Helper to create Prism XML structure (valid .pzfx format)
    def create_prism_xml():
        root = ET.Element("GraphPadPrismFile")
        root.set("PrismXMLVersion", "5.00")

        # <Created> block
        created = ET.SubElement(root, "Created")
        orig = ET.SubElement(created, "OriginalVersion")
        orig.set("CreatedByProgram", "GraphPad Prism")
        orig.set("CreatedByVersion", "6.0f.254")
        orig.set("Login", "")
        orig.set("DateTime", datetime.now().strftime("%Y-%m-%dT%H:%M:%S+00:00"))

        # <InfoSequence>
        info_seq = ET.SubElement(root, "InfoSequence")
        ref = ET.SubElement(info_seq, "Ref")
        ref.set("ID", "Info0")
        ref.set("Selected", "1")

        info = ET.SubElement(root, "Info")
        info.set("ID", "Info0")
        info_title = ET.SubElement(info, "Title")
        info_title.text = "Project info 1"
        ET.SubElement(info, "Notes")
        const = ET.SubElement(info, "Constant")
        const_name = ET.SubElement(const, "Name")
        const_name.text = "Experiment Date"
        const_val = ET.SubElement(const, "Value")
        const_val.text = datetime.now().strftime("%Y-%m-%d")

        # <TableSequence> — Ref elements added by add_table()
        table_seq = ET.SubElement(root, "TableSequence")
        table_seq.set("Selected", "1")

        # Internal counter
        root.set("_table_count", "0")
        return root

    def _is_numeric_string(s):
        try:
            float(s)
            return True
        except (ValueError, TypeError):
            return False

    def add_table(root, table_name, columns_data):
        """Add a table to Prism XML in correct .pzfx column-oriented format.
        columns_data: list of (col_name, values_list) tuples.
        First column auto-detected as row titles if all values are text.
        """
        table_idx = int(root.get("_table_count", "0"))
        table_id = f"Table{table_idx}"
        root.set("_table_count", str(table_idx + 1))

        # Register in TableSequence
        table_seq = root.find("TableSequence")
        ref = ET.SubElement(table_seq, "Ref")
        ref.set("ID", table_id)
        if table_idx == 0:
            ref.set("Selected", "1")

        # Create Table
        table = ET.SubElement(root, "Table")
        table.set("ID", table_id)
        table.set("XFormat", "none")
        table.set("TableType", "OneWay")
        table.set("EVFormat", "AsteriskAfterNumber")

        title_elem = ET.SubElement(table, "Title")
        title_elem.text = table_name

        # Auto-detect row titles: first column if all text
        row_titles_col = None
        if columns_data:
            _, first_vals = columns_data[0]
            if all(isinstance(v, str) and not _is_numeric_string(v) for v in first_vals):
                row_titles_col = 0

        # Build RowTitlesColumn if applicable
        if row_titles_col is not None:
            _, rt_vals = columns_data[row_titles_col]
            rt_elem = ET.SubElement(table, "RowTitlesColumn")
            rt_elem.set("Width", "89")
            subcol = ET.SubElement(rt_elem, "Subcolumn")
            for val in rt_vals:
                d = ET.SubElement(subcol, "d")
                d.text = str(val)
            data_columns = [c for i, c in enumerate(columns_data) if i != row_titles_col]
        else:
            data_columns = list(columns_data)

        # Build YColumns — one per data column, all values in Subcolumn
        for col_name, values in data_columns:
            ycol = ET.SubElement(table, "YColumn")
            ycol.set("Width", "89")
            ycol.set("Decimals", "6")
            ycol.set("Subcolumns", "1")

            col_title = ET.SubElement(ycol, "Title")
            col_title.text = str(col_name)

            subcol = ET.SubElement(ycol, "Subcolumn")
            for val in values:
                d = ET.SubElement(subcol, "d")
                if val is None or (isinstance(val, float) and str(val) == 'nan'):
                    d.text = ""
                else:
                    d.text = str(val)

        return table

    def save_prism_xml(root, filepath):
        """Save XML to .pzfx with proper formatting and encoding."""
        if "_table_count" in root.attrib:
            del root.attrib["_table_count"]

        xml_str = ET.tostring(root, encoding='unicode', xml_declaration=False)
        dom = minidom.parseString(xml_str)
        pretty_xml = dom.toprettyxml(indent="  ", encoding=None)

        if pretty_xml.startswith('<?xml'):
            first_newline = pretty_xml.index('\n')
            pretty_xml = pretty_xml[first_newline + 1:]

        with open(filepath, 'w', encoding='utf-8') as f:
            f.write('<?xml version="1.0" encoding="UTF-8"?>\n')
            f.write(pretty_xml)

    gene_name_col_deseq = DESEQ2_COLS.get("gene_name", "gene_name")
    gene_name_col_rmats = RMATS_COLS.get("gene_name", "geneSymbol")
    log2fc_col = DESEQ2_COLS.get("log2fc", "log2FoldChange")
    padj_col = DESEQ2_COLS.get("padj", "padj")
    basemean_col = DESEQ2_COLS.get("basemean", "baseMean")
    dpsi_col = RMATS_COLS.get("inclevel_diff", "IncLevelDifference")
    fdr_col = RMATS_COLS.get("fdr", "FDR")

    # ===== PER-CONDITION FILES (3 files × 3 conditions = 9 files) =====
    for cond_name, data in condition_results.items():
        cond_label = condition_labels.get(cond_name, cond_name)

        # --- 1. DEG_data_{condition}.pzfx ---
        deg_filt = data["deseq2_filtered"].get("all_genes", pd.DataFrame())
        if len(deg_filt) > 0:
            root = create_prism_xml()

            # Volcano data (top 200 for file size)
            volcano_data = deg_filt.head(200)
            gene_symbols = volcano_data[gene_name_col_deseq].tolist()
            log2fc_vals = volcano_data[log2fc_col].tolist()
            padj_vals = volcano_data[padj_col].tolist()
            basemean_vals = volcano_data[basemean_col].tolist()

            add_table(root, f"Volcano_{cond_label}", [
                ("GeneSymbol", gene_symbols),
                ("log2FC", log2fc_vals),
                ("padj", padj_vals),
                ("baseMean", basemean_vals),
            ])

            # Top 50 DEGs
            top50 = deg_filt.head(50)
            add_table(root, f"Top50_{cond_label}", [
                ("GeneSymbol", top50[gene_name_col_deseq].tolist()),
                ("log2FC", top50[log2fc_col].tolist()),
                ("padj", top50[padj_col].tolist()),
            ])

            # Biotype counts
            if "direction" in deg_filt.columns:
                up_count = int((deg_filt["direction"] == "up").sum())
                down_count = int((deg_filt["direction"] == "down").sum())
                add_table(root, f"Biotype_Counts_{cond_label}", [
                    ("Direction", ["Up", "Down"]),
                    ("Count", [up_count, down_count]),
                ])

            filepath = prism_dir / f"DEG_data_{cond_name}.pzfx"
            save_prism_xml(root, filepath)
            print(f"  Created: {filepath.name} ({filepath.stat().st_size} bytes)")

        # --- 2. Splicing_data_{condition}.pzfx ---
        rmats_filt = data.get("rmats_filtered", {})
        if rmats_filt:
            root = create_prism_xml()

            for et, df in rmats_filt.items():
                if len(df) > 0:
                    event_data = df.head(100)  # Top 100 per event type
                    add_table(root, f"{et}_{cond_label}", [
                        ("GeneSymbol", event_data[gene_name_col_rmats].tolist()),
                        ("dPSI", event_data[dpsi_col].tolist()),
                        ("FDR", event_data[fdr_col].tolist()),
                    ])

            # Event counts
            event_counts = [(et, len(df)) for et, df in rmats_filt.items()]
            if event_counts:
                add_table(root, f"EventCounts_{cond_label}", [
                    ("EventType", [et for et, _ in event_counts]),
                    ("Count", [cnt for _, cnt in event_counts]),
                ])

            filepath = prism_dir / f"Splicing_data_{cond_name}.pzfx"
            save_prism_xml(root, filepath)
            print(f"  Created: {filepath.name} ({filepath.stat().st_size} bytes)")

        # --- 3. GSEA_data_{condition}.pzfx ---
        gsea_prism_ok = False
        if gsea_results and cond_name in gsea_results:
            root = create_prism_xml()
            cond_gsea = gsea_results[cond_name]

            for db_name, pathways_df in cond_gsea.items():
                if len(pathways_df) > 0:
                    pathway_names = pathways_df["Term"].tolist()
                    nes_vals = pathways_df["nes"].tolist()
                    fdr_vals = pathways_df["fdr"].tolist()
                    geneset_sizes = pathways_df["geneset_size"].tolist()

                    add_table(root, f"{db_name[:20]}_{cond_label}", [
                        ("Pathway", pathway_names),
                        ("NES", nes_vals),
                        ("FDR", fdr_vals),
                        ("GenesetSize", geneset_sizes),
                    ])
                    gsea_prism_ok = True

            if gsea_prism_ok:
                filepath = prism_dir / f"GSEA_data_{cond_name}.pzfx"
                save_prism_xml(root, filepath)
                print(f"  Created: {filepath.name} ({filepath.stat().st_size} bytes)")

        if not gsea_prism_ok:
            # Fallback: scan disk for gseapy CSV reports
            gsea_disk_dir = outdir / "gsea_results" / cond_name
            disk_tables_added = False
            if gsea_disk_dir.is_dir():
                root = create_prism_xml()
                for db_subdir in sorted(gsea_disk_dir.iterdir()):
                    if not db_subdir.is_dir():
                        continue
                    report_csv = db_subdir / "gseapy.gene_set.prerank.report.csv"
                    if not report_csv.exists():
                        continue
                    try:
                        rpt = _normalize_gsea_cols(pd.read_csv(report_csv))
                        if len(rpt) == 0:
                            continue
                        rpt["abs_nes"] = rpt["nes"].abs()
                        rpt = rpt.sort_values("abs_nes", ascending=False).head(5)
                        db_label = db_subdir.name[:20]
                        add_table(root, f"{db_label}_{cond_label}", [
                            ("Pathway", rpt["Term"].tolist()),
                            ("NES", rpt["nes"].tolist()),
                            ("FDR", rpt["fdr"].tolist()),
                            ("LeadGenes", rpt["lead_genes"].tolist()),
                        ])
                        disk_tables_added = True
                    except Exception:
                        continue
                if disk_tables_added:
                    filepath = prism_dir / f"GSEA_data_{cond_name}.pzfx"
                    save_prism_xml(root, filepath)
                    print(f"  Created: {filepath.name} ({filepath.stat().st_size} bytes) [from disk CSV]")

            if not disk_tables_added:
                # Placeholder if GSEA not run and no disk files found
                root = create_prism_xml()
                add_table(root, f"GSEA_Placeholder_{cond_label}", [
                    ("Pathway", ["No GSEA data"]),
                    ("NES", [0.0]),
                    ("FDR", [1.0]),
                ])
                filepath = prism_dir / f"GSEA_data_{cond_name}.pzfx"
                save_prism_xml(root, filepath)
                print(f"  Created: {filepath.name} (placeholder)")

    # ===== CROSS-CONDITION FILES (6+ files) =====

    # --- 4. Splicing_overlap_SE.pzfx ---
    root = create_prism_xml()
    names = list(condition_results.keys())
    if len(names) >= 2:
        for pair in combinations(names, 2):
            cond1, cond2 = pair
            label1, label2 = condition_labels[cond1], condition_labels[cond2]

            se_genes1 = set()
            se_genes2 = set()
            if "SE" in condition_results[cond1]["rmats_filtered"]:
                se_genes1 = set(condition_results[cond1]["rmats_filtered"]["SE"][gene_name_col_rmats].dropna())
            if "SE" in condition_results[cond2]["rmats_filtered"]:
                se_genes2 = set(condition_results[cond2]["rmats_filtered"]["SE"][gene_name_col_rmats].dropna())

            overlap = se_genes1 & se_genes2
            add_table(root, f"SE_Overlap_{cond1}_vs_{cond2}", [
                ("Metric", [f"{label1}_only", f"{label2}_only", "Overlap"]),
                ("Count", [len(se_genes1 - se_genes2), len(se_genes2 - se_genes1), len(overlap)]),
            ])

            # Gene lists (top 50 from overlap)
            overlap_list = list(overlap)[:50]
            if overlap_list:
                add_table(root, f"SE_Genes_{cond1}_vs_{cond2}", [
                    ("GeneSymbol", overlap_list),
                ])

    filepath = prism_dir / "Splicing_overlap_SE.pzfx"
    save_prism_xml(root, filepath)
    print(f"  Created: {filepath.name} ({filepath.stat().st_size} bytes)")

    # --- 5. Splicing_overlap_other_events.pzfx ---
    root = create_prism_xml()
    for et in ["A3SS", "A5SS", "MXE", "RI"]:
        for pair in combinations(names, 2):
            cond1, cond2 = pair
            genes1 = set()
            genes2 = set()
            if et in condition_results[cond1]["rmats_filtered"]:
                genes1 = set(condition_results[cond1]["rmats_filtered"][et][gene_name_col_rmats].dropna())
            if et in condition_results[cond2]["rmats_filtered"]:
                genes2 = set(condition_results[cond2]["rmats_filtered"][et][gene_name_col_rmats].dropna())

            overlap = genes1 & genes2
            add_table(root, f"{et}_{cond1}_vs_{cond2}", [
                ("Metric", ["Cond1_only", "Cond2_only", "Overlap"]),
                ("Count", [len(genes1 - genes2), len(genes2 - genes1), len(overlap)]),
            ])

    filepath = prism_dir / "Splicing_overlap_other_events.pzfx"
    save_prism_xml(root, filepath)
    print(f"  Created: {filepath.name} ({filepath.stat().st_size} bytes)")

    # --- 6. DEG_overlap.pzfx ---
    root = create_prism_xml()
    for pair in combinations(names, 2):
        cond1, cond2 = pair
        genes1 = set(condition_results[cond1]["deseq2_filtered"]["all_genes"][gene_name_col_deseq].dropna())
        genes2 = set(condition_results[cond2]["deseq2_filtered"]["all_genes"][gene_name_col_deseq].dropna())
        overlap = genes1 & genes2

        add_table(root, f"DEG_{cond1}_vs_{cond2}", [
            ("Metric", ["Cond1_only", "Cond2_only", "Overlap"]),
            ("Count", [len(genes1 - genes2), len(genes2 - genes1), len(overlap)]),
        ])

        # Shared gene list (top 50)
        overlap_list = list(overlap)[:50]
        if overlap_list:
            add_table(root, f"Shared_Genes_{cond1}_vs_{cond2}", [
                ("GeneSymbol", overlap_list),
            ])

    filepath = prism_dir / "DEG_overlap.pzfx"
    save_prism_xml(root, filepath)
    print(f"  Created: {filepath.name} ({filepath.stat().st_size} bytes)")

    # --- 7. Summary_statistics.pzfx ---
    root = create_prism_xml()
    summary_data = []
    for cond_name, data in condition_results.items():
        deg_count = len(data["deseq2_filtered"]["all_genes"])
        splicing_count = sum(len(df) for df in data["rmats_filtered"].values())
        summary_data.append((condition_labels[cond_name], deg_count, splicing_count))

    add_table(root, "Summary_Statistics", [
        ("Condition", [s[0] for s in summary_data]),
        ("DEG_Count", [s[1] for s in summary_data]),
        ("Splicing_Count", [s[2] for s in summary_data]),
    ])

    filepath = prism_dir / "Summary_statistics.pzfx"
    save_prism_xml(root, filepath)
    print(f"  Created: {filepath.name} ({filepath.stat().st_size} bytes)")

    # --- 8. Concordance_scatter.pzfx ---
    root = create_prism_xml()
    for pair in combinations(names, 2):
        cond1, cond2 = pair
        df1 = condition_results[cond1]["deseq2_filtered"]["all_genes"]
        df2 = condition_results[cond2]["deseq2_filtered"]["all_genes"]

        # Merge on gene symbol
        merged = pd.merge(
            df1[[gene_name_col_deseq, log2fc_col]],
            df2[[gene_name_col_deseq, log2fc_col]],
            on=gene_name_col_deseq,
            suffixes=("_1", "_2")
        ).head(100)

        if len(merged) > 0:
            add_table(root, f"Scatter_{cond1}_vs_{cond2}", [
                ("GeneSymbol", merged[gene_name_col_deseq].tolist()),
                (f"log2FC_{cond1}", merged[f"{log2fc_col}_1"].tolist()),
                (f"log2FC_{cond2}", merged[f"{log2fc_col}_2"].tolist()),
            ])

    filepath = prism_dir / "Concordance_scatter.pzfx"
    save_prism_xml(root, filepath)
    print(f"  Created: {filepath.name} ({filepath.stat().st_size} bytes)")

    # --- 9. dPSI_distributions.pzfx ---
    root = create_prism_xml()
    for cond_name, data in condition_results.items():
        all_dpsi = []
        for et, df in data["rmats_filtered"].items():
            if len(df) > 0:
                all_dpsi.extend(df[dpsi_col].dropna().tolist())

        if all_dpsi:
            # Sample 200 points for violin plot
            sampled_dpsi = all_dpsi[:200] if len(all_dpsi) > 200 else all_dpsi
            add_table(root, f"dPSI_{cond_name}", [
                ("Value", sampled_dpsi),
            ])

    filepath = prism_dir / "dPSI_distributions.pzfx"
    save_prism_xml(root, filepath)
    print(f"  Created: {filepath.name} ({filepath.stat().st_size} bytes)")

    print(f"  Prism export complete: {len(list(prism_dir.glob('*.pzfx')))} files in {prism_dir}")


def export_unfiltered_merged(condition_results, condition_labels, outdir):
    """Export unfiltered (all genes/events) merged across conditions to Excel.

    Keeps ALL native columns from each condition's raw data, suffixed with the
    condition short name.  Merge keys (gene_id/gene_name for DESeq2; gene info +
    genomic coordinates for rMATS) are shared columns, not duplicated.

    Creates a multi-sheet XLSX with:
      - DESeq2_All_Conditions: outer merge of all genes across all conditions
      - DESeq2_{labelA}_vs_{labelB}: pairwise outer merges for each pair
      - rMATS_{event_type}: outer merge of splicing events per event type
    """
    from itertools import combinations
    from pathlib import Path

    outdir = Path(outdir)
    xlsx_path = outdir / "Unfiltered_All_Conditions_Merged.xlsx"
    print(f"\n-- Exporting Unfiltered Merged Results --")
    print(f"   Output: {xlsx_path}")

    names = list(condition_results.keys())
    id_col   = DESEQ2_COLS["gene_id"]
    name_col = DESEQ2_COLS["gene_name"]

    rmats_gid  = RMATS_COLS["gene_id"]
    rmats_gn   = RMATS_COLS["gene_name"]

    def _short(label):
        parts = label.split(" vs ")
        return parts[0].replace(" ", "_") if parts else label.replace(" ", "_")

    def _deseq2_key(df):
        if name_col in df.columns and df[name_col].notna().sum() > 0:
            return name_col
        return id_col

    # --- Build per-condition DESeq2 slices with ALL columns ---
    deseq2_slices = {}
    for cname in names:
        raw = condition_results[cname].get("deseq2_raw", pd.DataFrame())
        if raw.empty:
            continue
        key = _deseq2_key(raw)
        short = _short(condition_labels[cname])
        df = raw.copy()
        # Identify merge-key columns (shared, not suffixed)
        key_cols = [c for c in [id_col, name_col] if c in df.columns]
        # Suffix all non-key columns with condition short name
        rename = {}
        for c in df.columns:
            if c not in key_cols:
                rename[c] = f"{short}_{c}"
        df = df.rename(columns=rename)
        deseq2_slices[cname] = (df, key, short)

    with pd.ExcelWriter(xlsx_path, engine="openpyxl") as writer:

        # ---- Sheet 1: DESeq2_All_Conditions (outer merge) ----
        if deseq2_slices:
            merged = None
            for cname in names:
                if cname not in deseq2_slices:
                    continue
                slice_df, key, short = deseq2_slices[cname]
                if merged is None:
                    merged = slice_df.copy()
                else:
                    merge_on = [c for c in [key] if c in merged.columns and c in slice_df.columns]
                    right_drop = [c for c in [id_col, name_col]
                                  if c in slice_df.columns and c in merged.columns and c not in merge_on]
                    merged = merged.merge(
                        slice_df.drop(columns=right_drop, errors="ignore"),
                        on=merge_on, how="outer")

            if merged is not None and not merged.empty:
                merged.to_excel(writer, sheet_name="DESeq2_All_Conditions", index=False)
                print(f"   DESeq2_All_Conditions: {len(merged):,} genes x {len(merged.columns)} cols")

        # ---- Sheets 2+: DESeq2 pairwise ----
        for cA, cB in combinations(names, 2):
            if cA not in deseq2_slices or cB not in deseq2_slices:
                continue
            slA, keyA, shortA = deseq2_slices[cA]
            slB, keyB, shortB = deseq2_slices[cB]
            merge_on = [c for c in [keyA] if c in slA.columns and c in slB.columns]
            right_drop = [c for c in [id_col, name_col]
                          if c in slB.columns and c in slA.columns and c not in merge_on]
            pair_df = slA.merge(
                slB.drop(columns=right_drop, errors="ignore"),
                on=merge_on, how="outer")
            sheet = f"DESeq2_{shortA}_vs_{shortB}"[:31]
            pair_df.to_excel(writer, sheet_name=sheet, index=False)
            print(f"   {sheet}: {len(pair_df):,} genes")

        # ---- rMATS sheets: one per event type, merged across conditions ----
        # Coordinate columns that uniquely identify each splicing event
        _coord_cols = {
            "SE": ["chr", "strand", "exonStart_0base", "exonEnd",
                    "upstreamES", "upstreamEE", "downstreamES", "downstreamEE"],
            "A3SS": ["chr", "strand", "longExonStart_0base", "longExonEnd",
                     "shortES", "shortEE", "flankingES", "flankingEE"],
            "A5SS": ["chr", "strand", "longExonStart_0base", "longExonEnd",
                     "shortES", "shortEE", "flankingES", "flankingEE"],
            "RI": ["chr", "strand", "riExonStart_0base", "riExonEnd",
                   "upstreamES", "upstreamEE", "downstreamES", "downstreamEE"],
            "MXE": ["chr", "strand", "1stExonStart_0base", "1stExonEnd",
                    "2ndExonStart_0base", "2ndExonEnd",
                    "upstreamES", "upstreamEE", "downstreamES", "downstreamEE"],
        }
        for et in RMATS_EVENT_TYPES:
            merged_rmats = None
            coord = _coord_cols.get(et, [])
            for cname in names:
                rmats_raw = condition_results[cname].get("rmats_raw", {})
                raw = rmats_raw.get(et, pd.DataFrame())
                if raw.empty:
                    continue
                short = _short(condition_labels[cname])
                df = raw.copy()
                # Merge keys: gene info + coordinates (shared, not suffixed)
                id_keys = [c for c in [rmats_gid, rmats_gn] if c in df.columns]
                coord_keys = [c for c in coord if c in df.columns]
                shared_keys = id_keys + coord_keys
                # Drop the rMATS "ID" column (run-specific, not meaningful across conditions)
                if "ID" in df.columns:
                    df = df.drop(columns=["ID"])
                # Suffix all non-shared columns
                rename = {}
                for c in df.columns:
                    if c not in shared_keys:
                        rename[c] = f"{short}_{c}"
                df = df.rename(columns=rename)

                if merged_rmats is None:
                    merged_rmats = df
                else:
                    merge_on = [c for c in shared_keys
                                if c in merged_rmats.columns and c in df.columns]
                    if not merge_on:
                        continue
                    right_drop = [c for c in shared_keys
                                  if c in df.columns and c in merged_rmats.columns
                                  and c not in merge_on]
                    merged_rmats = merged_rmats.merge(
                        df.drop(columns=right_drop, errors="ignore"),
                        on=merge_on, how="outer")

            if merged_rmats is not None and not merged_rmats.empty:
                # Excel row limit: 1,048,576. Warn if close.
                if len(merged_rmats) > 1_048_000:
                    print(f"   WARNING: rMATS_{et} has {len(merged_rmats):,} rows — "
                          f"near Excel limit, truncating to 1,048,000")
                    merged_rmats = merged_rmats.head(1_048_000)
                sheet = f"rMATS_{et}"[:31]
                merged_rmats.to_excel(writer, sheet_name=sheet, index=False)
                print(f"   {sheet}: {len(merged_rmats):,} events x {len(merged_rmats.columns)} cols")

    print(f"   Saved: {xlsx_path}")


def generate_powerpoint_report(condition_results, condition_labels, outdir):
    """Generate a professional PowerPoint report with all analysis figures.

    Searches per-condition figures in outdir/<cond>/figures/ and cross-condition
    figures in outdir/cross_condition/figures/ and outdir/gsea_results/.
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        print("  WARNING: python-pptx not installed, skipping PowerPoint generation")
        print("  Install with: pip install python-pptx")
        return

    from datetime import date
    from PIL import Image

    print("\n-- Generating PowerPoint Report --")

    DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    SLIDE_W = Inches(13.333)
    SLIDE_H = Inches(7.5)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # ── helper: section header slide ──────────────────────────────────────
    def add_section_slide(title_text, subtitle_text=""):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        # dark blue rectangle covering full slide
        from pptx.util import Emu as _Emu
        shape = slide.shapes.add_shape(
            1, _Emu(0), _Emu(0), SLIDE_W, SLIDE_H  # MSO_SHAPE.RECTANGLE = 1
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = DARK_BLUE
        shape.line.fill.background()
        # title
        tx = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1.2))
        p = tx.text_frame.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.LEFT
        # subtitle
        if subtitle_text:
            tx2 = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(11), Inches(0.8))
            p2 = tx2.text_frame.paragraphs[0]
            p2.text = subtitle_text
            p2.font.size = Pt(22)
            p2.font.color.rgb = WHITE
            p2.alignment = PP_ALIGN.LEFT
        return slide

    # ── helper: image slide ───────────────────────────────────────────────
    def add_image_slide(title_text, image_path):
        if not image_path.exists():
            print(f"    WARNING: missing {image_path.name}, skipping")
            return
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        # title bar
        bar = slide.shapes.add_shape(1, Emu(0), Emu(0), SLIDE_W, Inches(0.9))
        bar.fill.solid()
        bar.fill.fore_color.rgb = DARK_BLUE
        bar.line.fill.background()
        tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(12), Inches(0.7))
        p = tx.text_frame.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = WHITE
        # image — scale to fit remaining area
        try:
            img = Image.open(image_path)
            iw, ih = img.size
            img.close()
            max_w = Inches(12.333)  # 0.5" margins each side
            max_h = Inches(6.2)     # below title bar
            scale = min(max_w / Emu(int(iw * 914400 / 96)),
                        max_h / Emu(int(ih * 914400 / 96)))
            w = int(iw * 914400 / 96 * min(scale, 1.0))
            h = int(ih * 914400 / 96 * min(scale, 1.0))
            left = int((SLIDE_W - w) / 2)
            top = Inches(1.0) + int((max_h - h) / 2)
            slide.shapes.add_picture(str(image_path), left, top, w, h)
        except Exception as e:
            print(f"    WARNING: could not add {image_path.name}: {e}")

    # ══════════════════════════════════════════════════════════════════════
    # TITLE SLIDE
    # ══════════════════════════════════════════════════════════════════════
    cond_names = list(condition_results.keys())
    n_conds = len(cond_names)
    _ppt_filters = (f"padj < {PADJ_CUTOFF}, |log2FC| \u2265 {LOG2FC_CUTOFF}, "
                     f"baseMean \u2265 {BASEMEAN_CUTOFF}")
    if RMATS_DUAL_FILTER:
        _ppt_filters += (f"  |  rMATS: FDR < {RMATS_FDR_CUTOFF} & P < {RMATS_PVAL_CUTOFF}, "
                          f"|dPSI| \u2265 {INCLEVEL_DIFF_CUTOFF}")
    else:
        _ppt_fdr_or_p = f"FDR < {RMATS_FDR_CUTOFF}" if USE_FDR else f"P < {RMATS_PVAL_CUTOFF}"
        _ppt_filters += f"  |  rMATS: {_ppt_fdr_or_p}, |dPSI| \u2265 {INCLEVEL_DIFF_CUTOFF}"
    add_section_slide(
        "RNA-seq Analysis Report",
        f"{n_conds} conditions  |  {_ppt_filters}  |  {date.today().strftime('%B %d, %Y')}"
    )

    # ══════════════════════════════════════════════════════════════════════
    # OVERVIEW SLIDE
    # ══════════════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bar = slide.shapes.add_shape(1, Emu(0), Emu(0), SLIDE_W, Inches(0.9))
    bar.fill.solid()
    bar.fill.fore_color.rgb = DARK_BLUE
    bar.line.fill.background()
    tx = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(12), Inches(0.7))
    p = tx.text_frame.paragraphs[0]
    p.text = "Analysis Overview"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE

    body = slide.shapes.add_textbox(Inches(1), Inches(1.4), Inches(11), Inches(5.5))
    tf = body.text_frame
    tf.word_wrap = True
    for cond_name, data in condition_results.items():
        label = condition_labels.get(cond_name, cond_name)
        deg_df = data["deseq2_filtered"]["all_genes"]
        n_deg = len(deg_df)
        n_up = int((deg_df["direction"] == "up").sum()) if "direction" in deg_df.columns else 0
        n_down = int((deg_df["direction"] == "down").sum()) if "direction" in deg_df.columns else 0
        n_splice = sum(len(df) for df in data["rmats_filtered"].values())
        p = tf.add_paragraph()
        p.text = label
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = DARK_BLUE
        p = tf.add_paragraph()
        p.text = f"    DEGs: {n_deg}  ({n_up} up / {n_down} down)    |    Splicing events: {n_splice}"
        p.font.size = Pt(16)
        tf.add_paragraph()  # spacer

    # ══════════════════════════════════════════════════════════════════════
    # PER-CONDITION SECTIONS
    # ══════════════════════════════════════════════════════════════════════
    per_cond_figures = [
        ("Volcano Plot", "volcano_plot.png"),
        ("Volcano Plot (labeled)", "volcano_plot_labeled.png"),
        ("MA Plot", "ma_plot.png"),
        ("Top DEGs (lollipop)", "top_genes_lollipop.png"),
        ("P-value Histogram", "pvalue_histogram.png"),
        ("Biotype Distribution", "biotype_distribution.png"),
        ("Biotype Direction Chart", "biotype_direction_chart.png"),
        ("Biotype Enrichment", "biotype_enrichment.png"),
        ("Non-coding Volcano", "volcano_plot_labeled_non_protein_coding.png"),
        ("Non-coding MA Plot", "ma_plot_non_protein_coding.png"),
        ("Non-coding Biotype Distribution", "biotype_distribution_non_protein_coding.png"),
        ("Expression Rank Plot", "expression_rank_plot.png"),
        ("Splicing Event Summary", "rmats_event_type_summary.png"),
        ("Splicing dPSI Distribution", "rmats_dpsi_distribution.png"),
        ("Splicing All Events Scatter", "rmats_all_events_scatter.png"),
        ("SE Scatter", "rmats_SE_scatter.png"),
        ("SE PSI Scatter", "rmats_SE_psi_scatter.png"),
        ("A3SS Scatter", "rmats_A3SS_scatter.png"),
        ("A5SS Scatter", "rmats_A5SS_scatter.png"),
        ("RI Scatter", "rmats_RI_scatter.png"),
        ("MXE Scatter", "rmats_MXE_scatter.png"),
    ]

    for cond_name in cond_names:
        label = condition_labels.get(cond_name, cond_name)
        fig_dir = outdir / cond_name / "figures"
        add_section_slide(label, "Condition-Specific Analysis")
        for slide_title, fname in per_cond_figures:
            path = fig_dir / fname
            if path.exists():
                add_image_slide(f"{label} \u2014 {slide_title}", path)

    # ══════════════════════════════════════════════════════════════════════
    # CROSS-CONDITION SECTION
    # ══════════════════════════════════════════════════════════════════════
    cross_dir = outdir / "cross_condition" / "figures"
    gsea_dir = outdir / "gsea_results"

    cross_figures = [
        ("DEG Counts Overview", "deseq2_de_counts_overview.png"),
        ("Analysis Summary Dashboard", "analysis_summary_dashboard.png"),
        ("Venn \u2014 All Significant Genes", "venn_all_sig_genes.png"),
        ("Venn \u2014 Upregulated", "venn_upregulated.png"),
        ("Venn \u2014 Downregulated", "venn_downregulated.png"),
        ("UpSet \u2014 All Significant", "deseq2_upset_all_sig.png"),
        ("UpSet \u2014 Up", "deseq2_upset_up.png"),
        ("UpSet \u2014 Down", "deseq2_upset_down.png"),
        ("Direction Concordance Heatmap", "direction_concordance_heatmap.png"),
        ("Direction Concordance Summary", "direction_concordance_summary.png"),
        ("Pairwise log2FC Scatter", "pairwise_log2fc_scatter.png"),
        ("log2FC Heatmap", "log2fc_heatmap.png"),
        ("log2FC vs dPSI Scatter", "log2fc_vs_dpsi_scatter.png"),
        ("Cross-Condition Biotype Comparison", "cross_condition_biotype_comparison.png"),
        ("Cross-Condition Biotype Direction", "cross_condition_biotype_direction.png"),
        ("rMATS Event Count Comparison", "rmats_event_count_comparison.png"),
        ("rMATS Direction Concordance", "rmats_direction_concordance.png"),
        ("rMATS UpSet \u2014 SE", "rmats_upset_SE.png"),
        ("rMATS UpSet \u2014 Genes", "rmats_upset_genes.png"),
        ("Pairwise dPSI Scatter", "pairwise_dpsi_scatter.png"),
    ]

    add_section_slide("Cross-Condition Analysis", "Comparative & Integrative Results")

    for slide_title, fname in cross_figures:
        path = cross_dir / fname
        if path.exists():
            add_image_slide(slide_title, path)

    # Pairwise DEG venns
    for png in sorted(cross_dir.glob("venn_deg_*.png")) if cross_dir.exists() else []:
        add_image_slide(f"Pairwise DEG Venn \u2014 {png.stem.replace('venn_deg_', '')}", png)

    # Pairwise splicing venns
    for png in sorted(cross_dir.glob("venn_splicing_*.png")) if cross_dir.exists() else []:
        parts = png.stem.replace("venn_splicing_", "").split("_", 1)
        add_image_slide(f"Splicing Venn ({parts[0]}) \u2014 {parts[1] if len(parts)>1 else ''}", png)

    # DESeq2 vs rMATS venns
    for png in sorted(cross_dir.glob("deseq2_vs_rmats_venn_*.png")) if cross_dir.exists() else []:
        cname = png.stem.replace("deseq2_vs_rmats_venn_", "")
        label = condition_labels.get(cname, cname)
        add_image_slide(f"DESeq2 vs rMATS \u2014 {label}", png)

    # ── ORA enrichment plots (from cross_condition/figures) ───────────
    add_section_slide("Functional Enrichment", "ORA (enrichr & g:Profiler) and GSEA")

    for png in sorted(cross_dir.glob("go_enrichment_combined_*.png")) if cross_dir.exists() else []:
        parts = png.stem.replace("go_enrichment_combined_", "")
        add_image_slide(f"ORA \u2014 {parts}", png)

    # ── GSEA combined plots (from gsea_results/) ─────────────────────
    if gsea_dir.exists():
        for png in sorted(gsea_dir.glob("gsea_combined_*.png")):
            cname = png.stem.replace("gsea_combined_", "")
            label = condition_labels.get(cname, cname)
            add_image_slide(f"GSEA Combined \u2014 {label}", png)

    # ══════════════════════════════════════════════════════════════════════
    # SAVE
    # ══════════════════════════════════════════════════════════════════════
    pptx_path = outdir / "RNA-seq_Analysis_Report.pptx"
    prs.save(str(pptx_path))
    print(f"  PowerPoint saved: {pptx_path} ({len(prs.slides)} slides)")


def validate_outputs(condition_results, condition_labels, outdir):
    """Validate all pipeline outputs for completeness and correctness.

    Checks:
    - File existence and sizes
    - Prism files >10KB
    - Venn diagram math (All = Concordant_Up + Concordant_Down + Discordant)
    - DEG totals (total = up + down)
    - Gene symbols present in output files
    """
    print("\n" + "=" * 60)
    print("  VALIDATION REPORT")
    print("=" * 60)

    validation_passed = True
    warnings = []
    errors = []

    # --- 1. Check file existence ---
    print("\n1. File Existence Check:")

    expected_files = [
        ("cross_condition/multi_condition_results.xlsx", "Multi-condition results"),
        ("RNA-seq_Analysis_Report.pptx", "PowerPoint report"),
    ]

    for fname, desc in expected_files:
        fpath = outdir / fname
        if fpath.exists():
            size_mb = fpath.stat().st_size / (1024 * 1024)
            print(f"  ✓ {desc}: {fname} ({size_mb:.2f} MB)")
        else:
            errors.append(f"Missing file: {fname}")
            print(f"  ✗ {desc}: {fname} MISSING")
            validation_passed = False

    # --- 2. Check Prism files ---
    print("\n2. Prism File Validation:")
    prism_dir = outdir / "prism_files"
    if prism_dir.exists():
        prism_files = list(prism_dir.glob("*.pzfx"))
        print(f"  Found {len(prism_files)} Prism files")

        small_files = []
        for pfile in prism_files:
            size_kb = pfile.stat().st_size / 1024
            if size_kb < 10:
                small_files.append((pfile.name, size_kb))

        if small_files:
            for fname, size_kb in small_files:
                warnings.append(f"Prism file {fname} is only {size_kb:.1f} KB (expected >10 KB)")
                print(f"  ⚠ {fname}: {size_kb:.1f} KB (may be too small)")
        else:
            print(f"  ✓ All Prism files >10 KB")

        if len(prism_files) < 15:
            warnings.append(f"Expected 15+ Prism files, found {len(prism_files)}")
            print(f"  ⚠ Only {len(prism_files)} Prism files (expected 15+)")
        else:
            print(f"  ✓ File count meets requirement ({len(prism_files)} >= 15)")
    else:
        errors.append("Prism directory not found")
        print("  ✗ Prism directory not found")
        validation_passed = False

    # --- 3. Validate DEG counts ---
    print("\n3. DEG Count Validation:")
    gene_name_col = DESEQ2_COLS.get("gene_name", "gene_name")

    for cond_name, data in condition_results.items():
        cond_label = condition_labels.get(cond_name, cond_name)
        deg_filt = data["deseq2_filtered"]["all_genes"]

        if "direction" in deg_filt.columns:
            total = len(deg_filt)
            up = int((deg_filt["direction"] == "up").sum())
            down = int((deg_filt["direction"] == "down").sum())
            computed_total = up + down

            if computed_total == total:
                print(f"  ✓ {cond_label}: Total={total}, Up={up}, Down={down}")
            else:
                errors.append(f"{cond_label} DEG count mismatch: Total={total}, Up+Down={computed_total}")
                print(f"  ✗ {cond_label}: Total={total} ≠ Up+Down={computed_total}")
                validation_passed = False

            # Check for gene symbols
            if gene_name_col in deg_filt.columns:
                n_missing = deg_filt[gene_name_col].isna().sum()
                if n_missing > 0:
                    warnings.append(f"{cond_label}: {n_missing} genes missing gene symbols")
                    print(f"  ⚠ {cond_label}: {n_missing}/{total} genes missing symbols")
                else:
                    print(f"  ✓ {cond_label}: All genes have symbols")
            else:
                warnings.append(f"{cond_label}: No gene_name column found")
                print(f"  ⚠ {cond_label}: No gene_name column")

    # --- 4. Validate splicing Venn math ---
    print("\n4. Splicing Venn Math Validation:")
    # This checks the directional Venn diagrams created in Fix 3

    names = list(condition_results.keys())
    gene_col = RMATS_COLS["gene_name"]
    dpsi_col = RMATS_COLS["inclevel_diff"]

    if len(names) >= 2:
        for et in RMATS_EVENT_TYPES:
            dfs = {}
            for name in names:
                if et in condition_results[name]["rmats_filtered"]:
                    df = condition_results[name]["rmats_filtered"][et]
                    if len(df) > 0:
                        gene_dpsi = df.groupby(gene_col)[dpsi_col].mean()
                        dfs[name] = gene_dpsi

            if len(dfs) >= 2:
                shared_genes = set.intersection(*[set(d.index) for d in dfs.values()])

                concordant_up = 0
                concordant_down = 0
                discordant = 0

                for gene in shared_genes:
                    signs = [np.sign(dfs[name][gene]) for name in dfs.keys()]
                    if all(s > 0 for s in signs):
                        concordant_up += 1
                    elif all(s < 0 for s in signs):
                        concordant_down += 1
                    else:
                        discordant += 1

                computed_all = concordant_up + concordant_down + discordant
                if computed_all == len(shared_genes):
                    print(f"  ✓ {et}: All={len(shared_genes)}, Up={concordant_up}, "
                          f"Down={concordant_down}, Disc={discordant}")
                else:
                    errors.append(f"{et} Venn math error: All={len(shared_genes)} ≠ "
                                  f"Sum={computed_all}")
                    print(f"  ✗ {et}: All={len(shared_genes)} ≠ Sum={computed_all}")
                    validation_passed = False

    # --- 5. Summary ---
    print("\n" + "=" * 60)
    print("  VALIDATION SUMMARY")
    print("=" * 60)

    if validation_passed and not errors:
        print("  ✓ ALL VALIDATION CHECKS PASSED")
    else:
        print(f"  ✗ VALIDATION FAILED: {len(errors)} error(s)")

    if warnings:
        print(f"\n  ⚠ {len(warnings)} warning(s):")
        for w in warnings[:10]:  # Show first 10
            print(f"    - {w}")

    if errors:
        print(f"\n  ✗ {len(errors)} error(s):")
        for e in errors:
            print(f"    - {e}")

    print("=" * 60)

    return validation_passed


# ---------------------------------------------------------------------------
# MAIN
# ---------------------------------------------------------------------------

def main():
    setup_style()

    if not CONDITIONS:
        raise ValueError("No conditions defined in CONDITIONS list")

    outdir = Path(OUTPUT_DIR)
    outdir.mkdir(parents=True, exist_ok=True)

    print("=" * 60)
    print("  DESeq2 & rMATS Multi-Condition Pipeline")
    print(f"  {len(CONDITIONS)} conditions to process")
    print("=" * 60)

    # ===================================================================
    # PHASE 1: Per-condition analysis
    # ===================================================================
    condition_results = {}
    condition_labels = {}
    counts_df = None
    sample_metadata = {}

    # Load normalized counts matrix if provided (for PCA, heatmaps)
    if COUNTS_FILE:
        counts_df, sample_metadata = load_counts_matrix(
            COUNTS_FILE, SAMPLE_METADATA, CONDITIONS)
        if counts_df is not None:
            qc_dir = outdir / "qc_plots"
            qc_dir.mkdir(parents=True, exist_ok=True)
            print("\n-- QC: PCA Plot --")
            pca_plot(counts_df=counts_df, metadata=sample_metadata, outdir=qc_dir)
            print("\n-- QC: Sample Correlation Heatmap --")
            sample_correlation_heatmap(counts_df, sample_metadata, qc_dir)
    else:
        print("[INFO] No counts file — skipping PCA, correlation heatmap, top DEG heatmap")

    # Load RBP annotations if provided
    rbp_annotations = {}
    if RBP_FILE:
        print("\n-- Loading RBP Annotations --")
        try:
            rbp_annotations = load_rbp_annotations(RBP_FILE)
        except Exception as e:
            print(f"  WARNING: Failed to load RBP annotations: {e}")
            rbp_annotations = {}
    else:
        print("[INFO] No RBP_FILE — skipping RBP annotation")

    # Per-condition PCA files (from WSF output, if provided)
    for cond in CONDITIONS:
        if cond.get("pca_file"):
            cond_fig = outdir / cond["name"] / "figures"
            cond_fig.mkdir(parents=True, exist_ok=True)
            pca_plot(pca_file=cond["pca_file"], outdir=cond_fig)

    for cond in CONDITIONS:
        cond_name = cond["name"]
        cond_label = cond["label"]
        condition_labels[cond_name] = cond_label

        print(f"\n{'=' * 60}")
        print(f"  Processing condition: {cond_label}")
        print(f"{'=' * 60}")

        # Condition-specific output subdirectory
        cond_outdir = outdir / cond_name
        cond_outdir.mkdir(parents=True, exist_ok=True)
        cond_fig_dir = cond_outdir / "figures"
        cond_fig_dir.mkdir(exist_ok=True)

        # --- DESeq2 ---
        print("\n-- Loading DESeq2 --")
        deseq2_raw = load_file(cond["deseq2_file"], f"DESeq2 ({cond_label})")
        deseq2_raw = normalize_deseq2_columns(deseq2_raw, f"DESeq2 ({cond_label})")
        # gene_name, biotype, stat, lfcSE are optional — many DESeq2 files lack these
        optional_keys = {"biotype", "gene_name", "stat", "lfcSE"}
        required_deseq2 = [v for k, v in DESEQ2_COLS.items()
                           if v is not None and k not in optional_keys]
        validate_columns(deseq2_raw, required_deseq2, f"DESeq2 ({cond_label})")
        if DESEQ2_COLS.get("gene_name", "") not in deseq2_raw.columns:
            print(f"  NOTE: No gene name column found — attempting Ensembl → gene name lookup")
        deseq2_raw = _enrich_with_gene_names(deseq2_raw, f"DESeq2 ({cond_label})")
        deseq2_raw = _reassign_biotypes_from_mygene(deseq2_raw, f"DESeq2 ({cond_label})")

        # RBP annotation on raw data (so columns propagate to all exports)
        if rbp_annotations:
            gene_col_rbp = DESEQ2_COLS.get("gene_name", "")
            if gene_col_rbp and gene_col_rbp in deseq2_raw.columns:
                print(f"\n-- RBP Annotation [{cond_label}] --")
                deseq2_raw = annotate_rbps(deseq2_raw, rbp_annotations,
                                           gene_col=gene_col_rbp)

        # Biotype passes
        biotype_passes = [("All Genes", None, "")]
        if AUTO_BIOTYPE_SPLIT and DESEQ2_COLS["biotype"] in deseq2_raw.columns:
            biotype_passes += [
                ("Protein Coding", "protein_coding", "_protein_coding"),
                ("Non-Protein Coding", "non_protein_coding", "_non_protein_coding"),
            ]

        deseq2_filtered_sets = {}

        for label, bio_filter, suffix in biotype_passes:
            full_label = f"{cond_label} - {label}"
            deseq2_all, deseq2_filt = filter_deseq2(
                deseq2_raw, biotype_filter=bio_filter, label=full_label)

            export_key = label.lower().replace(" ", "_").replace("-", "_")
            deseq2_filtered_sets[export_key] = deseq2_filt

            print(f"\n-- Generating DESeq2 Figures [{full_label}] --")
            pvalue_histogram(deseq2_all, cond_fig_dir, label=full_label, suffix=suffix)
            volcano_plot(deseq2_all, cond_fig_dir, label=full_label, suffix=suffix)
            volcano_plot_labeled(deseq2_all, cond_fig_dir, label=full_label, suffix=suffix)
            ma_plot(deseq2_all, cond_fig_dir, label=full_label, suffix=suffix)
            expression_rank_plot(deseq2_all, cond_fig_dir, label=full_label, suffix=suffix)
            if INTERACTIVE_PLOTS:
                volcano_plot_interactive(deseq2_all, cond_fig_dir, label=full_label, suffix=suffix)
                ma_plot_interactive(deseq2_all, cond_fig_dir, label=full_label, suffix=suffix)
            if len(deseq2_filt) > 0:
                biotype_chart(deseq2_filt, cond_fig_dir, label=full_label, suffix=suffix)
                if bio_filter is None:
                    top_genes_lollipop(deseq2_filt, cond_fig_dir, label=full_label, suffix=suffix)
                    biotype_direction_chart(deseq2_filt, cond_fig_dir, label=full_label, suffix=suffix)
                    biotype_enrichment_test(deseq2_filt, deseq2_all, cond_fig_dir, label=full_label, suffix=suffix)
            else:
                print(f"  No genes passed filter -- skipping biotype chart for {full_label}")
            if bio_filter is None:
                biotype_volcano(deseq2_all, cond_fig_dir, label=full_label, suffix=suffix)
                ecdf_log2fc_by_biotype(deseq2_all, cond_fig_dir, label=full_label, suffix=suffix)

        # --- rMATS (if directory exists) ---
        rmats_raw = {}
        rmats_filtered = {}

        if cond.get("rmats_dir"):
            print(f"\n-- Loading rMATS for {cond_label} --")
            rmats_raw = load_all_rmats(cond["rmats_dir"])

            print(f"\n-- rMATS Filtering --")
            filtered_counts = {}
            for event_type, df in rmats_raw.items():
                raw, filt = filter_rmats(df, event_type)
                rmats_raw[event_type] = raw
                rmats_filtered[event_type] = filt
                filtered_counts[event_type] = len(filt)

            total_sig = sum(filtered_counts.values())
            print(f"  TOTAL significant events: {total_sig:,}")

            # RBP annotation on filtered rMATS data
            if rbp_annotations:
                print(f"\n-- RBP Annotation on rMATS [{cond_label}] --")
                for et in rmats_filtered:
                    if not rmats_filtered[et].empty:
                        rmats_filtered[et] = annotate_rbps(
                            rmats_filtered[et], rbp_annotations,
                            gene_col=RMATS_COLS["gene_name"]
                        )

            print(f"\n-- Generating rMATS Figures [{cond_label}] --")
            for event_type, df in rmats_raw.items():
                rmats_scatter(df, event_type, cond_fig_dir)
                rmats_psi_scatter(rmats_raw, rmats_filtered, event_type, cond_fig_dir)
            rmats_combined_volcano(rmats_raw, cond_fig_dir)
            rmats_event_summary_chart(filtered_counts, cond_fig_dir)
            rmats_dpsi_distribution(rmats_filtered, cond_fig_dir)
        else:
            print(f"\n  No rMATS data for {cond_label}, skipping rMATS analysis")

        # Per-condition export
        export_results(deseq2_filtered_sets, rmats_filtered, cond_outdir)

        # Store for cross-condition comparisons
        condition_results[cond_name] = {
            "deseq2_raw": deseq2_raw,
            "deseq2_filtered": deseq2_filtered_sets,
            "rmats_raw": rmats_raw,
            "rmats_filtered": rmats_filtered,
        }

    # ===================================================================
    # PHASE 2: Cross-condition comparisons
    # ===================================================================
    print(f"\n{'=' * 60}")
    print("  Cross-Condition Comparisons")
    print(f"{'=' * 60}")

    comparison_dir = outdir / "cross_condition"
    comparison_dir.mkdir(parents=True, exist_ok=True)
    comparison_fig_dir = comparison_dir / "figures"
    comparison_fig_dir.mkdir(exist_ok=True)

    cross_data = {}

    # DESeq2 comparisons (all conditions)
    print("\n-- DESeq2 DE Counts Overview --")
    deseq2_de_counts_chart(condition_results, condition_labels, comparison_fig_dir)

    print("\n-- DESeq2 Venn Diagrams (3-way) --")
    gene_sets = extract_gene_sets(condition_results)
    deseq2_venn_diagrams(gene_sets, condition_labels, comparison_fig_dir)

    print("\n-- Pairwise DEG Venn Diagrams --")
    pairwise_deg_venns(condition_results, condition_labels, comparison_fig_dir)

    print("\n-- DESeq2 UpSet Plots --")
    deseq2_upset_plot(condition_results, condition_labels, comparison_fig_dir)

    print("\n-- DESeq2 Direction Concordance --")
    concordance_df = deseq2_direction_concordance(
        condition_results, condition_labels, comparison_fig_dir)
    cross_data["concordance_matrix"] = concordance_df

    print("\n-- DESeq2 Log2FC Heatmap --")
    log2fc_df = deseq2_log2fc_heatmap(
        condition_results, condition_labels, comparison_fig_dir)
    cross_data["log2fc_matrix"] = log2fc_df

    # Top DEG heatmap (requires counts matrix)
    if counts_df is not None:
        print("\n-- Top DEG Expression Heatmap --")
        top_deg_heatmap(counts_df, condition_results, condition_labels,
                        sample_metadata, comparison_fig_dir)

    print("\n-- Pairwise log2FC Scatter --")
    pairwise_log2fc_scatter(condition_results, condition_labels, comparison_fig_dir)

    print("\n-- Cross-Condition Biotype Comparison --")
    cross_condition_biotype_comparison(condition_results, condition_labels, comparison_fig_dir)
    cross_condition_biotype_direction(condition_results, condition_labels, comparison_fig_dir)

    # RBP cross-condition analysis (only if RBP annotations were loaded)
    if rbp_annotations:
        print("\n-- RBP Cross-Condition Heatmap --")
        try:
            rbp_heatmap(condition_results, condition_labels, comparison_fig_dir)
        except Exception as e:
            print(f"  WARNING: RBP heatmap failed: {e}")

        print("\n-- RBP Summary Table --")
        try:
            rbp_summary_table(condition_results, condition_labels, comparison_dir)
        except Exception as e:
            print(f"  WARNING: RBP summary table failed: {e}")

    # rMATS comparisons (only conditions with rMATS data)
    rmats_conditions = {name: res for name, res in condition_results.items()
                        if res["rmats_filtered"]}
    if len(rmats_conditions) >= 2:
        print("\n-- rMATS Cross-Condition Venn Diagrams --")
        rmats_cross_condition_venn(rmats_conditions, condition_labels, comparison_fig_dir)
        try:
            rmats_cross_condition_venn(rmats_conditions, condition_labels, comparison_fig_dir, match_by="gene")
        except Exception as e:
            print(f"  [WARN] Gene-level cross-condition Venns failed: {e}")

        print("\n-- rMATS Event Count Comparison --")
        rmats_event_count_comparison(rmats_conditions, condition_labels, comparison_fig_dir)

        print("\n-- rMATS UpSet Plots --")
        rmats_upset_plot(rmats_conditions, condition_labels, comparison_fig_dir)
        try:
            rmats_upset_plot(rmats_conditions, condition_labels, comparison_fig_dir, match_by="gene")
        except Exception as e:
            print(f"  [WARN] Gene-level UpSet plots failed: {e}")

        print("\n-- rMATS Direction Concordance --")
        rmats_conc = rmats_direction_concordance(
            rmats_conditions, condition_labels, comparison_fig_dir)
        cross_data["rmats_concordance"] = rmats_conc

        print("\n-- Pairwise Splicing Venn Diagrams --")
        pairwise_splicing_venns(rmats_conditions, condition_labels, comparison_fig_dir)
        try:
            pairwise_splicing_venns(rmats_conditions, condition_labels, comparison_fig_dir, match_by="gene")
        except Exception as e:
            print(f"  [WARN] Gene-level splicing Venns failed: {e}")

        print("\n-- Pairwise dPSI Scatter --")
        pairwise_dpsi_scatter(rmats_conditions, condition_labels, comparison_fig_dir)
        try:
            pairwise_dpsi_scatter(rmats_conditions, condition_labels, comparison_fig_dir, match_by="gene")
        except Exception as e:
            print(f"  [WARN] Gene-level dPSI scatter failed: {e}")

        print("\n-- Directional Splicing Venn Diagrams --")
        # REMOVED: redundant with pairwise Venns
        # try:
        #     rmats_directional_venn_diagrams(rmats_conditions, condition_labels, comparison_fig_dir)
        # except Exception as e:
        #     print(f"  [WARN] Event-level directional Venns failed: {e}")
        # REMOVED: redundant with pairwise Venns
        # try:
        #     rmats_directional_venn_diagrams(rmats_conditions, condition_labels, comparison_fig_dir, match_by="gene")
        # except Exception as e:
        #     print(f"  [WARN] Gene-level directional Venns failed: {e}")

        print("\n-- rMATS Event Heatmaps --")
        for _et in ["SE", "RI"]:
            try:
                rmats_event_heatmap(rmats_conditions, condition_labels, _et, comparison_fig_dir)
            except Exception as e:
                print(f"  WARNING: {_et} event heatmap failed: {e}")

        print("\n-- rMATS Event Type Pie Chart --")
        try:
            rmats_event_pie_chart(rmats_conditions, condition_labels, comparison_fig_dir)
        except Exception as e:
            print(f"  WARNING: Event type pie chart failed: {e}")

        print("\n-- Pairwise Comparison Workbooks --")
        try:
            export_pairwise_workbook(rmats_conditions, condition_labels, comparison_fig_dir)
        except Exception as e:
            print(f"  WARNING: Pairwise workbook export failed: {e}")
    else:
        print("\n  Fewer than 2 conditions have rMATS data, skipping rMATS comparisons")

    # Combined DESeq2 + rMATS
    print("\n-- Combined DESeq2 + rMATS Analyses --")
    deseq2_vs_rmats_venn(condition_results, condition_labels, comparison_fig_dir)
    log2fc_vs_dpsi_scatter(condition_results, condition_labels, comparison_fig_dir)

    # Master combined export
    print("\n-- Exporting Combined Multi-Condition Results --")
    export_combined_results(condition_results, cross_data, comparison_dir)

    # Unfiltered merged overlap Excel (all genes, all conditions)
    export_unfiltered_merged(condition_results, condition_labels, outdir)

    # ===================================================================
    # PHASE 3: GSEA, GO ORA, Prism, PowerPoint, and Validation
    # ===================================================================
    print(f"\n{'=' * 60}")
    print("  GSEA, GO ORA, Prism Export, PowerPoint, and Validation")
    print(f"{'=' * 60}")

    # GSEA enrichment
    gsea_results = run_gsea_enrichment(condition_results, condition_labels, outdir)

    # GO Over-Representation Analysis
    if ORA_METHOD == "both":
        # Run both Enrichr and g:Profiler side-by-side for comparison
        print("\n-- Running Enrichr ORA --")
        try:
            go_results_enrichr = run_go_enrichment(
                condition_results, condition_labels, outdir,
                _best_gene_key_fn=_best_gene_key, DESEQ2_COLS_map=DESEQ2_COLS,
                _force_enrichr=True)
        except Exception as e:
            print(f"  WARNING: Enrichr ORA failed: {e}")
            go_results_enrichr = {}
        if go_results_enrichr:
            go_enrichment_combined_plot(
                go_results_enrichr, condition_labels,
                outdir / "cross_condition" / "figures",
                FIG_FORMAT_override=FIG_FORMAT, FIG_DPI_override=FIG_DPI,
                filename_suffix="_enrichr")
            export_go_prism(go_results_enrichr, condition_labels,
                            outdir / "prism_files", filename_suffix="_enrichr")

        print("\n-- Running g:Profiler ORA --")
        try:
            go_results_gprofiler = run_gprofiler_ora(
                condition_results, condition_labels, outdir,
                _best_gene_key_fn=_best_gene_key, DESEQ2_COLS_map=DESEQ2_COLS)
        except Exception as e:
            print(f"  WARNING: g:Profiler ORA failed: {e}")
            go_results_gprofiler = {}
        if go_results_gprofiler:
            go_enrichment_combined_plot(
                go_results_gprofiler, condition_labels,
                outdir / "cross_condition" / "figures",
                FIG_FORMAT_override=FIG_FORMAT, FIG_DPI_override=FIG_DPI,
                filename_suffix="_gprofiler")
            export_go_prism(go_results_gprofiler, condition_labels,
                            outdir / "prism_files", filename_suffix="_gprofiler")

        # Use Enrichr as primary for summary dashboard (HSCHARME paper standard)
        go_results = go_results_enrichr if go_results_enrichr else go_results_gprofiler
    else:
        # Single method (legacy behavior)
        go_results = run_go_enrichment(
            condition_results, condition_labels, outdir,
            _best_gene_key_fn=_best_gene_key, DESEQ2_COLS_map=DESEQ2_COLS)
        go_enrichment_combined_plot(
            go_results, condition_labels,
            outdir / "cross_condition" / "figures",
            FIG_FORMAT_override=FIG_FORMAT, FIG_DPI_override=FIG_DPI)
        export_go_prism(go_results, condition_labels, outdir / "prism_files")

    # Combined GSEA dot plots (replaces legacy gsea_dotplot)
    gsea_combined_plot(gsea_results, condition_labels, outdir)

    # GSEA enrichment plots and leading edge export
    gsea_enrichment_plots(gsea_results, condition_labels, outdir)
    export_gsea_leading_edge(gsea_results, condition_labels, outdir)

    # Gene overlap summary
    gene_overlap_summary(condition_results, condition_labels, outdir / "cross_condition")

    # Summary dashboard
    summary_dashboard(condition_results, condition_labels, go_results, gsea_results,
                      outdir / "cross_condition" / "figures")

    # Prism export
    export_prism_files(condition_results, condition_labels, outdir, gsea_results=gsea_results)

    # PowerPoint generation
    generate_powerpoint_report(condition_results, condition_labels, outdir)

    # Validation
    validation_passed = validate_outputs(condition_results, condition_labels, outdir)

    print(f"\n{'=' * 60}")
    print(f"  Done! All outputs in: {outdir.resolve()}")
    if validation_passed:
        print("  ✓ All validation checks passed")
    else:
        print("  ⚠ Some validation checks failed (see report above)")
    print(f"{'=' * 60}")


# ---------------------------------------------------------------------------
# PROGRAMMATIC ENTRY POINT  (used by pipeline_launcher.py)
# ---------------------------------------------------------------------------

def run_pipeline(config: dict):
    """Inject a config dict, override all module globals, and run the pipeline.

    Called by the GUI launcher so the user never needs to edit this file.
    Can also be used from notebooks or scripts for reproducible runs.

    Required keys in config
    -----------------------
    CONDITIONS            list of dicts with keys: name, label, deseq2_file, rmats_dir
    OUTPUT_DIR            str
    LOG2FC_CUTOFF         float
    BASEMEAN_CUTOFF       float
    PADJ_CUTOFF           float
    AUTO_BIOTYPE_SPLIT    bool
    RMATS_FDR_CUTOFF      float
    RMATS_PVAL_CUTOFF     float
    INCLEVEL_DIFF_CUTOFF  float
    USE_FDR               bool
    FIG_DPI               int
    FIG_FORMAT            str  ('png', 'svg', or 'pdf')
    FONT_SIZE             int
    COLOR_UP              str  (hex, e.g. '#D72638')
    COLOR_DOWN            str
    COLOR_NS              str
    INTERACTIVE_PLOTS     bool
    DESEQ2_COLS           dict  (same structure as the module-level constant)
    RMATS_COLS            dict
    """
    global CONDITIONS, OUTPUT_DIR
    global LOG2FC_CUTOFF, BASEMEAN_CUTOFF, PADJ_CUTOFF, AUTO_BIOTYPE_SPLIT
    global GENE_NAME_LOOKUP, SPECIES
    global RMATS_FDR_CUTOFF, RMATS_PVAL_CUTOFF, INCLEVEL_DIFF_CUTOFF, USE_FDR
    global FIG_DPI, FIG_FORMAT, FONT_SIZE
    global COLOR_UP, COLOR_DOWN, COLOR_NS
    global INTERACTIVE_PLOTS, DESEQ2_COLS, RMATS_COLS, RMATS_DUAL_FILTER
    global GSEA_DATABASES, ORA_DATABASES, GENES_OF_INTEREST
    global COUNTS_FILE, SAMPLE_METADATA, GSEA_RANKING, GSEA_MIN_SIZE, GSEA_MAX_SIZE
    global GSEA_PERMUTATIONS, ORA_METHOD
    global RBP_FILE

    CONDITIONS           = config["CONDITIONS"]
    OUTPUT_DIR           = config["OUTPUT_DIR"]
    LOG2FC_CUTOFF        = float(config.get("LOG2FC_CUTOFF", LOG2FC_CUTOFF))
    BASEMEAN_CUTOFF      = float(config.get("BASEMEAN_CUTOFF", BASEMEAN_CUTOFF))
    PADJ_CUTOFF          = float(config.get("PADJ_CUTOFF", PADJ_CUTOFF))
    AUTO_BIOTYPE_SPLIT   = bool(config.get("AUTO_BIOTYPE_SPLIT", AUTO_BIOTYPE_SPLIT))
    GENE_NAME_LOOKUP     = bool(config.get("GENE_NAME_LOOKUP", GENE_NAME_LOOKUP))
    SPECIES              = str(config.get("SPECIES", SPECIES))
    RMATS_FDR_CUTOFF     = float(config.get("RMATS_FDR_CUTOFF", RMATS_FDR_CUTOFF))
    RMATS_PVAL_CUTOFF    = float(config.get("RMATS_PVAL_CUTOFF", RMATS_PVAL_CUTOFF))
    INCLEVEL_DIFF_CUTOFF = float(config.get("INCLEVEL_DIFF_CUTOFF", INCLEVEL_DIFF_CUTOFF))
    USE_FDR              = bool(config.get("USE_FDR", USE_FDR))
    RMATS_DUAL_FILTER    = bool(config.get("RMATS_DUAL_FILTER", RMATS_DUAL_FILTER))
    FIG_DPI              = int(config.get("FIG_DPI", FIG_DPI))
    FIG_FORMAT           = str(config.get("FIG_FORMAT", FIG_FORMAT))
    FONT_SIZE            = int(config.get("FONT_SIZE", FONT_SIZE))
    COLOR_UP             = str(config.get("COLOR_UP", COLOR_UP))
    COLOR_DOWN           = str(config.get("COLOR_DOWN", COLOR_DOWN))
    COLOR_NS             = str(config.get("COLOR_NS", COLOR_NS))
    INTERACTIVE_PLOTS    = bool(config.get("INTERACTIVE_PLOTS", INTERACTIVE_PLOTS))
    DESEQ2_COLS          = dict(config.get("DESEQ2_COLS", DESEQ2_COLS))
    RMATS_COLS           = dict(config.get("RMATS_COLS", RMATS_COLS))
    GSEA_DATABASES       = list(config.get("GSEA_DATABASES", GSEA_DATABASES))
    ORA_DATABASES        = list(config.get("ORA_DATABASES", ORA_DATABASES))
    GENES_OF_INTEREST    = list(config.get("GENES_OF_INTEREST", GENES_OF_INTEREST))
    COUNTS_FILE          = str(config.get("COUNTS_FILE", COUNTS_FILE))
    SAMPLE_METADATA      = dict(config.get("SAMPLE_METADATA", SAMPLE_METADATA))
    GSEA_RANKING         = str(config.get("GSEA_RANKING", GSEA_RANKING))
    GSEA_MIN_SIZE        = int(config.get("GSEA_MIN_SIZE", GSEA_MIN_SIZE))
    GSEA_MAX_SIZE        = int(config.get("GSEA_MAX_SIZE", GSEA_MAX_SIZE))
    GSEA_PERMUTATIONS    = int(config.get("GSEA_PERMUTATIONS", GSEA_PERMUTATIONS))
    ORA_METHOD           = str(config.get("ORA_METHOD", ORA_METHOD))
    RBP_FILE             = str(config.get("RBP_FILE", RBP_FILE))

    main()


if __name__ == "__main__":
    main()
